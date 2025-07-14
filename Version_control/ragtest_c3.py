# --- START OF FILE ragtest_c3_query_adapt_vision.py --- # Renamed for clarity

import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd # Added for CSV/XLSX
import faiss
# import streamlit as st # Import only if needed / within UI code
import time
from sentence_transformers import SentenceTransformer
# from transformers import T5ForConditionalGeneration, T5Tokenizer # Keep if needed elsewhere
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage # Added SystemMessage
from dotenv import load_dotenv
from pptx import Presentation # Added for PPTX
from pptx.enum.shapes import MSO_SHAPE_TYPE # Needed for placeholder check
import sys # Import sys for checking streamlit context

# --- Imports for Image Captioning ---
import io
import base64
try:
    from PIL import Image
    # Suppress DecompressionBombWarning for large images if needed, use cautiously
    Image.MAX_IMAGE_PIXELS = None
except ImportError:
    print("Pillow library not found. Image captioning will be disabled.")
    print("Please install it using: pip install Pillow")
    Image = None # Set Image to None if import fails
# --- End Imports for Image Captioning ---


# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted from RAG2.py) ---

# Configure logging (optional for Streamlit, but good for debugging)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system_vision.log"), # Changed log filename
        # logging.StreamHandler() # Can be noisy in Streamlit console
    ]
)
logger = logging.getLogger(__name__)

# Check if Pillow is available after trying to import
if Image is None:
    logger.warning("Pillow library not installed. PDF image captioning feature is disabled.")


class RAGConfig:
    """Configuration class for RAG system."""
    # (This class remains unchanged from ragtest_c3_query_adapt.py)
    def __init__(self, config_path="config.ini"):
        """Initialize configuration from config file or defaults."""
        self.config = configparser.ConfigParser()
        self.defaults = {
            "PATHS": {"data_dir": "./Data/", "index_dir": "./Faiss_index/", "log_file": "rag_system_vision.log"}, # Updated log file default
            "MODELS": {"encoder_model": "sentence-transformers/all-MiniLM-L6-v2", "llm_model": "meta-llama/llama-4-scout-17b-16e-instruct", "device": "auto"}, # Ensure llm_model supports vision if needed
            "PARAMETERS": {"chunk_size": "200", "overlap": "50", "k_retrieval": "5", "temperature": "0.2", "max_context_tokens": "4000", "max_chars_per_element": "1000", "pptx_merge_threshold_words": "50"},
            "SUPPORTED_EXTENSIONS": {"extensions": ".pdf, .xlsx, .csv, .pptx"}
        }
        if os.path.exists(config_path):
            try: self.config.read(config_path); logger.info(f"Loaded configuration from {config_path}"); self._ensure_defaults()
            except Exception as e: logger.error(f"Error reading config file {config_path}: {e}"); self._set_defaults()
        else: logger.warning(f"Config file {config_path} not found, using defaults"); self._set_defaults()
    def _set_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section)
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value)
    def _ensure_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section); logger.info(f"Added missing section [{section}] from defaults.")
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value); logger.info(f"Added missing option '{option} = {value}' to section [{section}] from defaults.")
    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    @property
    def llm_model(self): return self.config.get("MODELS", "llm_model")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto")
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=200)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=50)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=5)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.2)
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=4000)
    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1000)
    @property
    def pptx_merge_threshold_words(self): return self.config.getint("PARAMETERS", "pptx_merge_threshold_words", fallback=50)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf")
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])


class RAGSystem:
    """Retrieval-Augmented Generation system for various document types."""
    # (Init remains largely unchanged, added check for Pillow)
    def __init__(self, config_path="config.ini"):
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system. Device: {self.config.device}, Supported files: {self.config.supported_extensions}")
        if Image is None:
            logger.warning("Pillow library not available. PDF image captioning feature is disabled.")
            self.image_captioning_enabled = False
        else:
             self.image_captioning_enabled = True # Assume enabled if Pillow loaded

        self._is_streamlit = "streamlit" in sys.modules
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)
        try:
            if self._is_streamlit:
                import streamlit as st
                with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."): self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            else: logger.info(f"Loading embedding model ({self.config.encoder_model})..."); self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            logger.info("Local embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}"); st.stop()
            else: raise
        self.file_chunks = {}
        self.faiss_indexes = {}
        self.processed_files = set()
        # Initialize Groq LLM client once if needed frequently
        self._llm = None


    def _get_llm(self):
        """Initializes and returns the Groq LLM client."""
        if self._llm is None:
            api_key = os.getenv("GROQ_API_KEY")
            if not api_key:
                logger.error("GROQ_API_KEY not found.")
                if self._is_streamlit:
                    import streamlit as st
                    st.error("Error: GROQ_API_KEY is not set.")
                raise ValueError("GROQ_API_KEY not configured.")
            try:
                # *** IMPORTANT: Ensure self.config.llm_model is a VISION-capable model ***
                # If not, image captioning will fail when calling invoke.
                self._llm = ChatGroq(
                    temperature=self.config.temperature,
                    model_name=self.config.llm_model,
                    groq_api_key=api_key,
                )
                logger.info(f"Groq LLM client initialized ({self.config.llm_model})")
            except Exception as e:
                 logger.error(f"Failed to initialize Groq client: {e}", exc_info=True)
                 if self._is_streamlit:
                     import streamlit as st
                     st.error(f"Failed to initialize LLM. Error: {e}")
                 raise
        return self._llm


    def clean_text(self, text):
        """Clean and normalize text."""
        # (This function remains unchanged)
        if not isinstance(text, str): text = str(text)
        text = re.sub(r"\(cid:.*?\)", "", text)
        text = re.sub(r"\s+", " ", text).strip()
        text = text.replace('\n', ' ').replace('\r', '')
        return text

    # --- Specific Content Extractors ---

    # ==========================================================================
    # START OF NEW HELPER METHOD FOR IMAGE CAPTIONING
    # ==========================================================================
    def _caption_image(self, image_bytes, page_num, img_idx):
        """Generates a caption for an image using the multimodal LLM."""
        if not self.image_captioning_enabled:
            logger.debug("Image captioning disabled (Pillow not loaded). Skipping.")
            return None

        try:
            # 1. Validate and Standardize Image Format (using Pillow)
            try:
                img = Image.open(io.BytesIO(image_bytes))
                img.verify() # Check if image data is corrupt
                # Re-open after verify
                img = Image.open(io.BytesIO(image_bytes))
                # Convert to RGB if it's not (common requirement for models)
                if img.mode != 'RGB':
                    logger.debug(f"Converting image {img_idx} on page {page_num} from {img.mode} to RGB")
                    img = img.convert('RGB')

                # Resize large images to prevent potential issues (optional, adjust size as needed)
                max_size = (1024, 1024)
                if img.size[0] > max_size[0] or img.size[1] > max_size[1]:
                    logger.debug(f"Resizing large image {img_idx} on page {page_num} from {img.size}")
                    img.thumbnail(max_size, Image.Resampling.LANCZOS)

                # Standardize to PNG for base64 encoding
                buffer = io.BytesIO()
                img.save(buffer, format="PNG")
                image_bytes_standardized = buffer.getvalue()
                image_format = "png"
            except Exception as e_img:
                logger.warning(f"Could not process image {img_idx} on page {page_num}. Skipping captioning. Error: {e_img}")
                return None # Skip if image is invalid/unsupported

            # 2. Encode image to base64
            base64_image = base64.b64encode(image_bytes_standardized).decode('utf-8')
            data_url = f"data:image/{image_format};base64,{base64_image}"

            # 3. Prepare Prompt for LLM
            caption_prompt = "Describe this image concisely, focusing on elements relevant to a business proposal (e.g., diagrams, charts, product visuals, key text visible within the image). If it's decorative or lacks substantive content, say 'Decorative image'."

            # 4. Call LLM (ensure your Groq model supports vision)
            llm = self._get_llm() # Get the initialized LLM client

            # Construct multimodal message (using Langchain core messages format)
            message = HumanMessage(
                content=[
                    {"type": "text", "text": caption_prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": data_url},
                    },
                ]
            )

            logger.info(f"Requesting caption from LLM for image {img_idx} on page {page_num}...")
            start_time = time.time()
            # Use invoke for single call with Langchain messages
            response = llm.invoke([message]) # Pass as a list
            end_time = time.time()
            caption = response.content.strip() if hasattr(response, 'content') else str(response).strip()

            if caption:
                logger.info(f"Caption received for image {img_idx} on page {page_num} in {end_time - start_time:.2f}s: '{caption[:100]}...'")
                # Optional: Filter out very short/uninformative captions if desired
                if caption.lower() == "decorative image." or len(caption) < 15:
                     logger.info(f"Skipping trivial caption for image {img_idx} page {page_num}.")
                     return None
                return caption
            else:
                logger.warning(f"Received empty caption for image {img_idx} on page {page_num}.")
                return "[Image present, but caption could not be generated]"

        except Exception as e:
            logger.error(f"Failed to generate caption for image {img_idx} on page {page_num}: {e}", exc_info=True)
            # Check if the error is related to the model not supporting images
            if "vision" in str(e).lower() or "image" in str(e).lower():
                 logger.error(f"Potential Issue: The configured LLM '{self.config.llm_model}' might not support image input.")
            return "[Image captioning failed due to error]" # Placeholder on failure

    # ==========================================================================
    # END OF NEW HELPER METHOD FOR IMAGE CAPTIONING
    # ==========================================================================


    # ==========================================================================
    # START OF MODIFIED _extract_pdf METHOD
    # ==========================================================================
    def _extract_pdf(self, file_path):
        all_content = []
        base_filename = os.path.basename(file_path)
        try:
            # Using pdfplumber as in the original code
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Extracting from PDF: {base_filename} ({total_pages} pages) - Image Captioning {'Enabled' if self.image_captioning_enabled else 'Disabled'}")
                page_bar = None
                # Show progress bar sooner if captioning might take time
                show_progress_threshold = 5 if self.image_captioning_enabled else 10
                if self._is_streamlit and total_pages > show_progress_threshold:
                    import streamlit as st
                    prog_text = f"Processing pages & images from {base_filename}..." if self.image_captioning_enabled else f"Extracting pages from {base_filename}..."
                    page_bar = st.progress(0, text=prog_text)

                for page_num, page in enumerate(pdf.pages):
                    current_page_num = page_num + 1
                    logger.debug(f"Processing Page {current_page_num}/{total_pages} of {base_filename}")

                    # --- Existing Text Extraction ---
                    try:
                        text = page.extract_text(layout="normal") or ""
                        cleaned_text = self.clean_text(text)
                        if cleaned_text:
                            all_content.append({
                                "type": "text",
                                "content": cleaned_text,
                                "source_info": {"page": current_page_num},
                                "file_type": "pdf"
                            })
                            logger.debug(f" P{current_page_num}: Added text block ({len(cleaned_text)} chars)")
                    except Exception as e_text:
                         logger.warning(f"Could not extract text on page {current_page_num} in {base_filename}: {e_text}")


                    # --- Existing Table Extraction ---
                    try:
                        page_tables = page.extract_tables()
                        if page_tables: # Check if list is not empty
                            for table_idx, table in enumerate(page_tables):
                                if table:
                                    table_df = pd.DataFrame(table).fillna('')
                                    table_string = table_df.to_string(index=False, header=True)
                                    cleaned_table_string = self.clean_text(table_string)
                                    if cleaned_table_string:
                                        all_content.append({
                                            "type": "table",
                                            "content": cleaned_table_string,
                                            "source_info": {"page": current_page_num, "table_index": table_idx},
                                            "file_type": "pdf"
                                        })
                                        logger.debug(f" P{current_page_num}: Added table block {table_idx} ({len(cleaned_table_string)} chars)")
                    except Exception as e_table:
                        logger.warning(f"Could not extract table(s) on page {current_page_num} in {base_filename}: {e_table}")

                    # --- *** NEW: Image Extraction and Captioning *** ---
                    if self.image_captioning_enabled:
                        try:
                            page_images = page.images
                            if page_images: # Check if list is not empty
                                logger.debug(f" P{current_page_num}: Found {len(page_images)} potential image object(s).")
                                for img_idx, img_props in enumerate(page_images):
                                    logger.debug(f" P{current_page_num}: Processing image object {img_idx}")
                                    image_bytes = None
                                    try:
                                        # Access image data. pdfplumber structure can vary.
                                        # Primary way: Access the stream associated via 'stream' key if present
                                        if 'stream' in img_props and hasattr(img_props['stream'], 'get_data'):
                                            image_bytes = img_props['stream'].get_data()
                                            logger.debug(f"  P{current_page_num}-Img{img_idx}: Got data from stream ({len(image_bytes)} bytes)")
                                        # Fallback: Check 'obj' if 'stream' isn't there or lacks get_data
                                        elif 'obj' in img_props and hasattr(img_props['obj'], 'get_data'):
                                            image_bytes = img_props['obj'].get_data()
                                            logger.debug(f"  P{current_page_num}-Img{img_idx}: Got data from obj ({len(image_bytes)} bytes)")
                                        else:
                                             logger.warning(f"  P{current_page_num}-Img{img_idx}: Could not find get_data method in image properties.")
                                             # Uncomment to log full properties if debugging is needed
                                             # logger.debug(f"  Image Properties: {img_props}")


                                    except Exception as e_img_data:
                                        logger.warning(f" P{current_page_num}: Could not get image data for image {img_idx}. Error: {e_img_data}")
                                        continue # Skip to next image if data extraction fails

                                    if image_bytes:
                                        # Call the captioning function
                                        caption = self._caption_image(image_bytes, current_page_num, img_idx)
                                        if caption:
                                            all_content.append({
                                                "type": "image_caption", # New type
                                                # Prefix caption for clarity in context
                                                "content": f"[Image Caption]: {caption}",
                                                "source_info": {"page": current_page_num, "image_index": img_idx},
                                                "file_type": "pdf"
                                            })
                                            logger.debug(f" P{current_page_num}: Added image caption block {img_idx} ({len(caption)} chars)")
                                    else:
                                        logger.debug(f" P{current_page_num}: No image bytes found or extracted for image {img_idx}")

                        except Exception as e_img_loop:
                             logger.error(f" P{current_page_num}: Error processing images on page: {e_img_loop}", exc_info=True)
                    # --- *** End of Image Processing *** ---

                    if page_bar:
                        progress_percent = min(1.0, (page_num + 1) / total_pages)
                        prog_text = f"Processing page {current_page_num}/{total_pages} (Text, Tables, Images)..." if self.image_captioning_enabled else f"Extracting page {current_page_num}/{total_pages}..."
                        page_bar.progress(progress_percent, text=prog_text)

                if page_bar: page_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks (incl. text, tables, captions) from PDF: {base_filename}")
            return all_content
        except pdfplumber.pdfminer.pdfdocument.PDFEncryptionError as e_encrypt:
             logger.error(f"Cannot process encrypted PDF: {base_filename}. Skipping. Error: {e_encrypt}")
             if self._is_streamlit: import streamlit as st; st.error(f"Cannot process encrypted PDF: {base_filename}")
             return []
        except Exception as e:
            logger.error(f"Error extracting content from PDF {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    # ==========================================================================
    # END OF MODIFIED _extract_pdf METHOD
    # ==========================================================================


    # (_extract_xlsx, _extract_csv, _extract_pptx remain unchanged)
    def _extract_xlsx(self, file_path):
        all_content = []
        base_filename = os.path.basename(file_path)
        max_chars = self.config.max_chars_per_element
        try:
            excel_file = pd.ExcelFile(file_path); logger.info(f"Extracting from XLSX: {base_filename} (Sheets: {', '.join(excel_file.sheet_names)})")
            sheet_bar = None; total_sheets = len(excel_file.sheet_names)
            if self._is_streamlit and total_sheets > 1: import streamlit as st; sheet_bar = st.progress(0, text=f"Extracting sheets from {base_filename}...")
            for i, sheet_name in enumerate(excel_file.sheet_names):
                try:
                    df = excel_file.parse(sheet_name); df = df.fillna(''); sheet_string = df.to_string(index=False, header=True, na_rep=''); cleaned_content = self.clean_text(sheet_string)
                    if cleaned_content:
                        if len(cleaned_content) > max_chars * 10: logger.warning(f"Sheet '{sheet_name}' in {base_filename} is very large, truncating."); cleaned_content = cleaned_content[:max_chars*10] + "... (truncated sheet)"
                        all_content.append({"type": "sheet_data", "content": cleaned_content, "source_info": {"sheet": sheet_name, "rows": f"1-{len(df)}"}, "file_type": "xlsx"})
                    logger.debug(f"Extracted content from sheet '{sheet_name}' in {base_filename}")
                except Exception as e_sheet: logger.warning(f"Could not process sheet '{sheet_name}' in {base_filename}: {e_sheet}")
                if sheet_bar: sheet_bar.progress(min(1.0, (i + 1) / total_sheets), text=f"Extracting sheet {i+1}/{total_sheets} ('{sheet_name}') from {base_filename}...")
            if sheet_bar: sheet_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks (sheets) from XLSX: {base_filename}"); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from XLSX {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def _extract_csv(self, file_path):
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element
        try:
            encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']; df = None
            for enc in encodings_to_try:
                 try: df = pd.read_csv(file_path, encoding=enc, low_memory=False); logger.info(f"Read CSV {base_filename} using encoding: {enc}"); break
                 except UnicodeDecodeError: logger.debug(f"CSV {base_filename} encoding failed with {enc}"); continue
                 except Exception as e_read: logger.warning(f"Error reading CSV {base_filename} with encoding {enc}: {e_read}")
            if df is None: logger.error(f"Could not read CSV {base_filename} with any tried encoding."); return []
            df = df.fillna(''); csv_string = df.to_string(index=False, header=True, na_rep=''); cleaned_content = self.clean_text(csv_string)
            if cleaned_content:
                if len(cleaned_content) > max_chars * 20: logger.warning(f"CSV {base_filename} is very large, truncating."); cleaned_content = cleaned_content[:max_chars*20] + "... (truncated CSV)"
                all_content.append({"type": "csv_data", "content": cleaned_content, "source_info": {"rows": f"1-{len(df)}"}, "file_type": "csv"})
            logger.info(f"Extracted content block from CSV: {base_filename} ({len(df)} rows)"); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from CSV {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def _extract_pptx(self, file_path):
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element; merge_threshold_words = self.config.pptx_merge_threshold_words
        try:
            prs = Presentation(file_path); logger.info(f"Extracting from PPTX: {base_filename} ({len(prs.slides)} slides) using merge strategy (threshold: {merge_threshold_words} words).")
            slide_bar = None; total_slides = len(prs.slides)
            if self._is_streamlit and total_slides > 1: import streamlit as st; slide_bar = st.progress(0, text=f"Extracting slides from {base_filename}...")
            pending_title_slide_data = None
            for i, slide in enumerate(prs.slides):
                current_slide_number = i + 1; current_slide_title_text = ""; current_slide_other_texts = []; current_slide_has_title_placeholder = False; title_shape = None
                try:
                    if slide.shapes.title: title_shape = slide.shapes.title; current_slide_has_title_placeholder = True
                except AttributeError: pass # Slide might not have a title placeholder
                if title_shape and title_shape.has_text_frame: cleaned_title = self.clean_text(title_shape.text_frame.text); current_slide_title_text = cleaned_title if cleaned_title else ""
                for shape in slide.shapes:
                    if shape == title_shape: continue # Already processed title
                    is_placeholder = shape.is_placeholder; is_body_placeholder = False
                    if is_placeholder:
                        try: ph_type = shape.placeholder_format.type; is_body_placeholder = ph_type in [MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT, MSO_SHAPE_TYPE.SUBTITLE] # Check specific types
                        except AttributeError: pass # Shape might not have placeholder_format
                    if shape.has_text_frame:
                        text = shape.text_frame.text; cleaned = self.clean_text(text)
                        if cleaned:
                             if len(cleaned) > max_chars: cleaned = cleaned[:max_chars] + "...(truncated shape)"
                             prefix = "[Body Placeholder]: " if is_body_placeholder else "[Text Box]: " if not is_placeholder else "[Other Placeholder]: " # Add prefix for context
                             current_slide_other_texts.append(prefix + cleaned)
                # Notes Slide Extraction
                if slide.has_notes_slide:
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text; cleaned_notes = self.clean_text(notes_text)
                        if cleaned_notes:
                            if len(cleaned_notes) > max_chars * 2: cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated notes)"
                            current_slide_other_texts.append(f"[Slide Notes]: {cleaned_notes}")
                    except Exception as e_notes: logger.warning(f"Could not extract notes from slide {current_slide_number}: {e_notes}")

                # Combine title and other texts for the current slide
                current_slide_full_content = current_slide_title_text
                if current_slide_other_texts: current_slide_full_content += ("\n" if current_slide_full_content else "") + "\n".join(current_slide_other_texts)
                current_slide_full_content = current_slide_full_content.strip()
                other_text_word_count = sum(len(s.split()) for s in current_slide_other_texts) # Word count excluding title

                # Merge Logic
                merged_content_block = None; should_merge = False
                # Condition: If there's a pending slide, it had a title, few words, and current slide has content
                if pending_title_slide_data:
                    if pending_title_slide_data['has_title'] and pending_title_slide_data['other_words'] <= merge_threshold_words and current_slide_full_content:
                        should_merge = True
                        logger.info(f"Merging slide {pending_title_slide_data['number']} (Title: '{pending_title_slide_data['title'][:30]}...') with slide {current_slide_number}.")

                if should_merge:
                    # Create merged text block
                    merged_text = f"[Content from Slide {pending_title_slide_data['number']} - Title: {pending_title_slide_data['title']}]:\n{pending_title_slide_data['content']}\n\n---\n\n[Content from Slide {current_slide_number}]:\n{current_slide_full_content}"
                    merged_content_block = {"type": "slide_text_merged", "content": merged_text.strip(), "source_info": {"slide": current_slide_number, "merged_from": pending_title_slide_data['number']}, "file_type": "pptx"}
                    all_content.append(merged_content_block)
                    pending_title_slide_data = None # Clear pending data after merging
                else:
                    # If not merging, process the pending slide (if any) first
                    if pending_title_slide_data:
                         if pending_title_slide_data['content']: # Only add if it had content
                             all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
                             logger.debug(f"Adding non-merged pending slide {pending_title_slide_data['number']}.")
                         pending_title_slide_data = None # Clear pending data

                    # Now decide about the *current* slide
                    # Is it a potential *new* title slide to hold? (Has title placeholder, few other words, and some content)
                    if current_slide_has_title_placeholder and other_text_word_count <= merge_threshold_words and current_slide_full_content:
                        logger.debug(f"Slide {current_slide_number} is potential title slide. Holding. (Title: '{current_slide_title_text[:30]}...', Words: {other_text_word_count})");
                        pending_title_slide_data = {
                            "content": current_slide_full_content,
                            "number": current_slide_number,
                            "has_title": current_slide_has_title_placeholder,
                            "other_words": other_text_word_count,
                            "title": current_slide_title_text # Store title for merge context
                        }
                    else:
                        # Otherwise, just add the current slide's content if it exists
                        if current_slide_full_content:
                            all_content.append({"type": "slide_text", "content": current_slide_full_content, "source_info": {"slide": current_slide_number}, "file_type": "pptx"})
                            logger.debug(f"Adding regular slide {current_slide_number}.")

                if slide_bar: slide_bar.progress(min(1.0, (i + 1) / total_slides), text=f"Processing slide {current_slide_number}/{total_slides}...")

            # Handle any remaining pending slide at the end of the presentation
            if pending_title_slide_data:
                logger.debug(f"Processing pending title slide {pending_title_slide_data['number']} at end.")
                if pending_title_slide_data['content']:
                    all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})

            if slide_bar: slide_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PPTX {base_filename} (Merge strategy applied)."); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PPTX {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []


    def extract_content(self, file_path):
        """Extract content based on file extension."""
        # (This function remains unchanged)
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf': return self._extract_pdf(file_path)
        elif extension == '.xlsx': return self._extract_xlsx(file_path)
        elif extension == '.csv': return self._extract_csv(file_path)
        elif extension == '.pptx': return self._extract_pptx(file_path)
        else: logger.warning(f"Unsupported file type skipped: {file_path}"); return []

    def chunk_content(self, all_content):
        """Split extracted content into chunks with overlap, preserving metadata."""
        # (This function remains unchanged - handles text content from any source type)
        chunks = []; total_words_estimate = 0; chunk_bar = None
        if not all_content: return chunks
        try: total_words_estimate = sum(len(str(item.get('content', '')).split()) for item in all_content)
        except Exception: pass
        # Show progress bar for chunking if estimate is large
        if self._is_streamlit and total_words_estimate > 5000:
             import streamlit as st; chunk_bar = st.progress(0, text=f"Chunking content ({len(all_content)} blocks)...")
        words_processed = 0
        logger.info(f"Starting chunking process for {len(all_content)} content blocks (estimated {total_words_estimate} words)...")
        for item_index, item in enumerate(all_content):
            content = item.get('content', '')
            source_info = item.get('source_info', {})
            file_type = item.get('file_type', 'unknown')
            content_type = item.get('type', 'unknown') # e.g., text, table, image_caption, slide_text
            if not isinstance(content, str): content = str(content) # Ensure content is string

            words = content.split()
            if not words:
                logger.debug(f"Skipping empty content block {item_index} (Type: {content_type}, Source: {source_info})")
                continue

            item_chunks_created = 0
            # Use stride to create overlapping chunks
            stride = self.config.chunk_size - self.config.overlap
            for i in range(0, len(words), stride):
                chunk_words = words[i : i + self.config.chunk_size]
                chunk_text = " ".join(chunk_words)
                if chunk_text: # Ensure chunk is not empty
                    chunks.append({
                        "content": chunk_text,
                        "source_info": source_info, # Carry over metadata
                        "file_type": file_type,
                        "type": content_type
                    })
                    item_chunks_created += 1

            logger.debug(f" Block {item_index} (Type: {content_type}, {len(words)} words) -> {item_chunks_created} chunks created.")

            if chunk_bar:
                words_processed += len(words)
                progress_percent = min(1.0, words_processed / total_words_estimate) if total_words_estimate > 0 else 0
                chunk_bar.progress(progress_percent, text=f"Chunking content... ({len(chunks)} chunks total)")

        if chunk_bar: chunk_bar.empty()
        logger.info(f"Created {len(chunks)} chunks from {len(all_content)} content blocks."); return chunks

    # --- Path Helpers --- (Unchanged)
    def _get_safe_filename(self, file_name): base_name = os.path.splitext(file_name)[0]; return re.sub(r'[^\w\.-]', '_', base_name)
    def get_index_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.index")
    def get_embedding_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.npy")
    def get_chunks_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.json")

    # --- FAISS Index and Chunk Loading/Saving --- (Unchanged)
    def load_faiss_index(self, file_name, embedding_dim):
        index_path = self.get_index_path(file_name)
        if os.path.exists(index_path):
            try: index = faiss.read_index(index_path); logger.info(f"Loaded FAISS index for {file_name} ({index.ntotal} vectors)"); return index
            except Exception as e: logger.error(f"Error reading FAISS index {index_path}: {e}. Creating new index.")
        logger.info(f"FAISS index not found at {index_path}. Will create new."); return faiss.IndexFlatL2(embedding_dim)
    def save_chunks(self, file_name, chunks):
        chunks_path = self.get_chunks_path(file_name)
        try:
            with open(chunks_path, 'w', encoding='utf-8') as f: json.dump(chunks, f, indent=2); logger.info(f"Saved {len(chunks)} chunks to {chunks_path}")
        except Exception as e: logger.error(f"Error saving chunks for {file_name} to {chunks_path}: {e}");
    def load_chunks(self, file_name):
        chunks_path = self.get_chunks_path(file_name)
        if os.path.exists(chunks_path):
            try:
                with open(chunks_path, 'r', encoding='utf-8') as f: chunks = json.load(f); logger.info(f"Loaded {len(chunks)} chunks from {chunks_path}"); return chunks
            except Exception as e: logger.error(f"Error loading/decoding chunks from {chunks_path}: {e}. Will re-process."); return None
        logger.info(f"Chunks file not found at {chunks_path}"); return None

    # --- Document Processing --- (Unchanged, relies on extract_content)
    def process_files(self, progress_callback=None):
        self.file_chunks = {}; self.faiss_indexes = {}; self.processed_files = set(); data_dir = self.config.data_dir; supported_ext = self.config.supported_extensions
        try: all_files = os.listdir(data_dir); process_list = sorted([f for f in all_files if f.lower().endswith(supported_ext)])
        except Exception as e: logger.error(f"Error listing files in data directory {data_dir}: {e}", exc_info=True); return False
        if not process_list: logger.warning(f"No supported files found in {data_dir}"); return True
        logger.info(f"Processing {len(process_list)} supported files from {data_dir}"); embedding_dim = self.encoder_model.get_sentence_embedding_dimension(); total_files = len(process_list)
        if progress_callback: progress_callback(f"Found {total_files} file(s). Starting processing...", current_step=0, total_steps=total_files)
        for idx, file_name in enumerate(process_list):
            current_file_msg = f"Processing ({idx+1}/{total_files}): {file_name}"; logger.info(f"--- {current_file_msg} ---")
            if progress_callback: progress_callback(current_file_msg, current_step=idx, total_steps=total_files)
            file_path = os.path.join(data_dir, file_name); index_path = self.get_index_path(file_name); emb_path = self.get_embedding_path(file_name); chunks_path = self.get_chunks_path(file_name)
            try:
                chunks = self.load_chunks(file_name)
                if chunks is None:
                    if progress_callback: progress_callback(f"{current_file_msg} - Extracting...", stage="Extracting")
                    start_extract = time.time()
                    all_content = self.extract_content(file_path); # This now potentially includes image captions
                    end_extract = time.time()
                    logger.info(f"Extraction for {file_name} took {end_extract - start_extract:.2f}s. Found {len(all_content)} content blocks.")
                    if not all_content: logger.warning(f"No content extracted from {file_name}, skipping."); continue

                    if progress_callback: progress_callback(f"{current_file_msg} - Chunking...", stage="Chunking")
                    start_chunk = time.time()
                    chunks = self.chunk_content(all_content);
                    end_chunk = time.time()
                    logger.info(f"Chunking for {file_name} took {end_chunk - start_chunk:.2f}s. Created {len(chunks)} chunks.")
                    if not chunks: logger.warning(f"No chunks created for {file_name}, skipping."); continue
                    self.save_chunks(file_name, chunks)

                self.file_chunks[file_name] = chunks; logger.debug(f"Stored {len(chunks)} chunks for {file_name}")
                faiss_index = None; regenerate_embeddings = False
                if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                    if progress_callback: progress_callback(f"{current_file_msg} - Verifying...", stage="Verifying")
                    try:
                        embeddings = np.load(emb_path)
                        if embeddings.ndim != 2 or embeddings.shape[1] != embedding_dim or embeddings.shape[0] != len(chunks): logger.warning(f"Data mismatch for {file_name} (Embeddings: {embeddings.shape}, Chunks: {len(chunks)}). Regenerating embeddings."); regenerate_embeddings = True
                        else: faiss_index = self.load_faiss_index(file_name, embedding_dim); logger.info(f"Verified existing data for {file_name}")
                    except Exception as e: logger.error(f"Error verifying existing data for {file_name}: {e}. Regenerating...", exc_info=True); regenerate_embeddings = True; faiss_index = None
                else: logger.info(f"Index/embeddings/chunks missing for {file_name}. Will generate."); regenerate_embeddings = True

                if regenerate_embeddings or faiss_index is None:
                    logger.info(f"Generating embeddings/index for {file_name}...")
                    if progress_callback: progress_callback(f"{current_file_msg} - Embedding...", stage="Embedding")
                    content_list = [chunk['content'] for chunk in chunks];
                    if not content_list: logger.warning(f"No content to embed for {file_name}."); continue
                    start_embed = time.time()
                    # Note: show_progress_bar=True can be noisy if running many files non-interactively
                    embeddings = self.encoder_model.encode(content_list, batch_size=64, show_progress_bar=False, convert_to_numpy=True).astype('float32')
                    end_embed = time.time()
                    logger.info(f"Embedding for {file_name} took {end_embed - start_embed:.2f}s.")
                    if embeddings.shape[0] == 0: logger.warning(f"Embedding yielded no vectors for {file_name}."); continue
                    np.save(emb_path, embeddings); logger.info(f"Saved {embeddings.shape[0]} embeddings to {emb_path}")

                    if progress_callback: progress_callback(f"{current_file_msg} - Indexing...", stage="Indexing")
                    start_index = time.time()
                    faiss_index = faiss.IndexFlatL2(embedding_dim); faiss_index.add(embeddings); faiss.write_index(faiss_index, index_path);
                    end_index = time.time()
                    logger.info(f"Indexing for {file_name} took {end_index - start_index:.2f}s.")
                    logger.info(f"Saved FAISS index ({faiss_index.ntotal} vectors) to {index_path}")

                if faiss_index is not None and faiss_index.ntotal > 0:
                     self.faiss_indexes[file_name] = faiss_index; self.processed_files.add(file_name); logger.debug(f"Stored FAISS index for {file_name}")
                else: logger.warning(f"No valid index created or loaded for {file_name}.")

            except Exception as e:
                logger.error(f"Failed to process {file_name}: {e}", exc_info=True)
                if progress_callback: progress_callback(f"❌ Error processing {file_name}: {e}", is_error=True)
                # Attempt cleanup on failure
                try:
                    for p in [index_path, emb_path, chunks_path]:
                        if os.path.exists(p): os.remove(p); logger.debug(f"Cleaned up {p} for failed file {file_name}")
                except OSError as e_clean: logger.error(f"Error cleaning up files for {file_name}: {e_clean}")
                # Remove from in-memory structures if partially processed
                if file_name in self.file_chunks: del self.file_chunks[file_name]
                if file_name in self.faiss_indexes: del self.faiss_indexes[file_name]

        final_indexed_count = len(self.faiss_indexes)
        if final_indexed_count > 0: logger.info(f"--- Processing Complete. Indexed {final_indexed_count}/{total_files} documents. ---"); return True
        else: logger.warning(f"--- Processing Complete. No documents indexed. ---"); return True


    # --- Querying Logic ---

    # (query_files remains unchanged)
    def query_files(self, query):
        """Query FAISS index for a single query string."""
        if not self.faiss_indexes:
            logger.warning("No indexes available for querying.")
            return {} # Return empty results, let caller handle

        query_results = {}
        try:
            logger.info(f"Encoding query for Faiss: '{query[:100]}...'")
            query_embedding = self.encoder_model.encode(query, convert_to_numpy=True).astype("float32")
            query_embedding = np.array([query_embedding])
            if query_embedding.ndim != 2: raise ValueError("Query embedding shape error")

            # Search across all indexes for this single query
            for file_name, index in self.faiss_indexes.items():
                if index is None or index.ntotal == 0:
                    logger.debug(f"Skipping empty or non-existent index for {file_name}")
                    continue
                try:
                    k_search = min(self.config.k_retrieval, index.ntotal)
                    logger.debug(f"Searching index {file_name} (size {index.ntotal}) for k={k_search} nearest neighbors.")
                    D, I = index.search(query_embedding, k=k_search) # D=distances (L2 squared), I=indices
                    indices, distances = I[0], D[0]

                    current_file_chunks = self.file_chunks.get(file_name)
                    if not current_file_chunks: logger.error(f"CRITICAL: Chunks missing for {file_name} during query!"); continue

                    file_results = []
                    processed_indices = set()
                    for i, idx in enumerate(indices):
                        # FAISS returns -1 for indices if k > ntotal or for padding
                        if idx == -1 or idx in processed_indices: continue
                        # Basic sanity check on index validity
                        if not (0 <= idx < len(current_file_chunks)):
                            logger.warning(f"Invalid index {idx} returned for file {file_name} (max: {len(current_file_chunks)-1}). Skipping.")
                            continue

                        processed_indices.add(idx)
                        chunk = current_file_chunks[idx]
                        # Make sure essential keys exist, provide defaults if not
                        result_item = {
                            "source_info": chunk.get('source_info', {}),
                            "file_type": chunk.get('file_type', 'unknown'),
                            "content": chunk.get('content', ''),
                            "score": round(float(distances[i]), 4), # L2 distance, lower is better
                            "type": chunk.get('type', 'unknown') # e.g., text, table, image_caption
                        }
                        file_results.append(result_item)
                        logger.debug(f"  Match in {file_name}: Index {idx}, Score {result_item['score']:.4f}, Type '{result_item['type']}', Content: '{result_item['content'][:50]}...'")


                    # Sort results for this file by score (lower is better)
                    file_results.sort(key=lambda x: x['score'])
                    if file_results:
                        # Store results under the filename
                        if file_name not in query_results: query_results[file_name] = []
                        query_results[file_name].extend(file_results)
                        logger.debug(f"Found {len(file_results)} relevant chunks in {file_name}.")
                        # Optional: Limit per file here if needed, but aggregation handles global limit later
                        # query_results[file_name] = query_results[file_name][:self.config.k_retrieval]

                except Exception as e_search:
                    logger.error(f"Error searching index for {file_name} with query '{query[:50]}...': {e_search}", exc_info=True)

            logger.debug(f"Faiss search complete for query '{query[:50]}...'. Found results in {len(query_results)} files.")
            return query_results

        except Exception as e_query:
            logger.error(f"Error during Faiss query processing for '{query[:50]}...': {e_query}", exc_info=True)
            return {} # Return empty on error

    # (aggregate_context remains unchanged, handles different chunk types via metadata)
    def aggregate_context(self, query_results, strategy="top_k"):
        all_context = {}
        max_chars = self.config.max_context_tokens * 3 # Allow more chars initially, LLM might truncate
        logger.info(f"Aggregating context using strategy '{strategy}', aiming for ~{self.config.max_context_tokens} tokens (approx. {max_chars} chars)")
        if not query_results:
             logger.warning("No query results provided for context aggregation.")
             return all_context

        # 1. Flatten results and sort globally by score (lower is better)
        flat_results = []
        for file_name, results in query_results.items():
            for res in results:
                flat_results.append({**res, "file_name": file_name}) # Add file_name to each result

        if not flat_results:
             logger.warning("Query results were empty after flattening.")
             return all_context

        flat_results.sort(key=lambda x: x['score'])
        logger.debug(f"Total unique chunks retrieved across all files: {len(flat_results)}")

        # 2. Build context string respecting max_chars and strategy
        aggregated_context_str = ""
        total_aggregated_chars = 0
        added_chunks_count = 0
        context_sources = set() # Keep track of source files included

        # Determine limit based on strategy
        limit = self.config.k_retrieval if strategy == "top_k" else len(flat_results)
        logger.debug(f"Aggregation limit set to {limit} chunks based on strategy '{strategy}'.")

        for i, res in enumerate(flat_results):
            if i >= limit:
                 logger.debug(f"Reached aggregation limit ({limit} chunks).")
                 break

            # Extract details for formatting
            source_info = res.get('source_info', {})
            file_name = res['file_name']
            file_type = res.get('file_type', 'unknown')
            content_type = res.get('type', 'unknown') # text, table, image_caption, etc.
            score = res['score']
            content_body = res['content'] # The actual chunk text

            # Build source string dynamically based on available info
            source_parts = [f"SourceFile: {file_name}"]
            if file_type == 'pptx':
                slide_info = source_info.get('slide', 'N/A')
                merged_from = source_info.get('merged_from')
                slide_display = f"Slide: {slide_info}" + (f" (merged from {merged_from})" if merged_from else "")
                source_parts.append(slide_display)
            elif file_type == 'pdf':
                source_parts.append(f"Page: {source_info.get('page', 'N/A')}")
                # Include image index if it's a caption
                if content_type == 'image_caption':
                     source_parts.append(f"ImageIndex: {source_info.get('image_index', 'N/A')}")
                elif content_type == 'table':
                     source_parts.append(f"TableIndex: {source_info.get('table_index', 'N/A')}")
            elif file_type == 'xlsx':
                source_parts.append(f"Sheet: {source_info.get('sheet', 'N/A')}")
                # Could add row info here if chunking was row-based, but it's sheet-based
            elif file_type == 'csv':
                source_parts.append(f"CSV (Rows: {source_info.get('rows', 'all')})") # Row info might be less useful if chunked

            source_parts.append(f"ContentType: {content_type}")
            source_parts.append(f"RetrievalScore: {score:.4f}") # Add score for context

            # Format the context block
            source_str = ", ".join(source_parts)
            content_header = f"--- Context Block {i+1} from {source_str} ---\n"
            content_to_add = content_header + content_body + "\n\n" # Add double newline separator
            content_chars = len(content_to_add)

            # Check if adding this block exceeds the character limit
            if total_aggregated_chars + content_chars <= max_chars:
                aggregated_context_str += content_to_add
                total_aggregated_chars += content_chars
                added_chunks_count += 1
                context_sources.add(file_name)
            else:
                # If even the *first* chunk is too large, try to truncate it
                 if added_chunks_count == 0:
                     logger.warning(f"First context block (Block {i+1}, {content_chars} chars) alone exceeds max_chars ({max_chars}). Attempting truncation.")
                     # Calculate remaining space, accounting for header and truncation marker
                     remaining_chars = max_chars - len(content_header) - len("\n[...TRUNCATED CONTEXT BLOCK...]\n\n")
                     if remaining_chars > 50: # Only add if there's meaningful space
                         truncated_body = content_body[:remaining_chars]
                         aggregated_context_str += content_header + truncated_body + "\n[...TRUNCATED CONTEXT BLOCK...]\n\n"
                         total_aggregated_chars = len(aggregated_context_str) # Update total chars accurately
                         added_chunks_count += 1
                         context_sources.add(file_name)
                         logger.warning(f"Context truncated to fit {total_aggregated_chars} chars.")
                     else:
                         logger.error(f"First context block too large to fit even truncated ({remaining_chars} available). Cannot create context.")
                 # Stop adding more chunks if limit is reached
                 logger.info(f"Stopping context aggregation at {added_chunks_count} chunks ({total_aggregated_chars}/{max_chars} chars). Limit reached.")
                 break

        final_context = aggregated_context_str.strip()
        if final_context:
            all_context = {
                "combined_context": final_context,
                "source_files": sorted(list(context_sources))
            }
            logger.info(f"Aggregated {total_aggregated_chars} chars from {added_chunks_count} chunks across {len(context_sources)} file(s).")
        else:
            logger.warning("No context aggregated within character limits.")

        return all_context


    # (query_llm remains unchanged)
    def query_llm(self, query, context_data, retry_count=1):
        combined_context = context_data.get("combined_context", "")
        source_files = context_data.get("source_files", [])
        source_file_str = ", ".join(source_files) if source_files else "the provided documents"

        if not combined_context:
            logger.warning(f"LLM query attempted with no context for query: '{query[:100]}...'")
            return f"Could not generate answer: No relevant context was found or aggregated from {source_file_str}."

        try:
            llm = self._get_llm() # Use helper to get/init LLM

            # Construct the prompt using System and Human messages
            # System prompt sets the persona and constraints
            system_prompt = f"""You are an AI assistant answering questions based ONLY on the provided context from document(s): '{source_file_str}'.
Strictly adhere to the following rules:
1. Use ONLY information presented in the '--- START CONTEXT ---' to '--- END CONTEXT ---' section below.
2. Do NOT use any prior knowledge or external information.
3. If the answer is explicitly stated in the context, provide it directly.
4. If the answer can be inferred or calculated *solely* from the provided context, do so.
5. If the context contains conflicting information, state that the context provides conflicting details.
6. If the answer is not found in the context, state clearly: "Based on the provided context from {source_file_str}, I cannot answer this question."
7. Mention the source (e.g., Page number, Sheet name, Slide number, Image caption) from the context headers if it helps clarify the answer, but keep it concise.
8. Do not mention 'The context states...' unless necessary to highlight ambiguity or lack of information."""

            # Human prompt provides the context and the user's question
            human_prompt = f"""--- START CONTEXT ---

{combined_context}

--- END CONTEXT ---

User Question: {query}

Answer based ONLY on the provided context:"""

            full_prompt_messages = [
                SystemMessage(content=system_prompt),
                HumanMessage(content=human_prompt)
            ]

            logger.info(f"Querying Groq model {self.config.llm_model} with {len(combined_context)} context chars from {len(source_files)} file(s)...")
            # logger.debug(f"Full prompt (human part):\n{human_prompt}") # Careful logging full context

            answer = f"Error: LLM query failed after retries." # Default error message
            for attempt in range(retry_count + 1):
                try:
                    start_time = time.time()
                    response = llm.invoke(full_prompt_messages)
                    end_time = time.time()
                    answer = response.content.strip() if hasattr(response, 'content') else str(response).strip()
                    logger.info(f"Groq response received (attempt {attempt+1}/{retry_count+1}) in {end_time - start_time:.2f}s. Length: {len(answer)}")

                    # Basic check for non-empty response. If empty, retry.
                    if not answer and attempt < retry_count:
                        logger.warning("LLM returned an empty response, retrying...")
                        time.sleep(1 * (attempt + 1)) # Exponential backoff basic
                        continue
                    # If response received (even if empty on last try), break the loop
                    return answer or f"Received empty response from LLM after {retry_count+1} attempts for context from {source_file_str}."

                except Exception as e_api:
                    logger.warning(f"Groq API call attempt {attempt+1} failed: {e_api}")
                    # Stop retrying if it's likely a persistent issue (e.g., auth, model error)
                    if "authentication" in str(e_api).lower() or "invalid model" in str(e_api).lower():
                         logger.error("Unrecoverable API error detected, stopping retries.")
                         answer = f"Error: Failed to get answer from LLM. Unrecoverable API error: {e_api}"
                         break
                    if attempt < retry_count:
                        wait_time = 1.5 ** attempt
                        logger.info(f"Retrying Groq query in {wait_time:.1f} seconds...")
                        time.sleep(wait_time)
                    else:
                        # Log final failure after all retries
                        logger.error(f"LLM query failed after {retry_count+1} attempts.")
                        answer = f"Error: Failed to get answer from LLM after {retry_count+1} attempts. Last API Error: {e_api}"
            return answer # Return the result after loop (success, empty, or error)

        except Exception as e_setup:
            logger.error(f"Error setting up or calling Groq API: {e_setup}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"LLM query failed: {e_setup}")
            return f"Error: Could not query LLM due to setup error: {e_setup}"

    # ==========================================================================
    # Query Analysis Helpers (_call_llm_for_analysis, _classify_query, _decompose_query)
    # (These remain unchanged as they operate on the text query itself)
    # ==========================================================================
    def _call_llm_for_analysis(self, prompt_messages, task_description):
        """Helper to call LLM for internal analysis tasks like classification/decomposition."""
        try:
            # Use a potentially different/cheaper/faster model for analysis if desired
            # analysis_llm = ChatGroq(...) # Or reuse self._get_llm()
            analysis_llm = self._get_llm() # Reuse main LLM for simplicity here
            logger.info(f"Calling LLM ({analysis_llm.model_name}) for {task_description}...")
            response = analysis_llm.invoke(prompt_messages)
            content = response.content.strip() if hasattr(response, 'content') else str(response).strip()
            logger.info(f"LLM response for {task_description}: '{content[:100]}...'")
            return content
        except Exception as e:
            logger.error(f"LLM call failed during {task_description}: {e}", exc_info=True)
            return None # Indicate failure

    def _classify_query(self, query):
        """Classify the query type using the configured LLM."""
        # Consider making the prompt more robust or providing examples
        system_prompt = """You are an expert query analyzer. Classify the user query into one category:
'Simple Retrieval': Asks for specific facts, definitions, summaries, or information likely found directly in text or simple tables.
'Complex/Reasoning': Requires calculations (sum, average, count), comparisons across multiple data points/rows/sheets/slides, multi-step lookups, or analysis of potentially structured data (like complex tables or financial reports).

Respond ONLY with 'Simple Retrieval' or 'Complex/Reasoning'."""
        human_prompt = f"Classify the following user query: \"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        classification = self._call_llm_for_analysis(messages, "query classification")

        # Normalize response slightly
        if classification:
            classification = classification.strip().replace("'", "").replace('"', '')

        if classification in ["Simple Retrieval", "Complex/Reasoning"]:
            return classification
        else:
            logger.warning(f"Unexpected classification result: '{classification}'. Defaulting to 'Simple Retrieval'.")
            # Fallback heuristic: if query contains math keywords, lean towards complex
            math_keywords = ['sum', 'average', 'total', 'count', 'compare', 'difference', 'maximum', 'minimum', 'trend']
            if any(keyword in query.lower() for keyword in math_keywords):
                 logger.info("Fallback heuristic: Query contains math keywords, classifying as Complex/Reasoning.")
                 return "Complex/Reasoning"
            return "Simple Retrieval" # Default fallback

    def _decompose_query(self, query):
        """Decompose a complex query into simpler sub-queries using the configured LLM."""
        # Prompt engineering is key here for good decomposition
        system_prompt = """You are an expert query decomposer. Break down the complex user query below into a series of simple, factual sub-queries. Each sub-query should aim to retrieve a *single specific piece of information* or a small, well-defined set of related data points needed to ultimately answer the original complex query.
Focus on extracting the necessary raw data points or simple facts. Avoid asking for calculations, comparisons, or aggregations in the sub-queries themselves; these should be done based on the retrieved results.
If the query is already simple and cannot be decomposed, return just the original query.
Format the output as a JSON list of strings, where each string is a sub-query. Example: ["Find the total revenue for 2023", "Find the total expenses for 2023"]

Query to decompose:
\"\"\"
{query}
\"\"\"

Decomposed sub-queries (JSON list of strings):"""
        human_prompt = f"Decompose the following complex query into simple, factual sub-queries:\n\"{query}\"" # Kept simple human prompt, system prompt is more detailed now.

        # Using a prompt template might be better for complex system prompts
        # from langchain_core.prompts import ChatPromptTemplate
        # prompt_template = ChatPromptTemplate.from_messages([
        #     ("system", system_prompt),
        #     # No human message needed if query is in system prompt template
        # ])
        # messages = prompt_template.format_messages(query=query)

        # Sticking to simpler System/Human messages for now:
        messages = [SystemMessage(content=system_prompt.format(query=query))] # Put query directly in system prompt

        decomposition_str = self._call_llm_for_analysis(messages, "query decomposition")

        if decomposition_str:
            # Attempt to parse the response as JSON
            try:
                # Clean potential markdown code blocks ```json ... ```
                if decomposition_str.startswith("```json"):
                    decomposition_str = decomposition_str.strip()[7:-3].strip()
                elif decomposition_str.startswith("```"):
                     decomposition_str = decomposition_str.strip()[3:-3].strip()

                sub_queries = json.loads(decomposition_str)
                if isinstance(sub_queries, list) and all(isinstance(sq, str) for sq in sub_queries) and sub_queries:
                    # Basic validation passed
                    logger.info(f"Decomposed query into: {sub_queries}")
                    # Further filter out empty strings just in case
                    sub_queries = [sq.strip() for sq in sub_queries if sq.strip()]
                    if sub_queries:
                        return sub_queries
                    else:
                         logger.warning("Decomposition resulted in empty list after stripping. Using original query.")
                         return [query]
                else:
                    logger.warning(f"LLM decomposition response was not a valid JSON list of strings: '{decomposition_str}'. Using original query.")
                    return [query] # Fallback
            except json.JSONDecodeError as e:
                 logger.warning(f"Failed to parse LLM decomposition response as JSON: {e}. Response was: '{decomposition_str}'. Using original query.")
                 # Fallback: Try simple newline splitting if JSON fails
                 sub_queries_fallback = [re.sub(r"^\d+\.\s*", "", line).strip() for line in decomposition_str.split('\n') if line.strip()]
                 if len(sub_queries_fallback) > 1:
                      logger.info(f"Using fallback newline split for decomposition: {sub_queries_fallback}")
                      return sub_queries_fallback
                 return [query] # Fallback to original query if parsing fails
        else:
            logger.warning(f"Decomposition failed (LLM call returned None). Using original query.")
            return [query] # Fallback to original query


    # ==========================================================================
    # Modified run_query Function (remains structurally the same, relies on updated sub-components)
    # ==========================================================================
    def run_query(self, query, context_strategy="top_k"):
        """Complete query pipeline: classify, (decompose), retrieve, aggregate, query LLM."""
        start_run_time = time.time()
        logger.info(f"--- Starting query execution (v3-vision): '{query[:100]}...' ---")
        final_results = {
            "query": query,
            "classification": None,
            "sub_queries": None,
            "retrieval_results": {}, # Stores results per file AFTER de-duplication
            "aggregated_context_data": {},
            "answer": "",
            "status": "Started",
            "error": None,
            "timings": {} # To store timing info
        }
        is_streamlit = self._is_streamlit
        if is_streamlit: import streamlit as st

        try:
            # === 1. Classify Query ===
            step_start_time = time.time()
            spinner_msg = "Analyzing query..."
            if is_streamlit:
                with st.spinner(spinner_msg): classification = self._classify_query(query)
            else: logger.info(spinner_msg); classification = self._classify_query(query)
            final_results["classification"] = classification
            final_results["timings"]["1_classify"] = time.time() - step_start_time
            logger.info(f"Query classified as: {classification} (took {final_results['timings']['1_classify']:.2f}s)")

            # === 2. Decompose if Complex ===
            step_start_time = time.time()
            queries_to_retrieve = [query] # Default: use original query
            if classification == "Complex/Reasoning":
                spinner_msg = "Decomposing complex query..."
                if is_streamlit:
                     with st.spinner(spinner_msg): sub_queries = self._decompose_query(query)
                else: logger.info(spinner_msg); sub_queries = self._decompose_query(query)

                # Check if decomposition actually happened and produced multiple queries
                if sub_queries and len(sub_queries) > 1 and sub_queries != [query]:
                     queries_to_retrieve = sub_queries
                     final_results["sub_queries"] = sub_queries
                     logger.info(f"Using {len(sub_queries)} sub-queries for retrieval: {sub_queries}")
                elif sub_queries and len(sub_queries) == 1 and sub_queries != [query]:
                     # If decomposition returned a single, different query
                     queries_to_retrieve = sub_queries
                     final_results["sub_queries"] = sub_queries # Log it even if it's one
                     logger.info(f"Decomposition returned a single modified query: {sub_queries}")
                else:
                     logger.info("Decomposition didn't yield distinct sub-queries or failed, using original query for retrieval.")
                     # Keep queries_to_retrieve = [query]

            final_results["timings"]["2_decompose"] = time.time() - step_start_time
            logger.info(f"Decomposition step took {final_results['timings']['2_decompose']:.2f}s")


            # === 3. Retrieve relevant chunks for determined queries ===
            step_start_time = time.time()
            spinner_msg = f"Searching documents for {len(queries_to_retrieve)} query part(s)..."
            all_query_results_raw = {} # Temp storage before de-duplication

            if is_streamlit:
                with st.spinner(spinner_msg):
                    for i, q in enumerate(queries_to_retrieve):
                        logger.info(f"Retrieving for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                        results_for_q = self.query_files(q) # Calls Faiss search
                        # Merge results, simple aggregation for now
                        for file, res_list in results_for_q.items():
                            if file not in all_query_results_raw: all_query_results_raw[file] = []
                            all_query_results_raw[file].extend(res_list)
            else:
                logger.info(spinner_msg)
                for i, q in enumerate(queries_to_retrieve):
                     logger.info(f"Retrieving for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                     results_for_q = self.query_files(q)
                     for file, res_list in results_for_q.items():
                            if file not in all_query_results_raw: all_query_results_raw[file] = []
                            all_query_results_raw[file].extend(res_list)

            # --- De-duplicate and sort results PER FILE *after* merging from sub-queries ---
            final_retrieval_results_deduped = {}
            total_raw_chunks = 0
            total_deduped_chunks = 0
            for file, res_list in all_query_results_raw.items():
                 total_raw_chunks += len(res_list)
                 # Use content + precise source_info tuple as key for uniqueness
                 # Sort source_info items to ensure consistent key order
                 unique_content = {}
                 for res in res_list:
                     source_key = tuple(sorted(res.get('source_info',{}).items()))
                     key = (res['content'], source_key)
                     # Keep the one with the *lowest* score (best match) if duplicate found
                     if key not in unique_content or res['score'] < unique_content[key]['score']:
                         unique_content[key] = res
                 # Sort the unique results for this file by score
                 sorted_unique_results = sorted(unique_content.values(), key=lambda x: x['score'])
                 final_retrieval_results_deduped[file] = sorted_unique_results
                 total_deduped_chunks += len(sorted_unique_results)

            logger.info(f"Retrieval raw chunk count: {total_raw_chunks}, after de-duplication: {total_deduped_chunks}")
            final_results["retrieval_results"] = final_retrieval_results_deduped # Store final, de-duped results
            final_results["timings"]["3_retrieve_dedupe"] = time.time() - step_start_time
            logger.info(f"Retrieval & De-duplication step took {final_results['timings']['3_retrieve_dedupe']:.2f}s")


            if not final_retrieval_results_deduped:
                logger.warning("No relevant chunks found across any documents for the query/sub-queries after de-duplication.")
                final_results["answer"] = "Based on the scanned documents, I could not find any relevant information to answer this question."
                final_results["status"] = "Completed: No relevant information found."
                if is_streamlit: st.info("Could not find relevant information in the documents.")
                # Record total time even on early exit
                final_results["timings"]["total_runtime"] = time.time() - start_run_time
                return final_results

            # === 4. Aggregate context (uses combined & de-duplicated results) ===
            step_start_time = time.time()
            spinner_msg = "Gathering context..."
            if is_streamlit:
                with st.spinner(spinner_msg): aggregated_context_data = self.aggregate_context(final_retrieval_results_deduped, strategy=context_strategy)
            else: logger.info(spinner_msg); aggregated_context_data = self.aggregate_context(final_retrieval_results_deduped, strategy=context_strategy)

            final_results["aggregated_context_data"] = aggregated_context_data
            final_results["timings"]["4_aggregate"] = time.time() - step_start_time
            logger.info(f"Context aggregation step took {final_results['timings']['4_aggregate']:.2f}s")


            if not aggregated_context_data or not aggregated_context_data.get("combined_context"):
                 logger.warning(f"Context aggregation failed or yielded empty context for query: '{query[:100]}...'")
                 final_results["answer"] = "Could not generate an answer because the relevant information found could not be prepared for the AI."
                 final_results["status"] = "Completed: Context aggregation failed or empty."
                 final_results["error"] = "Context aggregation failed or yielded empty context."
                 if is_streamlit: st.warning("Found some relevant chunks, but couldn't prepare the final context for answering.")
                 # Record total time even on early exit
                 final_results["timings"]["total_runtime"] = time.time() - start_run_time
                 return final_results

            # === 5. Query LLM with aggregated context BUT use the ORIGINAL query ===
            step_start_time = time.time()
            spinner_msg = "Generating final answer..."
            # Use the ORIGINAL query for the final answer generation
            original_query_for_llm = query
            if is_streamlit:
                 with st.spinner(spinner_msg): answer = self.query_llm(original_query_for_llm, aggregated_context_data)
            else: logger.info(spinner_msg); answer = self.query_llm(original_query_for_llm, aggregated_context_data)

            final_results["answer"] = answer
            final_results["timings"]["5_llm_query"] = time.time() - step_start_time
            logger.info(f"LLM query step took {final_results['timings']['5_llm_query']:.2f}s")

            final_results["status"] = "Completed Successfully."


        except Exception as e:
            logger.error(f"Unexpected error during run_query (v3-vision): {e}", exc_info=True)
            final_results["status"] = "Failed"
            final_results["error"] = str(e)
            final_results["answer"] = f"An unexpected error occurred while processing the query: {e}" # Provide error in answer
            if is_streamlit: st.error(f"An unexpected error occurred: {e}")

        # Record total time
        final_results["timings"]["total_runtime"] = time.time() - start_run_time
        logger.info(f"--- Finished query execution (v3-vision) for: '{query[:100]}...' in {final_results['timings']['total_runtime']:.2f}s ---")
        return final_results
    # ==========================================================================
    # END OF MODIFIED run_query FUNCTION
    # ==========================================================================


# --- START OF STREAMLIT UI CODE ---
# (Streamlit UI code remains largely identical, but updated titles/captions and added timing display)
if __name__ == "__main__":
    try: import streamlit as st
    except ImportError: logger.error("Streamlit not installed."); sys.exit(1)

    st.set_page_config(layout="wide", page_title="Multi-Document Q&A with RAG v3+Vision")
    st.title("📄 Multi-Document Question Answering System (v3 + Vision)")
    st.caption("Query PDFs (incl. Images via Captioning), XLSX, CSV, PPTX using Embeddings, FAISS, Groq LLM (with Query Analysis)")

    @st.cache_resource
    def load_rag_system(config_path="config.ini"):
        # Check for Groq key early in Streamlit context
        if not os.getenv("GROQ_API_KEY"):
            st.error("🚨 Environment variable `GROQ_API_KEY` is not set! The application cannot contact the LLM.")
            st.info("Please set the `GROQ_API_KEY` in your environment or a `.env` file.")
            st.stop()
        # Check for Pillow if not already checked
        if Image is None and "checked_pillow" not in st.session_state:
             st.warning("⚠️ Python library `Pillow` not found. PDF image captioning is disabled. Install with `pip install Pillow` and restart.")
             st.session_state.checked_pillow = True # Prevent repeated warnings

        try:
            with st.spinner("Initializing RAG System and loading models..."):
                 rag_system = RAGSystem(config_path=config_path)
            # Check if LLM init worked (it tries during __init__)
            try: rag_system._get_llm()
            except ValueError as e: # Catch API key error from _get_llm
                 st.error(f"🚨 {e}")
                 st.stop()
            except Exception as e_llm: # Catch other LLM init errors
                 st.error(f"🚨 Failed to initialize the LLM client: {e_llm}")
                 st.info("Check your `config.ini` model name and network connection.")
                 st.stop()
            return rag_system
        except Exception as e:
            st.error(f"Fatal RAG System Initialization Error: {e}")
            logger.error("Fatal RAG System Initialization Error:", exc_info=True)
            st.stop()

    # Function to process documents, uses cache_data for result, depends on rag_system resource
    @st.cache_data
    def process_documents(_rag_system):
        status_messages = []
        status_placeholder = st.empty() # Placeholder for dynamic updates

        # Callback function for process_files to update Streamlit UI
        def streamlit_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
            log_prefix = "✅ " if is_done else "❌ " if is_error else "⚠️ " if is_warning else "⏳ "
            progress_val = 1.0 if is_done else 0.0
            stage_text = f" ({stage})" if stage else ""

            if total_steps and current_step is not None:
                # More granular progress based on typical stages
                stage_progress_map = {"Extracting": 0.1, "Chunking": 0.4, "Verifying": 0.5, "Embedding": 0.6, "Indexing": 0.9}
                base_progress = current_step / total_steps
                stage_weight = (1 / total_steps) # Distribute progress over the file
                stage_prog_val = stage_progress_map.get(stage, 0.0) * stage_weight
                progress_val = min(1.0, base_progress + stage_prog_val)

            full_message = f"{log_prefix}{message}{stage_text}"
            status_messages.append(full_message) # Keep log of messages

            # Update placeholder using try/except for robustness in Streamlit callbacks
            try:
                # Only show progress bar if progress > 0 and not done/error/warn
                if 0 < progress_val < 1.0 and not (is_error or is_warning or is_done):
                    status_placeholder.progress(progress_val)
                else:
                    status_placeholder.empty() # Clear progress bar otherwise

                # Display message based on status
                if is_error: status_placeholder.error(full_message)
                elif is_warning: status_placeholder.warning(full_message)
                elif is_done: status_placeholder.success(full_message)
                else: status_placeholder.info(full_message) # Use info for ongoing status
            except Exception as e_ui:
                 # Log UI update errors but don't crash the backend process
                 logger.warning(f"Streamlit UI update error during callback: {e_ui}")

        st.header("📚 Document Processing")
        st.caption(f"Data Source: `{os.path.abspath(_rag_system.config.data_dir)}` | Index Storage: `{os.path.abspath(_rag_system.config.index_dir)}`")
        processing_successful = False
        with st.spinner("Checking documents and processing if needed... This may take time for large files or images."):
            try:
                start_proc_time = time.time()
                processing_successful = _rag_system.process_files(progress_callback=streamlit_progress_callback)
                end_proc_time = time.time()
                proc_duration = end_proc_time - start_proc_time
                indexed_count = len(_rag_system.faiss_indexes)

                if indexed_count > 0:
                    status_placeholder.success(f"✅ Ready! Indexed {indexed_count} document(s). (Processing took {proc_duration:.2f}s)")
                elif processing_successful: # Finished but indexed 0
                    status_placeholder.warning(f"⚠️ Processing finished, but no documents were successfully indexed. Check logs. (Took {proc_duration:.2f}s)")
                elif not status_messages or not any("❌" in msg for msg in status_messages):
                    # If process_files returned False but no specific error was reported via callback
                    status_placeholder.error(f"❌ Processing finished with errors, but no documents indexed. Check logs. (Took {proc_duration:.2f}s)")
                # If errors were reported via callback, they should already be displayed
            except Exception as e:
                logger.error(f"Fatal error during document processing call: {e}", exc_info=True)
                status_placeholder.error(f"❌ Fatal error during processing: {e}. Check logs.")

        # Return success status and list of indexed files
        return len(_rag_system.faiss_indexes) > 0, sorted(list(_rag_system.faiss_indexes.keys()))

    # --- Main App Logic ---
    try:
        # Load the RAG system (cached)
        rag_sys = load_rag_system()

        # Button to force reprocessing
        if st.button("🔄 Re-process All Documents", help="Clears cache and re-runs extraction, chunking, embedding, and indexing for all files in the data directory."):
            # Clear specific caches related to document processing
            st.cache_data.clear()
            # We might not need to clear the resource cache if only data changed,
            # but clear it too for a full reset. Be cautious if resource loading is slow.
            st.cache_resource.clear()
            st.info("Cache cleared. Re-processing documents...")
            st.rerun() # Rerun the script to trigger reprocessing

        # Process documents and get status/file list
        is_ready, indexed_files = process_documents(rag_sys)

        # Sidebar Information
        if is_ready and indexed_files:
            st.sidebar.success(f"System Ready ({len(indexed_files)} Doc(s) Indexed)")
            with st.sidebar.expander("Show Indexed Files"):
                 for fname in indexed_files: st.caption(f"- {fname}")
        elif not indexed_files:
             st.sidebar.warning("System Ready (No Docs Indexed)")
        else:
             st.sidebar.error("System Not Ready")

        st.sidebar.markdown("---")
        st.sidebar.subheader("Configuration")
        st.sidebar.info(f"LLM: `{rag_sys.config.llm_model}`")
        st.sidebar.caption(f"(Vision Capable: {'Yes' if rag_sys.image_captioning_enabled else 'No (Pillow missing)'})")
        st.sidebar.info(f"Encoder: `{rag_sys.config.encoder_model.split('/')[-1]}`")
        st.sidebar.info(f"Retrieval K: `{rag_sys.config.k_retrieval}`")
        st.sidebar.info(f"Chunk Size: `{rag_sys.config.chunk_size}`")
        st.sidebar.info(f"Overlap: `{rag_sys.config.overlap}`")
        st.sidebar.info(f"Max Context Tokens (Approx): `{rag_sys.config.max_context_tokens}`")


        # Main Query Interface (only if system is ready with indexed files)
        if is_ready and indexed_files:
            st.header("💬 Ask a Question")
            user_query = st.text_input("Enter your query about the indexed documents:", key="query_input", placeholder="e.g., What solutions were proposed for Client X? Summarize the key benefits.")

            if user_query: # Only show button if query is entered
                if st.button("Get Answer", key="submit_query", type="primary"):
                    query_start_time = time.time()
                    # Call the main query execution function
                    results_data = rag_sys.run_query(user_query, context_strategy="top_k")
                    query_end_time = time.time()
                    query_duration = query_end_time - query_start_time

                    st.markdown("---") # Separator before results

                    # Display Answer
                    st.subheader("💡 Answer")
                    answer = results_data.get("answer")
                    status = results_data.get("status")
                    error = results_data.get("error")

                    if "Completed" in status and answer:
                        st.markdown(answer) # Render markdown for better formatting
                    elif error:
                        st.error(f"Error generating answer: {error}")
                        if answer: st.info(f"LLM Raw Output (if any): {answer}") # Show raw if error occurred but answer exists
                    elif not answer:
                         st.warning("Could not generate an answer. The LLM may have returned an empty response or context might have been insufficient.")
                    else: # Catch-all for unexpected status
                        st.warning(f"Query status: {status}. Answer: {answer}")

                    # Display Query Analysis Details
                    with st.expander("📊 Query Analysis & Performance", expanded=False):
                        col1, col2 = st.columns(2)
                        with col1:
                            st.metric(label="Query Classification", value=results_data.get("classification", "N/A"))
                            sub_queries = results_data.get("sub_queries")
                            if sub_queries:
                                st.caption("Sub-queries used for retrieval:")
                                for i, sq in enumerate(sub_queries):
                                    st.code(f"{i+1}. {sq}", language=None)
                            else:
                                 st.caption("Original query used directly for retrieval.")
                        with col2:
                             st.metric(label="Total Query Time", value=f"{query_duration:.2f} s")
                             timings = results_data.get("timings", {})
                             st.caption(f"Analysis: {timings.get('1_classify', 0):.2f}s | Decompose: {timings.get('2_decompose', 0):.2f}s")
                             st.caption(f"Retrieve/Dedupe: {timings.get('3_retrieve_dedupe', 0):.2f}s | Aggregate: {timings.get('4_aggregate', 0):.2f}s")
                             st.caption(f"LLM Query: {timings.get('5_llm_query', 0):.2f}s")


                    # Display Supporting Evidence (Context + Retrieval)
                    st.subheader("🔍 Supporting Evidence")
                    agg_context_data = results_data.get("aggregated_context_data", {})
                    combined_context = agg_context_data.get("combined_context", "")
                    context_source_files = agg_context_data.get("source_files", [])

                    if combined_context:
                         with st.expander(f"Context Used by LLM (from {len(context_source_files)} file(s))", expanded=False):
                             st.text_area("Combined Context Sent to LLM", combined_context, height=300, key="context_display", help="This is the exact context provided to the LLM to generate the answer.")
                    else:
                         st.info("No context was aggregated or sent to the LLM.")

                    retrieval_results = results_data.get("retrieval_results", {}) # This is the dict of {filename: [results]}
                    if retrieval_results:
                        # Flatten results for display, respecting original file source
                        display_chunks = []
                        for file_name, res_list in retrieval_results.items():
                             for res in res_list:
                                  display_chunks.append({**res, 'file_name': file_name})
                        # Sort by score globally for display order (lower is better)
                        display_chunks.sort(key=lambda x: x['score'])
                        num_chunks_to_display = rag_sys.config.k_retrieval # Show top K overall

                        with st.expander(f"Retrieved Chunks (Top {num_chunks_to_display} Overall)", expanded=False):
                            st.caption("These are the most relevant text/table/caption chunks found across all documents based on the query/sub-queries, sorted by relevance score (lower is better).")
                            if not display_chunks:
                                 st.info("No chunks were retrieved.")

                            for i, res in enumerate(display_chunks[:num_chunks_to_display]):
                                source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown'); content_type = res.get('type', 'unknown'); score = res['score']

                                # Build display string for source
                                source_display_parts = [f"File: **{file_name}**"]
                                if file_type == 'pptx':
                                     slide_info = source_info.get('slide', 'N/A'); merged = source_info.get('merged_from')
                                     source_display_parts.append(f"Slide: {slide_info}{f' (merged from {merged})' if merged else ''}")
                                elif file_type == 'pdf':
                                     source_display_parts.append(f"Page: {source_info.get('page', 'N/A')}")
                                     if content_type == 'image_caption': source_display_parts.append(f"Image: {source_info.get('image_index', 'N/A')}")
                                     if content_type == 'table': source_display_parts.append(f"Table: {source_info.get('table_index', 'N/A')}")
                                elif file_type == 'xlsx': source_display_parts.append(f"Sheet: {source_info.get('sheet', 'N/A')}")
                                elif file_type == 'csv': source_display_parts.append(f"CSV (Rows: {source_info.get('rows', 'all')})")
                                source_display_parts.append(f"Type: `{content_type}`")

                                st.markdown(f"**Chunk {i+1} (Score: {score:.4f})**")
                                st.caption(" | ".join(source_display_parts))
                                # Use st.text or st.code for raw content display
                                st.text_area(f"Chunk {i+1} Content", res['content'], height=150, key=f"chunk_{i}", disabled=True)
                                # st.divider() # Removed divider for cleaner look with text_area border

                    else:
                         st.info("No specific document chunks were retrieved for this query.")

        elif not indexed_files and not is_ready: # If processing failed or yielded nothing
            st.warning(f"System not ready. Document processing may have failed or found no indexable documents.")
            st.info("Check the 'Document Processing' status above and the log file (`rag_system_vision.log`). Add supported files to the data directory and click 'Re-process All Documents'.")
        elif not indexed_files and is_ready: # Processing finished okay, but nothing indexed
             st.info("Document processing completed, but no files were successfully indexed.")
             st.info("Ensure supported files (.pdf, .xlsx, .csv, .pptx) exist in the data directory and are not corrupted or empty. Then click 'Re-process All Documents'.")

    except Exception as e:
        st.error(f"An error occurred in the Streamlit application: {e}")
        logger.exception("Streamlit application error:")
        st.info("Check the console output or the `rag_system_vision.log` file for more details.")

# --- END OF STREAMLIT UI CODE ---

# --- END OF FILE ragtest_c3_query_adapt_vision.py ---
