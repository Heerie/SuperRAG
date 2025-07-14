# --- START OF FILE ragtest_c2-6.py ---

# --- START OF FILE ragtest_c2-4_sql.py --- # Renamed

# --- START OF FILE ragtest_c4_dataframe.py --- # (Original comments kept for context)
# --- START OF FILE ragtest_c3_query_adapt.py --- #
# --- START OF FILE ragtest_c.py --- #

import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd
import faiss
# import streamlit as st # Import only if needed / within UI code
import time
from sentence_transformers import SentenceTransformer
# from transformers import T5ForConditionalGeneration, T5Tokenizer # Keep if needed elsewhere
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import io
from contextlib import redirect_stdout
import sqlite3 # Added for SQL database
import requests # Added for DeepSeek API

# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted from RAG2.py + SQL Integration) ---

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system_sql.log"), # Log to a different file maybe
        # logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


class RAGConfig:
    """Configuration class for RAG system."""
    def __init__(self, config_path="config.ini"):
        self.config = configparser.ConfigParser()
        # Added DEEPSEEK_API_KEY reference and SQLITE_DB path
        self.defaults = {
            "PATHS": {
                "data_dir": "./Data/",
                "index_dir": "./Faiss_index/",
                "log_file": "rag_system_sql.log",
                "sqlite_db_path": ":memory:" # Default to in-memory SQLite DB
            },
            "MODELS": {
                "encoder_model": "sentence-transformers/all-MiniLM-L6-v2",
                "llm_model": "llama3-8b-8192", # Use a known good Groq model
                "device": "auto",
                "deepseek_model": "deepseek-reasoner" # Specify DeepSeek model for metadata
            },
            "PARAMETERS": {
                "chunk_size": "200",
                "overlap": "50",
                "k_retrieval": "5",
                "temperature": "0.1",
                "max_context_tokens": "4000",
                "max_chars_per_element": "1000",
                "pptx_merge_threshold_words": "50",
                "dataframe_query_confidence_threshold": "0.75" # Slightly lower threshold maybe
            },
            "SUPPORTED_EXTENSIONS": {
                "extensions": ".pdf, .xlsx, .csv, .pptx"
            },
            "API_KEYS": { # Section for API keys (optional, can still rely purely on .env)
                 "deepseek_api_key_config": "" # Can be set here or via .env
            }
        }
        # --- Load/Save logic (unchanged) ---
        if os.path.exists(config_path):
            try: self.config.read(config_path); logger.info(f"Loaded configuration from {config_path}"); self._ensure_defaults()
            except Exception as e: logger.error(f"Error reading config file {config_path}: {e}"); self._set_defaults()
        else: logger.warning(f"Config file {config_path} not found, using defaults"); self._set_defaults()
    def _set_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section)
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value)
    def _ensure_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section); logger.info(f"Added missing section [{section}] from defaults.")
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value); logger.info(f"Added missing option '{option} = {value}' to section [{section}] from defaults.")

    # --- Property definitions (updated/added) ---
    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def sqlite_db_path(self): return self.config.get("PATHS", "sqlite_db_path", fallback=":memory:")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    @property
    def llm_model(self): return self.config.get("MODELS", "llm_model")
    @property
    def deepseek_model(self): return self.config.get("MODELS", "deepseek_model", fallback="deepseek-reasoner")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto")
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=200)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=50)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=5)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.1)
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=4000)
    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1000)
    @property
    def pptx_merge_threshold_words(self): return self.config.getint("PARAMETERS", "pptx_merge_threshold_words", fallback=50)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf")
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])
    @property
    def dataframe_query_confidence_threshold(self):
        return self.config.getfloat("PARAMETERS", "dataframe_query_confidence_threshold", fallback=0.75)
    @property
    def deepseek_api_key(self):
        # Prioritize environment variable, then config file
        key = os.getenv("DEEPSEEK_API_KEY")
        if not key:
            key = self.config.get("API_KEYS", "deepseek_api_key_config", fallback=None)
            if key: logger.info("Using DeepSeek API key from config file.")
        # else: logger.info("Using DeepSeek API key from environment variable.") # Can be noisy
        return key


class RAGSystem:
    """Retrieval-Augmented Generation system supporting text RAG and SQL querying."""
    def __init__(self, config_path="config.ini"):
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system (SQL Version). Device: {self.config.device}, Supported files: {self.config.supported_extensions}")
        self._is_streamlit = "streamlit" in sys.modules
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)

        # --- Embedding Model Loading (unchanged) ---
        try:
            if self._is_streamlit:
                import streamlit as st
                with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."): self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            else: logger.info(f"Loading embedding model ({self.config.encoder_model})..."); self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            logger.info("Local embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}"); st.stop()
            else: raise

        # --- State Variables ---
        self.file_chunks = {} # Stores text chunks for FAISS
        self.faiss_indexes = {} # Stores FAISS indexes
        self.processed_files = set() # Tracks files processed for text RAG *or* SQL

        # --- SQL Database and Metadata Stores ---
        self.db_conn = None
        self.db_cursor = None
        self.table_metadata = {} # { "table_name": {"columns": [...], "source_file": "...", "source_sheet": "...", "row_count": N} }
        self.file_to_table_map = {} # { "filename.xlsx": ["table_sheet1", "table_sheet2"], ... }

        self._llm = None # Groq LLM client cache

        # --- API Key Checks ---
        self.groq_api_key = os.getenv("GROQ_API_KEY")
        self.deepseek_api_key = self.config.deepseek_api_key # Get key via config property
        if not self.groq_api_key: logger.warning("GROQ_API_KEY not found in environment variables."); # UI will show error later if needed
        if not self.deepseek_api_key: logger.warning("DEEPSEEK_API_KEY not found in environment variables or config file. Spreadsheet metadata generation will fail.")

        # --- Connect to SQLite DB ---
        self._connect_db()

    def _connect_db(self):
        """Establish connection to the SQLite database."""
        try:
            db_path = self.config.sqlite_db_path
            self.db_conn = sqlite3.connect(db_path, check_same_thread=False) # Allow access from different threads (e.g., Streamlit)
            self.db_cursor = self.db_conn.cursor()
            logger.info(f"Connected to SQLite database at '{db_path}'")
        except sqlite3.Error as e:
            logger.error(f"Error connecting to SQLite database at '{self.config.sqlite_db_path}': {e}", exc_info=True)
            self.db_conn = None
            self.db_cursor = None
            if self._is_streamlit:
                import streamlit as st
                st.error(f"Fatal Error: Failed to connect to SQLite DB. Check logs. Error: {e}")
                st.stop()
            else:
                raise ConnectionError(f"Failed to connect to SQLite DB: {e}") from e

    def _close_db(self):
        """Close the SQLite database connection."""
        if self.db_conn:
            try:
                self.db_conn.commit() # Commit any pending changes
                self.db_conn.close()
                logger.info("Closed SQLite database connection.")
                self.db_conn = None
                self.db_cursor = None
            except sqlite3.Error as e:
                logger.error(f"Error closing SQLite connection: {e}", exc_info=True)

    # Make sure to close the DB when the object is destroyed
    def __del__(self):
        self._close_db()

    # --- LLM Initialization (Groq) --- (Unchanged)
    def _get_llm(self):
        if self._llm is None:
            if not self.groq_api_key:
                logger.error("GROQ_API_KEY not found.")
                if self._is_streamlit: import streamlit as st; st.error("Error: GROQ_API_KEY is not set.")
                raise ValueError("GROQ_API_KEY not configured.")
            try:
                self._llm = ChatGroq(
                    temperature=self.config.temperature,
                    model_name=self.config.llm_model,
                    groq_api_key=self.groq_api_key,
                )
                logger.info(f"Groq LLM client initialized ({self.config.llm_model})")
            except Exception as e:
                 logger.error(f"Failed to initialize Groq client: {e}", exc_info=True)
                 if self._is_streamlit: import streamlit as st; st.error(f"Failed to initialize LLM. Error: {e}")
                 raise
        return self._llm

    # --- Text Cleaning and Basic File Handling --- (Unchanged)
    def clean_text(self, text):
        if not isinstance(text, str): text = str(text)
        text = re.sub(r"\(cid:.*?\)", "", text)
        text = re.sub(r"\s+", " ", text).strip()
        text = text.replace('\n', ' ').replace('\r', '')
        return text

    def _get_safe_filename(self, file_name):
        base_name = os.path.splitext(file_name)[0]
        return re.sub(r'[^\w\.-]', '_', base_name)

    def _get_safe_table_name(self, file_name, sheet_name=None):
        """Creates a safe, unique name for an SQL table."""
        base = self._get_safe_filename(file_name)
        if sheet_name:
             safe_sheet = re.sub(r'[^\w]', '_', sheet_name)
             name = f"{base}_{safe_sheet}"
        else:
             name = base
        # Ensure it starts with a letter or underscore if needed by some DBs, and truncate length
        if not re.match(r"^[a-zA-Z_]", name): name = "_" + name
        name = name[:60] # Limit length
        return name.lower() # Often best to use lowercase table names

    def _table_exists(self, table_name):
        """Check if a table exists in the database."""
        if not self.db_cursor: return False
        try:
            self.db_cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name=?;", (table_name,))
            return self.db_cursor.fetchone() is not None
        except sqlite3.Error as e:
            logger.error(f"Error checking if table '{table_name}' exists: {e}"); return False

    def get_index_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.index")
    def get_embedding_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.npy")
    def get_chunks_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.json")


    # --- Content Extractors ---
    # PDF and PPTX remain unchanged, focused on text extraction for FAISS
    def _extract_pdf(self, file_path):
        # (Unchanged from previous version - extracts text/table text for Faiss)
        all_content = []
        base_filename = os.path.basename(file_path)
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages); logger.info(f"Extracting text/tables from PDF: {base_filename} ({total_pages} pages)")
                page_bar = None
                if self._is_streamlit and total_pages > 10: import streamlit as st; page_bar = st.progress(0, text=f"Extracting pages from {base_filename}...")
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text(layout="normal") or ""; cleaned_text = self.clean_text(text)
                    if cleaned_text: all_content.append({"type": "text", "content": cleaned_text, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                    try:
                        for table in page.extract_tables():
                             if table:
                                 table_df = pd.DataFrame(table).fillna(''); table_string = table_df.to_string(index=False, header=True); cleaned_table_string = self.clean_text(table_string)
                                 if cleaned_table_string: all_content.append({"type": "table", "content": cleaned_table_string, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                    except Exception as e_table: logger.warning(f"Could not extract table on page {page_num + 1} in {base_filename}: {e_table}")
                    if page_bar: progress_percent = min(1.0, (page_num + 1) / total_pages); page_bar.progress(progress_percent, text=f"Extracting page {page_num + 1}/{total_pages} from {base_filename}...")
                if page_bar: page_bar.empty()
            logger.info(f"Extracted {len(all_content)} text/table content blocks from PDF: {base_filename}"); return all_content
        except Exception as e: logger.error(f"Error extracting content from PDF {base_filename}: {e}", exc_info=True); return []

    def _extract_pptx(self, file_path):
        # (Unchanged from previous version - extracts text for Faiss)
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element; merge_threshold_words = self.config.pptx_merge_threshold_words
        try:
            prs = Presentation(file_path); logger.info(f"Extracting text from PPTX: {base_filename} ({len(prs.slides)} slides)")
            slide_bar = None; total_slides = len(prs.slides)
            if self._is_streamlit and total_slides > 1: import streamlit as st; slide_bar = st.progress(0, text=f"Extracting slides from {base_filename}...")
            pending_title_slide_data = None
            for i, slide in enumerate(prs.slides):
                 current_slide_number = i + 1; current_slide_title_text = ""; current_slide_other_texts = []; current_slide_has_title_placeholder = False; title_shape = None
                 try:
                    if slide.shapes.title: title_shape = slide.shapes.title; current_slide_has_title_placeholder = True
                 except AttributeError: pass
                 if title_shape and title_shape.has_text_frame: cleaned_title = self.clean_text(title_shape.text_frame.text); current_slide_title_text = cleaned_title if cleaned_title else ""
                 for shape in slide.shapes:
                     if shape == title_shape: continue
                     is_placeholder = shape.is_placeholder; is_body_placeholder = False
                     if is_placeholder:
                         try: ph_type = shape.placeholder_format.type; is_body_placeholder = ph_type in [MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT, MSO_SHAPE_TYPE.SUBTITLE]
                         except AttributeError: pass
                     if shape.has_text_frame:
                         text = shape.text_frame.text; cleaned = self.clean_text(text)
                         if cleaned:
                              if len(cleaned) > max_chars: cleaned = cleaned[:max_chars] + "...(truncated shape)"
                              prefix = "[Body]: " if is_body_placeholder else ""; current_slide_other_texts.append(prefix + cleaned)
                 if slide.has_notes_slide:
                    try:
                         notes_text = slide.notes_slide.notes_text_frame.text; cleaned_notes = self.clean_text(notes_text)
                         if cleaned_notes:
                             if len(cleaned_notes) > max_chars * 2: cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated notes)"
                             current_slide_other_texts.append(f"[Notes]: {cleaned_notes}")
                    except Exception as e_notes: logger.warning(f"Could not extract notes from slide {current_slide_number}: {e_notes}")
                 current_slide_full_content = current_slide_title_text
                 if current_slide_other_texts: current_slide_full_content += ("\n" if current_slide_full_content else "") + "\n".join(current_slide_other_texts)
                 current_slide_full_content = current_slide_full_content.strip(); other_text_word_count = sum(len(s.split()) for s in current_slide_other_texts)
                 merged_content_block = None; should_merge = False
                 if pending_title_slide_data:
                     if pending_title_slide_data['has_title'] and pending_title_slide_data['other_words'] <= merge_threshold_words and current_slide_full_content:
                         should_merge = True; logger.info(f"Merging slide {pending_title_slide_data['number']} with slide {current_slide_number}.")
                 if should_merge:
                    merged_text = f"[Content from Slide {pending_title_slide_data['number']}]:\n{pending_title_slide_data['content']}\n\n---\n\n[Content from Slide {current_slide_number}]:\n{current_slide_full_content}"
                    merged_content_block = {"type": "slide_text_merged", "content": merged_text.strip(), "source_info": {"slide": current_slide_number, "merged_from": pending_title_slide_data['number']}, "file_type": "pptx"}
                    all_content.append(merged_content_block); pending_title_slide_data = None
                 else:
                    if pending_title_slide_data:
                         if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
                         pending_title_slide_data = None
                    if current_slide_has_title_placeholder and other_text_word_count <= merge_threshold_words and current_slide_full_content:
                        logger.debug(f"Slide {current_slide_number} is potential title slide. Holding."); pending_title_slide_data = {"content": current_slide_full_content, "number": current_slide_number, "has_title": current_slide_has_title_placeholder, "other_words": other_text_word_count, "title": current_slide_title_text}
                    else:
                        if current_slide_full_content: all_content.append({"type": "slide_text", "content": current_slide_full_content, "source_info": {"slide": current_slide_number}, "file_type": "pptx"})
                 if slide_bar: slide_bar.progress(min(1.0, (i + 1) / total_slides), text=f"Extracting slide {current_slide_number}/{total_slides}...")
            if pending_title_slide_data:
                 logger.debug(f"Processing pending title slide {pending_title_slide_data['number']} at end.")
                 if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
            if slide_bar: slide_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PPTX {base_filename} (Merge strategy applied)."); return all_content
        except Exception as e: logger.error(f"Error extracting content from PPTX {base_filename}: {e}", exc_info=True); return []


    # --- Spreadsheet Handling (Now involves loading DF, getting metadata, loading to SQL) ---
    def _extract_and_load_xlsx(self, file_path, file_name, progress_callback=None):
        """Loads XLSX sheets into DataFrames, gets metadata, and loads into SQL."""
        base_filename = os.path.basename(file_path)
        sheet_load_success_count = 0
        sheet_metadata_success_count = 0
        sheet_sql_load_success_count = 0
        processed_sheet_tables = []

        try:
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            logger.info(f"Processing XLSX: {base_filename} (Sheets: {len(sheet_names)})")
            if not sheet_names: logger.warning(f"No sheets found in {base_filename}"); return False, []

            sheet_bar = None; total_sheets = len(sheet_names)
            if self._is_streamlit and total_sheets > 1: import streamlit as st; sheet_bar = st.progress(0, text=f"Processing sheets in {base_filename}...")

            for i, sheet_name in enumerate(sheet_names):
                stage_msg = f"Processing sheet '{sheet_name}' ({i+1}/{total_sheets}) in {base_filename}"
                logger.info(stage_msg)
                if progress_callback: progress_callback(f"{stage_msg} - Loading...", stage="Loading Sheet")
                if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.1) / total_sheets), text=f"Loading '{sheet_name}'...")

                try:
                    df = excel_file.parse(sheet_name)
                    df = df.fillna('') # Basic NA handling
                    # Clean column names for SQL compatibility
                    df.columns = [re.sub(r'[^a-zA-Z0-9_]', '_', str(col)).strip('_') for col in df.columns]
                    # Ensure column names are unique
                    cols = pd.Series(df.columns)
                    for dup in cols[cols.duplicated()].unique():
                        cols[cols[cols == dup].index.values.tolist()] = [dup + '_' + str(j) if j != 0 else dup for j in range(sum(cols == dup))]
                    df.columns = cols

                    if df.empty: logger.warning(f"Sheet '{sheet_name}' in {base_filename} is empty. Skipping."); continue
                    sheet_load_success_count += 1
                    logger.info(f"Loaded DataFrame from sheet '{sheet_name}' (Rows: {len(df)}, Cols: {len(df.columns)})")

                    # Generate Table Name
                    table_name = self._get_safe_table_name(file_name, sheet_name)
                    if not table_name: logger.error(f"Could not generate table name for sheet '{sheet_name}'. Skipping SQL load."); continue

                    # Get Column Metadata via DeepSeek
                    if progress_callback: progress_callback(f"{stage_msg} - Analyzing Columns (DeepSeek)...", stage="Analyzing Columns")
                    if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.4) / total_sheets), text=f"Analyzing '{sheet_name}'...")

                    column_metadata = self._get_column_metadata_deepseek(df, file_name, sheet_name)
                    if not column_metadata:
                        logger.warning(f"Failed to get column metadata for sheet '{sheet_name}'. Proceeding without descriptions.")
                        # Create basic metadata if DeepSeek fails
                        column_metadata = [{"name": col, "type": self._infer_sql_type(df[col]), "description": "N/A"} for col in df.columns]
                    else:
                        sheet_metadata_success_count +=1

                    # Load DataFrame to SQL
                    if progress_callback: progress_callback(f"{stage_msg} - Loading to SQL Database...", stage="Loading SQL")
                    if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.8) / total_sheets), text=f"Storing '{sheet_name}'...")

                    sql_loaded = self._load_df_to_sql(df, table_name, column_metadata, file_name, sheet_name)
                    if sql_loaded:
                        sheet_sql_load_success_count += 1
                        processed_sheet_tables.append(table_name)
                        logger.info(f"Successfully loaded sheet '{sheet_name}' into SQL table '{table_name}'")
                    else:
                         logger.error(f"Failed to load sheet '{sheet_name}' into SQL.")

                except Exception as e_sheet:
                    logger.error(f"Error processing sheet '{sheet_name}' in {base_filename}: {e_sheet}", exc_info=True)
                    if progress_callback: progress_callback(f"⚠️ Error processing sheet '{sheet_name}': {e_sheet}", is_warning=True)

                if sheet_bar: sheet_bar.progress(min(1.0, (i + 1) / total_sheets), text=f"Finished sheet {i+1}/{total_sheets}...")

            if sheet_bar: sheet_bar.empty()
            logger.info(f"Finished processing XLSX {base_filename}. Sheets Loaded: {sheet_load_success_count}/{total_sheets}, Metadata OK: {sheet_metadata_success_count}, SQL OK: {sheet_sql_load_success_count}")
            return sheet_sql_load_success_count > 0, processed_sheet_tables

        except Exception as e:
            logger.error(f"Error opening or processing XLSX file {base_filename}: {e}", exc_info=True)
            if progress_callback: progress_callback(f"❌ Error processing XLSX {base_filename}: {e}", is_error=True)
            return False, []

    def _extract_and_load_csv(self, file_path, file_name, progress_callback=None):
        """Loads CSV into a DataFrame, gets metadata, and loads into SQL."""
        base_filename = os.path.basename(file_path)
        df = None
        load_success = False
        metadata_success = False
        sql_load_success = False
        table_name = None

        try:
            stage_msg = f"Processing CSV: {base_filename}"
            logger.info(stage_msg)
            if progress_callback: progress_callback(f"{stage_msg} - Loading...", stage="Loading CSV")

            encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
            for enc in encodings_to_try:
                try:
                    df = pd.read_csv(file_path, encoding=enc, low_memory=False)
                    df = df.fillna('') # Basic NA handling
                    # Clean column names
                    df.columns = [re.sub(r'[^a-zA-Z0-9_]', '_', str(col)).strip('_') for col in df.columns]
                    # Ensure column names are unique
                    cols = pd.Series(df.columns)
                    for dup in cols[cols.duplicated()].unique():
                         cols[cols[cols == dup].index.values.tolist()] = [dup + '_' + str(j) if j != 0 else dup for j in range(sum(cols == dup))]
                    df.columns = cols

                    logger.info(f"Read CSV {base_filename} using encoding: {enc}. (Rows: {len(df)}, Cols: {len(df.columns)})")
                    load_success = True
                    break
                except UnicodeDecodeError: continue
                except Exception as e_read: logger.warning(f"Pandas read_csv error for {base_filename} with encoding {enc}: {e_read}")

            if df is None or df.empty:
                logger.error(f"Could not read CSV {base_filename} or it is empty."); return False, []
            if not load_success: return False, [] # Should be caught above, but safety check

            # Generate Table Name
            table_name = self._get_safe_table_name(file_name) # No sheet name for CSV
            if not table_name: logger.error("Could not generate table name for CSV. Skipping SQL load."); return False, []

            # Get Column Metadata via DeepSeek
            if progress_callback: progress_callback(f"{stage_msg} - Analyzing Columns (DeepSeek)...", stage="Analyzing Columns")
            column_metadata = self._get_column_metadata_deepseek(df, file_name)
            if not column_metadata:
                logger.warning("Failed to get column metadata for CSV. Proceeding without descriptions.")
                column_metadata = [{"name": col, "type": self._infer_sql_type(df[col]), "description": "N/A"} for col in df.columns]
            else:
                metadata_success = True

            # Load DataFrame to SQL
            if progress_callback: progress_callback(f"{stage_msg} - Loading to SQL Database...", stage="Loading SQL")
            sql_loaded = self._load_df_to_sql(df, table_name, column_metadata, file_name)
            if sql_loaded:
                sql_load_success = True
                logger.info(f"Successfully loaded CSV {base_filename} into SQL table '{table_name}'")
            else:
                logger.error(f"Failed to load CSV {base_filename} into SQL.")

            return sql_load_success, ([table_name] if sql_load_success else [])

        except Exception as e:
            logger.error(f"Error processing CSV file {base_filename}: {e}", exc_info=True)
            if progress_callback: progress_callback(f"❌ Error processing CSV {base_filename}: {e}", is_error=True)
            return False, []


    def _infer_sql_type(self, series):
        """Basic inference of SQL type from Pandas Series dtype."""
        if pd.api.types.is_integer_dtype(series.dtype): return "INTEGER"
        if pd.api.types.is_float_dtype(series.dtype): return "REAL" # Use REAL for floats in SQLite
        if pd.api.types.is_bool_dtype(series.dtype): return "INTEGER" # Store bools as 0/1
        if pd.api.types.is_datetime64_any_dtype(series.dtype): return "TEXT" # Store datetimes as ISO text
        # pd.api.types.is_string_dtype(series.dtype) # This often includes 'object'
        return "TEXT" # Default to TEXT

    def _get_column_metadata_deepseek(self, df, file_name, sheet_name=None):
        """Uses DeepSeek API to get column descriptions and inferred SQL types."""
        if not self.deepseek_api_key:
            logger.warning("DeepSeek API key not available. Cannot generate column metadata.")
            return None

        api_url = "https://api.deepseek.com/v1/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.deepseek_api_key}"
        }

        # Prepare sample data (first 5 rows)
        sample_data_str = df.head(5).to_string(index=False)
        column_names = list(df.columns)

        prompt = f"""Analyze the following table columns from a spreadsheet ('{file_name}'{f", sheet '{sheet_name}'" if sheet_name else ""}).
The column names are: {column_names}
Here are the first few rows of data:
{sample_data_str}

For each column name, provide:
1. A concise one-sentence description of what the data in the column represents.
2. The most appropriate SQL data type (choose from: TEXT, INTEGER, REAL, DATE, DATETIME - use TEXT for dates/datetimes for SQLite compatibility).

Respond ONLY with a valid JSON object where keys are the exact column names and values are objects containing 'description' and 'sql_type'. Example:
{{
  "ColumnA": {{ "description": "Unique identifier for each record.", "sql_type": "INTEGER" }},
  "ColumnB": {{ "description": "Name of the customer.", "sql_type": "TEXT" }},
  "TransactionDate": {{ "description": "Date when the transaction occurred.", "sql_type": "TEXT" }}
}}
"""
        payload = {
            "model": self.config.deepseek_model,
            "messages": [
                {"role": "system", "content": "You are an expert data analyst specializing in understanding spreadsheet columns and assigning SQL types."},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.1, # Low temp for consistent JSON output
            "max_tokens": 1500, # Adjust as needed based on column count
            "stream": False
        }

        try:
            response = requests.post(api_url, headers=headers, json=payload, timeout=60) # 60 second timeout
            response.raise_for_status() # Raise HTTPError for bad responses (4xx or 5xx)

            result = response.json()
            raw_content = result['choices'][0]['message']['content']
            logger.debug(f"Raw DeepSeek response for metadata: {raw_content}")

            # Attempt to parse the JSON (sometimes LLMs add extra text)
            try:
                # Find the first '{' and last '}' to extract the JSON part
                json_start = raw_content.find('{')
                json_end = raw_content.rfind('}')
                if json_start != -1 and json_end != -1 and json_end > json_start:
                    json_str = raw_content[json_start:json_end+1]
                    metadata_dict = json.loads(json_str)
                else:
                    raise json.JSONDecodeError("Could not find JSON object delimiters", raw_content, 0)

                # Validate and format into list
                formatted_metadata = []
                missing_cols = set(column_names)
                for col_name, meta in metadata_dict.items():
                    if col_name in column_names and isinstance(meta, dict) and 'description' in meta and 'sql_type' in meta:
                         # Basic SQL type validation
                         sql_type = meta['sql_type'].upper()
                         if sql_type not in ["TEXT", "INTEGER", "REAL"]:
                              logger.warning(f"DeepSeek proposed invalid SQL type '{sql_type}' for column '{col_name}'. Defaulting to TEXT.")
                              sql_type = "TEXT" # Default to TEXT for safety
                         formatted_metadata.append({
                            "name": col_name,
                            "description": str(meta['description']).strip(),
                            "type": sql_type
                         })
                         missing_cols.discard(col_name)
                    else:
                        logger.warning(f"Invalid metadata format received from DeepSeek for column '{col_name}'. Skipping.")

                # Add basic metadata for columns DeepSeek might have missed
                if missing_cols:
                     logger.warning(f"DeepSeek metadata missing for columns: {missing_cols}. Adding basic entries.")
                     for col_name in missing_cols:
                          formatted_metadata.append({
                            "name": col_name,
                            "description": "N/A - Metadata generation failed.",
                            "type": self._infer_sql_type(df[col_name]) # Infer type locally
                          })
                     # Reorder to match original df column order might be good practice
                     original_order_map = {name: i for i, name in enumerate(column_names)}
                     formatted_metadata.sort(key=lambda x: original_order_map.get(x['name'], 999))


                logger.info(f"Successfully generated metadata for {len(formatted_metadata)} columns using DeepSeek for '{file_name}'{f'/{sheet_name}' if sheet_name else ''}.")
                return formatted_metadata

            except json.JSONDecodeError as e_json:
                logger.error(f"Failed to parse JSON metadata from DeepSeek: {e_json}. Raw response: {raw_content}")
                return None

        except requests.exceptions.RequestException as e_req:
            logger.error(f"DeepSeek API request failed: {e_req}", exc_info=True)
            return None
        except Exception as e:
            logger.error(f"Error getting column metadata from DeepSeek: {e}", exc_info=True)
            return None

    def _load_df_to_sql(self, df, table_name, column_metadata, file_name, sheet_name=None):
        """Loads a Pandas DataFrame into the SQLite database."""
        if not self.db_conn or not self.db_cursor:
            logger.error("Database connection not available. Cannot load table.")
            return False
        if not table_name:
             logger.error("Invalid table name provided. Cannot load table.")
             return False
        if self._table_exists(table_name):
             logger.warning(f"Table '{table_name}' already exists. Replacing.")
             try: self.db_cursor.execute(f"DROP TABLE \"{table_name}\"") # Use quotes for safety
             except sqlite3.Error as e_drop: logger.error(f"Error dropping existing table {table_name}: {e_drop}"); return False

        try:
            # Use df.to_sql - it handles schema creation and data insertion
            logger.info(f"Loading DataFrame into SQL table '{table_name}'...")
            df.to_sql(table_name, self.db_conn, if_exists='replace', index=False) # Replace if exists, don't write index
            self.db_conn.commit()
            logger.info(f"DataFrame loaded into table '{table_name}' ({len(df)} rows).")

            # Store metadata
            self.table_metadata[table_name] = {
                "columns": column_metadata,
                "source_file": file_name,
                "source_sheet": sheet_name, # Will be None for CSV
                "row_count": len(df)
            }
            # Update file_to_table_map
            if file_name not in self.file_to_table_map: self.file_to_table_map[file_name] = []
            if table_name not in self.file_to_table_map[file_name]: self.file_to_table_map[file_name].append(table_name)

            return True
        except sqlite3.Error as e_sql:
            logger.error(f"SQLite error loading DataFrame to table '{table_name}': {e_sql}", exc_info=True)
            try: self.db_conn.rollback()
            except: pass
            return False
        except Exception as e:
            logger.error(f"Unexpected error loading DataFrame to table '{table_name}': {e}", exc_info=True)
            try: self.db_conn.rollback()
            except: pass
            return False

    # --- Text Chunking and FAISS Indexing (Largely unchanged, operates on text content) ---
    def extract_content(self, file_path):
        """Extract content for text RAG based on file extension. Skips CSV/XLSX here."""
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf': return self._extract_pdf(file_path)
        elif extension == '.pptx': return self._extract_pptx(file_path)
        elif extension in ['.csv', '.xlsx']:
             logger.debug(f"Skipping text extraction for {os.path.basename(file_path)} in extract_content, handled separately for SQL.")
             return [] # Handled by _extract_and_load_xlsx/csv
        else: logger.warning(f"Unsupported file type for text extraction: {file_path}"); return []

    def chunk_content(self, all_content):
        # (Unchanged - chunks text content for FAISS)
        chunks = []; total_words_estimate = 0; chunk_bar = None
        if not all_content: return chunks
        try: total_words_estimate = sum(len(str(item.get('content', '')).split()) for item in all_content)
        except Exception: pass
        if self._is_streamlit and total_words_estimate > 5000: import streamlit as st; chunk_bar = st.progress(0, text=f"Chunking text content...")
        words_processed = 0
        for item_index, item in enumerate(all_content):
            content = item.get('content', ''); source_info = item.get('source_info', {}); file_type = item.get('file_type', 'unknown'); content_type = item.get('type', 'unknown')
            if not isinstance(content, str): content = str(content)
            words = content.split()
            if not words: logger.debug(f"Skipping empty text content block {item_index}"); continue
            item_chunks_created = 0
            for i in range(0, len(words), self.config.chunk_size - self.config.overlap):
                chunk_words = words[i:i + self.config.chunk_size]; chunk_text = " ".join(chunk_words)
                if chunk_text: chunks.append({"content": chunk_text, "source_info": source_info, "file_type": file_type, "type": content_type}); item_chunks_created += 1
            if chunk_bar: words_processed += len(words); progress_percent = min(1.0, words_processed / total_words_estimate) if total_words_estimate > 0 else 0; chunk_bar.progress(progress_percent, text=f"Chunking text content... ({len(chunks)} chunks created)")
        if chunk_bar: chunk_bar.empty()
        logger.info(f"Created {len(chunks)} text chunks from {len(all_content)} content blocks."); return chunks

    def load_faiss_index(self, file_name, embedding_dim):
        # (Unchanged)
        index_path = self.get_index_path(file_name)
        if os.path.exists(index_path):
            try: index = faiss.read_index(index_path); logger.info(f"Loaded FAISS index for {file_name} ({index.ntotal} vectors)"); return index
            except Exception as e: logger.error(f"Error reading FAISS index {index_path}: {e}. Creating new index.")
        logger.info(f"FAISS index not found at {index_path}. Will create new."); return faiss.IndexFlatL2(embedding_dim)
    def save_chunks(self, file_name, chunks):
        # (Unchanged)
        chunks_path = self.get_chunks_path(file_name)
        try:
            with open(chunks_path, 'w', encoding='utf-8') as f: json.dump(chunks, f, indent=2); logger.info(f"Saved {len(chunks)} text chunks to {chunks_path}")
        except Exception as e: logger.error(f"Error saving text chunks for {file_name} to {chunks_path}: {e}");
    def load_chunks(self, file_name):
        # (Unchanged)
        chunks_path = self.get_chunks_path(file_name)
        if os.path.exists(chunks_path):
            try:
                with open(chunks_path, 'r', encoding='utf-8') as f: chunks = json.load(f); logger.info(f"Loaded {len(chunks)} text chunks from {chunks_path}"); return chunks
            except Exception as e: logger.error(f"Error loading/decoding text chunks from {chunks_path}: {e}. Will re-process."); return None
        logger.info(f"Text chunks file not found at {chunks_path}"); return None

    # --- Document Processing (Modified to handle SQL loading) ---
    def process_files(self, progress_callback=None):
        # Reset states
        self.file_chunks = {}; self.faiss_indexes = {}; self.processed_files = set()
        self.table_metadata = {}; self.file_to_table_map = {}
        # Clear existing tables in the DB if it's in-memory or configured to reset
        if self.config.sqlite_db_path == ":memory:" and self.db_cursor:
            logger.info("Clearing existing tables from in-memory database...")
            try:
                self.db_cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = self.db_cursor.fetchall()
                for table in tables: self.db_cursor.execute(f"DROP TABLE IF EXISTS \"{table[0]}\"")
                self.db_conn.commit()
                logger.info(f"Dropped {len(tables)} tables.")
            except sqlite3.Error as e_drop: logger.error(f"Error dropping tables: {e_drop}")

        data_dir = self.config.data_dir; supported_ext = self.config.supported_extensions
        try: all_files = os.listdir(data_dir); process_list = sorted([f for f in all_files if f.lower().endswith(supported_ext)])
        except Exception as e: logger.error(f"Error listing files in data directory {data_dir}: {e}", exc_info=True); return False
        if not process_list: logger.warning(f"No supported files found in {data_dir}"); return True # Completed successfully, but nothing to process
        logger.info(f"Processing {len(process_list)} supported files from {data_dir}"); embedding_dim = self.encoder_model.get_sentence_embedding_dimension(); total_files = len(process_list)
        if progress_callback: progress_callback(f"Found {total_files} file(s). Starting processing...", current_step=0, total_steps=total_files)

        any_success = False
        for idx, file_name in enumerate(process_list):
            current_file_msg = f"Processing ({idx+1}/{total_files}): {file_name}"; logger.info(f"--- {current_file_msg} ---")
            if progress_callback: progress_callback(current_file_msg, current_step=idx, total_steps=total_files)
            file_path = os.path.join(data_dir, file_name); index_path = self.get_index_path(file_name); emb_path = self.get_embedding_path(file_name); chunks_path = self.get_chunks_path(file_name)
            file_extension = os.path.splitext(file_name)[1].lower()
            file_processed_ok = False

            # --- Spreadsheet Processing (SQL Loading) ---
            if file_extension in ['.csv', '.xlsx']:
                sql_success = False
                processed_tables = []
                if file_extension == '.xlsx':
                     sql_success, processed_tables = self._extract_and_load_xlsx(file_path, file_name, progress_callback)
                elif file_extension == '.csv':
                     sql_success, processed_tables = self._extract_and_load_csv(file_path, file_name, progress_callback)

                if sql_success:
                    logger.info(f"Successfully loaded data from {file_name} into SQL table(s): {processed_tables}")
                    self.processed_files.add(file_name) # Mark file as processed if SQL load worked
                    file_processed_ok = True
                    any_success = True
                else:
                    logger.error(f"Failed to load data from {file_name} into SQL database.")
                    if progress_callback: progress_callback(f"❌ Failed SQL processing for {file_name}", is_error=True)
                # No separate text indexing needed for these files using the SQL approach

            # --- Text Processing (Chunking/Indexing) - PDF/PPTX ---
            elif file_extension in ['.pdf', '.pptx']:
                try:
                    chunks = self.load_chunks(file_name)
                    if chunks is None:
                        if progress_callback: progress_callback(f"{current_file_msg} - Extracting Text...", stage="Extracting Text")
                        all_content = self.extract_content(file_path); # Extracts text representation
                        if not all_content: logger.warning(f"No text content extracted from {file_name}. Skipping text indexing."); continue # Skip to next file
                        else:
                            if progress_callback: progress_callback(f"{current_file_msg} - Chunking Text...", stage="Chunking Text")
                            chunks = self.chunk_content(all_content);
                            if not chunks: logger.warning(f"No text chunks generated for {file_name}."); continue
                            else: self.save_chunks(file_name, chunks)

                    # Proceed with embedding/indexing only if chunks exist
                    if chunks:
                        self.file_chunks[file_name] = chunks; logger.debug(f"Stored {len(chunks)} text chunks for {file_name}")
                        faiss_index = None; regenerate_embeddings = False
                        if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                            if progress_callback: progress_callback(f"{current_file_msg} - Verifying Index...", stage="Verifying Index")
                            try:
                                embeddings = np.load(emb_path)
                                if embeddings.ndim != 2 or embeddings.shape[1] != embedding_dim or embeddings.shape[0] != len(chunks): logger.warning(f"Index mismatch for {file_name}. Regenerating."); regenerate_embeddings = True
                                else: faiss_index = self.load_faiss_index(file_name, embedding_dim); logger.info(f"Verified existing text index data for {file_name}")
                            except Exception as e: logger.error(f"Error verifying index for {file_name}: {e}. Regenerating...", exc_info=True); regenerate_embeddings = True; faiss_index = None
                        else: regenerate_embeddings = True; logger.info(f"No complete index found for {file_name}, will generate.")

                        if regenerate_embeddings or faiss_index is None:
                            logger.info(f"Generating text embeddings/index for {file_name}...")
                            if progress_callback: progress_callback(f"{current_file_msg} - Embedding Text...", stage="Embedding Text")
                            content_list = [chunk['content'] for chunk in chunks];
                            if not content_list: logger.warning(f"No text content to embed for {file_name}.");
                            else:
                                embeddings = self.encoder_model.encode(content_list, batch_size=64, show_progress_bar=False, convert_to_numpy=True).astype('float32')
                                if embeddings.shape[0] == 0: logger.warning(f"Text embedding yielded no vectors for {file_name}.");
                                else:
                                    np.save(emb_path, embeddings); logger.info(f"Saved {embeddings.shape[0]} text embeddings for {file_name}")
                                    if progress_callback: progress_callback(f"{current_file_msg} - Indexing Text...", stage="Indexing Text")
                                    faiss_index = faiss.IndexFlatL2(embedding_dim); faiss_index.add(embeddings); faiss.write_index(faiss_index, index_path); logger.info(f"Saved FAISS text index for {file_name} ({faiss_index.ntotal} vectors)")

                        if faiss_index is not None and faiss_index.ntotal > 0:
                             self.faiss_indexes[file_name] = faiss_index; self.processed_files.add(file_name); logger.debug(f"Stored FAISS text index for {file_name}")
                             file_processed_ok = True
                             any_success = True
                        else: logger.warning(f"No valid text index created for {file_name}.")
                    else:
                        logger.info(f"No text chunks for {file_name}, skipping text indexing.")

                except Exception as e_text:
                    logger.error(f"Failed text processing for {file_name}: {e_text}", exc_info=True)
                    if progress_callback: progress_callback(f"❌ Error during text processing for {file_name}: {e_text}", is_error=True)
                    try: # Clean up potentially inconsistent text index files
                        for p in [index_path, emb_path, chunks_path]:
                            if os.path.exists(p): os.remove(p); logger.debug(f"Removed potentially inconsistent text file {p}")
                    except OSError as e_clean: logger.error(f"Error cleaning up text files for {file_name}: {e_clean}")
                    if file_name in self.file_chunks: del self.file_chunks[file_name]
                    if file_name in self.faiss_indexes: del self.faiss_indexes[file_name]
            else:
                 logger.warning(f"Skipping file {file_name} due to unsupported extension {file_extension} in main processing loop.")


        # Final summary
        final_text_index_count = len(self.faiss_indexes)
        final_sql_table_count = len(self.table_metadata)
        total_processed_ok = len(self.processed_files)

        if any_success:
            logger.info(f"--- Processing Complete. Processed {total_processed_ok}/{total_files} files. Text Indices: {final_text_index_count}. SQL Tables: {final_sql_table_count}. ---")
        else:
            logger.warning(f"--- Processing Complete. No documents successfully processed (no text index or SQL table). ---")

        # Return True if the process finished, even if nothing was loaded. The UI can check counts.
        return True


    # --- Querying Logic (Includes Faiss search, Query Classification, SQL Generation/Execution, Text RAG) ---

    # Faiss search for text documents (unchanged)
    def query_files(self, query):
        if not self.faiss_indexes:
            logger.warning("No text indexes available for Faiss querying.")
            return {}
        query_results = {}; # Store results per file
        try:
            logger.info(f"Encoding query for Faiss text search: '{query[:100]}...'")
            query_embedding = self.encoder_model.encode(query, convert_to_numpy=True).astype("float32")
            query_embedding = np.array([query_embedding])
            if query_embedding.ndim != 2: raise ValueError("Query embedding shape error")

            for file_name, index in self.faiss_indexes.items():
                if index is None or index.ntotal == 0: logger.debug(f"Skipping empty text index for {file_name}"); continue
                try:
                    k_search = min(self.config.k_retrieval, index.ntotal)
                    D, I = index.search(query_embedding, k=k_search) # D=distances (L2), I=indices
                    indices, distances = I[0], D[0]

                    current_file_chunks = self.file_chunks.get(file_name)
                    if not current_file_chunks: logger.error(f"Text chunks missing for {file_name} during query!"); continue

                    file_results = []
                    processed_indices = set()
                    for i, idx in enumerate(indices):
                        if idx == -1 or idx in processed_indices or not (0 <= idx < len(current_file_chunks)): continue
                        processed_indices.add(idx)
                        chunk = current_file_chunks[idx]
                        file_results.append({
                            "source_info": chunk.get('source_info', {}),
                            "file_type": chunk.get('file_type', 'unknown'),
                            "content": chunk.get('content', ''),
                            "score": round(float(distances[i]), 4), # L2 distance, lower is better
                            "type": chunk.get('type', 'unknown')
                        })

                    file_results.sort(key=lambda x: x['score']) # Ensure sorted by score ascending
                    if file_results:
                        if file_name not in query_results: query_results[file_name] = []
                        query_results[file_name].extend(file_results)

                except Exception as e_search:
                    logger.error(f"Error searching text index {file_name} for query '{query[:50]}...': {e_search}", exc_info=True)

            logger.debug(f"Faiss text search complete for query '{query[:50]}...'. Found results in {len(query_results)} files.")
            return query_results

        except Exception as e_query:
            logger.error(f"Error during Faiss text query processing for '{query[:50]}...': {e_query}", exc_info=True)
            return {}

    # Aggregate text context (unchanged, uses Faiss results)
    def aggregate_context(self, query_results, strategy="top_k"):
        all_context = {}; max_chars = self.config.max_context_tokens * 3 # Estimate based on tokens
        logger.info(f"Aggregating text context using strategy '{strategy}', max_chars ~{max_chars}")
        if not query_results: return all_context
        flat_results = []
        for file_name, results in query_results.items():
            for res in results: flat_results.append({**res, "file_name": file_name})
        flat_results.sort(key=lambda x: x['score']) # Sort globally by score (lower L2 is better)
        aggregated_context_str = ""; total_aggregated_chars = 0; added_chunks_count = 0; context_sources = set()
        limit = self.config.k_retrieval if strategy == "top_k" else len(flat_results)
        for i, res in enumerate(flat_results):
            if added_chunks_count >= limit: logger.debug(f"Reached global top-k limit ({limit}) for text context aggregation."); break
            source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown'); content_type = res.get('type', 'unknown'); score = res['score']; content_body = res['content']
            source_parts = [f"Source: {file_name}"]
            if file_type == 'pptx': source_parts.append(f"Slide: {source_info.get('slide', 'N/A')}" + (f" (merged from {source_info.get('merged_from')})" if source_info.get('merged_from') else ""))
            elif file_type == 'pdf': source_parts.append(f"Page: {source_info.get('page', 'N/A')}")
            source_parts.append(f"Type: {content_type}"); source_parts.append(f"Score: {score:.4f}")
            source_str = ", ".join(source_parts); content_header = f"--- Context from {source_str} ---\n"; content_to_add = content_header + content_body + "\n\n"; content_chars = len(content_to_add)
            if total_aggregated_chars + content_chars <= max_chars:
                aggregated_context_str += content_to_add; total_aggregated_chars += content_chars; added_chunks_count += 1; context_sources.add(file_name)
            else:
                if added_chunks_count == 0: # Try to add truncated first chunk if too large
                     remaining_chars = max_chars - total_aggregated_chars - len(content_header) - 20 # Headroom
                     if remaining_chars > 50:
                         truncated_body = content_body[:remaining_chars]; aggregated_context_str += content_header + truncated_body + "\n[...TRUNCATED CONTEXT...]\n\n"; total_aggregated_chars += len(aggregated_context_str); added_chunks_count += 1; context_sources.add(file_name); logger.warning("Text context truncated (first chunk too large).")
                     else: logger.warning("First text chunk too large to fit even truncated, skipping.")
                logger.info(f"Stopping text context aggregation at {added_chunks_count} chunks ({total_aggregated_chars}/{max_chars} chars). Limit reached."); break
        final_context = aggregated_context_str.strip()
        if final_context: all_context = {"combined_context": final_context, "source_files": sorted(list(context_sources))}; logger.info(f"Aggregated {total_aggregated_chars} chars from {added_chunks_count} text chunks across {len(context_sources)} files.")
        else: logger.warning("No text context aggregated within limits.")
        return all_context

    # Query LLM with text context (unchanged)
    def query_llm(self, query, context_data, retry_count=1):
        combined_context = context_data.get("combined_context", ""); source_files = context_data.get("source_files", []); source_file_str = ", ".join(source_files) if source_files else "the provided text documents"
        if not combined_context: logger.warning("No text context provided for LLM query."); return f"Could not generate answer: No relevant text context found in {source_file_str}."
        try:
            llm = self._get_llm()
            system_prompt = f"""You are an AI assistant answering questions based ONLY on the provided text context from document(s): '{source_file_str}'. Use ONLY information presented between '--- START CONTEXT ---' and '--- END CONTEXT ---'. Do not use any prior knowledge. If the answer cannot be found in the text, state: "Based on the provided text context from {source_file_str}, I cannot answer this question." Do NOT perform calculations unless the calculation is explicitly shown in the text."""
            human_prompt = f"""Context from document(s) '{source_file_str}':\n--- START CONTEXT ---\n{combined_context}\n--- END CONTEXT ---\n\nUser Question: {query}\n\nAnswer based ONLY on the provided text context:"""
            full_prompt_messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
            logger.info(f"Querying Groq model {self.config.llm_model} using text context from {source_file_str}...")
            answer = f"Error: LLM text query failed after retries."
            for attempt in range(retry_count + 1):
                try:
                    response = llm.invoke(full_prompt_messages)
                    answer = response.content.strip() if hasattr(response, 'content') else str(response).strip()
                    logger.info(f"Groq response received for text RAG (attempt {attempt+1}). Length: {len(answer)}")
                    if not answer and attempt < retry_count: logger.warning("Empty text RAG response, retrying..."); time.sleep(1); continue
                    return answer or f"Received empty response from LLM for text context from {source_file_str}."
                except Exception as e_api:
                    logger.warning(f"Groq API attempt {attempt+1} for text RAG failed: {e_api}")
                    if attempt < retry_count: time.sleep(1.5 ** attempt); logger.info("Retrying Groq text RAG query...")
                    else: answer = f"Error: Failed to get answer from LLM for text RAG after {retry_count+1} attempts. (API Error: {e_api})"
            return answer
        except Exception as e_setup:
            logger.error(f"Error setting up/calling Groq API for text RAG: {e_setup}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"LLM text query failed: {e_setup}")
            return f"Error: Could not query LLM for text RAG due to setup error: {e_setup}"

    # --- Query Analysis Helpers ---
    # LLM call helper (unchanged)
    def _call_llm_for_analysis(self, prompt_messages, task_description):
        try:
            llm = self._get_llm()
            logger.info(f"Calling Groq LLM ({self.config.llm_model}) for internal task: {task_description}...")
            response = llm.invoke(prompt_messages)
            content = response.content.strip() if hasattr(response, 'content') else str(response).strip()
            logger.info(f"LLM response for {task_description}: '{content[:150]}...'")
            return content
        except Exception as e:
            logger.error(f"LLM call failed during {task_description}: {e}", exc_info=True)
            return None

    # Query Classification (Revised JSON parsing)
    def _classify_query(self, query):
        system_prompt = """You are an expert query analyzer. Classify the user query into ONE of the following categories:
1.  'Simple Retrieval': Asking for specific facts, definitions, or simple summaries directly extractable from text.
2.  'Complex/Reasoning (Text)': Requires combining information from multiple text passages, summarization across sections, or reasoning based *only* on the text provided. Does not involve calculations on tabular data.
3.  'Structured Query (SQL)': Requires operations on structured data (like CSV or Excel tables stored in a SQL database), such as calculations (sum, average, count), filtering based on values, grouping, sorting, or comparisons across rows/columns.

Analyze the query carefully. Respond ONLY in JSON format with two keys: "classification" (the category string) and "confidence" (a float between 0.0 and 1.0 indicating your certainty, e.g., 0.95).

Examples:
Query: "What is the definition of RAG?" -> {"classification": "Simple Retrieval", "confidence": 0.99}
Query: "Summarize the main challenges discussed in the report." -> {"classification": "Complex/Reasoning (Text)", "confidence": 0.9}
Query: "What is the total revenue for product X in Q3?" -> {"classification": "Structured Query (SQL)", "confidence": 0.98}
Query: "Compare the findings section with the conclusion." -> {"classification": "Complex/Reasoning (Text)", "confidence": 0.85}
Query: "List all employees hired after 2022." -> {"classification": "Structured Query (SQL)", "confidence": 0.96}
Query: "Tell me about the company history." -> {"classification": "Simple Retrieval", "confidence": 0.9}
"""
        human_prompt = f"Classify the following user query:\n\"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        raw_response = self._call_llm_for_analysis(messages, "query classification")

        if not raw_response:
             logger.warning("Query classification failed (LLM call unsuccessful). Defaulting.")
             return "Simple Retrieval", 0.5

        try:
            # Attempt to find JSON within the response more robustly
            json_match = re.search(r"\{.*\}", raw_response, re.DOTALL)
            if json_match:
                json_str = json_match.group(0)
                result = json.loads(json_str)
                classification = result.get("classification")
                confidence = float(result.get("confidence", 0.0))
                valid_classifications = ["Simple Retrieval", "Complex/Reasoning (Text)", "Structured Query (SQL)"]
                if classification in valid_classifications and 0.0 <= confidence <= 1.0:
                    logger.info(f"Query classified as '{classification}' with confidence {confidence:.2f}")
                    return classification, confidence
                else:
                    logger.warning(f"Invalid classification or confidence in JSON: {json_str}. Defaulting.")
            else:
                 logger.warning(f"No JSON object found in classification response: {raw_response}. Defaulting.")

            # Fallback if JSON parsing failed or content was invalid
            return "Simple Retrieval", 0.5

        except (json.JSONDecodeError, TypeError, AttributeError, ValueError) as e:
            logger.warning(f"Failed to parse classification JSON: '{raw_response}'. Error: {e}. Defaulting.")
            return "Simple Retrieval", 0.5


    # Text Query Decomposition (unchanged, uses LLM)
    def _decompose_query(self, query):
        system_prompt = """You are an expert query decomposer. Break down a complex user query needing TEXT-based information into simpler, factual sub-queries. Each sub-query should aim to retrieve a specific piece of information from text. Do NOT decompose queries asking for calculations or filtering on tables. Format as a numbered list. If the query is simple or cannot be decomposed, return the original query prefixed with '1. '."""
        human_prompt = f"Decompose the following complex text query:\n\"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        decomposition = self._call_llm_for_analysis(messages, "text query decomposition")
        if decomposition:
            sub_queries = [re.sub(r"^\d+\.\s*", "", line).strip() for line in decomposition.split('\n') if line.strip()]
            # Ensure decomposition actually changed the query
            if sub_queries and not (len(sub_queries) == 1 and sub_queries[0].strip() == query.strip()):
                logger.info(f"Decomposed text query into: {sub_queries}")
                return sub_queries
            else: logger.info("Decomposition didn't yield distinct sub-queries for text. Using original."); return [query]
        else: logger.warning("Text query decomposition failed. Using original query."); return [query]

    # --- [NEW] SQL Target Table Identification ---
    def _identify_target_sql_table(self, query):
        """Uses LLM to identify the most relevant SQL table based on query and available schemas."""
        if not self.table_metadata:
            logger.warning("No table metadata available to identify target SQL table.")
            return None

        # Prepare schema overview for the LLM
        schema_overview = "Available SQL Tables:\n"
        table_options = []
        for table_name, meta in self.table_metadata.items():
            table_options.append(table_name)
            col_descs = []
            for col in meta.get('columns', []):
                # Safely get description, default to "N/A"
                desc = col.get('description', 'N/A')
                # Limit description length for the prompt if necessary
                desc_short = (desc[:100] + '...') if len(desc) > 103 else desc
                col_descs.append(f"- `{col.get('name', 'UNKNOWN_COL')}` (Type: {col.get('type', 'UNKNOWN_TYPE')}, Desc: {desc_short})")

            sheet_info = f" (Sheet: '{meta.get('source_sheet')}')" if meta.get('source_sheet') else ""
            schema_overview += f"\nTable: `{table_name}` (From File: '{meta.get('source_file', 'N/A')}'{sheet_info}, Rows: {meta.get('row_count', 'N/A')})\n"
            schema_overview += "\n".join(col_descs) + "\n"

        if not table_options:
             logger.warning("No tables found in metadata during identification.")
             return None

        system_prompt = f"""You are an expert data analyst. Your task is to choose the single most relevant SQL table to answer the user's query, based on the table names, source files, and column descriptions provided.

Available Table Schemas:
{schema_overview}

Analyze the user's query and the schemas. Respond ONLY with the exact name of the single most appropriate table from the list: {table_options}.
If no single table seems clearly relevant or sufficient to answer the query, respond ONLY with the word "NONE".
Do not add any explanation or introductory text.
"""
        human_prompt = f"User Query: \"{query}\"\n\nMost relevant table name:"
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        target_table_name = self._call_llm_for_analysis(messages, "SQL target table identification")

        if not target_table_name:
            logger.error("SQL target table identification failed (LLM call unsuccessful).")
            return None

        target_table_name = target_table_name.strip().strip('`"') # Clean potential markdown/quotes

        if target_table_name == "NONE":
            logger.info("LLM indicated no single relevant SQL table found.")
            return None
        elif target_table_name in self.table_metadata:
            logger.info(f"LLM identified '{target_table_name}' as the target SQL table.")
            return target_table_name
        else:
            logger.warning(f"LLM returned an invalid table name '{target_table_name}'. Valid options were: {table_options}")
            # Optional: Add fuzzy matching here if needed
            return None

    # --- SQL Query Generation (Improved Sanitization) ---
    def _generate_sql_query(self, query, table_name, table_meta):
        """Generate SQL query using LLM based on user query and table schema/metadata."""
        if not table_meta or 'columns' not in table_meta:
            logger.error(f"Metadata missing for table '{table_name}'. Cannot generate SQL.")
            return None

        # Format schema with descriptions for the LLM
        schema_description = f"Table Name: `{table_name}`\nColumns:\n"
        for col in table_meta['columns']:
             safe_col_name = f'"{col["name"]}"' # Ensure column names are quoted
             col_type = col.get('type', 'TEXT') # Default type if missing
             col_desc = col.get('description', 'N/A') # Default description if missing
             schema_description += f"- {safe_col_name} (Type: {col_type}, Description: {col_desc})\n"

        system_prompt = f"""You are an expert SQL query writer specializing in SQLite. Given a user query and the schema of a table (including column descriptions), write a concise and valid SQLite query to answer the user's question.

Table Schema:
{schema_description}

Instructions:
- Write ONLY a single, valid SQLite query.
- Use the exact column names provided in the schema, ensuring they are correctly quoted (e.g., "Column Name With Spaces").
- Pay close attention to the column descriptions to understand the data's meaning.
- Handle potential data types correctly in your query (e.g., use appropriate functions for dates stored as TEXT if needed, ensure numeric comparisons are valid). Use `CAST(column AS REAL)` or `CAST(column AS INTEGER)` if needed for comparisons or calculations on TEXT columns containing numbers.
- Output ONLY the raw SQL query. No explanations, comments, markdown backticks (```sql ... ```), or introductory text like "Here is the SQL query:".
- If the query cannot be answered accurately with the given table and columns, output ONLY the exact text: QUERY_CANNOT_BE_ANSWERED
"""
        human_prompt = f"User Query: \"{query}\"\n\nSQLite Query:" # Explicitly ask for SQLite
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        sql_query = self._call_llm_for_analysis(messages, f"SQL query generation for table {table_name}")

        if not sql_query:
            logger.error("SQL query generation failed (LLM call unsuccessful).")
            return None

        # Basic validation and cleanup
        sql_query = sql_query.strip()
        if sql_query == "QUERY_CANNOT_BE_ANSWERED":
             logger.warning(f"LLM indicated query cannot be answered for table {table_name}.")
             return None

        # Remove potential markdown backticks more robustly
        sql_query = re.sub(r"^```(?:sql)?\s*", "", sql_query, flags=re.IGNORECASE | re.DOTALL)
        sql_query = re.sub(r"\s*```$", "", sql_query, flags=re.DOTALL)
        sql_query = sql_query.strip() # Strip again after removing backticks

        # Allow only SELECT statements for safety in this context
        if not sql_query.upper().startswith("SELECT"):
             logger.warning(f"Generated SQL is not a SELECT statement: '{sql_query[:100]}...' - Discarding for safety.")
             return None

        # Remove trailing semicolon if present
        if sql_query.endswith(';'): sql_query = sql_query[:-1]

        logger.info(f"Generated SQL query for table '{table_name}':\n{sql_query}")
        return sql_query

    # --- SQL Execution (Improved Error Logging and Result Formatting) ---
    def _execute_sql_query(self, sql_query):
        """Executes a SQL query against the database and returns results."""
        if not self.db_cursor:
            logger.error("No database cursor available for SQL execution.")
            return None, "Database connection error."
        if not sql_query:
            logger.error("No SQL query provided for execution.")
            return None, "No SQL query generated."

        try:
            logger.info(f"Executing SQL: {sql_query}")
            start_time = time.time()
            # Use pandas read_sql_query for potentially better type handling? Stick to cursor for now.
            self.db_cursor.execute(sql_query)
            results = self.db_cursor.fetchall() # List of tuples
            column_names = [desc[0] for desc in self.db_cursor.description] if self.db_cursor.description else []

            exec_time = time.time() - start_time
            logger.info(f"SQL query executed successfully in {exec_time:.3f}s. Fetched {len(results)} rows.")

            # Format results into a readable string
            if not results:
                return "No results found.", None # Success, but no data

            # Use pandas to create a formatted string (easier for tables)
            try:
                 df_results = pd.DataFrame(results, columns=column_names)
                 max_rows_preview = 20 # Limit preview size in string output
                 output_str = df_results.to_string(index=False, max_rows=max_rows_preview)
                 if len(df_results) > max_rows_preview:
                      output_str += f"\n... (truncated, {len(df_results)} total rows)"
                 # Add header for clarity
                 output_str = f"Query Result ({len(results)} row(s)):\n" + output_str

            except Exception as e_pd:
                 logger.warning(f"Pandas formatting of SQL results failed: {e_pd}. Falling back to basic string format.")
                 # Fallback basic formatting
                 max_rows_preview = 20
                 output_str = f"Query Result ({len(results)} row(s)):\n"
                 header = ", ".join(column_names)
                 output_str += header + "\n"
                 output_str += "-" * len(header) + "\n" # Separator line
                 for i, row in enumerate(results):
                     if i >= max_rows_preview:
                         output_str += f"... (truncated, {len(results) - max_rows_preview} more rows)\n"
                         break
                     # Convert each item in row to string safely
                     row_str = ", ".join(map(lambda x: str(x) if x is not None else "NULL", row))
                     output_str += row_str + "\n"

            return output_str.strip(), None # Return formatted string and no error

        except sqlite3.Error as e_sql:
            # Provide more context in the error message returned to the main loop
            err_msg = f"SQLite execution error: {e_sql} (Query: {sql_query[:200]}...)"
            logger.error(err_msg, exc_info=False) # Avoid full traceback for common SQL errors
            try: self.db_conn.rollback()
            except: pass
            return None, f"Database error during execution: {e_sql}" # Return cleaner error message
        except Exception as e:
            err_msg = f"Unexpected error during SQL execution: {e}"
            logger.error(f"{err_msg} for query: {sql_query}", exc_info=True)
            try: self.db_conn.rollback()
            except: pass
            return None, err_msg

    # --- Final Answer Synthesis from SQL Result (Improved Handling of No Results) ---
    def _synthesize_answer_from_sql(self, user_query, sql_result_str, table_name, table_meta):
        """Uses the main LLM to create a natural language answer from SQL results."""
        source_desc = f"the table '{table_name}'"
        if table_meta:
            source_desc += f" from file '{table_meta.get('source_file', 'N/A')}'"
            if table_meta.get('source_sheet'): source_desc += f" (sheet '{table_meta.get('source_sheet')}')"

        # Handle "No results found" case directly before calling LLM if desired
        if sql_result_str == "No results found.":
             logger.info("SQL query returned no results. Synthesizing direct 'not found' answer.")
             # Optional: Ask LLM to phrase this nicely, or just return a standard message
             # return f"Based on the data in {source_desc}, I could not find any information matching your query: '{user_query}'"
             # Let's keep using LLM for consistency, but provide specific instruction
             system_prompt = f"""You are an AI assistant. You were asked a question by a user, and a SQL query was run against {source_desc}. The query returned "No results found.".
Your task is to formulate a polite, natural language response indicating that the requested information could not be found in the specified data source.
Do not mention the SQL query. Start the answer directly.
"""
             human_prompt = f"""Original User Question: "{user_query}"

Data retrieved from {source_desc}: No results found.

Formulate the final "not found" answer:"""
        else:
            # Original prompt for when results *are* found
            system_prompt = f"""You are an AI assistant. You were asked a question by a user, and relevant data was retrieved by executing a SQL query against {source_desc}.
Your task is to formulate a clear, concise, and natural language answer to the original user question based *only* on the provided SQL query results.

Do NOT just repeat the raw data. Summarize or present the key information needed to answer the question directly and conversationally.
If the data clearly answers the question, state the answer directly.
If the data seems incomplete or doesn't directly answer the question, state what you found and mention that it might not be the full answer.
Do not mention the SQL query itself or the table name unless the user specifically asked how the data was obtained.
Start your answer directly, without introductory phrases like "Based on the data...".
"""
            human_prompt = f"""Original User Question: "{user_query}"

Data retrieved from {source_desc}:
--- START SQL RESULT DATA ---
{sql_result_str}
--- END SQL RESULT DATA ---

Formulate the final answer based ONLY on the provided data:"""

        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        final_answer = self._call_llm_for_analysis(messages, "SQL result synthesis")

        if not final_answer:
            logger.error("Final answer synthesis from SQL result failed (LLM call unsuccessful). Returning structured info.")
            if sql_result_str == "No results found.":
                 return f"I searched in {source_desc} but couldn't find information for your query."
            else:
                 return f"I found the following data in {source_desc}, but had trouble summarizing it:\n```\n{sql_result_str}\n```"
        else:
            logger.info(f"Synthesized final answer from SQL result for query '{user_query[:50]}...'")
            return final_answer


    # ==========================================================================
    # REVISED run_query FUNCTION (v2 - Incorporating _identify_target_sql_table)
    # ==========================================================================
    def run_query(self, query, context_strategy="top_k"):
        """Complete query pipeline: classify, (SQL query OR text RAG), synthesize."""
        logger.info(f"--- Starting query execution (SQL v2): '{query[:100]}...' ---")
        final_results = {
            "query": query,
            "classification": None, "classification_confidence": None,
            "sub_queries": None, "retrieval_results": {}, # Text retrieval results (still useful)
            "target_table": None, "target_table_metadata": None, # SQL target info
            "sql_query_generated": None, # Generated SQL
            "sql_execution_result": None, # Raw result from SQL execution
            "aggregated_context_data": {}, # Text RAG context
            "answer": "", "answer_source": None, # 'SQL Query' or 'Text RAG'
            "status": "Started", "error": None
        }
        is_streamlit = self._is_streamlit
        if is_streamlit: import streamlit as st

        try:
            # === Run Faiss search EARLY (can be useful for context later or hybrid answers) ===
            spinner_msg_text = "Pre-searching text documents (if any)..."
            text_retrieval_results = {}
            if self.faiss_indexes:
                 if is_streamlit:
                     with st.spinner(spinner_msg_text):
                        text_retrieval_results = self.query_files(query)
                 else:
                     logger.info(spinner_msg_text)
                     text_retrieval_results = self.query_files(query)
                 final_results["retrieval_results"] = text_retrieval_results
            else:
                 logger.info("No text indexes (PDF/PPTX) found for pre-search.")


            # === 1. Classify Query ===
            spinner_msg_classify = "Analyzing query intent..."
            if is_streamlit:
                with st.spinner(spinner_msg_classify):
                    classification, confidence = self._classify_query(query)
            else:
                logger.info(spinner_msg_classify)
                classification, confidence = self._classify_query(query)

            final_results["classification"] = classification
            final_results["classification_confidence"] = confidence
            logger.info(f"Query classified as: {classification} (Confidence: {confidence:.2f})")

            # === Attempt SQL Query Path if classified and confident ===
            sql_attempted = False
            sql_succeeded = False
            target_table_name = None # Initialize target table name

            # Check if query is classified for SQL, confidence is high, AND tables exist
            if classification == "Structured Query (SQL)" and confidence >= self.config.dataframe_query_confidence_threshold and self.table_metadata:
                sql_attempted = True
                logger.info("Attempting SQL query path based on classification and confidence.")

                # === 2a. Identify Target Table using LLM and Metadata === <--- NEW LOGIC
                spinner_msg_identify = "Identifying relevant SQL table..."
                if is_streamlit:
                    with st.spinner(spinner_msg_identify):
                        target_table_name = self._identify_target_sql_table(query)
                else:
                    logger.info(spinner_msg_identify)
                    target_table_name = self._identify_target_sql_table(query)

                if target_table_name and target_table_name in self.table_metadata:
                    final_results["target_table"] = target_table_name
                    target_meta = self.table_metadata[target_table_name]
                    final_results["target_table_metadata"] = target_meta
                    logger.info(f"Proceeding with SQL query on identified table: '{target_table_name}'")

                    # === 3a. Generate SQL Query ===
                    spinner_msg_gen = f"Generating SQL query for '{target_table_name}'..."
                    sql_query = None
                    if is_streamlit:
                        with st.spinner(spinner_msg_gen):
                            sql_query = self._generate_sql_query(query, target_table_name, target_meta)
                    else:
                        logger.info(spinner_msg_gen)
                        sql_query = self._generate_sql_query(query, target_table_name, target_meta)

                    if sql_query:
                        final_results["sql_query_generated"] = sql_query

                        # === 4a. Execute SQL Query ===
                        spinner_msg_exec = f"Executing SQL query on '{target_table_name}'..."
                        sql_result_str, sql_error = None, None
                        if is_streamlit:
                             with st.spinner(spinner_msg_exec):
                                sql_result_str, sql_error = self._execute_sql_query(sql_query)
                        else:
                             logger.info(spinner_msg_exec)
                             sql_result_str, sql_error = self._execute_sql_query(sql_query)

                        if sql_error:
                            logger.error(f"SQL execution failed: {sql_error}")
                            final_results["status"] = "Failed during SQL execution"
                            # Make the error message more informative for the user/logs
                            final_results["error"] = f"SQL query execution failed: {sql_error}"
                            # Fall through to text RAG
                        else:
                            final_results["sql_execution_result"] = sql_result_str
                            logger.info(f"SQL execution successful. Raw result: {str(sql_result_str)[:200]}...") # Log start of result

                            # === 5a. Synthesize Final Answer from SQL Result ===
                            spinner_msg_synth = f"Generating final answer from '{target_table_name}' data..."
                            final_answer = None
                            if is_streamlit:
                                with st.spinner(spinner_msg_synth):
                                    final_answer = self._synthesize_answer_from_sql(query, sql_result_str, target_table_name, target_meta)
                            else:
                                logger.info(spinner_msg_synth)
                                final_answer = self._synthesize_answer_from_sql(query, sql_result_str, target_table_name, target_meta)

                            final_results["answer"] = final_answer
                            final_results["status"] = "Completed Successfully"
                            final_results["answer_source"] = "SQL Query"
                            sql_succeeded = True
                            final_results["error"] = None # Clear error on full SQL success

                    else: # SQL Generation Failed
                        logger.warning(f"Failed to generate SQL query for table '{target_table_name}'. Falling back to text search.")
                        final_results["status"] = "Failed SQL query generation"
                        final_results["error"] = f"Could not generate SQL query for table '{target_table_name}'."
                        # Fall through to text RAG
                else: # Target Table Identification Failed
                     if not target_table_name:
                         logger.warning("LLM could not identify a suitable SQL table. Falling back to text search.")
                         final_results["error"] = "Could not identify a relevant SQL table for the query."
                     else: # Table name returned but not in metadata (shouldn't happen often with new logic)
                          logger.error(f"Identified table '{target_table_name}' not found in metadata. Inconsistency?")
                          final_results["error"] = f"Identified table '{target_table_name}' not found in metadata."
                     final_results["status"] = "Failed SQL table identification"
                     # Fall through to text RAG

            # === Fallback to Text RAG Path ===
            if not sql_succeeded:
                 if sql_attempted: logger.info("SQL query path failed or was inconclusive. Proceeding with standard text RAG.")
                 else: logger.info("Query classified for text retrieval/reasoning or no suitable SQL path. Proceeding with text RAG.")

                 # Check if text indexes exist for Text RAG
                 if not self.faiss_indexes:
                      logger.warning("No text indexes available (PDF/PPTX). Cannot perform text RAG.")
                      final_results["status"] = "Failed: No suitable processing path"
                      # Preserve previous error if SQL was attempted
                      if not final_results.get("error"):
                           final_results["error"] = "Query requires text analysis, but no text documents (PDF/PPTX) are indexed."
                      # Construct answer indicating failure reason
                      fail_reason = final_results.get('error', 'No text documents available for analysis.')
                      final_results["answer"] = f"Could not process the query. Reason: {fail_reason}"
                      return final_results # Cannot proceed

                 final_results["answer_source"] = "Text RAG"

                 # === 2b. Decompose if Complex Text Query ===
                 queries_to_retrieve = [query]
                 if classification == "Complex/Reasoning (Text)":
                     spinner_msg_decompose = "Decomposing complex text query..."
                     sub_queries = []
                     if is_streamlit:
                         with st.spinner(spinner_msg_decompose):
                             sub_queries = self._decompose_query(query)
                     else:
                         logger.info(spinner_msg_decompose)
                         sub_queries = self._decompose_query(query)
                     if sub_queries and sub_queries != [query]:
                          queries_to_retrieve = sub_queries
                          final_results["sub_queries"] = sub_queries
                          logger.info(f"Using sub-queries for text retrieval: {sub_queries}")

                 # === 3b. Retrieve relevant text chunks (using pre-fetched or new results) ===
                 # Use results if already fetched, otherwise fetch now
                 current_retrieval_results = final_results.get("retrieval_results")
                 if not current_retrieval_results:
                     spinner_msg_retrieve = f"Searching text documents for {len(queries_to_retrieve)} query part(s)..."
                     all_query_results = {}
                     if is_streamlit:
                         with st.spinner(spinner_msg_retrieve):
                             for i, q in enumerate(queries_to_retrieve):
                                 logger.info(f"Retrieving text for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                                 results_for_q = self.query_files(q)
                                 for file, res_list in results_for_q.items():
                                     if file not in all_query_results: all_query_results[file] = []
                                     all_query_results[file].extend(res_list) # Combine results from sub-queries
                     else:
                         logger.info(spinner_msg_retrieve)
                         for i, q in enumerate(queries_to_retrieve):
                              logger.info(f"Retrieving text for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                              results_for_q = self.query_files(q)
                              for file, res_list in results_for_q.items():
                                     if file not in all_query_results: all_query_results[file] = []
                                     all_query_results[file].extend(res_list)

                     # De-duplicate and sort text results
                     final_retrieval_results = {}
                     for file, res_list in all_query_results.items():
                          unique_content = {}
                          for res in res_list:
                              # Use content and basic source info (page/slide) as key for uniqueness
                              source_key = tuple(sorted(res.get('source_info',{}).items()))
                              key = (res['content'], source_key)
                              # Keep the one with the best (lowest) score
                              if key not in unique_content or res['score'] < unique_content[key]['score']:
                                  unique_content[key] = res
                          # Sort unique results by score and take top K for this file
                          sorted_unique_results = sorted(unique_content.values(), key=lambda x: x['score'])
                          final_retrieval_results[file] = sorted_unique_results[:self.config.k_retrieval] # Apply K limit per file
                     final_results["retrieval_results"] = final_retrieval_results
                     current_retrieval_results = final_retrieval_results # Update for aggregation step
                 else:
                      logger.info("Reusing text retrieval results from pre-search step.")
                      # Ensure results are sorted and limited if reused
                      processed_retrieval_results = {}
                      for file, res_list in current_retrieval_results.items():
                            # Ensure sorting and limiting happens even if reused
                            res_list.sort(key=lambda x: x['score'])
                            processed_retrieval_results[file] = res_list[:self.config.k_retrieval]
                      final_results["retrieval_results"] = processed_retrieval_results # Store processed results
                      current_retrieval_results = processed_retrieval_results # Update for aggregation step


                 # Check if retrieval actually found something relevant
                 if not current_retrieval_results or not any(current_retrieval_results.values()):
                     logger.warning("No relevant text chunks found for the query/sub-queries.")
                     final_results["status"] = "Completed: No relevant information found."
                     final_results["answer"] = "Could not find relevant information in the text documents (PDF/PPTX)."
                     # Append previous error message if exists
                     if final_results["error"]: final_results["answer"] += f" (Previous attempt failed: {final_results['error']})"
                     return final_results

                 # === 4b. Aggregate text context ===
                 spinner_msg_agg = "Gathering text context..."
                 aggregated_context_data = {}
                 if is_streamlit:
                     with st.spinner(spinner_msg_agg):
                         aggregated_context_data = self.aggregate_context(current_retrieval_results, strategy=context_strategy)
                 else:
                     logger.info(spinner_msg_agg)
                     aggregated_context_data = self.aggregate_context(current_retrieval_results, strategy=context_strategy)
                 final_results["aggregated_context_data"] = aggregated_context_data

                 if not aggregated_context_data or not aggregated_context_data.get("combined_context"):
                      logger.warning(f"Text context aggregation failed or yielded empty context for query: '{query[:100]}...'")
                      final_results["status"] = "Completed: Text context aggregation failed or empty."
                      prev_error = final_results.get("error")
                      final_results["error"] = "Text context aggregation failed or empty."
                      if prev_error: final_results["error"] += f" (Previous error: {prev_error})"
                      # Ensure answer reflects the error
                      final_results["answer"] = f"Could not process the query effectively. Failed to gather relevant text context. {('Previous error: ' + prev_error) if prev_error else ''}"
                      return final_results

                 # === 5b. Query LLM with aggregated text context (using ORIGINAL query) ===
                 spinner_msg_llm = "Generating final answer from text..."
                 answer = ""
                 if is_streamlit:
                      with st.spinner(spinner_msg_llm):
                         answer = self.query_llm(query, aggregated_context_data)
                 else:
                      logger.info(spinner_msg_llm)
                      answer = self.query_llm(query, aggregated_context_data)

                 final_results["answer"] = answer
                 # Only mark as fully successful if there wasn't a prior SQL error
                 if not final_results.get("error"):
                    final_results["status"] = "Completed Successfully"
                 else:
                    # Status reflects the fallback nature
                    final_results["status"] = "Completed via Text Fallback after SQL Issue"
                    logger.info("Text RAG fallback successful, but noting prior SQL path error.")


            # Final logging
            logger.info(f"--- Finished query execution (SQL v2) for: '{query[:100]}...'. Answer Source: {final_results.get('answer_source')} ---")

        except Exception as e:
            logger.error(f"Unexpected error during run_query (SQL v2): {e}", exc_info=True)
            final_results["status"] = "Failed"; final_results["error"] = f"Unexpected runtime error: {str(e)}"
            final_results["answer"] = f"An unexpected error occurred: {e}"
            # Error will be displayed in the UI based on the returned dict

        return final_results
    # ==========================================================================
    # END OF REVISED run_query FUNCTION (v2)
    # ==========================================================================


# --- START OF STREAMLIT UI CODE ---
# (Streamlit UI code remains largely the same, minor adjustments for clarity)
if __name__ == "__main__":
    try: import streamlit as st
    except ImportError: logger.error("Streamlit not installed."); sys.exit(1)

    st.set_page_config(layout="wide", page_title="Multi-Document Q&A with RAG+SQL")
    st.title("📄 Multi-Document Question Answering System (v5-SQL-Improved)") # Version bump
    st.caption("Query PDFs/PPTX (Text RAG) & XLSX/CSV (SQL Query / Text RAG Fallback) using Embeddings, FAISS, SQLite, DeepSeek, Groq LLM")

    @st.cache_resource # Cache the RAG system instance
    def load_rag_system(config_path="config.ini"):
        groq_key = os.getenv("GROQ_API_KEY")
        temp_config = RAGConfig(config_path)
        deepseek_key = temp_config.deepseek_api_key
        if not deepseek_key:
             st.warning("⚠️ DEEPSEEK_API_KEY not found. Spreadsheet metadata generation will be skipped (using basic types).")

        if not groq_key: st.error("🚨 GROQ_API_KEY not set in environment!"); st.stop()

        try:
            system = RAGSystem(config_path=config_path)
            if not system.db_conn:
                 st.error("🚨 Failed to establish SQLite database connection. Check logs.")
                 st.stop()
            return system
        except Exception as e:
            st.error(f"Fatal RAG Init Error: {e}")
            logger.error("RAG Init Error", exc_info=True)
            st.stop()

    def process_documents(_rag_system):
        status_messages = []; status_placeholder = st.empty()
        def streamlit_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
            log_prefix = "✅ " if is_done else "❌ " if is_error else "⚠️ " if is_warning else f"⏳ ({stage}) " if stage else "⏳ "
            progress_val = 1.0 if is_done else 0.0
            if total_steps and current_step is not None:
                 stage_prog = {"Extracting Text": 0.1, "Chunking Text": 0.15, "Verifying Index": 0.2, "Embedding Text": 0.3, "Indexing Text": 0.4,
                               "Loading Sheet": 0.5, "Analyzing Columns": 0.7, "Loading SQL": 0.9}.get(stage, 0.0)
                 progress_val = min(1.0, (current_step + stage_prog) / total_steps) if total_steps > 0 else 0.0
            full_message = f"{log_prefix}{message}"; status_messages.append(full_message)
            try:
                if progress_val > 0 and progress_val < 1.0 and not is_error and not is_warning : status_placeholder.progress(progress_val)
                status_placeholder.caption(full_message)
            except Exception as e: logger.warning(f"Streamlit UI update error in callback: {e}")

        st.header("📚 Document Processing")
        st.caption(f"Data: `{os.path.abspath(_rag_system.config.data_dir)}` | Index: `{os.path.abspath(_rag_system.config.index_dir)}` | DB: `{_rag_system.config.sqlite_db_path}`")
        processing_successful = False
        with st.spinner("Initializing and processing documents (including SQL setup)..."):
            try:
                processing_successful = _rag_system.process_files(progress_callback=streamlit_progress_callback)
                text_index_count = len(_rag_system.faiss_indexes)
                sql_table_count = len(_rag_system.table_metadata)
                processed_file_count = len(_rag_system.processed_files)

                final_message = f"Processed {processed_file_count} file(s). Text Indices: {text_index_count}, SQL Tables: {sql_table_count}."
                if processed_file_count > 0:
                    status_placeholder.success(f"✅ Ready! {final_message}")
                elif processing_successful:
                     status_placeholder.warning(f"⚠️ Processing finished, but no documents loaded. {final_message}")
                else:
                    status_placeholder.error(f"❌ Processing failed. Check logs.")
            except Exception as e:
                 logger.error(f"Fatal error during document processing call: {e}", exc_info=True); status_placeholder.error(f"❌ Fatal error: {e}. Check logs.")

        return len(_rag_system.processed_files) > 0, sorted(list(_rag_system.processed_files))


    try:
        rag_sys = load_rag_system()
        if st.button("🔄 Re-process Documents & Reload System"):
            st.cache_resource.clear()
            st.rerun()

        is_ready, processed_files_list = process_documents(rag_sys)

        if is_ready and processed_files_list:
            st.sidebar.success(f"Processed Assets ({len(processed_files_list)}):")
            with st.sidebar.expander("Show Processed Files & Assets"):
                 for fname in processed_files_list:
                     assets = []
                     if fname in rag_sys.faiss_indexes: assets.append("Text Index")
                     if fname in rag_sys.file_to_table_map:
                         tables = rag_sys.file_to_table_map.get(fname, [])
                         if tables: assets.append(f"SQL Table(s): {len(tables)}")
                     st.caption(f"- {fname} [{', '.join(assets)}]")

            st.sidebar.info(f"LLM: `{rag_sys.config.llm_model}`")
            st.sidebar.info(f"Encoder: `{rag_sys.config.encoder_model}`")
            st.sidebar.info(f"Metadata LLM: `{rag_sys.config.deepseek_model}`")
            st.sidebar.info(f"Retrieval K (Text): `{rag_sys.config.k_retrieval}`")
            st.sidebar.info(f"SQL Query Confidence: `{rag_sys.config.dataframe_query_confidence_threshold}`")


            st.header("💬 Ask a Question")
            user_query = st.text_input("Enter your query (about text content or structured data):", key="query_input")
            if user_query and st.button("Get Answer", key="submit_query"):
                query_start_time = time.time()
                results_data = rag_sys.run_query(user_query, context_strategy="top_k")
                query_end_time = time.time()

                answer_source = results_data.get('answer_source', 'N/A')
                st.subheader(f"💡 Answer (via {answer_source})")
                answer = results_data.get("answer")
                query_status = results_data.get("status", "Unknown")
                query_error = results_data.get("error")

                # Display answer or errors more clearly
                if "Completed" in query_status and answer: # More robust check for success/fallback
                    st.markdown(answer)
                    if query_error and "Fallback" in query_status: # Show warning only if fallback happened due to error
                        st.warning(f"Note: The primary query path encountered an issue, but an answer was generated via fallback. Details: {query_error}")
                    elif query_error and query_status != "Completed Successfully": # Show warning for other non-fatal errors if status indicates issue
                         st.warning(f"Note: An issue occurred during processing: {query_error}")
                else: # No answer or failed status
                     st.error(f"Could not generate an answer. Status: {query_status}")
                     if query_error: st.error(f"Details: {query_error}")


                # Display Query Analysis & Execution Info
                with st.expander("📊 Query Analysis & Execution Path", expanded=False):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.caption(f"Query Classification: `{results_data.get('classification', 'N/A')}`")
                        conf = results_data.get('classification_confidence')
                        st.caption(f"Confidence: `{conf:.2f}`" if conf is not None else "Confidence: N/A")
                        if results_data.get("sub_queries"):
                             st.caption("Text Sub-queries used:")
                             for i, sq in enumerate(results_data["sub_queries"]): st.code(f"{i+1}. {sq}", language=None)
                    with col2:
                         st.caption(f"Final Answer Source: `{results_data.get('answer_source', 'N/A')}`")
                         st.caption(f"Processing Status: `{results_data.get('status', 'N/A')}`")
                         if results_data.get("target_table"):
                             st.caption(f"Target SQL Table: `{results_data.get('target_table')}`")
                         if results_data.get("sql_query_generated"):
                             st.caption("SQL Query Generated:")
                             st.code(results_data["sql_query_generated"], language='sql')


                # Display Supporting Evidence (Context / SQL Info / Text Retrieval)
                with st.expander("🔍 Supporting Evidence / Context", expanded = (answer_source == "SQL Query") ):

                    # Show SQL Info if used
                    if answer_source == "SQL Query":
                         table_name = results_data.get("target_table")
                         metadata = results_data.get("target_table_metadata")
                         sql_result = results_data.get("sql_execution_result")

                         if table_name and metadata:
                              st.markdown(f"**Details for SQL Table Used: `{table_name}`**")
                              sheet_info = f" (Sheet: {metadata.get('source_sheet')})" if metadata.get('source_sheet') else ""
                              st.caption(f"Source File: {metadata.get('source_file', 'N/A')}{sheet_info}")
                              st.caption(f"Total Rows: {metadata.get('row_count', 'N/A')}")

                              st.caption("Columns (with descriptions used for SQL generation):")
                              column_data = metadata.get('columns', [])
                              if column_data and isinstance(column_data, list):
                                  try:
                                      # Display specific columns, handle missing 'name' gracefully
                                      display_data = []
                                      for col in column_data:
                                           display_data.append({
                                               'Column Name': col.get('name', 'N/A'),
                                               'Type': col.get('type', 'N/A'),
                                               'Description': col.get('description', 'N/A')
                                           })
                                      cols_df = pd.DataFrame(display_data)
                                      st.dataframe(cols_df.set_index('Column Name'), height=150)
                                  except Exception as e_df:
                                       logger.error(f"Error creating/displaying metadata DataFrame for {table_name}: {e_df}")
                                       st.warning(f"Could not display column metadata table: {e_df}")
                              else: st.caption("No column metadata available or format error.")


                         if sql_result:
                              st.markdown("**Raw SQL Query Result (Preview)**")
                              st.text_area("SQL Result", sql_result, height=150, key="sql_result_display")

                         if table_name or sql_result: st.divider()

                    # Show Text RAG Context if used
                    agg_context_data = results_data.get("aggregated_context_data", {})
                    combined_context = agg_context_data.get("combined_context", "")
                    context_source_files = agg_context_data.get("source_files", [])
                    if combined_context and answer_source == "Text RAG":
                         st.markdown(f"**Text Context Used from: {', '.join(context_source_files)}**")
                         st.text_area("Combined Text Context Sent to LLM", combined_context, height=200, key="context_display")
                         # Add divider only if SQL section wasn't shown before this
                         if not (answer_source == "SQL Query" and (results_data.get("target_table") or results_data.get("sql_execution_result"))):
                              st.divider()
                    elif answer_source == "SQL Query" and not combined_context: # Only show if SQL was used and no text context generated
                         st.info("Answer generated via SQL query. Raw SQL result (above) was used by the LLM to synthesize the final answer. No separate text context was needed.")


                    # Show Text Retrieval Results (useful regardless of final path)
                    retrieval_results = results_data.get("retrieval_results", {})
                    if retrieval_results and any(retrieval_results.values()): # Check if not empty dictionary and has values
                         st.markdown("**Retrieved Text Chunks (Top candidates from PDF/PPTX)**")
                         st.caption("(Used for Text RAG context or potentially for initial SQL relevance check)")
                         display_chunks = []
                         for file_name, res_list in retrieval_results.items():
                             # Add file_name to each result dict for easier processing
                             for res in res_list: display_chunks.append({**res, 'file_name': file_name})

                         # Sort globally by score (lower is better)
                         display_chunks.sort(key=lambda x: x['score'])

                         # Display top K overall retrieved text chunks
                         for i, res in enumerate(display_chunks[:rag_sys.config.k_retrieval]):
                             source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown'); content_type = res.get('type', 'unknown'); score = res['score']
                             source_display_parts = [f"File: {file_name}"]
                             if file_type == 'pptx': source_display_parts.append(f"Slide: {source_info.get('slide', 'N/A')}" + (f" (merged from {source_info.get('merged_from')})" if source_info.get('merged_from') else ""))
                             elif file_type == 'pdf': source_display_parts.append(f"Page: {source_info.get('page', 'N/A')}")
                             source_display_parts.append(f"Type: {content_type}")

                             st.markdown(f"**Chunk {i+1} (Score: {score:.4f})**")
                             st.caption(" | ".join(source_display_parts))
                             # Use st.text or st.code for potentially long/unformatted text
                             st.text(f"{res.get('content', 'N/A')[:300]}...") # Shorter preview

                    elif answer_source == "SQL Query": # Only show if SQL was used and no text results found/used
                         st.info("No relevant text documents (PDF/PPTX) were found or searched.")
                    # If Text RAG failed because of retrieval, the main error message should cover it.

                st.caption(f"Query processed in {query_end_time - query_start_time:.2f} seconds.")

        elif not processed_files_list and processing_successful: # Processing ran but nothing loaded
            st.warning(f"No documents processed successfully (No Text Index or SQL Table).")
            st.info(f"Add supported files ({', '.join(rag_sys.config.supported_extensions)}) to the '{rag_sys.config.data_dir}' directory and click 'Re-process Documents'.")
        elif not is_ready: # Processing failed to run or returned false
             st.error("System failed to process documents. Check logs.")

    except Exception as e:
        st.error(f"Streamlit App Error: {e}")
        logger.exception("Streamlit application error:")
        st.info("Check console or `rag_system_sql.log` for details.")

# --- END OF STREAMLIT UI CODE ---
