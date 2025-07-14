# --- START OF FILE ragtest_c.py ---

import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd # Added for CSV/XLSX
import faiss
import streamlit as st
import time
from sentence_transformers import SentenceTransformer
from transformers import T5ForConditionalGeneration, T5Tokenizer # Keep T5 if needed, but marked as likely removable later
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage # Keep for potential future use
from dotenv import load_dotenv
from pptx import Presentation # Added for PPTX

# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted from RAG2.py) ---

# Configure logging (optional for Streamlit, but good for debugging)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system.log"),
        # logging.StreamHandler() # Can be noisy in Streamlit console
    ]
)
logger = logging.getLogger(__name__)


class RAGConfig:
    """Configuration class for RAG system."""

    def __init__(self, config_path="config.ini"):
        """Initialize configuration from config file or defaults."""
        self.config = configparser.ConfigParser()

        # Default configuration
        self.defaults = {
            "PATHS": {
                # Renamed pdf_dir to data_dir
                "data_dir": "./Data/",
                "index_dir": "./Faiss_index/",
                "log_file": "rag_system.log"
            },
            "MODELS": {
                "encoder_model": "sentence-transformers/all-MiniLM-L6-v2",
                # Summarizer model is not used in the main query flow
                # "summarizer_model": "t5-small",
                "llm_model": "meta-llama/llama-4-scout-17b-16e-instruct",
                "device": "auto"
            },
            "PARAMETERS": {
                "chunk_size": "200",
                "overlap": "50",
                "k_retrieval": "5",
                "temperature": "0.2",
                # Max context tokens for LLM (approximate)
                "max_context_tokens": "4000",
                # Max characters per Excel cell/PPT shape to avoid huge chunks
                "max_chars_per_element": "1000"
            },
            "SUPPORTED_EXTENSIONS": {
                # Define supported extensions here
                "extensions": ".pdf, .xlsx, .csv, .pptx"
            }
        }

        if os.path.exists(config_path):
            try:
                self.config.read(config_path)
                logger.info(f"Loaded configuration from {config_path}")
                # Ensure all sections/options exist after loading, applying defaults if missing
                self._ensure_defaults()
            except Exception as e:
                logger.error(f"Error reading config file {config_path}: {e}")
                self._set_defaults()
        else:
            logger.warning(f"Config file {config_path} not found, using defaults")
            self._set_defaults()

    def _set_defaults(self):
        """Set default configuration values."""
        for section, options in self.defaults.items():
            if not self.config.has_section(section):
                self.config.add_section(section)
            for option, value in options.items():
                 # Check if the option exists in the specific section before setting
                if not self.config.has_option(section, option):
                    self.config.set(section, option, value)

    def _ensure_defaults(self):
        """Ensure all default sections and options exist in the loaded config."""
        for section, options in self.defaults.items():
            if not self.config.has_section(section):
                self.config.add_section(section)
                logger.info(f"Added missing section [{section}] from defaults.")
            for option, value in options.items():
                if not self.config.has_option(section, option):
                    self.config.set(section, option, value)
                    logger.info(f"Added missing option '{option} = {value}' to section [{section}] from defaults.")


    # --- Properties (getters) for configuration values ---
    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    # @property
    # def summarizer_model(self): return self.config.get("MODELS", "summarizer_model") # Removed as unused
    @property
    def llm_model(self): return self.config.get("MODELS", "llm_model")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto") # Add fallback
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=200)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=50)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=5)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.2)
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=4000)
    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1000)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf")
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])


class RAGSystem:
    """Retrieval-Augmented Generation system for various document types."""

    def __init__(self, config_path="config.ini"):
        """Initialize the RAG system with configuration."""
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system. Device: {self.config.device}, Supported files: {self.config.supported_extensions}")

        # Create necessary directories
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)

        # Initialize local models (encoder)
        try:
            with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."):
                 self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            logger.info("Local embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}")
            st.stop()

        # LLM (Groq) is initialized on-demand in query_llm

        # Initialize data containers
        self.file_chunks = {} # Renamed pdf_chunks to file_chunks
        self.faiss_indexes = {}
        self.processed_files = set()


    def clean_text(self, text):
        """Clean and normalize text."""
        if not isinstance(text, str): text = str(text) # Ensure text is string
        text = re.sub(r"\(cid:.*?\)", "", text) # Remove bad encoding chars
        text = re.sub(r"\s+", " ", text).strip() # Normalize whitespace
        text = text.replace('\n', ' ').replace('\r', '') # Remove newlines embedded within content
        # Add any other cleaning specific to data formats if needed
        return text


    # --- START: Specific Content Extractors ---

    def _extract_pdf(self, file_path):
        """Extract text and tables from PDF files."""
        all_content = []
        base_filename = os.path.basename(file_path)
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Extracting from PDF: {base_filename} ({total_pages} pages)")
                page_bar = None
                if total_pages > 10:
                    page_bar = st.progress(0, text=f"Extracting pages from {base_filename}...")

                for page_num, page in enumerate(pdf.pages):
                    # Extract text
                    text = page.extract_text(layout="normal") or ""
                    cleaned_text = self.clean_text(text)
                    if cleaned_text:
                        all_content.append({
                            "type": "text",
                            "content": cleaned_text,
                            "source_info": {"page": page_num + 1},
                            "file_type": "pdf"
                        })

                    # Extract tables
                    try:
                        for table in page.extract_tables():
                             if table:
                                 table_df = pd.DataFrame(table).fillna('')
                                 # Use markdown for better table representation if needed later by LLM
                                 # table_string = table_df.to_markdown(index=False)
                                 table_string = table_df.to_string(index=False, header=True)
                                 cleaned_table_string = self.clean_text(table_string)
                                 if cleaned_table_string:
                                     all_content.append({
                                         "type": "table",
                                         "content": cleaned_table_string,
                                         "source_info": {"page": page_num + 1},
                                         "file_type": "pdf"
                                     })
                    except Exception as e_table:
                         logger.warning(f"Could not extract table on page {page_num + 1} in {base_filename}: {e_table}")

                    if page_bar:
                         progress_percent = min(1.0, (page_num + 1) / total_pages)
                         page_bar.progress(progress_percent, text=f"Extracting page {page_num + 1}/{total_pages} from {base_filename}...")

                if page_bar: page_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PDF: {base_filename}")
            return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PDF {base_filename}: {e}", exc_info=True)
            st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []


    def _extract_xlsx(self, file_path):
        """Extract content from XLSX files, sheet by sheet."""
        all_content = []
        base_filename = os.path.basename(file_path)
        max_chars = self.config.max_chars_per_element
        try:
            excel_file = pd.ExcelFile(file_path)
            logger.info(f"Extracting from XLSX: {base_filename} (Sheets: {', '.join(excel_file.sheet_names)})")
            sheet_bar = st.progress(0, text=f"Extracting sheets from {base_filename}...")
            total_sheets = len(excel_file.sheet_names)

            for i, sheet_name in enumerate(excel_file.sheet_names):
                try:
                    df = excel_file.parse(sheet_name)
                    df = df.fillna('') # Replace NaN with empty strings

                    # Simple approach: Convert entire sheet to a string representation (e.g., CSV-like)
                    # Could be refined to process cell-by-cell or row-by-row for more granularity if needed
                    # Consider adding headers explicitly.
                    sheet_string = df.to_string(index=False, header=True, na_rep='')

                    # Alternative: Iterate through cells (might be too granular/slow, but captures position)
                    # sheet_text = ""
                    # for r_idx, row in df.iterrows():
                    #     row_texts = []
                    #     for c_idx, cell_value in enumerate(row):
                    #         cell_str = self.clean_text(str(cell_value))
                    #         if cell_str:
                    #              # Truncate long cell content
                    #              if len(cell_str) > max_chars:
                    #                  cell_str = cell_str[:max_chars] + "... (truncated)"
                    #              # Optional: Add cell coordinate? Might be too noisy.
                    #              # sheet_text += f"Row {r_idx+1}, Col {df.columns[c_idx]}: {cell_str}\n"
                    #              row_texts.append(cell_str)
                    #     if row_texts:
                    #         sheet_text += " | ".join(row_texts) + "\n" # Join cells in a row

                    cleaned_content = self.clean_text(sheet_string)
                    if cleaned_content:
                        # Truncate if the whole sheet as string is excessively long (adjust limit as needed)
                        # This is a fallback, ideally chunking handles this better
                        if len(cleaned_content) > max_chars * 10: # Example limit for whole sheet string
                             logger.warning(f"Sheet '{sheet_name}' in {base_filename} is very large, truncating representation.")
                             cleaned_content = cleaned_content[:max_chars*10] + "... (truncated sheet)"

                        all_content.append({
                            "type": "sheet_data",
                            "content": cleaned_content,
                            "source_info": {"sheet": sheet_name, "rows": f"1-{len(df)}"},
                            "file_type": "xlsx"
                        })
                    logger.debug(f"Extracted content from sheet '{sheet_name}' in {base_filename}")

                except Exception as e_sheet:
                    logger.warning(f"Could not process sheet '{sheet_name}' in {base_filename}: {e_sheet}")

                sheet_bar.progress(min(1.0, (i + 1) / total_sheets), text=f"Extracting sheet {i+1}/{total_sheets} ('{sheet_name}') from {base_filename}...")

            sheet_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks (sheets) from XLSX: {base_filename}")
            return all_content
        except Exception as e:
            logger.error(f"Error extracting content from XLSX {base_filename}: {e}", exc_info=True)
            st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []


    def _extract_csv(self, file_path):
        """Extract content from CSV files."""
        all_content = []
        base_filename = os.path.basename(file_path)
        max_chars = self.config.max_chars_per_element # Reuse config for consistency
        try:
            # Use pandas for robust CSV parsing (handles various delimiters, encodings etc.)
            # Try detecting encoding, fall back to utf-8 or latin-1
            encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
            df = None
            for enc in encodings_to_try:
                 try:
                     df = pd.read_csv(file_path, encoding=enc)
                     logger.info(f"Read CSV {base_filename} using encoding: {enc}")
                     break
                 except UnicodeDecodeError:
                     logger.debug(f"Failed to read {base_filename} with encoding {enc}")
                     continue
                 except Exception as e_read: # Catch other pandas read errors
                     logger.warning(f"Error reading CSV {base_filename} with encoding {enc}: {e_read}")
                     # Don't break, try next encoding

            if df is None:
                 logger.error(f"Could not read CSV {base_filename} with any attempted encoding.")
                 st.warning(f"Could not read CSV file {base_filename}. Check encoding or format.")
                 return []


            df = df.fillna('') # Replace NaN

            # Convert entire CSV DataFrame to string
            csv_string = df.to_string(index=False, header=True, na_rep='')
            cleaned_content = self.clean_text(csv_string)

            if cleaned_content:
                # Truncate if necessary (similar to XLSX sheet)
                if len(cleaned_content) > max_chars * 20: # Larger limit for potentially large CSVs
                     logger.warning(f"CSV {base_filename} is very large, truncating representation.")
                     cleaned_content = cleaned_content[:max_chars*20] + "... (truncated CSV)"

                all_content.append({
                    "type": "csv_data",
                    "content": cleaned_content,
                    "source_info": {"rows": f"1-{len(df)}"},
                    "file_type": "csv"
                })
            logger.info(f"Extracted content block from CSV: {base_filename} ({len(df)} rows)")
            return all_content
        except Exception as e:
            logger.error(f"Error extracting content from CSV {base_filename}: {e}", exc_info=True)
            st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def _extract_pptx(self, file_path):
        """Extract text from PPTX files (slides and notes)."""
        all_content = []
        base_filename = os.path.basename(file_path)
        max_chars = self.config.max_chars_per_element
        try:
            prs = Presentation(file_path)
            logger.info(f"Extracting from PPTX: {base_filename} ({len(prs.slides)} slides)")
            slide_bar = st.progress(0, text=f"Extracting slides from {base_filename}...")
            total_slides = len(prs.slides)

            for i, slide in enumerate(prs.slides):
                slide_texts = []
                # Extract text from shapes on the slide
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        cleaned = self.clean_text(text)
                        if cleaned:
                             # Truncate individual text elements if too long
                             if len(cleaned) > max_chars:
                                 cleaned = cleaned[:max_chars] + "...(truncated)"
                             slide_texts.append(cleaned)

                # Extract text from notes slide, if it exists
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    cleaned_notes = self.clean_text(notes_text)
                    if cleaned_notes:
                        if len(cleaned_notes) > max_chars * 2: # Allow notes to be a bit longer
                            cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated notes)"
                        slide_texts.append(f"Notes: {cleaned_notes}")

                # Combine all text from the slide into one content block
                slide_content = "\n".join(slide_texts).strip()

                if slide_content:
                    all_content.append({
                        "type": "slide_text",
                        "content": slide_content,
                        "source_info": {"slide": i + 1},
                        "file_type": "pptx"
                    })

                slide_bar.progress(min(1.0, (i + 1) / total_slides), text=f"Extracting slide {i+1}/{total_slides} from {base_filename}...")

            slide_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks (slides) from PPTX: {base_filename}")
            return all_content
        except Exception as e:
            # Catch specific pptx errors if needed, e.g., file corruption
            logger.error(f"Error extracting content from PPTX {base_filename}: {e}", exc_info=True)
            st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    # --- END: Specific Content Extractors ---


    def extract_content(self, file_path):
        """Extract content based on file extension."""
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf':
            return self._extract_pdf(file_path)
        elif extension == '.xlsx':
            return self._extract_xlsx(file_path)
        elif extension == '.csv':
            return self._extract_csv(file_path)
        elif extension == '.pptx':
            return self._extract_pptx(file_path)
        else:
            logger.warning(f"Unsupported file type skipped: {file_path}")
            return []


    def chunk_content(self, all_content):
        """Split extracted content into chunks with overlap, preserving metadata."""
        chunks = []
        if not all_content:
            return chunks

        # Estimate total words for progress bar (optional but nice)
        try:
            total_words_estimate = sum(len(str(item.get('content', '')).split()) for item in all_content)
        except Exception:
            total_words_estimate = 0 # Handle potential non-string content gracefully

        chunk_bar = None
        if total_words_estimate > 5000:
             chunk_bar = st.progress(0, text=f"Chunking content...")

        words_processed = 0
        for item in all_content:
            content = item.get('content', '')
            source_info = item.get('source_info', {}) # Get original source info
            file_type = item.get('file_type', 'unknown') # Get file type
            content_type = item.get('type', 'unknown') # Get content block type (text, table, slide etc)

            if not isinstance(content, str): content = str(content) # Ensure string
            words = content.split() # Simple whitespace split

            if not words: continue

            item_chunks_created = 0 # Track chunks from this item
            for i in range(0, len(words), self.config.chunk_size - self.config.overlap):
                chunk_words = words[i:i + self.config.chunk_size]
                chunk_text = " ".join(chunk_words)
                if chunk_text:
                    chunks.append({
                        "content": chunk_text,
                        "source_info": source_info, # Preserve metadata
                        "file_type": file_type,     # Preserve metadata
                        "type": content_type        # Preserve metadata
                    })
                    item_chunks_created += 1

            # Update progress bar based on words processed in this item
            if chunk_bar:
                words_processed += len(words)
                progress_percent = min(1.0, words_processed / total_words_estimate) if total_words_estimate > 0 else 0
                chunk_bar.progress(progress_percent, text=f"Chunking content... ({len(chunks)} chunks created)")

        if chunk_bar: chunk_bar.empty()

        logger.info(f"Created {len(chunks)} chunks from {len(all_content)} content blocks.")
        return chunks


    # --- Path Helpers ---
    def _get_safe_filename(self, file_name):
        """Generate a safe filename for index/cache files (removes extension)."""
        base_name = os.path.splitext(file_name)[0]
        return re.sub(r'[^\w\.-]', '_', base_name)

    def get_index_path(self, file_name):
        safe_name = self._get_safe_filename(file_name)
        return os.path.join(self.config.index_dir, f"{safe_name}.index")

    def get_embedding_path(self, file_name):
        safe_name = self._get_safe_filename(file_name)
        return os.path.join(self.config.index_dir, f"{safe_name}.npy")

    def get_chunks_path(self, file_name):
        safe_name = self._get_safe_filename(file_name)
        return os.path.join(self.config.index_dir, f"{safe_name}.json")

    # --- FAISS Index and Chunk Loading/Saving (Mostly Unchanged) ---
    def load_faiss_index(self, file_name, embedding_dim):
        """Load FAISS index from disk or create new one."""
        index_path = self.get_index_path(file_name)
        if os.path.exists(index_path):
            try:
                logger.debug(f"Attempting to load FAISS index from {index_path}")
                index = faiss.read_index(index_path)
                if index.d != embedding_dim:
                    logger.warning(f"Index dimension mismatch for {file_name} (Expected {embedding_dim}, Got {index.d}). Rebuilding.")
                    return faiss.IndexFlatL2(embedding_dim)
                logger.info(f"Loaded FAISS index for {file_name} ({index.ntotal} vectors)")
                return index
            except Exception as e:
                logger.error(f"Error reading FAISS index {index_path}: {e}. Creating new index.")
        logger.info(f"FAISS index not found at {index_path}. Will create a new one.")
        return faiss.IndexFlatL2(embedding_dim)

    def save_chunks(self, file_name, chunks):
        """Save chunks to disk."""
        chunks_path = self.get_chunks_path(file_name)
        try:
            with open(chunks_path, 'w', encoding='utf-8') as f:
                json.dump(chunks, f, indent=2)
            logger.info(f"Saved {len(chunks)} chunks to {chunks_path}")
        except Exception as e:
            logger.error(f"Error saving chunks for {file_name} to {chunks_path}: {e}")
            st.warning(f"Could not save chunks for {os.path.basename(file_name)} to disk.")

    def load_chunks(self, file_name):
        """Load chunks from disk."""
        chunks_path = self.get_chunks_path(file_name)
        if os.path.exists(chunks_path):
            try:
                with open(chunks_path, 'r', encoding='utf-8') as f:
                    chunks = json.load(f)
                logger.info(f"Loaded {len(chunks)} chunks from {chunks_path}")
                return chunks
            except json.JSONDecodeError as e:
                logger.error(f"Error decoding JSON from {chunks_path}: {e}. File might be corrupted. Will re-process.")
                return None
            except Exception as e:
                logger.error(f"Error loading chunks from {chunks_path}: {e}")
                return None
        logger.info(f"Chunks file not found at {chunks_path}")
        return None


    def process_files(self, progress_callback=None): # Renamed from process_pdfs
        """Process all supported files in directory, generating or loading embeddings and indexes."""
        self.file_chunks = {} # Reset
        self.faiss_indexes = {}
        self.processed_files = set()

        data_dir = self.config.data_dir
        supported_ext = self.config.supported_extensions

        try:
             all_files = os.listdir(data_dir)
             # Filter by supported extensions
             process_list = [f for f in all_files if f.lower().endswith(supported_ext)]
             process_list.sort()
        except FileNotFoundError:
             logger.error(f"Data directory not found: {data_dir}")
             if progress_callback: progress_callback(f"❌ Error: Data directory '{data_dir}' not found. Check config.ini.", is_error=True)
             return False
        except Exception as e:
             logger.error(f"Error listing files in data directory {data_dir}: {e}")
             if progress_callback: progress_callback(f"❌ Error accessing data directory: {e}", is_error=True)
             return False

        if not process_list:
            logger.warning(f"No supported files ({', '.join(supported_ext)}) found in {data_dir}")
            if progress_callback: progress_callback(f"⚠️ No supported files found in '{data_dir}'. Add files ({', '.join(supported_ext)}) to process.", is_warning=True)
            return True # Not an error, just nothing to do

        logger.info(f"Processing {len(process_list)} supported files from {data_dir}")
        embedding_dim = self.encoder_model.get_sentence_embedding_dimension()

        total_files = len(process_list)
        if progress_callback: progress_callback(f"Found {total_files} supported file(s). Starting processing...", current_step=0, total_steps=total_files)

        # --- Main Processing Loop ---
        for idx, file_name in enumerate(process_list):
            current_file_msg = f"Processing ({idx+1}/{total_files}): {file_name}"
            if progress_callback: progress_callback(current_file_msg, current_step=idx, total_steps=total_files)

            logger.info(f"--- {current_file_msg} ---")
            file_path = os.path.join(data_dir, file_name)
            index_path = self.get_index_path(file_name)
            emb_path = self.get_embedding_path(file_name)
            chunks_path = self.get_chunks_path(file_name)

            try:
                # 1. Load or Extract/Save Chunks
                chunks = self.load_chunks(file_name)
                if chunks is None:
                    if progress_callback: progress_callback(f"{current_file_msg} - Extracting content...", current_step=idx, total_steps=total_files, stage="Extracting")
                    logger.info(f"Extracting and chunking content for {file_name}...")
                    # Use the generalized extractor
                    all_content = self.extract_content(file_path)
                    if not all_content:
                         logger.warning(f"No content extracted from {file_name}. Skipping.")
                         if progress_callback: progress_callback(f"⚠️ No content extracted from {file_name}. Skipping.", is_warning=True, current_step=idx, total_steps=total_files)
                         continue

                    if progress_callback: progress_callback(f"{current_file_msg} - Chunking content...", current_step=idx, total_steps=total_files, stage="Chunking")
                    chunks = self.chunk_content(all_content)
                    if not chunks:
                         logger.warning(f"No chunks generated for {file_name} after extraction. Skipping.")
                         if progress_callback: progress_callback(f"⚠️ No chunks generated for {file_name}. Skipping.", is_warning=True, current_step=idx, total_steps=total_files)
                         continue
                    self.save_chunks(file_name, chunks) # Save newly created chunks

                # Store chunks in memory
                self.file_chunks[file_name] = chunks
                logger.debug(f"Stored {len(chunks)} chunks in memory for {file_name}")


                # 2. Load or Generate/Save Embeddings and Index (Logic remains largely the same)
                faiss_index = None
                regenerate_embeddings = False

                if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                    logger.info(f"Found existing index, embeddings, and chunks for {file_name}. Verifying...")
                    if progress_callback: progress_callback(f"{current_file_msg} - Verifying existing index...", current_step=idx, total_steps=total_files, stage="Verifying")
                    try:
                        embeddings = np.load(emb_path)
                        if embeddings.ndim != 2 or embeddings.shape[1] != embedding_dim:
                             logger.warning(f"Embeddings dimension mismatch for {file_name}. Regenerating.")
                             regenerate_embeddings = True
                        elif embeddings.shape[0] != len(chunks):
                             logger.warning(f"Chunk/Embedding count mismatch for {file_name}. Regenerating.")
                             regenerate_embeddings = True
                        else:
                             faiss_index = self.load_faiss_index(file_name, embedding_dim)
                             if faiss_index.ntotal != embeddings.shape[0]:
                                 logger.warning(f"FAISS index count doesn't match embeddings for {file_name}. Rebuilding index.")
                                 faiss_index = faiss.IndexFlatL2(embedding_dim)
                                 faiss_index.add(embeddings.astype('float32'))
                                 faiss.write_index(faiss_index, index_path)
                             else:
                                 logger.info(f"Successfully loaded and verified existing embeddings and index for {file_name}")

                    except Exception as e:
                        logger.error(f"Error loading/verifying pre-existing data for {file_name}: {e}. Regenerating...", exc_info=True)
                        regenerate_embeddings = True
                        faiss_index = None

                else:
                     logger.info(f"Required files missing for {file_name}. Generating embeddings and index.")
                     regenerate_embeddings = True


                if regenerate_embeddings or faiss_index is None:
                    logger.info(f"Generating embeddings and creating/updating index for {file_name}...")
                    if progress_callback: progress_callback(f"{current_file_msg} - Generating embeddings...", current_step=idx, total_steps=total_files, stage="Embedding")

                    content_list = [chunk['content'] for chunk in chunks]
                    if not content_list:
                         logger.warning(f"No content to embed for {file_name}. Skipping.")
                         continue

                    embeddings = self.encoder_model.encode(
                        content_list, batch_size=64, show_progress_bar=False, convert_to_numpy=True
                    ).astype('float32')

                    if embeddings.shape[0] == 0:
                        logger.warning(f"Embedding yielded no vectors for {file_name}. Skipping index.")
                        continue

                    np.save(emb_path, embeddings)
                    logger.info(f"Saved {embeddings.shape[0]} embeddings to {emb_path}")

                    if progress_callback: progress_callback(f"{current_file_msg} - Creating FAISS index...", current_step=idx, total_steps=total_files, stage="Indexing")
                    faiss_index = faiss.IndexFlatL2(embedding_dim)
                    faiss_index.add(embeddings)
                    faiss.write_index(faiss_index, index_path)
                    logger.info(f"Created and saved FAISS index with {faiss_index.ntotal} vectors to {index_path}")

                # Store final index in memory
                if faiss_index is not None and faiss_index.ntotal > 0:
                    self.faiss_indexes[file_name] = faiss_index
                    self.processed_files.add(file_name)
                    logger.debug(f"Stored FAISS index in memory for {file_name}")
                else:
                    logger.warning(f"No valid FAISS index generated or loaded for {file_name}.")
                    if progress_callback: progress_callback(f"⚠️ Could not create/load index for {file_name}.", is_warning=True, current_step=idx, total_steps=total_files)


            except Exception as e:
                logger.error(f"Failed to process {file_name}: {e}", exc_info=True)
                if progress_callback: progress_callback(f"❌ Error processing {file_name}: {e}", is_error=True, current_step=idx, total_steps=total_files)
                # Cleanup attempt
                try:
                    if os.path.exists(index_path): os.remove(index_path)
                    if os.path.exists(emb_path): os.remove(emb_path)
                    # if os.path.exists(chunks_path): os.remove(chunks_path) # Keep chunks?
                    logger.info(f"Cleaned up index/embedding files for failed {file_name}")
                except OSError as e_clean:
                    logger.error(f"Error cleaning up files for {file_name}: {e_clean}")
                if file_name in self.file_chunks: del self.file_chunks[file_name]
                if file_name in self.faiss_indexes: del self.faiss_indexes[file_name]


        # --- Post-processing Summary ---
        final_indexed_count = len(self.faiss_indexes)
        if final_indexed_count > 0:
             logger.info(f"--- File Processing Complete. Indexed {final_indexed_count}/{total_files} documents. ---")
             if progress_callback: progress_callback(f"✅ Processing Complete. Ready to query {final_indexed_count} document(s).", current_step=total_files, total_steps=total_files, is_done=True)
             return True
        else:
             logger.warning(f"--- File Processing Complete. No documents were successfully indexed. ---")
             if progress_callback: progress_callback(f"⚠️ Processing Finished. No documents could be indexed. Check logs and files.", is_warning=True, is_done=True)
             return False


    def query_files(self, query): # Renamed from query_pdfs
        """Query all processed files with a user query."""
        if not self.faiss_indexes:
            logger.warning("No indexes available. Run process_files() first.")
            st.warning("No documents seem to be indexed. Please process files first or check logs.")
            return {}

        all_results = {}
        try:
            logger.info(f"Encoding query: '{query[:100]}...'")
            query_embedding = self.encoder_model.encode(query, convert_to_numpy=True).astype("float32")
            query_embedding = np.array([query_embedding]) # Reshape for FAISS
            if query_embedding.ndim != 2:
                 logger.error(f"Query embedding error shape: {query_embedding.shape}")
                 st.error("Internal error: Could not prepare query embedding.")
                 return {}

            logger.info(f"Searching across {len(self.faiss_indexes)} indexed documents.")
            search_bar = None
            total_indexes = len(self.faiss_indexes)
            if total_indexes > 5:
                 search_bar = st.progress(0, text=f"Searching documents...")

            processed_index_count = 0
            for file_name, index in self.faiss_indexes.items():
                file_results = []
                if index is None or index.ntotal == 0:
                    logger.warning(f"Index for {file_name} is empty/invalid. Skipping.")
                    continue

                processed_index_count += 1
                if search_bar:
                     progress_percent = min(1.0, processed_index_count / total_indexes)
                     search_bar.progress(progress_percent, text=f"Searching {file_name} ({processed_index_count}/{total_indexes})...")

                try:
                    k_search = min(self.config.k_retrieval, index.ntotal)
                    logger.debug(f"Searching index for {file_name} (k={k_search}, size={index.ntotal})")
                    D, I = index.search(query_embedding, k=k_search)
                    indices = I[0]; distances = D[0]
                    logger.debug(f"Retrieved {len(indices)} results for {file_name}.")

                    current_file_chunks = self.file_chunks.get(file_name) # Use file_chunks
                    if not current_file_chunks:
                         logger.error(f"FATAL: Chunks for {file_name} not found in memory! Skipping.")
                         st.error(f"Internal inconsistency: Chunk data missing for {file_name}. Skipping.")
                         continue

                    processed_indices = set()
                    for i, idx in enumerate(indices):
                        if idx == -1 or idx in processed_indices: continue
                        if not (0 <= idx < len(current_file_chunks)):
                             logger.warning(f"Retrieved invalid index {idx} for {file_name}. Skipping.")
                             continue

                        processed_indices.add(idx)
                        chunk = current_file_chunks[idx]
                        distance = distances[i]
                        score = float(distance)

                        # Include the detailed source info and file type
                        file_results.append({
                            "source_info": chunk.get('source_info', {}),
                            "file_type": chunk.get('file_type', 'unknown'),
                            "content": chunk.get('content', ''),
                            "score": round(score, 4),
                            "type": chunk.get('type', 'unknown') # e.g., text, table, sheet_data, slide_text
                        })

                    file_results.sort(key=lambda x: x['score'])
                    if file_results:
                        all_results[file_name] = file_results
                        logger.debug(f"Found {len(file_results)} relevant chunks in {file_name}")
                    else:
                         logger.debug(f"No relevant chunks found in {file_name}")

                except Exception as e_search:
                    logger.error(f"Error searching index for {file_name}: {e_search}", exc_info=True)
                    st.warning(f"Could not search {os.path.basename(file_name)}. Error: {e_search}")

            if search_bar: search_bar.empty()
            logger.info(f"Search complete. Found relevant results in {len(all_results)} documents.")
            return all_results

        except Exception as e_query:
            logger.error(f"Error during query encoding/search: {e_query}", exc_info=True)
            st.error(f"An error occurred during the search: {e_query}")
            return {}


    def aggregate_context(self, query_results, strategy="top_k"):
        """Aggregate context from query results, respecting token limits and formatting source info."""
        all_context = {}
        # Use configured max tokens
        max_context_tokens = self.config.max_context_tokens
        # Estimate max chars: 1 token ~ 3-4 chars. Use 3 for safety margin.
        max_chars = max_context_tokens * 3
        logger.info(f"Aggregating context using strategy '{strategy}', max_chars ~{max_chars} (for ~{max_context_tokens} tokens)")

        if not query_results:
            logger.warning("No query results provided to aggregate context.")
            return all_context

        total_aggregated_chars = 0

        for file_name, results in query_results.items():
            if not results: continue

            context = ""
            current_chars = 0
            num_chunks_added = 0
            limit = self.config.k_retrieval if strategy == "top_k" else len(results)

            for i, res in enumerate(results[:limit]):
                source_info = res.get('source_info', {})
                file_type = res.get('file_type', 'unknown')
                content_type = res.get('type', 'unknown') # e.g. text, table, sheet_data
                score = res['score']
                content_body = res['content']

                # --- Format Source Information Dynamically ---
                source_parts = [f"Source: {file_name}"]
                if file_type == 'pdf':
                    source_parts.append(f"Page: {source_info.get('page', 'N/A')}")
                elif file_type == 'xlsx':
                    source_parts.append(f"Sheet: {source_info.get('sheet', 'N/A')}")
                    # Optional: Add row info if available and needed
                    # if 'rows' in source_info: source_parts.append(f"Rows: {source_info.get('rows')}")
                elif file_type == 'csv':
                     # CSV might just use filename, or row info if chunked by row
                     if 'rows' in source_info: source_parts.append(f"Rows: {source_info.get('rows')}")
                     else: source_parts.append("(CSV)")
                elif file_type == 'pptx':
                    source_parts.append(f"Slide: {source_info.get('slide', 'N/A')}")

                source_parts.append(f"Type: {content_type}")
                source_parts.append(f"Score: {score:.4f}")
                source_str = ", ".join(source_parts)
                # --- End Formatting ---

                content_header = f"[{source_str}]\n"
                content_to_add = content_header + content_body + "\n\n"
                content_chars = len(content_to_add)

                if current_chars + content_chars <= max_chars:
                    context += content_to_add
                    current_chars += content_chars
                    num_chunks_added += 1
                else:
                    if i == 0: # First chunk is too large
                         remaining_chars = max_chars - current_chars - len(content_header) - len("\n\n[...TRUNCATED...]\n\n")
                         if remaining_chars > 50:
                             truncated_body = content_body[:remaining_chars]
                             context += content_header + truncated_body + "\n[...TRUNCATED...]\n\n"
                             current_chars += len(content_header) + len(truncated_body) + len("\n[...TRUNCATED...]\n\n")
                             num_chunks_added += 1
                             logger.warning(f"Context for {file_name} truncated (first chunk too large).")
                         else:
                              logger.warning(f"First chunk for {file_name} too large to fit, skipping context.")
                    logger.info(f"Stopping context aggregation for {file_name} at {num_chunks_added} chunks ({current_chars}/{max_chars} chars).")
                    break

            final_context = context.strip()
            if final_context:
                 all_context[file_name] = final_context
                 total_aggregated_chars += len(final_context)
                 logger.debug(f"Aggregated {len(final_context)} chars of context for {file_name} from {num_chunks_added} chunks.")
            else:
                 logger.warning(f"No context aggregated for {file_name} within limits.")

        logger.info(f"Total aggregated context characters: {total_aggregated_chars}")
        return all_context


    def query_llm(self, query, context, source_file_name, retry_count=1): # Renamed pdf_file_source
        """Query the Groq API with context and handle errors."""
        if not context:
            logger.warning(f"No context for LLM query regarding {source_file_name}")
            return f"Could not generate answer for {source_file_name}: No relevant context found."

        api_key = os.getenv("GROQ_API_KEY")
        if not api_key:
            logger.error("GROQ_API_KEY not found.")
            st.error("Error: GROQ_API_KEY is not set.")
            return "Error: Groq API key not configured."

        try:
            llm = ChatGroq(
                temperature=self.config.temperature,
                model_name=self.config.llm_model,
                groq_api_key=api_key,
            )

            # --- Generalized System Prompt ---
            system_prompt = f"""
            You are an AI assistant specialized in answering questions based *only* on the provided context from various source documents. The current context is derived from the document '{source_file_name}'.
Your task is to synthesize information found *solely* within the given text excerpts.
- Answer the user's question accurately using *only* the information present in the context below.
- If the context does not contain the information needed to answer the question, you MUST state: "Based on the provided context from '{source_file_name}', I cannot answer this question."
- Do NOT use any external knowledge, make assumptions, or infer information not explicitly stated.
- Be concise and directly answer the question.
- Reference the source document name ('{source_file_name}') if you cannot answer.
"""
            human_prompt = f"""Context from document '{source_file_name}':
--- START CONTEXT ---
{context}
--- END CONTEXT ---

User Question: {query}

Answer based *only* on the provided context:"""

            full_prompt_messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": human_prompt}
            ]

            logger.info(f"Querying Groq model {self.config.llm_model} for {source_file_name}...")
            answer = f"Error: LLM query failed for {source_file_name} after retries."
            for attempt in range(retry_count + 1):
                try:
                    # Use standard invoke which expects list of dicts or Message objects
                    response = llm.invoke(full_prompt_messages)

                    if hasattr(response, 'content'):
                         answer = response.content.strip()
                    elif isinstance(response, str): # Fallback if it returns a string directly
                         answer = response.strip()
                    else:
                         logger.error(f"Unexpected response type from Groq: {type(response)}")
                         answer = "Error: Received unexpected response format from LLM."
                         break # Don't retry on format error

                    logger.info(f"Groq response received for {source_file_name} (attempt {attempt+1}). Length: {len(answer)}")
                    if not answer:
                         logger.warning(f"Received empty response from Groq for {source_file_name} on attempt {attempt+1}")
                         answer = f"Received an empty response from the language model for {source_file_name}."
                         # Continue retrying if empty, maybe temporary issue
                         if attempt < retry_count: time.sleep(1); continue
                         else: break # Stop if still empty after retries

                    # If valid answer (even if it's "cannot answer"), return it.
                    return answer

                except Exception as e_api:
                    logger.warning(f"Groq API query attempt {attempt+1} for {source_file_name} failed: {e_api}")
                    if attempt < retry_count:
                        wait_time = 1.5 ** attempt
                        logger.info(f"Retrying Groq query for {source_file_name} in {wait_time:.1f}s...")
                        time.sleep(wait_time)
                    else:
                        logger.error(f"All Groq API query attempts failed for {source_file_name}: {e_api}", exc_info=True)
                        answer = f"Error: Failed to get answer from language model for {source_file_name} after {retry_count+1} attempts. (API Error: {e_api})"
                        # Return the error instead of retrying further
                        return answer

            # Should only be reached if all retries failed (e.g., all empty responses)
            return answer

        except Exception as e_setup:
            logger.error(f"Error setting up/calling Groq API for {source_file_name}: {e_setup}", exc_info=True)
            st.error(f"Failed to query the language model for {source_file_name}. Error: {e_setup}")
            return f"Error: Could not query language model for {source_file_name} due to setup error: {e_setup}"


    def run_query(self, query, context_strategy="top_k"):
        """Complete query pipeline: retrieve, aggregate context, query LLM."""
        logger.info(f"--- Starting query execution: '{query[:100]}...' ---")
        final_results = {
            "query": query,
            "retrieval_results": {},
            "aggregated_context": {},
            "answers": {},
            "status": "Started",
            "error": None
        }

        try:
            # 1. Retrieve relevant chunks (uses query_files)
            with st.spinner("Searching relevant documents..."):
                 query_results = self.query_files(query) # Use renamed method
                 final_results["retrieval_results"] = query_results

            if not query_results:
                logger.warning("No relevant chunks found across any documents.")
                final_results["status"] = "Completed: No relevant information found."
                st.info("Could not find relevant information in the indexed documents.")
                return final_results

            # 2. Aggregate context
            with st.spinner("Gathering context..."):
                contexts = self.aggregate_context(query_results, strategy=context_strategy)
                final_results["aggregated_context"] = contexts

            if not contexts:
                 logger.warning(f"Context aggregation failed for query: '{query[:100]}...'")
                 final_results["status"] = "Completed: Found relevant parts, but could not prepare context."
                 final_results["error"] = "Context aggregation failed."
                 st.warning("Found relevant document parts, but couldn't prepare context for answering.")
                 return final_results

            # 3. Query LLM for each context
            answers = {}
            num_contexts = len(contexts)
            llm_bar = st.progress(0, text=f"Generating answers (0/{num_contexts})...")
            processed_count = 0

            for file_name, context in contexts.items():
                processed_count += 1
                llm_bar.progress(min(1.0, processed_count / num_contexts), text=f"Generating answer based on {file_name} ({processed_count}/{num_contexts})...")
                logger.info(f"Generating answer for {file_name}...")
                # Pass the filename as source_file_name
                answer = self.query_llm(query, context, source_file_name=file_name)
                answers[file_name] = answer
                logger.info(f"Answer received for {file_name}.")

            final_results["answers"] = answers
            llm_bar.empty()

            final_results["status"] = "Completed Successfully."
            logger.info(f"--- Finished query execution for: '{query[:100]}...' ---")

        except Exception as e:
            logger.error(f"Unexpected error during run_query: {e}", exc_info=True)
            final_results["status"] = "Failed"
            final_results["error"] = str(e)
            st.error(f"An unexpected error occurred: {e}")

        return final_results

# --- END OF RAG SYSTEM CODE ---


# --- START OF STREAMLIT UI CODE ---

st.set_page_config(layout="wide", page_title="Multi-Document Q&A with RAG")

st.title("📄 Multi-Document Question Answering System")
st.caption("Query PDFs, XLSX, CSV, and PPTX files using Local Embeddings, FAISS, and Groq LLM")

# --- Initialization and Caching ---
@st.cache_resource
def load_rag_system(config_path="config.ini"):
    if not os.getenv("GROQ_API_KEY"):
        st.error("🚨 GROQ_API_KEY environment variable not set! Please set it in your .env file or environment.")
        st.stop()
    try:
        rag_system = RAGSystem(config_path=config_path)
        return rag_system
    except Exception as e:
        st.error(f"Fatal Error during RAG System Initialization: {e}")
        logger.error(f"Fatal Error during RAG System Initialization: {e}", exc_info=True)
        st.stop()


@st.cache_data # Cache the result of document processing status
def process_documents(_rag_system):
    """Processes documents and returns status. Caching prevents reprocessing unless files change."""
    status_messages = []
    processing_successful = False
    supported_ext_str = ", ".join(_rag_system.config.supported_extensions)

    # Define callback here or make it accessible
    status_placeholder = st.empty() # Placeholder needs to be defined before use in callback

    def streamlit_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
        log_prefix = ""
        if is_error: log_prefix = "❌ Error: "
        elif is_warning: log_prefix = "⚠️ Warning: "
        elif is_done: log_prefix = "✅ "

        full_message = f"{log_prefix}{message}"
        status_messages.append(full_message) # Keep history

        if total_steps and current_step is not None:
             progress_percent = min(1.0, (current_step + (0.5 if stage else 1) ) / total_steps) if total_steps > 0 else 0
             # Ensure placeholder exists before updating
             if status_placeholder:
                 status_placeholder.progress(progress_percent)
                 status_placeholder.caption(full_message)
        else:
             if status_placeholder:
                  status_placeholder.caption(full_message)

        if is_error: logger.error(message)
        elif is_warning: logger.warning(message)
        else: logger.info(message)


    st.header("📚 Document Processing")
    # status_placeholder defined above

    with st.spinner(f"Initializing and processing documents ({supported_ext_str})... This may take time."):
        try:
            # Use process_files (renamed)
            processing_successful = _rag_system.process_files(progress_callback=streamlit_progress_callback)

            # Final status update
            final_message = status_messages[-1] if status_messages else "Processing state unknown."
            indexed_count = len(_rag_system.faiss_indexes)
            if processing_successful and indexed_count > 0:
                status_placeholder.success(f"✅ Ready! Processed and indexed {indexed_count} document(s).")
            elif _rag_system.processed_files and indexed_count == 0 : # Processed some files but none indexed
                 status_placeholder.warning(f"⚠️ Processing finished, but no documents were successfully indexed. Check logs.")
            elif not _rag_system.processed_files: # No supported files found or processed at all
                 status_placeholder.warning(f"⚠️ No supported documents ({supported_ext_str}) found or processed in '{_rag_system.config.data_dir}'.")
            # Handle case where processing_successful is False (major error during processing)
            elif not processing_successful:
                 status_placeholder.error(f"❌ Processing failed. Check logs for errors.")


        except Exception as e:
            logger.error(f"Fatal error during document processing call: {e}", exc_info=True)
            status_placeholder.error(f"❌ A fatal error occurred during processing: {e}. Check logs.")
            processing_successful = False

    return processing_successful, list(_rag_system.faiss_indexes.keys())


# --- Main App Logic ---
try:
    rag_sys = load_rag_system()

    is_ready, indexed_files = process_documents(rag_sys) # This uses the cached function

    if is_ready and indexed_files:
        st.sidebar.success(f"Indexed Documents ({len(indexed_files)}):")
        with st.sidebar.expander("Show Indexed Files"):
             for fname in indexed_files:
                 st.caption(f"- {fname}")
        st.sidebar.info(f"LLM: `{rag_sys.config.llm_model}`")
        st.sidebar.info(f"Retrieval K: `{rag_sys.config.k_retrieval}`")
        st.sidebar.info(f"Max Context Tokens: `{rag_sys.config.max_context_tokens}`")

        st.header("💬 Ask a Question")
        user_query = st.text_input("Enter your query about the indexed documents:", key="query_input")

        if user_query:
            if st.button("Get Answer", key="submit_query"):
                with st.spinner("Thinking... (Retrieving & Querying LLM)"):
                    start_time = time.time()
                    results_data = rag_sys.run_query(user_query, context_strategy="top_k")
                    end_time = time.time()

                st.subheader("💡 Answer(s)")
                if results_data and results_data.get("answers"):
                    for file_name, answer in results_data["answers"].items():
                        with st.expander(f"Answer based on: **{file_name}**", expanded=True):
                            st.markdown(answer) # Display LLM answer
                else:
                    st.warning("Could not generate an answer. No relevant information found or an error occurred.")
                    if results_data.get("error"):
                        st.error(f"Error details: {results_data['error']}")


                # Display Retrieval Details with enhanced source info
                st.subheader("🔍 Retrieval Details (Supporting Evidence)")
                retrieval_results = results_data.get("retrieval_results", {})
                if retrieval_results:
                     for file_name, file_res_list in retrieval_results.items():
                         if file_res_list:
                             with st.expander(f"Top {len(file_res_list)} relevant chunks from: **{file_name}**"):
                                 for i, res in enumerate(file_res_list):
                                     source_info = res.get('source_info', {})
                                     file_type = res.get('file_type', 'unknown')
                                     content_type = res.get('type', 'unknown')
                                     score = res['score']

                                     # Format source string for display
                                     source_display_parts = []
                                     if file_type == 'pdf': source_display_parts.append(f"Page: {source_info.get('page', 'N/A')}")
                                     elif file_type == 'xlsx': source_display_parts.append(f"Sheet: {source_info.get('sheet', 'N/A')}")
                                     elif file_type == 'pptx': source_display_parts.append(f"Slide: {source_info.get('slide', 'N/A')}")
                                     elif file_type == 'csv': source_display_parts.append(f"CSV (Rows: {source_info.get('rows', 'all')})")
                                     source_display_parts.append(f"Type: {content_type}")

                                     st.markdown(f"**Chunk {i+1} (Score: {score:.4f})**")
                                     st.caption(" | ".join(source_display_parts))
                                     st.text(f"{res['content'][:350]}...") # Show preview
                                     st.divider()
                else:
                     st.info("No specific document chunks were retrieved for this query.")

                st.caption(f"Query processed in {end_time - start_time:.2f} seconds.")


    elif not indexed_files:
        supported_ext_str = ", ".join(rag_sys.config.supported_extensions)
        st.warning(f"No documents are currently indexed. Please add files ({supported_ext_str}) to the data directory and refresh.")
        st.info(f"Looking for files in: `{os.path.abspath(rag_sys.config.data_dir)}`")
        st.info(f"Storing index files in: `{os.path.abspath(rag_sys.config.index_dir)}`")
    else: # Not ready, processing likely failed
         st.error("The RAG system could not be initialized or documents could not be processed successfully. Please check the logs (`rag_system.log`) for details.")


except Exception as e:
    st.error(f"An unexpected error occurred in the Streamlit application: {e}")
    logger.error(f"Streamlit application error: {e}", exc_info=True)
    st.info("Please check the console or `rag_system.log` for more details.")


# --- END OF STREAMLIT UI CODE ---

# --- END OF FILE ragtest_c.py ---