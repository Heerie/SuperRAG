# --- START OF FILE ragtest_c2-4.py --- # MODIFIED (v5.1 - exec fix)

# --- START OF FILE ragtest_c4_dataframe.py ---

# --- START OF FILE ragtest_c3_query_adapt.py --- # Renamed for clarity
# --- START OF FILE ragtest_c.py --- # <-- Original comment

import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd # Added for CSV/XLSX
import faiss
# import streamlit as st # Import only if needed / within UI code
import time
from sentence_transformers import SentenceTransformer
# from transformers import T5ForConditionalGeneration, T5Tokenizer # Keep if needed elsewhere
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage # Added SystemMessage
from dotenv import load_dotenv
from pptx import Presentation # Added for PPTX
from pptx.enum.shapes import MSO_SHAPE_TYPE # Needed for placeholder check
import sys # Import sys for checking streamlit context
import io # Added for capturing print output from exec
from contextlib import redirect_stdout # Added for capturing print output

# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted from RAG2.py) ---

# Configure logging (optional for Streamlit, but good for debugging)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system.log"),
        # logging.StreamHandler() # Can be noisy in Streamlit console
    ]
)
logger = logging.getLogger(__name__)


class RAGConfig:
    """Configuration class for RAG system."""
    # (Unchanged from previous version)
    def __init__(self, config_path="config.ini"):
        """Initialize configuration from config file or defaults."""
        self.config = configparser.ConfigParser()
        self.defaults = {
            "PATHS": {"data_dir": "./Data/", "index_dir": "./Faiss_index/", "log_file": "rag_system.log"},
            "MODELS": {"encoder_model": "sentence-transformers/all-MiniLM-L6-v2", "llm_model": "llama3-70b-8192", "device": "auto"}, # Updated default LLM maybe
            "PARAMETERS": {"chunk_size": "200", "overlap": "50", "k_retrieval": "5", "temperature": "0.1", "max_context_tokens": "4000", "max_chars_per_element": "1000", "pptx_merge_threshold_words": "50", "dataframe_query_confidence_threshold": "0.8"}, # Added threshold
            "SUPPORTED_EXTENSIONS": {"extensions": ".pdf, .xlsx, .csv, .pptx"}
        }
        if os.path.exists(config_path):
            try: self.config.read(config_path); logger.info(f"Loaded configuration from {config_path}"); self._ensure_defaults()
            except Exception as e: logger.error(f"Error reading config file {config_path}: {e}"); self._set_defaults()
        else: logger.warning(f"Config file {config_path} not found, using defaults"); self._set_defaults()
    def _set_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section)
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value)
    def _ensure_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section); logger.info(f"Added missing section [{section}] from defaults.")
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value); logger.info(f"Added missing option '{option} = {value}' to section [{section}] from defaults.")
    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    @property
    def llm_model(self): return self.config.get("MODELS", "llm_model")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto")
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=200)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=50)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=5)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.1) # Lowered temp slightly for code gen?
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=4000)
    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1000)
    @property
    def pptx_merge_threshold_words(self): return self.config.getint("PARAMETERS", "pptx_merge_threshold_words", fallback=50)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf")
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])
    @property
    def dataframe_query_confidence_threshold(self): # New config for classification confidence
        return self.config.getfloat("PARAMETERS", "dataframe_query_confidence_threshold", fallback=0.8)


class RAGSystem:
    """Retrieval-Augmented Generation system supporting text and structured data querying."""
    # (Init slightly modified)
    def __init__(self, config_path="config.ini"):
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system. Device: {self.config.device}, Supported files: {self.config.supported_extensions}")
        self._is_streamlit = "streamlit" in sys.modules
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)
        try:
            if self._is_streamlit:
                import streamlit as st
                with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."): self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            else: logger.info(f"Loading embedding model ({self.config.encoder_model})..."); self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            logger.info("Local embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}"); st.stop()
            else: raise
        self.file_chunks = {} # Stores text chunks for FAISS (PDF, PPTX only)
        self.faiss_indexes = {} # Stores FAISS indexes (PDF, PPTX only)
        self.processed_files = set() # Tracks files processed (for Text RAG OR DataFrame)
        self.dataframes = {} # Stores loaded DataFrames {filename: pd.DataFrame} (CSV, XLSX only)
        self._llm = None


    def _get_llm(self):
        """Initializes and returns the Groq LLM client."""
        # (Unchanged)
        if self._llm is None:
            api_key = os.getenv("GROQ_API_KEY")
            if not api_key:
                logger.error("GROQ_API_KEY not found.")
                if self._is_streamlit:
                    import streamlit as st
                    st.error("Error: GROQ_API_KEY is not set.")
                raise ValueError("GROQ_API_KEY not configured.")
            try:
                self._llm = ChatGroq(
                    temperature=self.config.temperature,
                    model_name=self.config.llm_model,
                    groq_api_key=api_key,
                )
                logger.info(f"Groq LLM client initialized ({self.config.llm_model})")
            except Exception as e:
                 logger.error(f"Failed to initialize Groq client: {e}", exc_info=True)
                 if self._is_streamlit:
                     import streamlit as st
                     st.error(f"Failed to initialize LLM. Error: {e}")
                 raise
        return self._llm

    def clean_text(self, text):
        """Clean and normalize text."""
        # (Unchanged)
        if not isinstance(text, str): text = str(text)
        text = re.sub(r"\(cid:.*?\)", "", text)
        text = re.sub(r"\s+", " ", text).strip()
        text = text.replace('\n', ' ').replace('\r', '')
        return text

    # --- Specific Content Extractors ---
    # NOTE: _extract_xlsx and _extract_csv are NO LONGER USED for text extraction path
    # They might be used internally if needed for complex DF loading later, but not for RAG index.
    def _extract_pdf(self, file_path):
        # (Unchanged)
        all_content = []
        base_filename = os.path.basename(file_path)
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Extracting text/tables from PDF: {base_filename} ({total_pages} pages)")
                page_bar = None
                if self._is_streamlit and total_pages > 10: import streamlit as st; page_bar = st.progress(0, text=f"Extracting pages from {base_filename}...")
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text(layout="normal") or ""; cleaned_text = self.clean_text(text)
                    if cleaned_text: all_content.append({"type": "text", "content": cleaned_text, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                    try:
                        for table in page.extract_tables():
                             if table:
                                 table_df = pd.DataFrame(table).fillna(''); table_string = table_df.to_string(index=False, header=True); cleaned_table_string = self.clean_text(table_string)
                                 if cleaned_table_string: all_content.append({"type": "table", "content": cleaned_table_string, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                    except Exception as e_table: logger.warning(f"Could not extract table on page {page_num + 1} in {base_filename}: {e_table}")
                    if page_bar: progress_percent = min(1.0, (page_num + 1) / total_pages); page_bar.progress(progress_percent, text=f"Extracting page {page_num + 1}/{total_pages} from {base_filename}...")
                if page_bar: page_bar.empty()
            logger.info(f"Extracted {len(all_content)} text/table content blocks from PDF: {base_filename}"); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PDF {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    # --- _extract_xlsx REMOVED FROM MAIN TEXT PROCESSING FLOW ---
    # def _extract_xlsx(self, file_path): ... # NO LONGER CALLED by extract_content for text RAG

    # --- _extract_csv REMOVED FROM MAIN TEXT PROCESSING FLOW ---
    # def _extract_csv(self, file_path): ... # NO LONGER CALLED by extract_content for text RAG

    def _extract_pptx(self, file_path):
        # (Unchanged)
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element; merge_threshold_words = self.config.pptx_merge_threshold_words
        try:
            prs = Presentation(file_path); logger.info(f"Extracting from PPTX: {base_filename} ({len(prs.slides)} slides) using merge strategy (threshold: {merge_threshold_words} words).")
            slide_bar = None; total_slides = len(prs.slides)
            if self._is_streamlit and total_slides > 1: import streamlit as st; slide_bar = st.progress(0, text=f"Extracting slides from {base_filename}...")
            pending_title_slide_data = None
            for i, slide in enumerate(prs.slides):
                current_slide_number = i + 1; current_slide_title_text = ""; current_slide_other_texts = []; current_slide_has_title_placeholder = False; title_shape = None
                try:
                    if slide.shapes.title: title_shape = slide.shapes.title; current_slide_has_title_placeholder = True
                except AttributeError: pass
                if title_shape and title_shape.has_text_frame: cleaned_title = self.clean_text(title_shape.text_frame.text); current_slide_title_text = cleaned_title if cleaned_title else ""
                for shape in slide.shapes:
                    if shape == title_shape: continue
                    is_placeholder = shape.is_placeholder; is_body_placeholder = False
                    if is_placeholder:
                        try: ph_type = shape.placeholder_format.type; is_body_placeholder = ph_type in [MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT, MSO_SHAPE_TYPE.SUBTITLE]
                        except AttributeError: pass
                    if shape.has_text_frame:
                        text = shape.text_frame.text; cleaned = self.clean_text(text)
                        if cleaned:
                             if len(cleaned) > max_chars: cleaned = cleaned[:max_chars] + "...(truncated shape)"
                             prefix = "[Body]: " if is_body_placeholder else ""; current_slide_other_texts.append(prefix + cleaned)
                if slide.has_notes_slide:
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text; cleaned_notes = self.clean_text(notes_text)
                        if cleaned_notes:
                            if len(cleaned_notes) > max_chars * 2: cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated notes)"
                            current_slide_other_texts.append(f"[Notes]: {cleaned_notes}")
                    except Exception as e_notes: logger.warning(f"Could not extract notes from slide {current_slide_number}: {e_notes}")
                current_slide_full_content = current_slide_title_text
                if current_slide_other_texts: current_slide_full_content += ("\n" if current_slide_full_content else "") + "\n".join(current_slide_other_texts)
                current_slide_full_content = current_slide_full_content.strip(); other_text_word_count = sum(len(s.split()) for s in current_slide_other_texts)
                merged_content_block = None; should_merge = False
                if pending_title_slide_data:
                    if pending_title_slide_data['has_title'] and pending_title_slide_data['other_words'] <= merge_threshold_words and current_slide_full_content:
                        should_merge = True; logger.info(f"Merging slide {pending_title_slide_data['number']} with slide {current_slide_number}.")
                if should_merge:
                    merged_text = f"[Content from Slide {pending_title_slide_data['number']}]:\n{pending_title_slide_data['content']}\n\n---\n\n[Content from Slide {current_slide_number}]:\n{current_slide_full_content}"
                    merged_content_block = {"type": "slide_text_merged", "content": merged_text.strip(), "source_info": {"slide": current_slide_number, "merged_from": pending_title_slide_data['number']}, "file_type": "pptx"}
                    all_content.append(merged_content_block); pending_title_slide_data = None
                else:
                    if pending_title_slide_data:
                         if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
                         pending_title_slide_data = None
                    if current_slide_has_title_placeholder and other_text_word_count <= merge_threshold_words and current_slide_full_content:
                        logger.debug(f"Slide {current_slide_number} is potential title slide. Holding."); pending_title_slide_data = {"content": current_slide_full_content, "number": current_slide_number, "has_title": current_slide_has_title_placeholder, "other_words": other_text_word_count, "title": current_slide_title_text}
                    else:
                        if current_slide_full_content: all_content.append({"type": "slide_text", "content": current_slide_full_content, "source_info": {"slide": current_slide_number}, "file_type": "pptx"})
                if slide_bar: slide_bar.progress(min(1.0, (i + 1) / total_slides), text=f"Extracting slide {current_slide_number}/{total_slides}...")
            if pending_title_slide_data:
                logger.debug(f"Processing pending title slide {pending_title_slide_data['number']} at end.")
                if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
            if slide_bar: slide_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PPTX {base_filename} (Merge strategy applied)."); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PPTX {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def extract_content(self, file_path):
        """Extract content for text RAG based on file extension (PDF, PPTX ONLY)."""
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf': return self._extract_pdf(file_path)
        # elif extension == '.xlsx': return self._extract_xlsx(file_path) # SKIPPED
        # elif extension == '.csv': return self._extract_csv(file_path) # SKIPPED
        elif extension == '.pptx': return self._extract_pptx(file_path)
        else: logger.warning(f"Unsupported file type for text extraction: {file_path}"); return []

    def chunk_content(self, all_content):
        """Split extracted text content into chunks with overlap, preserving metadata (PDF, PPTX ONLY)."""
        # (Unchanged logic, but now only receives content from PDF/PPTX)
        chunks = []; total_words_estimate = 0; chunk_bar = None
        if not all_content: return chunks
        try: total_words_estimate = sum(len(str(item.get('content', '')).split()) for item in all_content)
        except Exception: pass
        if self._is_streamlit and total_words_estimate > 5000: import streamlit as st; chunk_bar = st.progress(0, text=f"Chunking content...")
        words_processed = 0
        for item_index, item in enumerate(all_content):
            content = item.get('content', ''); source_info = item.get('source_info', {}); file_type = item.get('file_type', 'unknown'); content_type = item.get('type', 'unknown')
            if not isinstance(content, str): content = str(content)
            words = content.split()
            if not words: logger.debug(f"Skipping empty content block {item_index}"); continue
            item_chunks_created = 0
            for i in range(0, len(words), self.config.chunk_size - self.config.overlap):
                chunk_words = words[i:i + self.config.chunk_size]; chunk_text = " ".join(chunk_words)
                if chunk_text: chunks.append({"content": chunk_text, "source_info": source_info, "file_type": file_type, "type": content_type}); item_chunks_created += 1
            if chunk_bar: words_processed += len(words); progress_percent = min(1.0, words_processed / total_words_estimate) if total_words_estimate > 0 else 0; chunk_bar.progress(progress_percent, text=f"Chunking content... ({len(chunks)} chunks created)")
        if chunk_bar: chunk_bar.empty()
        logger.info(f"Created {len(chunks)} text chunks from {len(all_content)} content blocks."); return chunks

    # --- Path Helpers --- (Unchanged)
    def _get_safe_filename(self, file_name): base_name = os.path.splitext(file_name)[0]; return re.sub(r'[^\w\.-]', '_', base_name)
    def get_index_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.index")
    def get_embedding_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.npy")
    def get_chunks_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.json")

    # --- FAISS Index and Chunk Loading/Saving --- (Unchanged logic, but only used for PDF/PPTX now)
    def load_faiss_index(self, file_name, embedding_dim):
        index_path = self.get_index_path(file_name)
        if os.path.exists(index_path):
            try: index = faiss.read_index(index_path); logger.info(f"Loaded FAISS index for {file_name} ({index.ntotal} vectors)"); return index
            except Exception as e: logger.error(f"Error reading FAISS index {index_path}: {e}. Creating new index.")
        logger.info(f"FAISS index not found at {index_path}. Will create new."); return faiss.IndexFlatL2(embedding_dim)
    def save_chunks(self, file_name, chunks):
        chunks_path = self.get_chunks_path(file_name)
        try:
            with open(chunks_path, 'w', encoding='utf-8') as f: json.dump(chunks, f, indent=2); logger.info(f"Saved {len(chunks)} text chunks to {chunks_path}")
        except Exception as e: logger.error(f"Error saving text chunks for {file_name} to {chunks_path}: {e}");
    def load_chunks(self, file_name):
        chunks_path = self.get_chunks_path(file_name)
        if os.path.exists(chunks_path):
            try:
                with open(chunks_path, 'r', encoding='utf-8') as f: chunks = json.load(f); logger.info(f"Loaded {len(chunks)} text chunks from {chunks_path}"); return chunks
            except Exception as e: logger.error(f"Error loading/decoding text chunks from {chunks_path}: {e}. Will re-process."); return None
        logger.info(f"Text chunks file not found at {chunks_path}"); return None

    # --- Document Processing --- (MODIFIED: Skip Text RAG steps for CSV/XLSX)
    def process_files(self, progress_callback=None):
        self.file_chunks = {}; self.faiss_indexes = {}; self.processed_files = set(); self.dataframes = {} # Reset all states
        data_dir = self.config.data_dir; supported_ext = self.config.supported_extensions
        try: all_files = os.listdir(data_dir); process_list = sorted([f for f in all_files if f.lower().endswith(supported_ext)])
        except Exception as e: logger.error(f"Error listing files in data directory {data_dir}: {e}", exc_info=True); return False
        if not process_list: logger.warning(f"No supported files found in {data_dir}"); return True
        logger.info(f"Processing {len(process_list)} supported files from {data_dir}"); embedding_dim = self.encoder_model.get_sentence_embedding_dimension(); total_files = len(process_list)
        if progress_callback: progress_callback(f"Found {total_files} file(s). Starting processing...", current_step=0, total_steps=total_files)

        for idx, file_name in enumerate(process_list):
            current_file_msg = f"Processing ({idx+1}/{total_files}): {file_name}"; logger.info(f"--- {current_file_msg} ---")
            if progress_callback: progress_callback(current_file_msg, current_step=idx, total_steps=total_files)
            file_path = os.path.join(data_dir, file_name); index_path = self.get_index_path(file_name); emb_path = self.get_embedding_path(file_name); chunks_path = self.get_chunks_path(file_name)
            file_extension = os.path.splitext(file_name)[1].lower()

            # --- Text Processing (Chunking/Indexing) - PDF/PPTX ONLY ---
            if file_extension in ['.pdf', '.pptx']:
                logger.info(f"Starting text processing for {file_name}")
                try:
                    chunks = self.load_chunks(file_name)
                    if chunks is None:
                        if progress_callback: progress_callback(f"{current_file_msg} - Extracting Text...", stage="Extracting")
                        all_content = self.extract_content(file_path); # Extracts text representation
                        if not all_content: logger.warning(f"No text content extracted from {file_name}. Skipping text indexing.");
                        else:
                            if progress_callback: progress_callback(f"{current_file_msg} - Chunking Text...", stage="Chunking")
                            chunks = self.chunk_content(all_content);
                            if not chunks: logger.warning(f"No text chunks generated for {file_name}.");
                            else: self.save_chunks(file_name, chunks)

                    # Proceed with embedding/indexing only if chunks exist
                    if chunks:
                        self.file_chunks[file_name] = chunks; logger.debug(f"Stored {len(chunks)} text chunks for {file_name}")
                        faiss_index = None; regenerate_embeddings = False
                        if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                            if progress_callback: progress_callback(f"{current_file_msg} - Verifying Index...", stage="Verifying")
                            try:
                                embeddings = np.load(emb_path)
                                # Check for consistency
                                loaded_chunks_len = len(self.load_chunks(file_name) or []) # Reload to ensure consistency check is valid
                                if embeddings.ndim != 2 or embeddings.shape[1] != embedding_dim or embeddings.shape[0] != loaded_chunks_len:
                                     logger.warning(f"Index/Chunk mismatch for {file_name} (Embeddings: {embeddings.shape}, Chunks: {loaded_chunks_len}). Regenerating."); regenerate_embeddings = True
                                else: faiss_index = self.load_faiss_index(file_name, embedding_dim); logger.info(f"Verified existing text index data for {file_name}")
                            except Exception as e: logger.error(f"Error verifying index for {file_name}: {e}. Regenerating...", exc_info=True); regenerate_embeddings = True; faiss_index = None; chunks = None # Invalidate chunks if verification fails badly
                        else: regenerate_embeddings = True; logger.info(f"Missing index/embedding/chunk files for {file_name}. Regenerating.")

                        # Regenerate if needed (ensure chunks is still valid)
                        if (regenerate_embeddings or faiss_index is None) and chunks:
                            logger.info(f"Generating text embeddings/index for {file_name}...")
                            if progress_callback: progress_callback(f"{current_file_msg} - Embedding Text...", stage="Embedding")
                            content_list = [chunk['content'] for chunk in chunks];
                            if not content_list: logger.warning("No text content to embed.");
                            else:
                                embeddings = self.encoder_model.encode(content_list, batch_size=64, show_progress_bar=False, convert_to_numpy=True).astype('float32')
                                if embeddings.shape[0] == 0: logger.warning("Text embedding yielded no vectors.");
                                else:
                                    np.save(emb_path, embeddings); logger.info(f"Saved {embeddings.shape[0]} text embeddings")
                                    if progress_callback: progress_callback(f"{current_file_msg} - Indexing Text...", stage="Indexing")
                                    faiss_index = faiss.IndexFlatL2(embedding_dim); faiss_index.add(embeddings); faiss.write_index(faiss_index, index_path); logger.info(f"Saved FAISS text index ({faiss_index.ntotal} vectors)")

                        if faiss_index is not None and faiss_index.ntotal > 0:
                             self.faiss_indexes[file_name] = faiss_index; self.processed_files.add(file_name); logger.debug(f"Stored FAISS text index for {file_name}")
                        else: logger.warning(f"No valid text index created for {file_name}.")
                    else:
                        logger.info(f"No text chunks for {file_name}, skipping text indexing.")

                except Exception as e_text:
                    logger.error(f"Failed text processing for {file_name}: {e_text}", exc_info=True)
                    if progress_callback: progress_callback(f"❌ Error during text processing for {file_name}: {e_text}", is_error=True)
                    # Clean up potentially inconsistent text index files
                    try:
                        for p in [index_path, emb_path, chunks_path]:
                            if os.path.exists(p): os.remove(p); logger.debug(f"Removed potentially inconsistent text file {p} for {file_name}")
                    except OSError as e_clean: logger.error(f"Error cleaning up text files for {file_name}: {e_clean}")
                    if file_name in self.file_chunks: del self.file_chunks[file_name]
                    if file_name in self.faiss_indexes: del self.faiss_indexes[file_name]
                    if file_name in self.processed_files: self.processed_files.remove(file_name) # Remove from processed if text part failed

            # --- DataFrame Loading (for CSV/XLSX) ---
            elif file_extension in ['.csv', '.xlsx']:
                logger.info(f"Skipping text indexing for {file_name}. Attempting DataFrame load.")
                df_load_msg = f"{current_file_msg} - Loading DataFrame..."
                if progress_callback: progress_callback(df_load_msg, stage="Loading DataFrame")
                try:
                    df = None
                    if file_extension == '.csv':
                        encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
                        for enc in encodings_to_try:
                            try:
                                df = pd.read_csv(file_path, encoding=enc, low_memory=False)
                                logger.info(f"Successfully loaded DataFrame from CSV {file_name} using encoding: {enc}")
                                break
                            except UnicodeDecodeError: continue
                            except Exception as e_read: logger.warning(f"Pandas read_csv error for {file_name} with encoding {enc}: {e_read}")
                        if df is None: raise ValueError(f"Could not read CSV {file_name} with any attempted encoding.")
                    elif file_extension == '.xlsx':
                        excel_file = pd.ExcelFile(file_path)
                        # Simplification: Load only the first sheet. Could be enhanced later.
                        if excel_file.sheet_names:
                            first_sheet_name = excel_file.sheet_names[0]
                            # Try parsing with header detection, fallback if needed
                            try:
                                df = excel_file.parse(first_sheet_name, header=0) # Assume header is first row
                            except Exception as e_parse_header:
                                logger.warning(f"Parsing sheet '{first_sheet_name}' with header failed, trying without: {e_parse_header}. Reloading without header.")
                                try:
                                     df = excel_file.parse(first_sheet_name, header=None)
                                     # Optionally rename columns to generic names if no header
                                     df.columns = [f'col_{i}' for i in range(len(df.columns))]
                                except Exception as e_parse_no_header:
                                     logger.error(f"Failed to parse sheet '{first_sheet_name}' even without header: {e_parse_no_header}")
                                     raise # Re-raise the error

                            logger.info(f"Successfully loaded DataFrame from XLSX {file_name}, sheet: '{first_sheet_name}' (Rows: {len(df)}, Cols: {len(df.columns)})")
                        else:
                            raise ValueError(f"XLSX file {file_name} has no sheets.")

                    if df is not None and not df.empty:
                        # Clean column names: replace non-alphanumeric with underscore
                        original_columns = df.columns
                        df.columns = [re.sub(r'[^0-9a-zA-Z_]+', '_', str(col).strip()) for col in df.columns]
                        cleaned_columns = df.columns
                        if list(original_columns) != list(cleaned_columns):
                            logger.info(f"Cleaned DataFrame column names for {file_name}. Original: {list(original_columns)}, Cleaned: {list(cleaned_columns)}")

                        self.dataframes[file_name] = df # Store the DataFrame
                        logger.info(f"Stored DataFrame for {file_name} (Rows: {len(df)}, Columns: {len(df.columns)})")
                        self.processed_files.add(file_name) # Mark as processed because DF loaded
                    else:
                         logger.warning(f"Loaded DataFrame for {file_name} is None or empty.")

                except Exception as e_df:
                    logger.error(f"Failed to load DataFrame for {file_name}: {e_df}", exc_info=True)
                    if progress_callback: progress_callback(f"⚠️ Warning: Could not load DataFrame for {file_name}: {e_df}", is_warning=True)
                    if file_name in self.dataframes: del self.dataframes[file_name] # Clean up if failed midway
                    if file_name in self.processed_files: self.processed_files.remove(file_name) # Remove if DF loading failed
            else:
                 logger.warning(f"File type {file_extension} not processed for text index or DataFrame: {file_name}")


        final_indexed_count = len(self.faiss_indexes) # Only PDF/PPTX
        final_dataframe_count = len(self.dataframes) # Only CSV/XLSX
        total_processed_ok = len(self.processed_files) # Files with either index OR dataframe successfully loaded

        if total_processed_ok > 0:
            logger.info(f"--- Processing Complete. Processed {total_processed_ok}/{total_files} files. Text Indices (PDF/PPTX): {final_indexed_count}. DataFrames (CSV/XLSX): {final_dataframe_count}. ---")
            return True
        else:
            logger.warning(f"--- Processing Complete. No documents successfully processed (no text index or DataFrame loaded). ---")
            return True # Still completed, just nothing loaded


    # --- Querying Logic ---

    def query_files(self, query):
        """Query FAISS index for text chunks relevant to a query string (PDF/PPTX ONLY)."""
        if not self.faiss_indexes:
            logger.warning("No text indexes (PDF/PPTX) available for querying.")
            return {} # Return empty results

        query_results = {}
        try:
            logger.info(f"Encoding query for Faiss text search (PDF/PPTX): '{query[:100]}...'")
            query_embedding = self.encoder_model.encode(query, convert_to_numpy=True).astype("float32")
            query_embedding = np.array([query_embedding])
            if query_embedding.ndim != 2: raise ValueError("Query embedding shape error")

            # Search across all available text indexes (PDF/PPTX)
            for file_name, index in self.faiss_indexes.items():
                if index is None or index.ntotal == 0: continue
                try:
                    k_search = min(self.config.k_retrieval, index.ntotal)
                    D, I = index.search(query_embedding, k=k_search)
                    indices, distances = I[0], D[0]

                    current_file_chunks = self.file_chunks.get(file_name)
                    if not current_file_chunks: logger.error(f"Text chunks missing for {file_name} during query!"); continue

                    file_results = []
                    processed_indices = set()
                    for i, idx in enumerate(indices):
                        if idx == -1 or idx in processed_indices: continue
                        if not (0 <= idx < len(current_file_chunks)): continue
                        processed_indices.add(idx)
                        chunk = current_file_chunks[idx]
                        file_results.append({
                            "source_info": chunk.get('source_info', {}),
                            "file_type": chunk.get('file_type', 'unknown'),
                            "content": chunk.get('content', ''),
                            "score": round(float(distances[i]), 4), # L2 distance, lower is better
                            "type": chunk.get('type', 'unknown')
                        })

                    file_results.sort(key=lambda x: x['score']) # Lower L2 is better
                    if file_results:
                        if file_name not in query_results: query_results[file_name] = []
                        query_results[file_name].extend(file_results)

                except Exception as e_search:
                    logger.error(f"Error searching text index {file_name} for query '{query[:50]}...': {e_search}", exc_info=True)

            logger.debug(f"Faiss text search complete for query '{query[:50]}...'. Found results in {len(query_results)} PDF/PPTX files.")
            return query_results

        except Exception as e_query:
            logger.error(f"Error during Faiss text query processing for '{query[:50]}...': {e_query}", exc_info=True)
            return {} # Return empty on error

    def aggregate_context(self, query_results, strategy="top_k"):
        """Aggregate text context from FAISS results (PDF/PPTX ONLY)."""
        # (Unchanged logic, but now only receives results from PDF/PPTX)
        all_context = {}; max_chars = self.config.max_context_tokens * 3
        logger.info(f"Aggregating text context (PDF/PPTX) using strategy '{strategy}', max_chars ~{max_chars}")
        if not query_results: return all_context
        flat_results = []
        for file_name, results in query_results.items():
            for res in results: flat_results.append({**res, "file_name": file_name})
        flat_results.sort(key=lambda x: x['score']) # Sort globally by score (lower is better)
        aggregated_context_str = ""; total_aggregated_chars = 0; added_chunks_count = 0; context_sources = set()
        limit = self.config.k_retrieval if strategy == "top_k" else len(flat_results)
        for i, res in enumerate(flat_results):
            if added_chunks_count >= limit:
                 logger.debug(f"Reached global top-k limit ({limit}) for text context aggregation.")
                 break
            source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown'); content_type = res.get('type', 'unknown'); score = res['score']; content_body = res['content']
            source_parts = [f"Source: {file_name}"]
            slide_info = source_info.get('slide', 'N/A'); merged_from = source_info.get('merged_from')
            if file_type == 'pptx': slide_display = f"Slide: {slide_info}" + (f" (merged from {merged_from})" if merged_from else ""); source_parts.append(slide_display)
            elif file_type == 'pdf': source_parts.append(f"Page: {source_info.get('page', 'N/A')}")
            # Note: xlsx/csv text chunk types should no longer appear here
            elif file_type in ['xlsx', 'csv']: logger.warning(f"Unexpected {file_type} chunk found during text aggregation."); continue # Skip just in case
            source_parts.append(f"Type: {content_type}"); source_parts.append(f"Score: {score:.4f}")
            source_str = ", ".join(source_parts); content_header = f"--- Context from {source_str} ---\n"; content_to_add = content_header + content_body + "\n\n"; content_chars = len(content_to_add)
            if total_aggregated_chars + content_chars <= max_chars:
                aggregated_context_str += content_to_add; total_aggregated_chars += content_chars; added_chunks_count += 1; context_sources.add(file_name)
            else:
                 if added_chunks_count == 0:
                     remaining_chars = max_chars - total_aggregated_chars - len(content_header) - 20 # Headroom
                     if remaining_chars > 50:
                         truncated_body = content_body[:remaining_chars]; aggregated_context_str += content_header + truncated_body + "\n[...TRUNCATED CONTEXT...]\n\n"; total_aggregated_chars += len(content_to_add); added_chunks_count += 1; context_sources.add(file_name); logger.warning("Context truncated (first chunk too large).")
                     else: logger.warning("First chunk too large to fit even truncated, skipping.")
                 logger.info(f"Stopping text context aggregation at {added_chunks_count} chunks ({total_aggregated_chars}/{max_chars} chars). Limit reached."); break
        final_context = aggregated_context_str.strip()
        if final_context: all_context = {"combined_context": final_context, "source_files": sorted(list(context_sources))}; logger.info(f"Aggregated {total_aggregated_chars} chars from {added_chunks_count} PDF/PPTX text chunks across {len(context_sources)} files.")
        else: logger.warning("No text context aggregated from PDF/PPTX within limits.")
        return all_context

    def query_llm(self, query, context_data, retry_count=1):
        """Query LLM with aggregated text context (PDF/PPTX ONLY)."""
        # (Unchanged logic, but context_data now only contains PDF/PPTX info)
        combined_context = context_data.get("combined_context", ""); source_files = context_data.get("source_files", []); source_file_str = ", ".join(source_files) if source_files else "the provided PDF/PPTX documents"
        if not combined_context: logger.warning("No text context provided for LLM query."); return f"Could not generate answer: No relevant text context found in {source_file_str}."
        try:
            llm = self._get_llm()
            system_prompt = f"""You are an AI assistant answering questions based ONLY on the provided text context from document(s): '{source_file_str}'. Use ONLY information in the 'START CONTEXT'/'END CONTEXT' section. If the answer isn't there, state: "Based on the provided text context from {source_file_str}, I cannot answer this question." Do NOT use external knowledge or perform calculations not explicitly shown in the text."""
            human_prompt = f"""Context from document(s) '{source_file_str}':\n--- START CONTEXT ---\n{combined_context}\n--- END CONTEXT ---\n\nUser Question: {query}\n\nAnswer based ONLY on the provided text context:"""
            full_prompt_messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
            logger.info(f"Querying Groq model {self.config.llm_model} using text context from {source_file_str}...")
            answer = f"Error: LLM query failed after retries."
            for attempt in range(retry_count + 1):
                try:
                    response = llm.invoke(full_prompt_messages)
                    answer = response.content.strip() if hasattr(response, 'content') else str(response).strip()
                    logger.info(f"Groq response received (attempt {attempt+1}). Length: {len(answer)}")
                    if not answer and attempt < retry_count: logger.warning("Empty response, retrying..."); time.sleep(1); continue
                    return answer or f"Received empty response from LLM for text context from {source_file_str}."
                except Exception as e_api:
                    logger.warning(f"Groq API attempt {attempt+1} failed: {e_api}")
                    if attempt < retry_count: time.sleep(1.5 ** attempt); logger.info("Retrying Groq query...")
                    else: answer = f"Error: Failed to get answer from LLM after {retry_count+1} attempts. (API Error: {e_api})"
            return answer
        except Exception as e_setup:
            logger.error(f"Error setting up/calling Groq API for text RAG: {e_setup}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"LLM text query failed: {e_setup}")
            return f"Error: Could not query LLM due to setup error: {e_setup}"


    # ==========================================================================
    # START OF NEW/MODIFIED HELPER METHODS FOR QUERY ANALYSIS & DATAFRAME OPS
    # ==========================================================================

    def _call_llm_for_analysis(self, prompt_messages, task_description):
        """Helper to call LLM for internal analysis tasks."""
        # (Unchanged)
        try:
            llm = self._get_llm()
            logger.info(f"Calling LLM for {task_description}...")
            # Consider adjusting temperature for analysis tasks if needed
            # llm.temperature = 0.0 # Example: Lower temp for deterministic tasks
            response = llm.invoke(prompt_messages)
            # llm.temperature = self.config.temperature # Reset if changed
            content = response.content.strip() if hasattr(response, 'content') else str(response).strip()
            logger.info(f"LLM response for {task_description}: '{content[:150]}...'")
            return content
        except Exception as e:
            logger.error(f"LLM call failed during {task_description}: {e}", exc_info=True)
            return None # Indicate failure

    def _classify_query(self, query):
        """
        Classify the query type: Simple Retrieval, Complex/Reasoning (Text), or Structured Query (DataFrame).
        Returns a tuple: (classification_string, confidence_score [0-1] or None)
        """
        # (Unchanged)
        system_prompt = """You are an expert query analyzer. Classify the user query into ONE of the following categories:
1.  'Simple Retrieval': Asking for specific facts, definitions, or simple summaries directly extractable from text (likely from PDF or PPTX).
2.  'Complex/Reasoning (Text)': Requires combining information from multiple text passages (likely from PDF or PPTX), summarization across sections, or reasoning based *only* on the text provided. Does not involve calculations on tabular data.
3.  'Structured Query (DataFrame)': Requires operations on structured data (like CSV or Excel tables), such as calculations (sum, average, count), filtering based on values, grouping, sorting, or comparisons across rows/columns.

Analyze the query carefully. Respond ONLY in JSON format with two keys: "classification" (the category string) and "confidence" (a float between 0.0 and 1.0 indicating your certainty, e.g., 0.95).

Examples:
Query: "What is the definition of RAG?" -> {"classification": "Simple Retrieval", "confidence": 0.99}
Query: "Summarize the main challenges discussed in the report." -> {"classification": "Complex/Reasoning (Text)", "confidence": 0.9}
Query: "What is the total revenue for product X in Q3?" -> {"classification": "Structured Query (DataFrame)", "confidence": 0.98}
Query: "Compare the findings section with the conclusion." -> {"classification": "Complex/Reasoning (Text)", "confidence": 0.85}
Query: "List all employees hired after 2022." -> {"classification": "Structured Query (DataFrame)", "confidence": 0.96}
Query: "Tell me about the company history." -> {"classification": "Simple Retrieval", "confidence": 0.9}
Query: "How many farms in the South region use organic fertilizer?" -> {"classification": "Structured Query (DataFrame)", "confidence": 0.97}
"""
        human_prompt = f"Classify the following user query:\n\"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        raw_response = self._call_llm_for_analysis(messages, "query classification")

        try:
            result = json.loads(raw_response)
            classification = result.get("classification")
            confidence = float(result.get("confidence", 0.0))
            valid_classifications = ["Simple Retrieval", "Complex/Reasoning (Text)", "Structured Query (DataFrame)"]
            if classification in valid_classifications and 0.0 <= confidence <= 1.0:
                logger.info(f"Query classified as '{classification}' with confidence {confidence:.2f}")
                return classification, confidence
            else:
                logger.warning(f"Invalid classification or confidence in response: {raw_response}. Defaulting.")
                return "Simple Retrieval", 0.5
        except (json.JSONDecodeError, TypeError, AttributeError, ValueError) as e:
            logger.warning(f"Failed to parse classification JSON response: '{raw_response}'. Error: {e}. Defaulting.")
            raw_response_cleaned = raw_response.strip().replace("'", '"')
            if raw_response_cleaned in ["Simple Retrieval", "Complex/Reasoning (Text)", "Structured Query (DataFrame)"]:
                 logger.info(f"Using raw response as classification: '{raw_response_cleaned}' (confidence unknown)")
                 return raw_response_cleaned, 0.6
            return "Simple Retrieval", 0.5

    def _decompose_query(self, query):
        """Decompose a complex TEXT query into simpler sub-queries (for PDF/PPTX)."""
        # (Unchanged - still relevant for Complex/Reasoning (Text) on PDF/PPTX)
        system_prompt = """You are an expert query decomposer. Your task is to break down a complex user query that requires combining information from TEXT passages (likely from PDF or PPTX) into a series of simpler, factual sub-queries. Each sub-query should aim to retrieve a specific piece of information from text needed to answer the original query.
Do NOT decompose queries asking for calculations or filtering on tables (those are handled differently). Focus on text-based information gathering steps.
Format the output as a numbered list, with each sub-query on a new line. If the query cannot be decomposed or is already simple, return the original query prefixed with '1. '."""
        human_prompt = f"Decompose the following complex text query into simple, factual sub-queries:\n\"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        decomposition = self._call_llm_for_analysis(messages, "query decomposition")

        if decomposition:
            sub_queries = [re.sub(r"^\d+\.\s*", "", line).strip() for line in decomposition.split('\n') if line.strip()]
            if sub_queries and not (len(sub_queries) == 1 and sub_queries[0] == query):
                logger.info(f"Decomposed text query into: {sub_queries}")
                return sub_queries
            else:
                 logger.info(f"Decomposition didn't yield distinct sub-queries for text. Using original query.")
                 return [query]
        else:
            logger.warning(f"Text query decomposition failed. Using original query.")
            return [query]

    def _generate_pandas_code(self, query, df_columns, df_name="df"):
        """Generate Python Pandas code to answer a query using a DataFrame."""
        # Added cleaning instructions for column names if needed
        system_prompt = f"""You are a Python Pandas expert. Given a user query and the list of columns in a Pandas DataFrame named `{df_name}`, write Python code using this DataFrame to answer the query.
Instructions:
- The DataFrame is already loaded and available as the variable `{df_name}`.
- Assume `pandas` is imported as `pd` and `numpy` as `np`.
- Column names in the DataFrame have been cleaned (non-alphanumeric characters replaced by underscores). Refer to the columns using these cleaned names: {df_columns}.
- Write code that calculates or retrieves the specific information asked in the query.
- The final line of your code MUST be a `print()` statement that outputs the computed result (e.g., a number, a list, a small summary DataFrame) or the relevant slice of the DataFrame. Aim for concise output.
- Output ONLY the Python code, nothing else (no explanations, comments, or markdown formatting like ```python).
- Handle potential errors gracefully (e.g., check if columns exist before using them, handle division by zero, check for empty results before calculations). If the query cannot be answered with the given columns, print an informative error message like "Cannot answer query with available columns".

DataFrame Columns (cleaned): {df_columns}
"""
        human_prompt = f"User Query: \"{query}\"\n\nGenerate the Python code to answer this query using the DataFrame `{df_name}`:"
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        code = self._call_llm_for_analysis(messages, "Pandas code generation")

        if code:
            # Basic cleanup: remove markdown code fences if present
            code = re.sub(r"^```python\n?", "", code, flags=re.MULTILINE)
            code = re.sub(r"\n?```$", "", code, flags=re.MULTILINE)
            code = code.strip() # Remove leading/trailing whitespace

            # Ensure it contains the df_name variable and a print statement
            if df_name not in code:
                 logger.warning(f"Generated code does not reference the DataFrame variable '{df_name}'. Discarding.")
                 return None
            if "print(" not in code:
                 logger.warning(f"Generated code does not contain a 'print()' statement. Appending print of last expression if possible, otherwise discarding.")
                 # Try to intelligently add print if it's just an expression on the last line
                 lines = code.splitlines()
                 if lines and not lines[-1].strip().startswith(('print', '#', 'import', 'def', 'class', 'for', 'while', 'if', 'with')):
                    code = "\n".join(lines[:-1]) + f"\nprint({lines[-1].strip()})"
                    logger.info(f"Appended print statement. New code:\n{code}")
                 else:
                     logger.error("Could not automatically add print statement. Discarding code.")
                     return None # Discard if no print

            logger.info(f"Generated Pandas code:\n{code}")
            return code
        else:
            logger.error("Pandas code generation failed (LLM call unsuccessful).")
            return None

    # --- NEW HELPER FOR FRAMING DATAFRAME RESULTS ---
    def _frame_dataframe_answer(self, query, dataframe_result_string, source_filename):
        """Use LLM to create a natural language answer from the DataFrame query result."""
        logger.info(f"Framing DataFrame result for query: '{query[:100]}...'")
        if not dataframe_result_string:
            logger.warning("Cannot frame answer: DataFrame result string is empty.")
            return f"The analysis of '{source_filename}' did not produce a result for your query."

        # Clean up potential truncation markers or excessive newlines from result string
        dataframe_result_string = re.sub(r'\n\.\.\.(\s*\n)*', '\n...\n', dataframe_result_string).strip()
        if len(dataframe_result_string) > 2000: # Limit context sent to framing LLM
            dataframe_result_string = dataframe_result_string[:2000] + "\n... (result truncated)"
            logger.warning("DataFrame result string truncated before sending for framing.")

        system_prompt = f"""You are an AI assistant. Your task is to present the result of a data analysis in a clear, natural language sentence or paragraph.
You will be given the original user query and the data result obtained by running code on a structured file ('{source_filename}').
Base your answer *solely* on the provided data result. Do not add external information or opinions.
If the result indicates an error or inability to answer, state that clearly.
Be concise and directly answer the user's query using the data."""

        human_prompt = f"""Original User Query: "{query}"

Data analysis result from '{source_filename}':
--- START DATA RESULT ---
{dataframe_result_string}
--- END DATA RESULT ---

Please provide a natural language answer to the user's query based *only* on the data result above:
"""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        framed_answer = self._call_llm_for_analysis(messages, "DataFrame answer framing")

        if framed_answer:
            logger.info("Successfully framed DataFrame answer.")
            return framed_answer
        else:
            logger.error("LLM call failed during DataFrame answer framing. Returning raw result.")
            # Fallback to returning the raw (potentially truncated) data string if framing fails
            return f"Analysis of '{source_filename}' produced the following result:\n```\n{dataframe_result_string}\n```"

    # ==========================================================================
    # END OF NEW/MODIFIED HELPER METHODS
    # ==========================================================================


    # ==========================================================================
    # START OF MODIFIED run_query FUNCTION (CSV/XLSX -> Structured Only)
    # ==========================================================================
    def run_query(self, query, context_strategy="top_k"):
        """Complete query pipeline: classify, then route to DataFrame query (CSV/XLSX) or Text RAG (PDF/PPTX)."""
        logger.info(f"--- Starting query execution (v5.1 - Structured Only for CSV/XLSX): '{query[:100]}...' ---")
        final_results = {
            "query": query,
            "classification": None,
            "classification_confidence": None,
            "sub_queries": None, # For text RAG path
            "retrieval_results": {}, # Text retrieval results (PDF/PPTX only for context)
            "target_dataframe": None, # DF used for structured query
            "pandas_code_used": None, # Code executed for structured query
            "dataframe_result": None, # Raw result from pandas execution
            "aggregated_context_data": {}, # For text RAG path
            "answer": "",
            "answer_source": None, # 'DataFrame (Framed)' or 'Text RAG'
            "status": "Started",
            "error": None
        }
        is_streamlit = self._is_streamlit
        if is_streamlit: import streamlit as st

        try:
            # === 1. Classify Query ===
            spinner_msg = "Analyzing query..."
            if is_streamlit:
                with st.spinner(spinner_msg): classification, confidence = self._classify_query(query)
            else: logger.info(spinner_msg); classification, confidence = self._classify_query(query)

            final_results["classification"] = classification
            final_results["classification_confidence"] = confidence
            logger.info(f"Query classified as: {classification} (Confidence: {confidence:.2f})")

            # === Route based on Classification ===

            # --- Path A: Structured Query (DataFrame - CSV/XLSX) ---
            if classification == "Structured Query (DataFrame)" and confidence >= self.config.dataframe_query_confidence_threshold and self.dataframes:
                logger.info("Attempting DataFrame query path.")
                final_results["answer_source"] = "DataFrame (Framed)" # Tentative

                # === 2a. Identify Target DataFrame ===
                target_df_name = None
                identified_by = "N/A"

                # Strategy 1: Use text search results (from PDF/PPTX) if available to guide selection
                spinner_msg = "Identifying relevant structured document (using text context if available)..."
                if is_streamlit:
                    with st.spinner(spinner_msg): text_retrieval_results = self.query_files(query) # Searches only PDF/PPTX indexes
                else: logger.info(spinner_msg); text_retrieval_results = self.query_files(query)
                final_results["retrieval_results"] = text_retrieval_results # Store for potential debug/info

                if text_retrieval_results:
                    best_score = float('inf')
                    best_file_from_text = None
                    for file_name, results in text_retrieval_results.items():
                         # Consider all results, not just first, maybe weighted? Simpler: best score overall.
                         if results:
                              current_best_score = min(r['score'] for r in results)
                              if current_best_score < best_score:
                                   best_score = current_best_score
                                   best_file_from_text = file_name # Store the file containing the best text chunk

                    # Now, we need to *infer* the target DataFrame based on the text file hit
                    # This is heuristic - assumes the structured file has a similar name or the query relates them
                    # Example: If query matches "Report.pdf" strongly, look for "Report.xlsx" or "Report_Data.csv"
                    if best_file_from_text:
                         logger.info(f"Best text match found in '{best_file_from_text}' (Score: {best_score:.4f}). Attempting to link to DataFrame.")
                         # Simple check: does a DF exist with the same base name?
                         base_name_text = os.path.splitext(best_file_from_text)[0]
                         potential_df_names = [dfn for dfn in self.dataframes.keys() if os.path.splitext(dfn)[0] == base_name_text]
                         if len(potential_df_names) == 1:
                             target_df_name = potential_df_names[0]
                             identified_by = f"Linked from text match in '{best_file_from_text}'"
                             logger.info(f"Heuristically linked to DataFrame: '{target_df_name}'")
                         elif len(potential_df_names) > 1 :
                              logger.warning(f"Ambiguous link: Text match '{best_file_from_text}' links to multiple DataFrames: {potential_df_names}. Cannot select target.")
                              final_results["error"] = f"Ambiguous: Cannot determine target DataFrame among {potential_df_names} based on text context."
                              # Proceed to check fallbacks

                # Strategy 2: If no target found via text link AND only one DataFrame exists, assume it's the target.
                if not target_df_name and len(self.dataframes) == 1:
                    target_df_name = list(self.dataframes.keys())[0]
                    identified_by = "Only one DataFrame loaded"
                    logger.info(f"Identified target DataFrame as it's the only one loaded: '{target_df_name}'")

                # Strategy 3: If still no target... fail.
                if not target_df_name:
                    if not self.dataframes: # Should have been caught by initial check, but defensive
                         logger.error("Classification indicated DataFrame, but no DataFrames are loaded.")
                         final_results["error"] = "Query requires structured data, but no CSV/XLSX files were successfully loaded."
                    elif len(self.dataframes) > 1:
                         logger.warning(f"Could not identify target DataFrame among multiple options: {list(self.dataframes.keys())}. Text context did not provide a clear link.")
                         final_results["error"] = f"Ambiguous: Query requires structured data, but couldn't determine target among multiple files ({list(self.dataframes.keys())}). Try referencing related text documents or simplifying the setup."
                    else: # Should not happen if len==1 case is handled
                         logger.error("Failed to identify target DataFrame.")
                         final_results["error"] = "Could not identify the target CSV/XLSX file for the query."

                    # --- EXIT DataFrame Path on Identification Failure ---
                    final_results["status"] = "Failed: Target DataFrame Identification"
                    final_results["answer"] = f"Error: {final_results['error']}"
                    logger.error(f"Stopping query execution: {final_results['error']}")
                    return final_results
                    # --- NO FALLBACK TO TEXT RAG ---

                # === 3a. Generate and Execute Pandas Code ===
                final_results["target_dataframe"] = target_df_name
                logger.info(f"Proceeding with DataFrame '{target_df_name}' (Identified by: {identified_by})")
                df_to_use = self.dataframes[target_df_name]
                # Use cleaned column names for code generation prompt
                df_columns = list(df_to_use.columns)

                spinner_msg = f"Generating analysis code for '{target_df_name}'..."
                if is_streamlit:
                    with st.spinner(spinner_msg): pandas_code = self._generate_pandas_code(query, df_columns)
                else: logger.info(spinner_msg); pandas_code = self._generate_pandas_code(query, df_columns)

                if pandas_code:
                    final_results["pandas_code_used"] = pandas_code
                    spinner_msg = f"Executing analysis code on '{target_df_name}'..."
                    logger.info(f"Executing Pandas code for query: '{query[:100]}...'")

                    execution_success = False
                    result_value = None
                    exec_error = None
                    try:
                        # Execute code and capture output
                        # Define the scope available to the executed code
                        execution_globals = {'df': df_to_use.copy(), 'pd': pd, 'np': np}
                        string_io = io.StringIO()
                        with redirect_stdout(string_io):
                            # !!! SECURITY RISK !!! Use with extreme caution / sandboxing in production
                            # Pass the scope as the globals dictionary. Locals will default to this.
                            exec(pandas_code, execution_globals) # <--- CORRECTED LINE (v5.1)
                        result_value = string_io.getvalue().strip()
                        final_results["dataframe_result"] = result_value
                        execution_success = True
                        logger.info(f"Pandas code execution successful. Raw Result Length: {len(result_value)}")
                        logger.debug(f"Raw Pandas Result:\n{result_value}")

                    except Exception as e_exec:
                        logger.error(f"Pandas code execution failed: {e_exec}", exc_info=True)
                        # Provide more specific error if possible
                        exec_error = f"Pandas code execution failed: {type(e_exec).__name__}: {e_exec}"
                        final_results["error"] = exec_error
                        final_results["status"] = "Failed: DataFrame Execution"
                        final_results["answer"] = f"Error executing analysis on '{target_df_name}': {e_exec}. Check the generated code and DataFrame columns."
                        # --- EXIT DataFrame Path on Execution Failure ---
                        logger.error(f"Stopping query execution: {final_results['error']}")
                        return final_results
                        # --- NO FALLBACK TO TEXT RAG ---

                    # === 4a. Frame Answer using LLM ===
                    if execution_success:
                        spinner_msg = f"Framing answer from '{target_df_name}' analysis..."
                        if is_streamlit:
                            with st.spinner(spinner_msg): framed_answer = self._frame_dataframe_answer(query, result_value, target_df_name)
                        else:
                            logger.info(spinner_msg); framed_answer = self._frame_dataframe_answer(query, result_value, target_df_name)

                        if framed_answer:
                            final_results["answer"] = framed_answer
                            final_results["status"] = "Completed Successfully"
                            final_results["error"] = None # Clear any previous non-fatal errors if framing succeeded
                        else:
                             # Framing failed, provide raw result as fallback (as handled in _frame_dataframe_answer)
                             final_results["answer"] = f"Failed to frame the result from '{target_df_name}'. Raw analysis output:\n```\n{result_value or '[No Output]'}\n```"
                             final_results["status"] = "Completed with Framing Failure"
                             final_results["error"] = "LLM failed to generate natural language answer from data."

                else: # Failed to generate pandas code
                    logger.warning("Failed to generate Pandas code.")
                    final_results["status"] = "Failed: Pandas Code Generation"
                    final_results["error"] = f"Could not generate analysis code for the query on '{target_df_name}'."
                    final_results["answer"] = f"Error: Unable to generate the analysis code needed to answer the query using '{target_df_name}'."
                    # --- EXIT DataFrame Path on Code Gen Failure ---
                    logger.error(f"Stopping query execution: {final_results['error']}")
                    return final_results
                    # --- NO FALLBACK TO TEXT RAG ---

            # --- Path B: Text RAG (PDF/PPTX) ---
            else:
                # Execute if:
                # - Classification is Simple Retrieval or Complex/Reasoning (Text)
                # - OR Classification is Structured but confidence too low
                # - OR Classification is Structured but no DataFrames were loaded
                if classification == "Structured Query (DataFrame)" and not self.dataframes:
                     logger.warning("Query classified as Structured, but no CSV/XLSX DataFrames loaded. Cannot proceed.")
                     final_results["status"] = "Failed: No DataFrames Loaded"
                     final_results["answer_source"] = "N/A"
                     final_results["error"] = "Query requires structured data (CSV/XLSX), but none were loaded."
                     final_results["answer"] = final_results["error"]
                     return final_results # Exit

                logger.info("Proceeding with Text RAG path (PDF/PPTX).")
                final_results["answer_source"] = "Text RAG"

                # === 2b. Decompose if Complex Text Query ===
                queries_to_retrieve = [query]
                if classification == "Complex/Reasoning (Text)":
                     spinner_msg = "Decomposing complex text query..."
                     if is_streamlit:
                          with st.spinner(spinner_msg): sub_queries = self._decompose_query(query)
                     else: logger.info(spinner_msg); sub_queries = self._decompose_query(query)
                     if sub_queries and sub_queries != [query]:
                          queries_to_retrieve = sub_queries
                          final_results["sub_queries"] = sub_queries
                          logger.info(f"Using sub-queries for text retrieval: {sub_queries}")

                # === 3b. Retrieve relevant text chunks (PDF/PPTX only) ===
                spinner_msg = f"Searching PDF/PPTX documents for {len(queries_to_retrieve)} query part(s)..."
                all_query_results = {}
                # Reuse results if already fetched during (failed) DF target identification?
                # Check if results exist and are non-empty. If so, maybe reuse.
                # For now, let's re-fetch to ensure consistency for the text path.
                # final_results["retrieval_results"] = {} # Reset if re-fetching
                if is_streamlit:
                    with st.spinner(spinner_msg):
                        for i, q in enumerate(queries_to_retrieve):
                            logger.info(f"Retrieving text for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                            results_for_q = self.query_files(q) # Calls Faiss search (PDF/PPTX only)
                            for file, res_list in results_for_q.items():
                                if file not in all_query_results: all_query_results[file] = []
                                all_query_results[file].extend(res_list)
                else:
                    logger.info(spinner_msg)
                    for i, q in enumerate(queries_to_retrieve):
                        logger.info(f"Retrieving text for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                        results_for_q = self.query_files(q)
                        for file, res_list in results_for_q.items():
                                if file not in all_query_results: all_query_results[file] = []
                                all_query_results[file].extend(res_list)

                # De-duplicate and sort text results per file
                final_retrieval_results = {}
                if all_query_results:
                    for file, res_list in all_query_results.items():
                        unique_content = {}
                        for res in res_list:
                            key = (res['content'], tuple(sorted(res.get('source_info',{}).items())))
                            # Use score to keep the best match if duplicates exist
                            if key not in unique_content or res['score'] < unique_content[key]['score']:
                                unique_content[key] = res
                        sorted_unique_results = sorted(unique_content.values(), key=lambda x: x['score'])
                        final_retrieval_results[file] = sorted_unique_results
                final_results["retrieval_results"] = final_retrieval_results

                # Check if text retrieval found anything
                if not final_results.get("retrieval_results"):
                     logger.warning("No relevant text chunks found in PDF/PPTX for the query/sub-queries.")
                     final_results["status"] = "Completed: No relevant information found"
                     final_results["answer"] = "Could not find relevant information in the PDF/PPTX documents."
                     if is_streamlit: st.info(final_results["answer"])
                     return final_results # Exit if no text chunks found

                # === 4b. Aggregate text context ===
                spinner_msg = "Gathering text context (PDF/PPTX)..."
                if is_streamlit:
                     with st.spinner(spinner_msg): aggregated_context_data = self.aggregate_context(final_results["retrieval_results"], strategy=context_strategy)
                else: logger.info(spinner_msg); aggregated_context_data = self.aggregate_context(final_results["retrieval_results"], strategy=context_strategy)
                final_results["aggregated_context_data"] = aggregated_context_data

                if not aggregated_context_data or not aggregated_context_data.get("combined_context"):
                      logger.warning(f"Text context aggregation failed or yielded empty context for query: '{query[:100]}...'")
                      final_results["status"] = "Completed: Text context aggregation failed or empty"
                      final_results["answer"] = "Could not prepare text context from PDF/PPTX for answering."
                      final_results["error"] = "Text context aggregation failed or empty."
                      if is_streamlit: st.warning(final_results["answer"])
                      return final_results

                # === 5b. Query LLM with aggregated text context ===
                spinner_msg = "Generating final answer from text (PDF/PPTX)..."
                if is_streamlit:
                      with st.spinner(spinner_msg): answer = self.query_llm(query, aggregated_context_data) # Pass ORIGINAL query
                else: logger.info(spinner_msg); answer = self.query_llm(query, aggregated_context_data)

                final_results["answer"] = answer
                final_results["status"] = "Completed Successfully"
                final_results["error"] = None # Clear any previous classification info/warnings

            # Final logging
            logger.info(f"--- Finished query execution (v5.1) for: '{query[:100]}...'. Answer Source: {final_results.get('answer_source')} ---")

        except Exception as e:
            logger.error(f"Unexpected error during run_query (v5.1): {e}", exc_info=True)
            final_results["status"] = "Failed: Unexpected Error"
            final_results["error"] = str(e)
            final_results["answer"] = f"An unexpected error occurred: {e}"
            if is_streamlit: st.error(f"An unexpected error occurred: {e}")

        return final_results
    # ==========================================================================
    # END OF MODIFIED run_query FUNCTION
    # ==========================================================================


# --- START OF STREAMLIT UI CODE ---
# (Streamlit UI code MODIFIED slightly for new answer source)
if __name__ == "__main__":
    try: import streamlit as st
    except ImportError: logger.error("Streamlit not installed."); sys.exit(1)

    st.set_page_config(layout="wide", page_title="Multi-Document Q&A with RAG v5.1")
    st.title("📄 Multi-Document Question Answering System (v5.1)")
    st.caption("Query PDFs/PPTX (Text RAG) & XLSX/CSV (Structured Query Only + LLM Framing) using Embeddings, FAISS, Pandas, Groq LLM")

    @st.cache_resource
    def load_rag_system(config_path="config.ini"):
        if not os.getenv("GROQ_API_KEY"): st.error("🚨 GROQ_API_KEY not set!"); st.stop()
        try: return RAGSystem(config_path=config_path)
        except Exception as e: st.error(f"Fatal RAG Init Error: {e}"); logger.error("RAG Init Error", exc_info=True); st.stop()

    def process_documents(_rag_system):
        status_messages = []; status_placeholder = st.empty()
        def streamlit_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
            log_prefix = "✅ " if is_done else "❌ " if is_error else "⚠️ " if is_warning else f"⏳ ({stage}) " if stage else "⏳ "; progress_val = 1.0 if is_done else 0.0
            if total_steps and current_step is not None:
                 # Adjust progress weights based on stages
                 stage_prog = {"Extracting Text": 0.1, "Chunking Text": 0.2, "Verifying Index": 0.3, "Embedding Text": 0.4, "Indexing Text": 0.6, "Loading DataFrame": 0.8}.get(stage, 0.0) # DataFrame load is later stage now
                 progress_val = min(1.0, (current_step + stage_prog) / total_steps) if total_steps > 0 else 0.0
            full_message = f"{log_prefix}{message}"; status_messages.append(full_message)
            try:
                if progress_val > 0 and not is_error and not is_warning : status_placeholder.progress(progress_val)
                status_placeholder.caption(full_message) # Always update caption
            except Exception as e: logger.warning(f"Streamlit UI update error: {e}")

        st.header("📚 Document Processing")
        st.caption(f"Data: `{os.path.abspath(_rag_system.config.data_dir)}` | Index: `{os.path.abspath(_rag_system.config.index_dir)}`")
        processing_successful = False
        with st.spinner("Initializing and processing documents..."):
            try:
                # We pass the callback directly to process_files
                processing_successful = _rag_system.process_files(progress_callback=streamlit_progress_callback)
                indexed_count = len(_rag_system.faiss_indexes) # PDF/PPTX only
                dataframe_count = len(_rag_system.dataframes) # CSV/XLSX only
                processed_count = len(_rag_system.processed_files) # Total files with either index or DF

                final_message = f"Processed {processed_count} file(s). Text Indices (PDF/PPTX): {indexed_count}, DataFrames (CSV/XLSX): {dataframe_count}."
                if processed_count > 0:
                    status_placeholder.success(f"✅ Ready! {final_message}")
                elif processing_successful: # Completed but nothing loaded
                     status_placeholder.warning(f"⚠️ Processing finished, but no documents loaded. {final_message}")
                else: # Should not happen if process_files returns bool correctly
                    status_placeholder.error(f"❌ Processing failed. Check logs.")
            except Exception as e:
                 logger.error(f"Fatal error during processing call: {e}", exc_info=True); status_placeholder.error(f"❌ Fatal error: {e}. Check logs.")

        return len(_rag_system.processed_files) > 0, sorted(list(_rag_system.processed_files))


    try:
        rag_sys = load_rag_system()

        if st.sidebar.button("🔄 Re-process Documents & Reload System"):
            st.cache_resource.clear()
            st.rerun()

        is_ready, processed_files_list = process_documents(rag_sys)

        if is_ready and processed_files_list:
            st.sidebar.success(f"Processed Documents ({len(processed_files_list)}):")
            with st.sidebar.expander("Show Processed Files & Type"):
                 for fname in processed_files_list:
                     tags = []
                     if fname in rag_sys.faiss_indexes: tags.append("Text Index") # PDF/PPTX
                     if fname in rag_sys.dataframes: tags.append("DataFrame") # CSV/XLSX
                     st.caption(f"- {fname} [{', '.join(tags)}]")

            st.sidebar.info(f"LLM: `{rag_sys.config.llm_model}`")
            st.sidebar.info(f"Retrieval K (Text): `{rag_sys.config.k_retrieval}`")
            st.sidebar.info(f"DataFrame Query Confidence: `{rag_sys.config.dataframe_query_confidence_threshold}`")


            st.header("💬 Ask a Question")
            user_query = st.text_input("Enter your query (about text in PDF/PPTX or structured data in CSV/XLSX):", key="query_input")

            if user_query and st.button("Get Answer", key="submit_query"):
                query_start_time = time.time()
                results_data = rag_sys.run_query(user_query, context_strategy="top_k") # Calls the modified run_query
                query_end_time = time.time()

                answer_source_display = results_data.get('answer_source', 'N/A')
                st.subheader(f"💡 Answer (via {answer_source_display})")
                answer = results_data.get("answer")

                if results_data.get("status") == "Completed Successfully":
                     st.markdown(answer) # Use markdown for potential formatting in framed answers
                elif results_data.get("status") == "Completed with Framing Failure":
                     st.warning("Could not generate a natural language answer, showing raw data result:")
                     st.markdown(answer) # Show the fallback raw result format
                elif "Failed" in results_data.get("status", ""):
                     st.error(f"Failed to generate answer. Error: {results_data.get('error', 'Unknown error')}")
                else: # e.g., No relevant info found
                     st.info(answer or "No answer generated.")


                # Display Query Analysis & Execution Info
                st.subheader("📊 Query Analysis & Execution Path")
                col1, col2 = st.columns(2)
                with col1:
                    st.caption(f"Query Classification: `{results_data.get('classification', 'N/A')}`")
                    st.caption(f"Confidence: `{results_data.get('classification_confidence', 0.0):.2f}`")
                    if results_data.get("sub_queries"):
                        with st.expander("Text Sub-queries used for retrieval (PDF/PPTX)"):
                            for i, sq in enumerate(results_data["sub_queries"]): st.code(f"{i+1}. {sq}", language=None)
                with col2:
                     st.caption(f"Final Answer Source: `{answer_source_display}`")
                     if results_data.get("target_dataframe"):
                         st.caption(f"Target DataFrame File: `{results_data.get('target_dataframe')}`")
                     # Expander for executed code is useful for DataFrame path
                     if results_data.get("pandas_code_used"):
                         with st.expander("Pandas Code Executed"):
                             st.code(results_data["pandas_code_used"], language='python')
                     # Expander for raw result might also be useful if framing was complex
                     if results_data.get("dataframe_result"):
                         with st.expander("Raw Result from Pandas Code"):
                             st.text(results_data["dataframe_result"])


                # Display Supporting Evidence (Context / Retrieval / DF Info)
                st.subheader("🔍 Supporting Evidence / Context")

                # Show DataFrame Info if used
                if results_data.get("answer_source", "").startswith("DataFrame") and results_data.get("target_dataframe"):
                     df_name = results_data.get("target_dataframe")
                     if df_name in rag_sys.dataframes:
                         df_preview = rag_sys.dataframes[df_name].head()
                         with st.expander(f"Preview of DataFrame Used: '{df_name}'"):
                             st.dataframe(df_preview, use_container_width=True)

                # Show Text RAG Context if used
                agg_context_data = results_data.get("aggregated_context_data", {})
                combined_context = agg_context_data.get("combined_context", "")
                context_source_files = agg_context_data.get("source_files", [])
                if combined_context and results_data.get("answer_source") == "Text RAG": # Only show if Text RAG was the final source
                     with st.expander(f"Text Context Used (PDF/PPTX) from: {', '.join(context_source_files)}"):
                         st.text_area("Combined Text Context Sent to LLM", combined_context, height=250, key="context_display")
                elif results_data.get("answer_source", "").startswith("DataFrame"):
                     st.info("Answer generated via direct DataFrame query. Text context (if any) was primarily used for target identification, not final answer generation.")


                # Show Text Retrieval Results (Useful for seeing what text was considered for DF identification or for Text RAG path)
                retrieval_results = results_data.get("retrieval_results", {})
                if retrieval_results:
                     expander_title = "Retrieved Text Chunks (PDF/PPTX)"
                     if results_data.get("answer_source", "").startswith("DataFrame"):
                          expander_title += " - Used for potential DataFrame identification"
                     else:
                           expander_title += " - Used for Text RAG context"

                     with st.expander(expander_title):
                         display_chunks = []
                         for file_name, res_list in retrieval_results.items():
                              for res in res_list:
                                   display_chunks.append({**res, 'file_name': file_name})
                         display_chunks.sort(key=lambda x: x['score'])

                         if not display_chunks:
                              st.caption("No relevant text chunks found in PDF/PPTX files.")
                         else:
                              for i, res in enumerate(display_chunks[:rag_sys.config.k_retrieval * 2]): # Show more for inspection
                                   source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown'); content_type = res.get('type', 'unknown'); score = res['score']
                                   source_display_parts = [f"File: {file_name}"]
                                   slide_info = source_info.get('slide', 'N/A'); merged_from = source_info.get('merged_from')
                                   if file_type == 'pptx': slide_display = f"Slide: {slide_info}" + (f" (merged from {merged_from})" if merged_from else ""); source_display_parts.append(slide_display)
                                   elif file_type == 'pdf': source_display_parts.append(f"Page: {source_info.get('page', 'N/A')}")
                                   # xlsx/csv should not appear here
                                   source_display_parts.append(f"Type: {content_type}")

                                   st.markdown(f"**Chunk {i+1} (Score: {score:.4f})**")
                                   st.caption(" | ".join(source_display_parts))
                                   st.text(f"{res['content'][:350]}...")
                                   if i < len(display_chunks) -1 : st.divider()
                elif results_data.get("answer_source") == "Text RAG":
                      st.info("No specific text document chunks were retrieved for this query.")


                st.caption(f"Query processed in {query_end_time - query_start_time:.2f} seconds.")

        elif not processed_files_list:
            st.warning(f"No documents processed successfully (No Text Index for PDF/PPTX or DataFrame for CSV/XLSX).")
            st.info("Add files to the 'Data' directory and click 'Re-process Documents'.")
        else: # is_ready is False but process_documents didn't raise error (shouldn't happen?)
             st.error("System processed files but is not marked as ready. Check logs.")

    except Exception as e:
        st.error(f"Streamlit App Error: {e}")
        logger.exception("Streamlit application error:")
        st.info("Check console or `rag_system.log`.")

# --- END OF STREAMLIT UI CODE ---