# --- START OF FILE ragtest_c2-7_modified.py ---

import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd
import faiss
# import streamlit as st # Import only if needed / within UI code (Done in __main__)
import time
from sentence_transformers import SentenceTransformer
# from transformers import T5ForConditionalGeneration, T5Tokenizer # Keep if needed elsewhere
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage # Added AIMessage
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder # For history management
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import io
from contextlib import redirect_stdout
import sqlite3 # Added for SQL database
import requests # Added for DeepSeek API

# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted from RAG2.py + SQL Integration + Chat History + PDF Structure Awareness) ---

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system_sql_chat_pdfstruct.log"), # Log to a different file maybe
        # logging.StreamHandler() # Uncomment for console logs
    ]
)
logger = logging.getLogger(__name__)


class RAGConfig:
    """Configuration class for RAG system."""
    def __init__(self, config_path="config.ini"):
        self.config = configparser.ConfigParser()
        # Added DEEPSEEK_API_KEY reference and SQLITE_DB path
        self.defaults = {
            "PATHS": {
                "data_dir": "./Data/",
                "index_dir": "./Faiss_index/",
                "log_file": "rag_system_sql_chat_pdfstruct.log",
                "sqlite_db_path": ":memory:" # Default to in-memory SQLite DB
            },
            "MODELS": {
                "encoder_model": "sentence-transformers/all-MiniLM-L6-v2",
                "llm_model": "meta-llama/llama-4-scout-17b-16e-instruct", # Updated known good Groq model
                "device": "auto",
                "deepseek_model": "deepseek-reasoner" # Updated DeepSeek model (verify availability)
            },
            "PARAMETERS": {
                "chunk_size": "250", # Slightly larger chunk for proposals?
                "overlap": "60",
                "k_retrieval": "6", # Retrieve slightly more context?
                "temperature": "0.1", # Lower temp for factual proposal answers
                "max_context_tokens": "6000", # Increase context size for LLAMA3-70b
                "max_chars_per_element": "1200",
                "pptx_merge_threshold_words": "50",
                "dataframe_query_confidence_threshold": "0.70", # Slightly lower maybe ok with better context
                "max_chat_history_turns": "5" # Max previous User/AI pairs to consider
            },
            "SUPPORTED_EXTENSIONS": {
                "extensions": ".pdf, .xlsx, .csv, .pptx"
            },
            "API_KEYS": { # Section for API keys (optional, can still rely purely on .env)
                 "deepseek_api_key_config": "" # Can be set here or via .env
            }
        }
        # --- Load/Save logic (unchanged) ---
        if os.path.exists(config_path):
            try: self.config.read(config_path); logger.info(f"Loaded configuration from {config_path}"); self._ensure_defaults()
            except Exception as e: logger.error(f"Error reading config file {config_path}: {e}"); self._set_defaults()
        else: logger.warning(f"Config file {config_path} not found, using defaults"); self._set_defaults()

    def _set_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section)
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value)

    def _ensure_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section); logger.info(f"Added missing section [{section}] from defaults.")
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value); logger.info(f"Added missing option '{option} = {value}' to section [{section}] from defaults.")

    # --- Property definitions (updated/added) ---
    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def sqlite_db_path(self): return self.config.get("PATHS", "sqlite_db_path", fallback=":memory:")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    @property
    def llm_model(self): return self.config.get("MODELS", "llm_model")
    @property
    def deepseek_model(self): return self.config.get("MODELS", "deepseek_model", fallback="deepseek-reasoner")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto")
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=250)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=60)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=6)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.1)
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=6000)
    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1200)
    @property
    def pptx_merge_threshold_words(self): return self.config.getint("PARAMETERS", "pptx_merge_threshold_words", fallback=50)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf, .xlsx, .csv, .pptx") # Added pptx to fallback
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])
    @property
    def dataframe_query_confidence_threshold(self):
        return self.config.getfloat("PARAMETERS", "dataframe_query_confidence_threshold", fallback=0.70)
    @property
    def max_chat_history_turns(self):
        return self.config.getint("PARAMETERS", "max_chat_history_turns", fallback=5)
    @property
    def deepseek_api_key(self):
        # Prioritize environment variable, then config file
        key = os.getenv("DEEPSEEK_API_KEY")
        if not key:
            key = self.config.get("API_KEYS", "deepseek_api_key_config", fallback=None)
            if key: logger.info("Using DeepSeek API key from config file.")
        # else: logger.info("Using DeepSeek API key from environment variable.") # Can be noisy
        return key


class RAGSystem:
    """Retrieval-Augmented Generation system supporting text RAG, SQL querying, chat history, and PDF structure awareness."""
    def __init__(self, config_path="config.ini"):
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system (SQL+Chat+PDFStruct Version). Device: {self.config.device}, Supported files: {self.config.supported_extensions}")
        self._is_streamlit = "streamlit" in sys.modules
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)

        # --- Embedding Model Loading ---
        try:
            # Check if running in streamlit before importing and using st
            if self._is_streamlit:
                import streamlit as st
                with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."):
                    self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            else:
                logger.info(f"Loading embedding model ({self.config.encoder_model})...")
                self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            logger.info("Local embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            if self._is_streamlit:
                import streamlit as st
                st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}")
                st.stop() # Stop Streamlit if embedding model fails critical load
            else: raise

        # --- State Variables ---
        self.file_chunks = {} # Stores text chunks for FAISS {filename: [chunk_dict, ...]}
        self.faiss_indexes = {} # Stores FAISS indexes {filename: faiss_index}
        self.processed_files = set() # Tracks files processed for text RAG *or* SQL

        # --- SQL Database and Metadata Stores ---
        self.db_conn = None
        self.db_cursor = None
        self.table_metadata = {} # { "table_name": {"columns": [...], "source_file": "...", "source_sheet": "...", "row_count": N} }
        self.file_to_table_map = {} # { "filename.xlsx": ["table_sheet1", "table_sheet2"], ... }

        self._llm = None # Groq LLM client cache

        # --- API Key Checks ---
        self.groq_api_key = os.getenv("GROQ_API_KEY")
        self.deepseek_api_key = self.config.deepseek_api_key # Get key via config property
        if not self.groq_api_key: logger.warning("GROQ_API_KEY not found in environment variables."); # UI will show error later if needed
        if not self.deepseek_api_key: logger.warning("DEEPSEEK_API_KEY not found in environment variables or config file. Spreadsheet metadata generation will fail.")

        # --- Connect to SQLite DB ---
        self._connect_db()

    # --- DB Connection Methods (Unchanged) ---
    def _connect_db(self):
        """Establish connection to the SQLite database."""
        try:
            db_path = self.config.sqlite_db_path
            self.db_conn = sqlite3.connect(db_path, check_same_thread=False) # Allow access from different threads (e.g., Streamlit)
            self.db_cursor = self.db_conn.cursor()
            logger.info(f"Connected to SQLite database at '{db_path}'")
        except sqlite3.Error as e:
            logger.error(f"Error connecting to SQLite database at '{self.config.sqlite_db_path}': {e}", exc_info=True)
            self.db_conn = None
            self.db_cursor = None
            if self._is_streamlit:
                import streamlit as st
                st.error(f"Fatal Error: Failed to connect to SQLite DB. Check logs. Error: {e}")
                st.stop() # Stop Streamlit if DB connection fails
            else:
                raise ConnectionError(f"Failed to connect to SQLite DB: {e}") from e

    def _close_db(self):
        """Close the SQLite database connection."""
        if self.db_conn:
            try:
                self.db_conn.commit() # Commit any pending changes
                self.db_conn.close()
                logger.info("Closed SQLite database connection.")
                self.db_conn = None
                self.db_cursor = None
            except sqlite3.Error as e:
                logger.error(f"Error closing SQLite connection: {e}", exc_info=True)

    # Make sure to close the DB when the object is destroyed
    def __del__(self):
        self._close_db()

    # --- LLM Initialization (Groq) --- (Unchanged)
    def _get_llm(self):
        if self._llm is None:
            if not self.groq_api_key:
                logger.error("GROQ_API_KEY not found.")
                if self._is_streamlit: import streamlit as st; st.error("Error: GROQ_API_KEY is not set.")
                raise ValueError("GROQ_API_KEY not configured.")
            try:
                self._llm = ChatGroq(
                    temperature=self.config.temperature,
                    model_name=self.config.llm_model,
                    groq_api_key=self.groq_api_key,
                )
                logger.info(f"Groq LLM client initialized ({self.config.llm_model})")
            except Exception as e:
                 logger.error(f"Failed to initialize Groq client: {e}", exc_info=True)
                 if self._is_streamlit: import streamlit as st; st.error(f"Failed to initialize LLM. Error: {e}")
                 raise
        return self._llm

    # --- Text Cleaning and File/Table Naming (Unchanged) ---
    def clean_text(self, text):
        if not isinstance(text, str): text = str(text)
        text = re.sub(r"\(cid:.*?\)", "", text) # Remove common PDF artifact
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text) # Remove most control characters except newline/tab
        text = re.sub(r"\s+", " ", text).strip() # Normalize whitespace
        # Keep newlines potentially for structure, but replace multiple newlines
        # text = text.replace('\r', '')
        # text = re.sub(r'\n\s*\n', '\n', text) # Replace multiple newlines with single
        return text # Keep original newlines for now, handled in chunking/context

    def _get_safe_filename(self, file_name):
        base_name = os.path.splitext(file_name)[0]
        return re.sub(r'[^\w\.-]', '_', base_name)

    def _get_safe_table_name(self, file_name, sheet_name=None):
        base = self._get_safe_filename(file_name)
        if sheet_name:
             safe_sheet = re.sub(r'[^\w]', '_', sheet_name)
             name = f"{base}_{safe_sheet}"
        else:
             name = base
        if not re.match(r"^[a-zA-Z_]", name): name = "_" + name
        # Limit length for SQLite compatibility
        name = name[:60]; return name.lower()

    def _table_exists(self, table_name):
        if not self.db_cursor: return False
        try: self.db_cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name=?;", (table_name,)); return self.db_cursor.fetchone() is not None
        except sqlite3.Error as e: logger.error(f"Error checking if table '{table_name}' exists: {e}"); return False

    def get_index_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.index")
    def get_embedding_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.npy")
    def get_chunks_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.json")

    # --- Content Extractors (PDF Modified, PPTX, XLSX, CSV Unchanged) ---

    # MODIFIED _extract_pdf for Structure Awareness
    def _extract_pdf(self, file_path):
        all_content = []
        base_filename = os.path.basename(file_path)
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Extracting text/tables from PDF: {base_filename} ({total_pages} pages) - Applying structural heuristics.")
                page_bar = None
                if self._is_streamlit and total_pages > 5: # Show progress bar for more than 5 pages
                    import streamlit as st
                    page_bar = st.progress(0, text=f"Extracting pages from {base_filename}...")

                # Pre-compile regex for efficiency
                toc_pattern = re.compile(r"(\.|\s){4,}\s*\d+\s*$") # Detects lines ending '.... page_num'
                intro_keywords = re.compile(r'\b(introduction|about sat ?sure|company profile|executive summary|overview|problem statement|client need|challenge)\b', re.IGNORECASE)
                closing_keywords = re.compile(r'\b(appendix|annexure|pricing|cost|timeline|schedule|next steps|conclusion)\b', re.IGNORECASE)

                for page_num, page in enumerate(pdf.pages):
                    current_page = page_num + 1
                    text = page.extract_text(layout="normal") or "" # Use layout preserving extraction
                    cleaned_text = self.clean_text(text) # Basic cleaning
                    page_word_count = len(cleaned_text.split())

                    # --- Structural Heuristics ---
                    potential_section = "main_content" # Default
                    is_first_page = (page_num == 0)
                    is_likely_toc_page = (page_num == 1 or page_num == 2) # Assume ToC often on page 2 or 3
                    is_near_end = (page_num >= total_pages - 3) # Heuristic for last few pages

                    if is_first_page:
                        potential_section = "title_page"
                        # Simple title extraction: longest line among the first few non-empty lines
                        potential_title = ""
                        lines = [line.strip() for line in text.split('\n') if line.strip()]
                        if lines:
                            first_few_lines = lines[:7] # Look at first 7 lines
                            if first_few_lines:
                                potential_title = max(first_few_lines, key=len)
                                # Basic check if it seems like a title (e.g., not too long, some caps?)
                                if len(potential_title) > 80 or potential_title.islower():
                                    potential_title = first_few_lines[0] # Fallback to first line
                        logger.debug(f"PDF Pg {current_page}: Marked as 'title_page'. Potential Title: '{potential_title[:50]}...'")

                    elif is_likely_toc_page:
                        # More robust ToC detection: look for multiple lines matching the pattern
                        lines = text.split('\n')
                        toc_lines_count = sum(1 for line in lines if toc_pattern.search(line.strip()))
                        if toc_lines_count >= 3 and page_word_count < 300 : # Require >=3 ToC-like lines and not too wordy page
                            potential_section = "table_of_contents"
                            logger.debug(f"PDF Pg {current_page}: Marked as 'table_of_contents' (Found {toc_lines_count} matching lines).")
                        # Fall through to check intro keywords if not clearly ToC

                    # Check for intro keywords if not title/ToC and early in doc
                    if potential_section == "main_content" and page_num < 5:
                        if intro_keywords.search(cleaned_text):
                            potential_section = "introduction_overview"
                            logger.debug(f"PDF Pg {current_page}: Marked as 'introduction_overview' based on keywords.")
                        elif page_num < 3: # Label pages 2/3 if not ToC/Intro
                             potential_section = "front_matter"

                    # Check for closing keywords if near the end
                    if potential_section == "main_content" and is_near_end:
                         if closing_keywords.search(cleaned_text):
                             potential_section = "closing_appendix"
                             logger.debug(f"PDF Pg {current_page}: Marked as 'closing_appendix' based on keywords.")
                         elif page_num == total_pages -1: # Mark very last page if not keyword identified
                             potential_section = "final_page"


                    # --- Store Content ---
                    source_info = {"page": current_page}
                    if is_first_page and potential_title:
                         source_info["potential_doc_title"] = potential_title

                    if cleaned_text:
                        all_content.append({
                            "type": "text",
                            "content": cleaned_text, # Store cleaned text
                            "source_info": source_info,
                            "file_type": "pdf",
                            "potential_section": potential_section # Store the determined section
                        })

                    # Extract tables (unchanged logic, but add section tag)
                    try:
                        for table_idx, table_content in enumerate(page.extract_tables()): # Changed variable name for clarity
                             if table_content:
                                 # Clean table representation
                                 table_df = pd.DataFrame(table_content)
                                 # Attempt to promote first row to header if it looks like one
                                 if not table_df.empty and not pd.api.types.is_numeric_dtype(table_df.iloc[0].dropna()):
                                     try:
                                         table_df.columns = table_df.iloc[0].fillna('').astype(str)
                                         table_df = table_df[1:]
                                     except Exception: pass # Ignore if header setting fails
                                 table_df = table_df.fillna('')
                                 table_string = table_df.to_string(index=False, header=True)
                                 cleaned_table_string = self.clean_text(table_string)

                                 if cleaned_table_string:
                                     all_content.append({
                                        "type": "table",
                                        "content": cleaned_table_string,
                                        "source_info": {**source_info, "table_index_on_page": table_idx + 1}, # Add table index
                                        "file_type": "pdf",
                                        "potential_section": potential_section # Assign same section as page
                                     })
                    except Exception as e_table:
                        logger.warning(f"Could not extract/process table on page {current_page} in {base_filename}: {e_table}")

                    if page_bar:
                        progress_percent = min(1.0, (page_num + 1) / total_pages)
                        page_bar.progress(progress_percent, text=f"Extracting page {current_page}/{total_pages} from {base_filename}...")

                if page_bar:
                    page_bar.empty() # Clear progress bar

            logger.info(f"Extracted {len(all_content)} text/table content blocks from PDF: {base_filename} with structural tagging.")
            return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PDF {base_filename}: {e}", exc_info=True)
            return []


    def _extract_pptx(self, file_path):
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element; merge_threshold_words = self.config.pptx_merge_threshold_words
        try:
            prs = Presentation(file_path); logger.info(f"Extracting text from PPTX: {base_filename} ({len(prs.slides)} slides)")
            slide_bar = None; total_slides = len(prs.slides)
            # Use self._is_streamlit check before using st
            if self._is_streamlit and total_slides > 1: import streamlit as st; slide_bar = st.progress(0, text=f"Extracting slides from {base_filename}...")
            pending_title_slide_data = None
            for i, slide in enumerate(prs.slides):
                 current_slide_number = i + 1; current_slide_title_text = ""; current_slide_other_texts = []; current_slide_has_title_placeholder = False; title_shape = None
                 try:
                    if slide.shapes.title: title_shape = slide.shapes.title; current_slide_has_title_placeholder = True
                 except AttributeError: pass # Some slides might not have a title placeholder by default
                 if title_shape and title_shape.has_text_frame: cleaned_title = self.clean_text(title_shape.text_frame.text); current_slide_title_text = cleaned_title if cleaned_title else ""
                 for shape in slide.shapes:
                     if shape == title_shape: continue # Already processed title
                     is_placeholder = shape.is_placeholder; is_body_placeholder = False
                     if is_placeholder:
                         try: ph_type = shape.placeholder_format.type; is_body_placeholder = ph_type in [MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT, MSO_SHAPE_TYPE.SUBTITLE, MSO_SHAPE_TYPE.CONTENT, MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.PICTURE] # Added more placeholder types that can contain text or imply body content
                         except AttributeError: pass
                     if shape.has_text_frame:
                         text = shape.text_frame.text; cleaned = self.clean_text(text)
                         if cleaned:
                              if len(cleaned) > max_chars: cleaned = cleaned[:max_chars] + "...(truncated shape)"
                              prefix = "[Body]: " if is_body_placeholder and not current_slide_title_text and not current_slide_other_texts else "" # Add prefix if it's the main body text and no other text collected yet from shapes
                              current_slide_other_texts.append(prefix + cleaned)
                 if slide.has_notes_slide:
                    try:
                         notes_text = slide.notes_slide.notes_text_frame.text; cleaned_notes = self.clean_text(notes_text)
                         if cleaned_notes:
                             if len(cleaned_notes) > max_chars * 2: cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated notes)"
                             current_slide_other_texts.append(f"[Notes]: {cleaned_notes}")
                    except Exception as e_notes: logger.warning(f"Could not extract notes from slide {current_slide_number}: {e_notes}")

                 current_slide_full_content = current_slide_title_text
                 if current_slide_other_texts: current_slide_full_content += ("\n" if current_slide_full_content else "") + "\n".join(current_slide_other_texts)
                 current_slide_full_content = current_slide_full_content.strip(); other_text_word_count = sum(len(s.split()) for s in current_slide_other_texts) # Word count of non-title texts on this slide

                 merged_content_block = None; should_merge = False
                 # Logic to decide if current slide should be merged with a PENDING title slide
                 if pending_title_slide_data:
                     # Condition: Pending was a title slide (had placeholder, its 'other_words' were low) AND current slide has content
                     if pending_title_slide_data['has_title'] and pending_title_slide_data['other_words'] <= merge_threshold_words and current_slide_full_content:
                         should_merge = True; logger.info(f"Merging slide {pending_title_slide_data['number']} (title-like) with content from slide {current_slide_number}.")

                 if should_merge:
                    # Create a single merged block
                    merged_text = f"[Title from Slide {pending_title_slide_data['number']}]: {pending_title_slide_data['title']}\n\n[Content from Slide {pending_title_slide_data['number']} (if any)]:\n{pending_title_slide_data['content_without_title']}\n\n---\n\n[Content from Slide {current_slide_number}]:\n{current_slide_full_content}"
                    merged_content_block = {"type": "slide_text_merged", "content": merged_text.strip(), "source_info": {"slide_title": pending_title_slide_data['number'], "slide_content": current_slide_number}, "file_type": "pptx"}
                    all_content.append(merged_content_block); pending_title_slide_data = None # Clear pending
                 else:
                    # If no merge, process any PENDING slide first
                    if pending_title_slide_data:
                         if pending_title_slide_data['content']: # If pending had any content
                             all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
                         pending_title_slide_data = None # Clear pending

                    # Now, evaluate current slide: IF it looks like a new title slide (has placeholder, low 'other_words'), make IT pending. ELSE, add it directly.
                    if current_slide_has_title_placeholder and current_slide_title_text and other_text_word_count <= merge_threshold_words and current_slide_full_content:
                        logger.debug(f"Slide {current_slide_number} ('{current_slide_title_text[:30]}...') is potential title slide for next content. Holding.");
                        content_without_title_parts = [txt for txt in current_slide_other_texts if not txt.startswith(current_slide_title_text)] # Attempt to get content other than title text if title was also in other_texts
                        content_w_o_title_str = "\n".join(content_without_title_parts).strip()

                        pending_title_slide_data = {
                            "content": current_slide_full_content, # Full content of this slide
                            "content_without_title": content_w_o_title_str,
                            "number": current_slide_number,
                            "has_title": current_slide_has_title_placeholder,
                            "other_words": other_text_word_count, # Words other than title on THIS slide
                            "title": current_slide_title_text
                        }
                    else: # Not a title slide, or has too much other content to be just a title for next
                        if current_slide_full_content:
                            all_content.append({"type": "slide_text", "content": current_slide_full_content, "source_info": {"slide": current_slide_number}, "file_type": "pptx"})

                 if slide_bar: slide_bar.progress(min(1.0, (i + 1) / total_slides), text=f"Extracting slide {current_slide_number}/{total_slides}...")
            # After loop, if there's still a pending title slide (e.g., last slide was a title slide)
            if pending_title_slide_data:
                 logger.debug(f"Processing final pending title slide {pending_title_slide_data['number']} at end.")
                 if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})

            if slide_bar: slide_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PPTX {base_filename} (Merge strategy applied)."); return all_content
        except Exception as e: logger.error(f"Error extracting content from PPTX {base_filename}: {e}", exc_info=True); return []

    # --- XLSX/CSV Extraction and Loading (Unchanged) ---
    def _extract_and_load_xlsx(self, file_path, file_name, progress_callback=None):
        base_filename = os.path.basename(file_path)
        sheet_load_success_count = 0; sheet_metadata_success_count = 0; sheet_sql_load_success_count = 0; processed_sheet_tables = []
        try:
            excel_file = pd.ExcelFile(file_path); sheet_names = excel_file.sheet_names
            logger.info(f"Processing XLSX: {base_filename} (Sheets: {len(sheet_names)})")
            if not sheet_names: logger.warning(f"No sheets found in {base_filename}"); return False, []
            sheet_bar = None; total_sheets = len(sheet_names)
            # Use self._is_streamlit check before using st
            if self._is_streamlit and total_sheets > 1: import streamlit as st; sheet_bar = st.progress(0, text=f"Processing sheets in {base_filename}...")
            for i, sheet_name in enumerate(sheet_names):
                stage_msg = f"Processing sheet '{sheet_name}' ({i+1}/{total_sheets}) in {base_filename}"
                logger.info(stage_msg)
                if progress_callback: progress_callback(f"{stage_msg} - Loading...", stage="Loading Sheet")
                if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.1) / total_sheets), text=f"Loading '{sheet_name}'...")
                try:
                    df = excel_file.parse(sheet_name); df = df.fillna('')
                    # Clean column names robustly
                    df.columns = [re.sub(r'[^a-zA-Z0-9_]', '_', str(col)).strip('_') for col in df.columns]
                    df.columns = [f"col_{j}" if not name else name for j, name in enumerate(df.columns)] # Ensure no empty names
                    # Handle duplicate column names
                    cols = pd.Series(df.columns)
                    for dup_idx, dup_val in cols[cols.duplicated()].items(): # Iterate over (index, value)
                         indices = cols[cols == dup_val].index.tolist()
                         for k_idx, col_idx in enumerate(indices):
                             if k_idx > 0: # only append suffix to duplicates, not the first one
                                 cols[col_idx] = f"{dup_val}_{k_idx}"
                    df.columns = cols

                    if df.empty: logger.warning(f"Sheet '{sheet_name}' in {base_filename} is empty. Skipping."); continue
                    sheet_load_success_count += 1; logger.info(f"Loaded DataFrame from sheet '{sheet_name}' (Rows: {len(df)}, Cols: {len(df.columns)})")
                    table_name = self._get_safe_table_name(file_name, sheet_name)
                    if not table_name: logger.error(f"Could not generate table name for sheet '{sheet_name}'. Skipping SQL load."); continue
                    if progress_callback: progress_callback(f"{stage_msg} - Analyzing Columns (DeepSeek)...", stage="Analyzing Columns")
                    if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.4) / total_sheets), text=f"Analyzing '{sheet_name}'...")
                    column_metadata = self._get_column_metadata_deepseek(df, file_name, sheet_name)
                    if not column_metadata: logger.warning(f"Failed to get column metadata for sheet '{sheet_name}'. Proceeding without descriptions."); column_metadata = [{"name": col, "type": self._infer_sql_type(df[col]), "description": "N/A"} for col in df.columns]
                    else: sheet_metadata_success_count +=1
                    if progress_callback: progress_callback(f"{stage_msg} - Loading to SQL Database...", stage="Loading SQL")
                    if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.8) / total_sheets), text=f"Storing '{sheet_name}'...")
                    sql_loaded = self._load_df_to_sql(df, table_name, column_metadata, file_name, sheet_name)
                    if sql_loaded: sheet_sql_load_success_count += 1; processed_sheet_tables.append(table_name); logger.info(f"Successfully loaded sheet '{sheet_name}' into SQL table '{table_name}'")
                    else: logger.error(f"Failed to load sheet '{sheet_name}' into SQL.")
                except Exception as e_sheet:
                    logger.error(f"Error processing sheet '{sheet_name}' in {base_filename}: {e_sheet}", exc_info=True);
                    if progress_callback: progress_callback(f"⚠️ Error processing sheet '{sheet_name}': {e_sheet}", is_warning=True) # Pass warning flag
                if sheet_bar: sheet_bar.progress(min(1.0, (i + 1) / total_sheets), text=f"Finished sheet {i+1}/{total_sheets}...")
            if sheet_bar: sheet_bar.empty()
            logger.info(f"Finished processing XLSX {base_filename}. Sheets Loaded: {sheet_load_success_count}/{total_sheets}, Metadata OK: {sheet_metadata_success_count}, SQL OK: {sheet_sql_load_success_count}")
            return sheet_sql_load_success_count > 0, processed_sheet_tables
        except Exception as e:
            logger.error(f"Error opening or processing XLSX file {base_filename}: {e}", exc_info=True);
            if progress_callback: progress_callback(f"❌ Error processing XLSX {base_filename}: {e}", is_error=True) # Pass error flag
            return False, []

    def _extract_and_load_csv(self, file_path, file_name, progress_callback=None):
        base_filename = os.path.basename(file_path); df = None; load_success = False; metadata_success = False; sql_load_success = False; table_name = None
        try:
            stage_msg = f"Processing CSV: {base_filename}"; logger.info(stage_msg)
            if progress_callback: progress_callback(f"{stage_msg} - Loading...", stage="Loading CSV")
            encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
            for enc in encodings_to_try:
                try:
                    # Try skipping bad lines if decode errors persist even with right encoding
                    df = pd.read_csv(file_path, encoding=enc, low_memory=False, on_bad_lines='warn')
                    df = df.fillna('')
                    # Clean column names robustly
                    df.columns = [re.sub(r'[^a-zA-Z0-9_]', '_', str(col)).strip('_') for col in df.columns]
                    df.columns = [f"col_{j}" if not name else name for j, name in enumerate(df.columns)] # Ensure no empty names
                    # Handle duplicate column names
                    cols = pd.Series(df.columns)
                    for dup_idx, dup_val in cols[cols.duplicated()].items(): # Iterate over (index, value)
                         indices = cols[cols == dup_val].index.tolist()
                         for k_idx, col_idx in enumerate(indices):
                             if k_idx > 0: # only append suffix to duplicates, not the first one
                                 cols[col_idx] = f"{dup_val}_{k_idx}"
                    df.columns = cols

                    logger.info(f"Read CSV {base_filename} using encoding: {enc}. (Rows: {len(df)}, Cols: {len(df.columns)})"); load_success = True; break
                except UnicodeDecodeError: continue
                except Exception as e_read: logger.warning(f"Pandas read_csv error for {base_filename} with encoding {enc}: {e_read}")
            if df is None: logger.error(f"Could not read CSV {base_filename} after trying encodings."); return False, []
            if df.empty: logger.warning(f"CSV {base_filename} is empty."); return False, [] # Changed from error to warning, still skip
            if not load_success: return False, [] # Should be caught by df is None, but safety
            table_name = self._get_safe_table_name(file_name)
            if not table_name: logger.error("Could not generate table name for CSV. Skipping SQL load."); return False, []
            if progress_callback: progress_callback(f"{stage_msg} - Analyzing Columns (DeepSeek)...", stage="Analyzing Columns")
            column_metadata = self._get_column_metadata_deepseek(df, file_name)
            if not column_metadata: logger.warning("Failed to get column metadata for CSV. Proceeding without descriptions."); column_metadata = [{"name": col, "type": self._infer_sql_type(df[col]), "description": "N/A"} for col in df.columns]
            else: metadata_success = True
            if progress_callback: progress_callback(f"{stage_msg} - Loading to SQL Database...", stage="Loading SQL")
            sql_loaded = self._load_df_to_sql(df, table_name, column_metadata, file_name)
            if sql_loaded: sql_load_success = True; logger.info(f"Successfully loaded CSV {base_filename} into SQL table '{table_name}'")
            else: logger.error(f"Failed to load CSV {base_filename} into SQL.")
            return sql_load_success, ([table_name] if sql_load_success else [])
        except Exception as e:
            logger.error(f"Error processing CSV file {base_filename}: {e}", exc_info=True);
            if progress_callback: progress_callback(f"❌ Error processing CSV {base_filename}: {e}", is_error=True) # Pass error flag
            return False, []

    # --- SQL Helper Methods (Unchanged) ---
    def _infer_sql_type(self, series):
        # Prioritize numeric types if possible, even if some strings exist
        num_numeric = pd.to_numeric(series, errors='coerce').notna().sum()
        num_total = len(series)
        if num_total > 0 and num_numeric / num_total > 0.9: # If >90% are numeric
            # Check if all *numeric* values are integers
            numeric_series = pd.to_numeric(series, errors='coerce').dropna()
            if (numeric_series == numeric_series.astype(np.int64)).all(): # Check if all numeric values are integers
                 return "INTEGER"
            return "REAL" # If some are float, treat as REAL

        # Check original dtype if not overwhelmingly numeric
        # This part needs to be careful if original series is object type due to mixed data
        if pd.api.types.is_integer_dtype(series.dtype): return "INTEGER"
        if pd.api.types.is_float_dtype(series.dtype): return "REAL"
        if pd.api.types.is_bool_dtype(series.dtype): return "INTEGER" # Represent bools as 0/1
        if pd.api.types.is_datetime64_any_dtype(series.dtype): return "TEXT" # Store dates/times as TEXT for SQLite compatibility
        return "TEXT" # Default

    def _get_column_metadata_deepseek(self, df, file_name, sheet_name=None):
        if not self.deepseek_api_key: logger.warning("DeepSeek API key not available. Cannot generate column metadata."); return None
        api_url = "https://api.deepseek.com/v1/chat/completions"; headers = {"Content-Type": "application/json", "Authorization": f"Bearer {self.deepseek_api_key}"}
        sample_data_str = df.head(5).to_string(index=False); column_names = list(df.columns)
        prompt = f"""Analyze the following table columns from a spreadsheet ('{file_name}'{f", sheet '{sheet_name}'" if sheet_name else ""}).
The column names are: {column_names}
Here are the first few rows of data:
{sample_data_str}

For each column name, provide:
1. A concise one-sentence description of what the data in the column represents.
2. The most appropriate SQL data type (choose from: TEXT, INTEGER, REAL - use TEXT for dates/datetimes or mixed types for SQLite compatibility).

Respond ONLY with a valid JSON object where keys are the exact column names and values are objects containing 'description' and 'sql_type'. Ensure the JSON is well-formed. Example:
{{
  "ColumnA": {{ "description": "Unique identifier for each record.", "sql_type": "INTEGER" }},
  "ColumnB": {{ "description": "Name of the customer.", "sql_type": "TEXT" }},
  "TransactionDate": {{ "description": "Date when the transaction occurred.", "sql_type": "TEXT" }}
}}
"""
        payload = {"model": self.config.deepseek_model, "messages": [{"role": "system", "content": "You are an expert data analyst specializing in understanding spreadsheet columns and assigning SQLite-compatible SQL types."}, {"role": "user", "content": prompt}], "temperature": 0.1, "max_tokens": 1500, "stream": False}
        try:
            response = requests.post(api_url, headers=headers, json=payload, timeout=90); response.raise_for_status() # Increased timeout
            result = response.json(); raw_content = result['choices'][0]['message']['content']; logger.debug(f"Raw DeepSeek response for metadata: {raw_content}")
            try:
                # Improved JSON extraction (handle potential markdown code blocks)
                json_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw_content, re.DOTALL | re.IGNORECASE)
                if not json_match: json_match = re.search(r"(\{.*?\})", raw_content, re.DOTALL) # Fallback to finding first {} block

                if json_match: json_str = json_match.group(1); metadata_dict = json.loads(json_str)
                else: raise json.JSONDecodeError("Could not find JSON object in response", raw_content, 0)

                formatted_metadata = []; missing_cols = set(column_names); processed_cols = set()

                # Use column_names to ensure order and handle case sensitivity differences
                for expected_col_name in column_names:
                    found_meta = None
                    # Try exact match first, then case-insensitive
                    if expected_col_name in metadata_dict:
                        found_meta = metadata_dict[expected_col_name]
                        del metadata_dict[expected_col_name] # Remove to handle case where LLM returns variations
                    else:
                        # Check remaining keys in metadata_dict case-insensitively
                        found_key_actual = None
                        for resp_col_name_key in list(metadata_dict.keys()): # Iterate over a copy of keys
                             if resp_col_name_key.lower() == expected_col_name.lower():
                                found_meta = metadata_dict[resp_col_name_key]
                                found_key_actual = resp_col_name_key
                                break
                        if found_key_actual:
                            del metadata_dict[found_key_actual] # Remove the found key

                    if found_meta and isinstance(found_meta, dict) and 'description' in found_meta and 'sql_type' in found_meta:
                         sql_type = str(found_meta['sql_type']).upper();
                         # Allow only TEXT, INTEGER, REAL for SQLite
                         if sql_type not in ["TEXT", "INTEGER", "REAL"]: logger.warning(f"DeepSeek proposed invalid SQL type '{sql_type}' for column '{expected_col_name}'. Defaulting to TEXT."); sql_type = "TEXT"
                         formatted_metadata.append({"name": expected_col_name, "description": str(found_meta['description']).strip(), "type": sql_type});
                         processed_cols.add(expected_col_name)
                    else:
                         logger.warning(f"Metadata missing or invalid format for column '{expected_col_name}' in DeepSeek response. Will add basic entry.")
                         # Add basic entry later if still missing

                missing_cols = set(column_names) - processed_cols
                if missing_cols:
                     logger.warning(f"DeepSeek metadata generation missed columns: {missing_cols}. Adding basic entries.")
                     # Add missing columns in the correct order
                     final_metadata_with_missing = []
                     metadata_map_by_name = {m['name']: m for m in formatted_metadata} # Use a map for quick lookup
                     for col_name_in_order in column_names: # Iterate through original df column order
                         if col_name_in_order in metadata_map_by_name:
                             final_metadata_with_missing.append(metadata_map_by_name[col_name_in_order])
                         else: # This column was missed or had issues
                             final_metadata_with_missing.append({"name": col_name_in_order, "description": "N/A - Metadata auto-generation failed or incomplete.", "type": self._infer_sql_type(df[col_name_in_order])})
                     formatted_metadata = final_metadata_with_missing

                logger.info(f"Successfully generated/handled metadata for {len(formatted_metadata)} columns using DeepSeek for '{file_name}'{f'/{sheet_name}' if sheet_name else ''}."); return formatted_metadata
            except json.JSONDecodeError as e_json: logger.error(f"Failed to parse JSON metadata from DeepSeek: {e_json}. Raw response: {raw_content}"); return None
        except requests.exceptions.RequestException as e_req: logger.error(f"DeepSeek API request failed: {e_req}", exc_info=True); return None
        except Exception as e: logger.error(f"Error getting column metadata from DeepSeek: {e}", exc_info=True); return None

    def _load_df_to_sql(self, df, table_name, column_metadata, file_name, sheet_name=None):
        if not self.db_conn or not self.db_cursor: logger.error("Database connection not available. Cannot load table."); return False
        if not table_name: logger.error("Invalid table name provided. Cannot load table."); return False
        if self._table_exists(table_name):
             logger.warning(f"Table '{table_name}' already exists. Replacing.")
             try: self.db_cursor.execute(f'DROP TABLE "{table_name}"') # Use quotes for safety
             except sqlite3.Error as e_drop: logger.error(f"Error dropping existing table {table_name}: {e_drop}"); return False
        try:
            logger.info(f"Loading DataFrame into SQL table '{table_name}'...")
            # Ensure column names in df match cleaned names used for metadata *exactly* before to_sql
            # This assumes metadata generation used the final cleaned names
            original_df_columns = list(df.columns)
            df.columns = [meta['name'] for meta in column_metadata]

            # Ensure data types are compatible before loading (especially for numerics)
            temp_df = df.copy()
            for i, meta in enumerate(column_metadata):
                col_name_meta = meta['name'] # Name from metadata (should be target SQL name)
                # col_name_original_df = original_df_columns[i] # Original name in df before renaming
                sql_type = meta['type']

                # Use the current column name in temp_df which should match col_name_meta
                current_col_in_temp_df = col_name_meta

                if sql_type in ["INTEGER", "REAL"]:
                    # Attempt conversion, set errors to None (which become NULL in SQL)
                    temp_df[current_col_in_temp_df] = pd.to_numeric(temp_df[current_col_in_temp_df], errors='coerce')
                    if sql_type == "INTEGER":
                        # Attempt conversion to nullable integer type if possible
                         try: temp_df[current_col_in_temp_df] = temp_df[current_col_in_temp_df].astype('Int64') # Pandas nullable integer
                         except Exception: pass # Keep as float if conversion fails (e.g., contains NaNs that couldn't be coerced to Int64)
                elif sql_type == "TEXT":
                     # Ensure text columns are actually strings to avoid issues with to_sql for mixed types
                     temp_df[current_col_in_temp_df] = temp_df[current_col_in_temp_df].astype(str)


            temp_df.to_sql(table_name, self.db_conn, if_exists='replace', index=False); self.db_conn.commit(); logger.info(f"DataFrame loaded into table '{table_name}' ({len(df)} rows).")

            # Store metadata after successful load
            self.table_metadata[table_name] = {"columns": column_metadata, "source_file": file_name, "source_sheet": sheet_name, "row_count": len(df)}
            if file_name not in self.file_to_table_map: self.file_to_table_map[file_name] = []
            if table_name not in self.file_to_table_map[file_name]: self.file_to_table_map[file_name].append(table_name)
            return True
        except sqlite3.Error as e_sql:
             logger.error(f"SQLite error loading DataFrame to table '{table_name}': {e_sql}", exc_info=True);
             try: self.db_conn.rollback()
             except: pass
             return False
        except Exception as e:
            logger.error(f"Unexpected error loading DataFrame to table '{table_name}': {e}", exc_info=True);
            try: self.db_conn.rollback()
            except: pass
            return False


    # --- Text Chunking and FAISS Indexing (MODIFIED for structure tag) ---
    def extract_content(self, file_path):
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf': return self._extract_pdf(file_path) # Calls modified version
        elif extension == '.pptx': return self._extract_pptx(file_path)
        elif extension in ['.csv', '.xlsx']: logger.debug(f"Skipping text extraction for {os.path.basename(file_path)} in extract_content, handled separately for SQL."); return []
        else: logger.warning(f"Unsupported file type for text extraction: {file_path}"); return []

    # MODIFIED chunk_content to pass through potential_section
    def chunk_content(self, all_content):
        chunks = []; total_words_estimate = 0; chunk_bar = None;
        if not all_content: return chunks
        try: total_words_estimate = sum(len(str(item.get('content', '')).split()) for item in all_content)
        except Exception: pass
        # Use self._is_streamlit check before using st
        if self._is_streamlit and total_words_estimate > 5000: import streamlit as st; chunk_bar = st.progress(0, text=f"Chunking text content...")
        words_processed = 0
        for item_index, item in enumerate(all_content):
            content = item.get('content', '');
            source_info = item.get('source_info', {}); # Contains page, potential title
            file_type = item.get('file_type', 'unknown');
            content_type = item.get('type', 'unknown')
            potential_section = item.get('potential_section', 'unknown') # Get the new tag

            if not isinstance(content, str): content = str(content)

            # Simple splitting by words for chunk size (could be improved with sentence splitting)
            # We split the *already cleaned* text from the item
            words = content.split();
            if not words: logger.debug(f"Skipping empty text content block {item_index}"); continue

            item_chunks_created = 0
            # Use a sliding window approach
            start_index = 0
            while start_index < len(words):
                end_index = start_index + self.config.chunk_size
                chunk_words = words[start_index:end_index]
                chunk_text = " ".join(chunk_words)

                if chunk_text:
                     chunks.append({
                         "content": chunk_text,
                         "source_info": source_info, # Pass original source info (page, etc.)
                         "file_type": file_type,
                         "type": content_type,
                         "potential_section": potential_section # << ADDED: Carry over section tag
                     })
                     item_chunks_created += 1

                # Move the window forward by chunk_size - overlap
                start_index += (self.config.chunk_size - self.config.overlap)
                if start_index >= len(words): break # Ensure we don't go past the end


            if chunk_bar:
                words_processed += len(words);
                progress_percent = min(1.0, words_processed / total_words_estimate) if total_words_estimate > 0 else 0;
                chunk_bar.progress(progress_percent, text=f"Chunking text content... ({len(chunks)} chunks created)")

        if chunk_bar: chunk_bar.empty()
        logger.info(f"Created {len(chunks)} text chunks from {len(all_content)} content blocks (including section tags).");
        return chunks

    # --- FAISS Indexing/Loading (Unchanged) ---
    def load_faiss_index(self, file_name, embedding_dim):
        index_path = self.get_index_path(file_name)
        if os.path.exists(index_path):
            try: index = faiss.read_index(index_path); logger.info(f"Loaded FAISS index for {file_name} ({index.ntotal} vectors)"); return index
            except Exception as e: logger.error(f"Error reading FAISS index {index_path}: {e}. Creating new index.")
        logger.info(f"FAISS index not found at {index_path}. Will create new."); return faiss.IndexFlatL2(embedding_dim)

    def save_chunks(self, file_name, chunks):
        chunks_path = self.get_chunks_path(file_name)
        try:
            with open(chunks_path, 'w', encoding='utf-8') as f: json.dump(chunks, f, indent=2); logger.info(f"Saved {len(chunks)} text chunks to {chunks_path}")
        except Exception as e: logger.error(f"Error saving text chunks for {file_name} to {chunks_path}: {e}");

    def load_chunks(self, file_name):
        chunks_path = self.get_chunks_path(file_name)
        if os.path.exists(chunks_path):
            try:
                with open(chunks_path, 'r', encoding='utf-8') as f: chunks = json.load(f); logger.info(f"Loaded {len(chunks)} text chunks from {chunks_path}"); return chunks
            except Exception as e: logger.error(f"Error loading/decoding text chunks from {chunks_path}: {e}. Will re-process."); return None
        logger.info(f"Text chunks file not found at {chunks_path}"); return None

    # --- Document Processing (process_files - Unchanged Logic, calls modified extractors) ---
    def process_files(self, progress_callback=None):
        # Reset states (partially, keep DB connection if not in-memory)
        self.file_chunks = {}; self.faiss_indexes = {}; self.processed_files = set()
        
        # Clear table metadata only if DB is in-memory or if a full re-process is implied
        # For now, always clear metadata and file_to_table_map to ensure consistency with current files
        self.table_metadata = {}; self.file_to_table_map = {}

        if self.config.sqlite_db_path == ":memory:" and self.db_cursor:
            logger.info("Clearing existing tables from in-memory database...")
            try:
                self.db_cursor.execute("SELECT name FROM sqlite_master WHERE type='table';"); tables = self.db_cursor.fetchall()
                for table_tuple in tables: # tables is a list of tuples
                    self.db_cursor.execute(f'DROP TABLE IF EXISTS "{table_tuple[0]}"') # Quote table names
                self.db_conn.commit(); logger.info(f"Dropped {len(tables)} tables from in-memory DB.")
            except sqlite3.Error as e_drop: logger.error(f"Error dropping tables from in-memory DB: {e_drop}")
        # If not in-memory, existing tables are kept unless explicitly overwritten by a new file with same name/sheet

        data_dir = self.config.data_dir; supported_ext = self.config.supported_extensions
        try: all_files = os.listdir(data_dir); process_list = sorted([f for f in all_files if f.lower().endswith(supported_ext)])
        except Exception as e: logger.error(f"Error listing files in data directory {data_dir}: {e}", exc_info=True); return False
        if not process_list: logger.warning(f"No supported files found in {data_dir}"); return True # Completed successfully, but nothing to process

        logger.info(f"Processing {len(process_list)} supported files from {data_dir}");
        embedding_dim = self.encoder_model.get_sentence_embedding_dimension();
        total_files = len(process_list)
        if progress_callback: progress_callback(f"Found {total_files} file(s). Starting processing...", current_step=0, total_steps=total_files)

        any_success = False
        for idx, file_name in enumerate(process_list):
            current_file_msg = f"Processing ({idx+1}/{total_files}): {file_name}"; logger.info(f"--- {current_file_msg} ---")
            if progress_callback: progress_callback(current_file_msg, current_step=idx, total_steps=total_files)
            file_path = os.path.join(data_dir, file_name); index_path = self.get_index_path(file_name); emb_path = self.get_embedding_path(file_name); chunks_path = self.get_chunks_path(file_name)
            file_extension = os.path.splitext(file_name)[1].lower(); file_processed_ok = False

            # --- Spreadsheet Processing (SQL Loading) ---
            if file_extension in ['.csv', '.xlsx']:
                sql_success = False; processed_tables = []
                # For non-in-memory DB, we might not want to drop tables unless file content changes.
                # Current logic for _load_df_to_sql is 'replace', which handles this.
                if file_extension == '.xlsx': sql_success, processed_tables = self._extract_and_load_xlsx(file_path, file_name, progress_callback)
                elif file_extension == '.csv': sql_success, processed_tables = self._extract_and_load_csv(file_path, file_name, progress_callback)

                if sql_success:
                    logger.info(f"Successfully loaded data from {file_name} into SQL table(s): {processed_tables}")
                    self.processed_files.add(file_name); file_processed_ok = True; any_success = True
                else:
                    logger.error(f"Failed to load data from {file_name} into SQL database.")
                    if progress_callback: progress_callback(f"❌ Failed SQL processing for {file_name}", is_error=True) # Pass error flag

            # --- Text Processing (Chunking/Indexing) - PDF/PPTX ---
            elif file_extension in ['.pdf', '.pptx']:
                try:
                    chunks = self.load_chunks(file_name)
                    if chunks is None: # If chunks file doesn't exist or fails to load
                        if progress_callback: progress_callback(f"{current_file_msg} - Extracting Content...", stage="Extracting Text")
                        all_content = self.extract_content(file_path); # Calls modified PDF extractor if .pdf
                        if not all_content: logger.warning(f"No text content extracted from {file_name}. Skipping text indexing."); continue # Skip to next file
                        else:
                            if progress_callback: progress_callback(f"{current_file_msg} - Chunking Content...", stage="Chunking Text")
                            chunks = self.chunk_content(all_content); # Calls modified chunker
                            if not chunks: logger.warning(f"No text chunks generated for {file_name}."); continue
                            else: self.save_chunks(file_name, chunks) # Save newly generated chunks

                    # Proceed with embedding/indexing only if chunks exist (either loaded or newly generated)
                    if chunks:
                        self.file_chunks[file_name] = chunks; logger.debug(f"Stored {len(chunks)} text chunks for {file_name}")
                        faiss_index = None; regenerate_embeddings = False
                        # Check if index, embeddings, AND chunks file exist for valid cache
                        if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                            if progress_callback: progress_callback(f"{current_file_msg} - Verifying Index...", stage="Verifying Index")
                            try:
                                embeddings_on_disk = np.load(emb_path)
                                # Verify consistency: num embeddings == num chunks, and correct dimension
                                if embeddings_on_disk.ndim != 2 or embeddings_on_disk.shape[1] != embedding_dim or embeddings_on_disk.shape[0] != len(chunks):
                                    logger.warning(f"Index/chunk mismatch for {file_name} (Embeddings: {embeddings_on_disk.shape}, Chunks: {len(chunks)}). Regenerating."); regenerate_embeddings = True
                                else:
                                    faiss_index = self.load_faiss_index(file_name, embedding_dim); # Tries to load existing index
                                    if faiss_index.ntotal != len(chunks): # Double check FAISS index size against chunks
                                        logger.warning(f"FAISS index size ({faiss_index.ntotal}) mismatch with chunk count ({len(chunks)}) for {file_name}. Regenerating.")
                                        regenerate_embeddings = True; faiss_index = None # Force regen
                                    else: logger.info(f"Verified existing text index data for {file_name}")
                            except Exception as e: logger.error(f"Error verifying index for {file_name}: {e}. Regenerating...", exc_info=True); regenerate_embeddings = True; faiss_index = None
                        else: regenerate_embeddings = True; logger.info(f"No complete index found for {file_name}, will generate.")

                        if regenerate_embeddings or faiss_index is None: # If index obj is None even after load attempt
                            logger.info(f"Generating text embeddings/index for {file_name}...")
                            if progress_callback: progress_callback(f"{current_file_msg} - Embedding Text...", stage="Embedding Text")
                            content_list = [chunk['content'] for chunk in chunks];
                            if not content_list: logger.warning(f"No text content to embed for {file_name}.");
                            else:
                                embeddings = self.encoder_model.encode(content_list, batch_size=32, show_progress_bar=False, convert_to_numpy=True).astype('float32')
                                if embeddings.shape[0] == 0: logger.warning(f"Text embedding yielded no vectors for {file_name}.");
                                else:
                                    np.save(emb_path, embeddings); logger.info(f"Saved {embeddings.shape[0]} text embeddings for {file_name}")
                                    if progress_callback: progress_callback(f"{current_file_msg} - Indexing Text...", stage="Indexing Text")
                                    faiss_index = faiss.IndexFlatL2(embedding_dim); faiss_index.add(embeddings); faiss.write_index(faiss_index, index_path); logger.info(f"Saved FAISS text index for {file_name} ({faiss_index.ntotal} vectors)")

                        if faiss_index is not None and faiss_index.ntotal > 0:
                             self.faiss_indexes[file_name] = faiss_index; self.processed_files.add(file_name); logger.debug(f"Stored FAISS text index for {file_name}"); file_processed_ok = True; any_success = True
                        elif faiss_index is not None and faiss_index.ntotal == 0 and len(chunks) > 0:
                             logger.warning(f"Index created for {file_name} but is empty, though chunks exist. Possible embedding issue.")
                        else: logger.warning(f"No valid text index created or loaded for {file_name}.")
                    else:
                        logger.info(f"No text chunks for {file_name}, skipping text indexing.")

                except Exception as e_text:
                    logger.error(f"Failed text processing for {file_name}: {e_text}", exc_info=True)
                    if progress_callback: progress_callback(f"❌ Error during text processing for {file_name}: {e_text}", is_error=True) # Pass error flag
                    try: # Clean up potentially inconsistent text index files
                        for p in [index_path, emb_path, chunks_path]:
                            if os.path.exists(p): os.remove(p); logger.debug(f"Removed potentially inconsistent text file {p}")
                    except OSError as e_clean: logger.error(f"Error cleaning up text files for {file_name}: {e_clean}")
                    if file_name in self.file_chunks: del self.file_chunks[file_name]
                    if file_name in self.faiss_indexes: del self.faiss_indexes[file_name]
            else:
                 logger.warning(f"Skipping file {file_name} due to unsupported extension {file_extension} in main processing loop.")

        # Final summary
        final_text_index_count = len(self.faiss_indexes)
        final_sql_table_count = len(self.table_metadata)
        total_processed_ok = len(self.processed_files)

        if any_success:
            logger.info(f"--- Processing Complete. Processed {total_processed_ok}/{total_files} files. Text Indices: {final_text_index_count}. SQL Tables: {final_sql_table_count}. ---")
        else:
            logger.warning(f"--- Processing Complete. No documents successfully processed (no text index or SQL table). ---")

        # Return True if the process finished, even if nothing was loaded. The UI can check counts.
        return True


    # --- Querying Logic (Faiss search MODIFIED, Aggregate MODIFIED, Query LLM MODIFIED, Analysis Helpers MODIFIED, SQL Unchanged) ---

    # MODIFIED Faiss search to include section tag in results
    def query_files(self, query):
        if not self.faiss_indexes: logger.warning("No text indexes available for Faiss querying."); return {}
        query_results = {};
        try:
            logger.info(f"Encoding query for Faiss text search: '{query[:100]}...'")
            query_embedding = self.encoder_model.encode(query, convert_to_numpy=True).astype("float32"); query_embedding = np.array([query_embedding])
            if query_embedding.ndim != 2: raise ValueError("Query embedding shape error")

            for file_name, index in self.faiss_indexes.items():
                if index is None or index.ntotal == 0: logger.debug(f"Skipping empty text index for {file_name}"); continue
                try:
                    k_search = min(self.config.k_retrieval, index.ntotal);
                    # logger.debug(f"Searching Faiss index for {file_name} (k={k_search}, N={index.ntotal})")
                    D, I = index.search(query_embedding, k=k_search); # D=distances, I=indices
                    indices, distances = I[0], D[0]

                    current_file_chunks = self.file_chunks.get(file_name)
                    if not current_file_chunks: logger.error(f"Text chunks missing for {file_name} during query!"); continue

                    file_results = []; processed_indices = set()
                    for i, idx in enumerate(indices):
                        if idx == -1 or idx in processed_indices or not (0 <= idx < len(current_file_chunks)): continue
                        processed_indices.add(idx); chunk = current_file_chunks[idx]

                        # Ensure all necessary fields are retrieved from the chunk
                        file_results.append({
                            "source_info": chunk.get('source_info', {}), # Page, maybe title
                            "file_type": chunk.get('file_type', 'unknown'),
                            "content": chunk.get('content', ''),
                            "score": round(float(distances[i]), 4), # L2 distance, lower is better
                            "type": chunk.get('type', 'unknown'),
                            "potential_section": chunk.get('potential_section', 'unknown') # << ADDED
                        })

                    file_results.sort(key=lambda x: x['score']) # Sort by distance (ascending)
                    if file_results:
                        if file_name not in query_results: query_results[file_name] = []
                        query_results[file_name].extend(file_results)
                        # logger.debug(f"Found {len(file_results)} relevant chunks in {file_name}.")

                except Exception as e_search: logger.error(f"Error searching text index {file_name} for query '{query[:50]}...': {e_search}", exc_info=True)

            logger.debug(f"Faiss text search complete for query '{query[:50]}...'. Found results in {len(query_results)} files.")
            return query_results
        except Exception as e_query: logger.error(f"Error during Faiss text query processing for '{query[:50]}...': {e_query}", exc_info=True); return {}

    # MODIFIED Aggregate text context to include section tag
    def aggregate_context(self, query_results, strategy="top_k"):
        all_context = {}; max_chars = self.config.max_context_tokens * 3 # Estimate based on tokens
        logger.info(f"Aggregating text context using strategy '{strategy}', max_chars ~{max_chars}")
        if not query_results: return all_context

        # Flatten results from all files and sort globally by score
        flat_results = []
        for file_name, results in query_results.items():
            for res in results: flat_results.append({**res, "file_name": file_name}) # Add filename for global sorting/source tracking
        flat_results.sort(key=lambda x: x['score']) # Sort globally by score (lower L2 distance is better)

        aggregated_context_str = ""; total_aggregated_chars = 0; added_chunks_count = 0; context_sources = set()

        # Limit by total chunks OR character limit, whichever comes first
        limit_k = self.config.k_retrieval # Use the config value as the max number of chunks

        for i, res in enumerate(flat_results):
            if added_chunks_count >= limit_k:
                logger.debug(f"Reached global top-k limit ({limit_k}) for text context aggregation.")
                break

            source_info = res.get('source_info', {});
            file_name = res['file_name'];
            file_type = res.get('file_type', 'unknown');
            content_type = res.get('type', 'unknown');
            score = res['score'];
            content_body = res['content']
            potential_section = res.get('potential_section', 'N/A') # << Get the section tag

            # Construct source string for context header
            source_parts = [f"Source: {file_name}"]
            if file_type == 'pptx':
                 slide_info_parts = []
                 if 'slide' in source_info: slide_info_parts.append(f"Slide: {source_info['slide']}")
                 if 'slide_title' in source_info : slide_info_parts.append(f"Title Slide: {source_info['slide_title']}") # For merged slides
                 if 'slide_content' in source_info : slide_info_parts.append(f"Content Slide: {source_info['slide_content']}") # For merged slides
                 if 'merged_from' in source_info: slide_info_parts.append(f"(merged from {source_info['merged_from']})") # Old merge style
                 source_parts.append(", ".join(filter(None, slide_info_parts)))

            elif file_type == 'pdf':
                 page_info = f"Page: {source_info.get('page', 'N/A')}"
                 if 'table_index_on_page' in source_info: page_info += f", Table {source_info['table_index_on_page']}"
                 source_parts.append(page_info)

            source_parts.append(f"Section: {potential_section}") # << ADDED section tag to header
            # source_parts.append(f"Type: {content_type}"); # Maybe redundant? Keep it concise.
            # source_parts.append(f"Score: {score:.4f}") # Score might not be useful for LLM, remove?

            source_str = ", ".join(filter(None, source_parts)); # Filter out empty parts before joining
            content_header = f"--- Context from {source_str} ---\n";
            content_to_add = content_header + content_body + "\n\n";
            content_chars = len(content_to_add)

            if total_aggregated_chars + content_chars <= max_chars:
                aggregated_context_str += content_to_add;
                total_aggregated_chars += content_chars;
                added_chunks_count += 1;
                context_sources.add(file_name)
            else:
                # Try to add truncated chunk if it's the first one and too large
                if added_chunks_count == 0:
                     remaining_chars = max_chars - total_aggregated_chars - len(content_header) - 20 # Headroom
                     if remaining_chars > 50:
                         truncated_body = content_body[:remaining_chars];
                         aggregated_context_str += content_header + truncated_body + "\n[...TRUNCATED CONTEXT...]\n\n";
                         total_aggregated_chars += len(aggregated_context_str); # Approx update
                         added_chunks_count += 1;
                         context_sources.add(file_name);
                         logger.warning("Text context truncated (first chunk too large).")
                     else: logger.warning("First text chunk too large to fit even truncated, skipping.")
                # Stop aggregation if character limit reached
                logger.info(f"Stopping text context aggregation at {added_chunks_count} chunks ({total_aggregated_chars}/{max_chars} chars). Character limit reached."); break

        final_context = aggregated_context_str.strip()
        if final_context:
            all_context = {"combined_context": final_context, "source_files": sorted(list(context_sources))};
            logger.info(f"Aggregated {total_aggregated_chars} chars from {added_chunks_count} text chunks across {len(context_sources)} files.")
        else: logger.warning("No text context aggregated within limits.")
        return all_context

    # --- LLM Call Helpers (Format History unchanged, Generic Call unchanged) ---

    def _format_chat_history_for_llm(self, chat_history):
        messages = []
        for msg in chat_history:
            role = msg.get("role")
            content = msg.get("content")
            if role == "user":
                messages.append(HumanMessage(content=content))
            elif role == "assistant":
                 # Remove footer from assistant message before adding to history
                 clean_content = re.sub(r"\n\n---\n\*Answer generated.*", "", content, flags=re.DOTALL).strip()
                 messages.append(AIMessage(content=clean_content))
            # Ignore system messages in history for simplicity, handled by prompts
        return messages

    def _call_llm_for_analysis(self, prompt_messages, task_description, chat_history=None):
        """Calls the LLM for internal analysis tasks, optionally considering chat history."""
        try:
            llm = self._get_llm()
            logger.info(f"Calling Groq LLM ({self.config.llm_model}) for internal task: {task_description}...")

            # History currently ignored for internal analysis (classification, SQL gen etc.)
            # Focus on using the refined query for these steps.
            if chat_history:
                logger.debug(f"Chat history provided for {task_description}, but currently ignored in this analysis step.")

            final_messages = prompt_messages # Pass only the task-specific messages

            response = llm.invoke(final_messages)
            content = response.content.strip() if hasattr(response, 'content') else str(response).strip()
            logger.info(f"LLM response for {task_description}: '{content[:150]}...'")
            return content
        except Exception as e:
            logger.error(f"LLM call failed during {task_description}: {e}", exc_info=True)
            return None

    # MODIFIED Query LLM with text context to include proposal/structure info and SatSure instructions
    def query_llm_with_history(self, query, context_data, chat_history, retry_count=1):
        combined_context = context_data.get("combined_context", "");
        source_files = context_data.get("source_files", []);
        source_file_str = ", ".join(source_files) if source_files else "the provided text documents"

        if not combined_context:
             logger.warning("No text context provided for LLM query.");
             # Return SatSure specific "don't know" response if no context
             return "I don't really know but my next versions will be able to answer this for sure! (Reason: No relevant text context found)"

        try:
            llm = self._get_llm()
            formatted_history = self._format_chat_history_for_llm(chat_history) # Format past messages

            # Create a prompt template that includes history, context, proposal info, and SatSure instructions
            prompt = ChatPromptTemplate.from_messages([
                SystemMessage(content=f"""You are the SatSure LLM, a helpful AI assistant answering questions based ONLY on the provided text context from SatSure project proposal documents and the ongoing conversation history. The proposals outline solutions for clients.
The context below comes from document(s): '{source_file_str}'. Context sections are marked with 'Source:', 'Page:', 'Table Index', 'Slide:', and 'Section:' (e.g., 'Section: title_page', 'Section: main_content').
Use ONLY information presented between '--- START CONTEXT ---' and '--- END CONTEXT ---' or previous messages.
This is sensitive data, there must be no hallucination. Do not use any prior knowledge.
If the answer cannot be found in the text or history, state EXACTLY: "I don't really know but my next versions will be able to answer this for sure!"
Respond conversationally, considering the flow of the chat. Do NOT perform calculations unless explicitly shown in the text."""),
                MessagesPlaceholder(variable_name="chat_history"), # Where history goes
                HumanMessage(content=f"""Retrieved Context from SatSure proposal(s):\n--- START CONTEXT ---\n{combined_context}\n--- END CONTEXT ---\n\nLatest User Question: {query}\n\nAnswer based ONLY on the provided text context and conversation history:""")
            ])

            chain = prompt | llm
            logger.info(f"Querying Groq model {self.config.llm_model} using text context from {source_file_str} and chat history (Proposal Aware)...")
            answer = f"Error: LLM text query failed after retries." # Default error

            for attempt in range(retry_count + 1):
                try:
                    response = chain.invoke({
                        "chat_history": formatted_history,
                        # Query is part of the last HumanMessage in the template
                    })
                    answer = response.content.strip() if hasattr(response, 'content') else str(response).strip()
                    logger.info(f"Groq response received for text RAG w/ history (attempt {attempt+1}). Length: {len(answer)}")
                    # Check if the answer is meaningful (not empty, not just an error message unless it's the specific SatSure one)
                    is_valid_answer = bool(answer) and answer.strip() != ""
                    is_api_error = answer.lower().startswith("error:")
                    is_satsure_dont_know = answer.strip() == "I don't really know but my next versions will be able to answer this for sure!"

                    if is_valid_answer and not is_api_error:
                         break # Got a valid answer (could be the "don't know" response, which is valid here)

                    # If empty or API error, retry if possible
                    if (not is_valid_answer or is_api_error) and attempt < retry_count:
                        logger.warning(f"Text RAG response invalid or API error (Attempt {attempt+1}), retrying..."); time.sleep(1.2 ** attempt)
                    elif attempt == retry_count: # Last attempt
                        if not is_valid_answer:
                           answer = f"Received empty response from LLM for text context from {source_file_str} (with history)."
                        # Keep the answer if it's an error message from the API itself or the SatSure "don't know" response
                        logger.warning(f"Final attempt for text RAG yielded: {answer}")

                except Exception as e_api:
                    logger.warning(f"Groq API attempt {attempt+1} for text RAG w/ history failed: {e_api}")
                    if attempt < retry_count: time.sleep(1.5 ** attempt); logger.info("Retrying Groq text RAG query...")
                    else: answer = f"Error: Failed to get answer from LLM for text RAG w/ history after {retry_count+1} attempts. (API Error: {e_api})"
            return answer
        except Exception as e_setup:
            logger.error(f"Error setting up/calling Groq API for text RAG w/ history: {e_setup}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"LLM text query failed: {e_setup}")
            # Return SatSure specific error format?
            return f"Error: Could not query LLM for text RAG w/ history due to setup error: {e_setup}"


    # --- Query Analysis Helpers (MODIFIED Prompts) ---
    # These operate best on a *standalone* query. We will refine the query *before* calling these.

    # MODIFIED classify prompt with SatSure context
    def _classify_query(self, query): # Takes refined query
        system_prompt = """You are the SatSure LLM, an expert query analyzer. Your task is to classify the user query related to SatSure project proposal documents (PDFs) and potentially related data (XLSX/CSV). This is sensitive data so there must be no hallucination during retrieval. If you don't know an answer, simple answer with "I don't really know but my next versions will be able to answer this for sure!"

Classify the user query into ONE of the following categories:
1.  'Simple Retrieval': Asking for specific facts, definitions, or simple summaries directly extractable from text within the proposals.
2.  'Complex/Reasoning (Text)': Requires combining information from multiple text passages across one or more proposals, summarization across sections, comparison between proposals, or reasoning based *only* on the text provided. Does not involve calculations on tabular data.
3.  'Structured Query (SQL)': Requires operations on structured data (like CSV or Excel tables stored in a SQL database), such as calculations (sum, average, count), filtering based on values, grouping, sorting, or comparisons across rows/columns.

Analyze the query carefully. Respond ONLY in JSON format with two keys: "classification" (the category string) and "confidence" (a float between 0.0 and 1.0 indicating your certainty, e.g., 0.95).

Examples:
Query: "What is the proposed satellite constellation for Project SkyWatch?" -> {"classification": "Simple Retrieval", "confidence": 0.99}
Query: "Compare the methodologies proposed in the AquaMonitor and TerraScan documents." -> {"classification": "Complex/Reasoning (Text)", "confidence": 0.9}
Query: "What is the total estimated cost for all Phase 1 deliverables listed in the attached spreadsheet?" -> {"classification": "Structured Query (SQL)", "confidence": 0.98}
Query: "Summarize the key risks mentioned across all proposals." -> {"classification": "Complex/Reasoning (Text)", "confidence": 0.85}
Query: "List all personnel requiring security clearance from the project_team.csv file." -> {"classification": "Structured Query (SQL)", "confidence": 0.96}
Query: "Tell me about SatSure's SAR capabilities mentioned in the proposals." -> {"classification": "Simple Retrieval", "confidence": 0.9}
"""
        human_prompt = f"Classify the following user query about SatSure proposals/data:\n\"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        raw_response = self._call_llm_for_analysis(messages, "query classification") # History not passed here
        if not raw_response: logger.warning("Query classification failed (LLM call unsuccessful). Defaulting."); return "Simple Retrieval", 0.5
        try:
            json_match = re.search(r"\{.*\}", raw_response, re.DOTALL)
            if json_match:
                json_str = json_match.group(0); result = json.loads(json_str)
                classification = result.get("classification"); confidence = float(result.get("confidence", 0.0))
                valid_classifications = ["Simple Retrieval", "Complex/Reasoning (Text)", "Structured Query (SQL)"]
                if classification in valid_classifications and 0.0 <= confidence <= 1.0: logger.info(f"Query classified as '{classification}' with confidence {confidence:.2f}"); return classification, confidence
                else: logger.warning(f"Invalid classification or confidence in JSON: {json_str}. Defaulting.")
            else: logger.warning(f"No JSON object found in classification response: {raw_response}. Defaulting.")
            return "Simple Retrieval", 0.5
        except (json.JSONDecodeError, TypeError, AttributeError, ValueError) as e: logger.warning(f"Failed to parse classification JSON: '{raw_response}'. Error: {e}. Defaulting."); return "Simple Retrieval", 0.5

    def _decompose_query(self, query): # Takes refined query
        # Prompt remains generic, decomposition is about breaking down complexity based on text needs.
        system_prompt = """You are an expert query decomposer. Break down a complex user query needing information potentially spanning multiple SatSure project proposal documents into simpler, factual sub-queries. Each sub-query should aim to retrieve a specific piece of information from the text content of the proposals. Do NOT decompose queries asking for calculations or filtering on tables. Format as a numbered list. If the query is simple or cannot be decomposed, return the original query prefixed with '1. '."""
        human_prompt = f"Decompose the following complex query about SatSure proposals:\n\"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        decomposition = self._call_llm_for_analysis(messages, "text query decomposition") # History not passed
        if decomposition:
            # Handle potential numbering or bullet points
            sub_queries = [re.sub(r"^\s*[\*\-\d]+\.?\s*", "", line).strip() for line in decomposition.split('\n') if line.strip()]
            # Ensure decomposition actually changed the query and didn't just return the number/bullet
            if sub_queries and not (len(sub_queries) == 1 and sub_queries[0].strip() == query.strip()):
                logger.info(f"Decomposed text query into: {sub_queries}")
                return sub_queries
            else: logger.info("Decomposition didn't yield distinct sub-queries for text. Using original."); return [query]
        else: logger.warning("Text query decomposition failed. Using original query."); return [query]

    def _identify_target_sql_table(self, query): # Takes refined query (Prompt unchanged)
        if not self.table_metadata: logger.warning("No table metadata available to identify target SQL table."); return None
        schema_overview = "Available SQL Tables (likely containing supplementary project data):\n"; table_options = []
        for table_name, meta in self.table_metadata.items():
            table_options.append(table_name); col_descs = []
            for col in meta.get('columns', []):
                desc = col.get('description', 'N/A'); desc_short = (desc[:100] + '...') if len(desc) > 103 else desc
                # Ensure column name is quoted for the prompt for clarity if it contains spaces or special chars
                col_name_display = f'`{col.get("name", "UNKNOWN_COL")}`'
                col_descs.append(f"- {col_name_display} (Type: {col.get('type', 'UNKNOWN_TYPE')}, Desc: {desc_short})")
            sheet_info = f" (Sheet: '{meta.get('source_sheet')}')" if meta.get('source_sheet') else ""
            schema_overview += f"\nTable: `{table_name}` (From File: '{meta.get('source_file', 'N/A')}'{sheet_info}, Rows: {meta.get('row_count', 'N/A')})\n"; schema_overview += "\n".join(col_descs) + "\n"
        if not table_options: logger.warning("No tables found in metadata during identification."); return None
        system_prompt = f"""You are an expert data analyst. Your task is to choose the single most relevant SQL table to answer the user's query, based on the table names, source files, and column descriptions provided. The user query likely relates to SatSure project proposals, and these tables might contain supporting data.

Available Table Schemas:\n{schema_overview}
Analyze the user's query and the schemas. Respond ONLY with the exact name of the single most appropriate table from the list: {table_options}.
If no single table seems clearly relevant or sufficient to answer the query, respond ONLY with the word "NONE". Do not add any explanation or introductory text."""
        human_prompt = f"User Query: \"{query}\"\n\nMost relevant table name:"
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        target_table_name = self._call_llm_for_analysis(messages, "SQL target table identification") # History not passed
        if not target_table_name: logger.error("SQL target table identification failed (LLM call unsuccessful)."); return None
        target_table_name = target_table_name.strip().strip('`"') # Clean potential markdown/quotes
        if target_table_name == "NONE": logger.info("LLM indicated no single relevant SQL table found."); return None
        elif target_table_name in self.table_metadata: logger.info(f"LLM identified '{target_table_name}' as the target SQL table."); return target_table_name
        else: logger.warning(f"LLM returned an invalid table name '{target_table_name}'. Valid options were: {table_options}"); return None

    # --- SQL Query Generation & Execution (Unchanged logic, prompt is fine) ---
    def _generate_sql_query(self, query, table_name, table_meta): # Takes refined query
        if not table_meta or 'columns' not in table_meta: logger.error(f"Metadata missing for table '{table_name}'. Cannot generate SQL."); return None
        schema_description = f"Table Name: `{table_name}` (This table likely contains data supporting a SatSure project proposal)\nColumns:\n"
        for col in table_meta['columns']:
             safe_col_name = f'"{col["name"]}"' # Double quotes for SQLite standard
             col_type = col.get('type', 'TEXT')
             col_desc = col.get('description', 'N/A')
             schema_description += f"- {safe_col_name} (Type: {col_type}, Description: {col_desc})\n"

        system_prompt = f"""You are an expert SQL query writer specializing in SQLite. Given a user query (likely related to SatSure project proposals) and the schema of a table (including column descriptions), write a concise and valid SQLite query to answer the user's question using ONLY this table.

Table Schema:\n{schema_description}
Instructions:
- Write ONLY a single, valid SQLite query.
- Use the exact column names provided in the schema, ensuring they are correctly quoted with double quotes (e.g., "Column Name With Spaces").
- Pay close attention to the column descriptions to understand the data's meaning.
- Handle potential data types correctly. Use `CAST(column AS REAL)` or `CAST(column AS INTEGER)` if needed for comparisons or calculations on TEXT columns containing numbers (e.g., `WHERE CAST("Numeric Column" AS REAL) > 100`). Remember to quote the column name inside CAST as well.
- Output ONLY the raw SQL query. No explanations, comments, markdown backticks (```sql ... ```), or introductory text.
- If the query cannot be answered accurately with the given table and columns, output ONLY the exact text: QUERY_CANNOT_BE_ANSWERED
"""
        human_prompt = f"User Query: \"{query}\"\n\nSQLite Query:"
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        sql_query = self._call_llm_for_analysis(messages, f"SQL query generation for table {table_name}") # History not passed
        if not sql_query: logger.error("SQL query generation failed (LLM call unsuccessful)."); return None
        sql_query = sql_query.strip();
        if sql_query == "QUERY_CANNOT_BE_ANSWERED": logger.warning(f"LLM indicated query cannot be answered for table {table_name}."); return None
        # Remove potential markdown and clean
        sql_query = re.sub(r"^```(?:sql)?\s*", "", sql_query, flags=re.IGNORECASE | re.DOTALL)
        sql_query = re.sub(r"\s*```$", "", sql_query, flags=re.DOTALL)
        sql_query = sql_query.strip()
        if not sql_query.upper().startswith("SELECT"): logger.warning(f"Generated SQL is not a SELECT statement: '{sql_query[:100]}...' - Discarding for safety."); return None
        if sql_query.endswith(';'): sql_query = sql_query[:-1]
        logger.info(f"Generated SQL query for table '{table_name}':\n{sql_query}"); return sql_query

    def _execute_sql_query(self, sql_query):
        if not self.db_cursor: logger.error("No database cursor available for SQL execution."); return None, "Database connection error."
        if not sql_query: logger.error("No SQL query provided for execution."); return None, "No SQL query generated."
        try:
            logger.info(f"Executing SQL: {sql_query}"); start_time = time.time()
            # Use pandas read_sql_query for potentially better type handling and easier formatting
            if self.db_conn:
                 try:
                      df_results = pd.read_sql_query(sql_query, self.db_conn)
                      results = df_results.values.tolist() # Get results as list of lists/tuples
                      column_names = df_results.columns.tolist()
                 except Exception as e_pd_read:
                      logger.warning(f"Pandas read_sql failed ({e_pd_read}), falling back to cursor execution.")
                      # Fallback to cursor execution if pandas fails
                      self.db_cursor.execute(sql_query)
                      results = self.db_cursor.fetchall() # List of tuples
                      column_names = [desc[0] for desc in self.db_cursor.description] if self.db_cursor.description else []
            else: # Should not happen if connected, but safety check
                 self.db_cursor.execute(sql_query)
                 results = self.db_cursor.fetchall()
                 column_names = [desc[0] for desc in self.db_cursor.description] if self.db_cursor.description else []

            exec_time = time.time() - start_time; logger.info(f"SQL query executed successfully in {exec_time:.3f}s. Fetched {len(results)} rows.")

            if not results: return "No results found.", None # Success, but no data

            # Use pandas to create a formatted string (easier for tables)
            try:
                 if 'df_results' not in locals(): df_results = pd.DataFrame(results, columns=column_names)
                 max_rows_preview = 25 # Limit preview size in string output
                 output_str = df_results.to_string(index=False, max_rows=max_rows_preview, na_rep='NULL') # Represent NaN as NULL
                 if len(df_results) > max_rows_preview:
                      output_str += f"\n... (truncated, {len(df_results)} total rows)"
                 output_str = f"Query Result ({len(results)} row(s)):\n" + output_str
            except Exception as e_pd_format:
                 logger.warning(f"Pandas formatting of SQL results failed: {e_pd_format}. Falling back to basic string format.")
                 # Fallback basic formatting
                 max_rows_preview = 25
                 output_str = f"Query Result ({len(results)} row(s)):\n"
                 header = ", ".join(map(str, column_names))
                 output_str += header + "\n" + "-" * len(header) + "\n"
                 for i, row in enumerate(results):
                     if i >= max_rows_preview:
                         output_str += f"... (truncated, {len(results) - max_rows_preview} more rows)\n"; break
                     row_str = ", ".join(map(lambda x: str(x) if x is not None else "NULL", row))
                     output_str += row_str + "\n"

            return output_str.strip(), None # Return formatted string and no error

        except sqlite3.Error as e_sql:
            err_msg = f"SQLite execution error: {e_sql} (Query: {sql_query[:200]}...)"
            logger.error(err_msg, exc_info=False)
            try: self.db_conn.rollback()
            except Exception as e_rb: logger.error(f"Rollback failed after SQL error: {e_rb}")
            return None, f"Database error during execution: {e_sql}"
        except Exception as e:
            err_msg = f"Unexpected error during SQL execution: {e}"
            logger.error(f"{err_msg} for query: {sql_query}", exc_info=True)
            try: self.db_conn.rollback()
            except Exception as e_rb: logger.error(f"Rollback failed after unexpected error: {e_rb}")
            return None, err_msg


    # --- Final Answer Synthesis from SQL Result (MODIFIED Prompt) ---
    # Added SatSure context and instructions to the prompt.
    def _synthesize_answer_from_sql_with_history(self, user_query, sql_result_str, table_name, table_meta, chat_history):
        source_desc = f"the table '{table_name}'"
        if table_meta:
            source_desc += f" from file '{table_meta.get('source_file', 'N/A')}'"
            if table_meta.get('source_sheet'): source_desc += f" (sheet '{table_meta.get('source_sheet')}')"
        source_desc += " (containing data potentially related to SatSure project proposals)"

        formatted_history = self._format_chat_history_for_llm(chat_history)
        system_content = ""
        human_content = ""

        if sql_result_str == "No results found.":
             logger.info("SQL query returned no results. Synthesizing direct 'not found' answer considering history.")
             system_content = f"""You are the SatSure LLM, an AI assistant engaging in a conversation about SatSure project proposals and related data. A search was performed in structured data ({source_desc}), but it returned "No results found." for the latest user query.
Your task is to formulate a polite, natural language response indicating that the requested information could not be found in the specified data source, considering the ongoing conversation history.
Do not mention the SQL query. Start the answer directly. If appropriate, use the SatSure specific phrase for unknown answers.
Remember: no hallucination. If the data truly wasn't found, state that clearly."""
             human_content = f"""Original User Question: "{user_query}"

Data retrieved from {source_desc}: No results found.

Formulate the final "not found" answer considering the conversation flow. If simply stating 'not found' is best, use "I looked in the relevant data table but couldn't find the specific information you asked for." Otherwise, you can use "I don't really know but my next versions will be able to answer this for sure!" if it fits the context better.:"""
        else:
            system_content = f"""You are the SatSure LLM, an AI assistant engaging in a conversation about SatSure project proposals and related data. Relevant data was retrieved by executing a SQL query against {source_desc}.
Your task is to formulate a clear, concise, and natural language answer to the *latest* user question based ONLY on the provided SQL query results and the conversation history.
Do NOT just repeat the raw data table. Summarize or present the key information needed to answer the question directly and conversationally.
If the data seems incomplete or doesn't directly answer the question fully, state what you found and mention that it might not be the full answer.
Do not mention the SQL query itself or the table name unless relevant to the conversation. Start your answer directly.
Remember: no hallucination. Base the answer strictly on the data provided below and the chat history. If the data doesn't support an answer, say so (perhaps using the SatSure 'don't know' phrase if appropriate)."""
            human_content = f"""Original User Question: "{user_query}"

Data retrieved from {source_desc}:
--- START SQL RESULT DATA ---
{sql_result_str}
--- END SQL RESULT DATA ---

Formulate the final answer based ONLY on the provided data and considering the conversation flow:"""

        # Use ChatPromptTemplate with MessagesPlaceholder
        prompt = ChatPromptTemplate.from_messages([
            SystemMessage(content=system_content),
            MessagesPlaceholder(variable_name="chat_history"),
            HumanMessage(content=human_content) # Contains the latest context/task
        ])

        chain = prompt | self._get_llm()
        final_answer = None
        try:
             logger.info("Synthesizing final answer from SQL result with history (SatSure context)...")
             response = chain.invoke({
                 "chat_history": formatted_history
             })
             final_answer = response.content.strip() if hasattr(response, 'content') else str(response).strip()
        except Exception as e:
            logger.error(f"LLM call failed during SQL result synthesis w/ history: {e}", exc_info=True)
            final_answer = None # Handled below

        if not final_answer or final_answer.strip() == "":
            logger.error("Final answer synthesis from SQL result w/ history failed or yielded empty. Returning structured info or SatSure 'don't know'.")
            if sql_result_str == "No results found.":
                 # Return the specific SatSure "don't know" phrase here as a fallback
                 return "I don't really know but my next versions will be able to answer this for sure! (Reason: SQL search found no results and synthesis failed)"
            else:
                 # Provide the raw data as a last resort if synthesis fails but data exists
                 return f"I found the following data in {source_desc} regarding our discussion, but had trouble summarizing it naturally:\n```\n{sql_result_str}\n```"
        else:
            logger.info(f"Synthesized final answer from SQL result w/ history for query '{user_query[:50]}...'")
            return final_answer

    # --- Query Refinement based on History (MODIFIED Prompt) ---
    # Added SatSure context to the refinement prompt.
    def _refine_query_with_history(self, current_query, chat_history):
        """Uses LLM to refine the current query based on chat history for standalone clarity, considering SatSure context."""
        if not chat_history: # No history, query is already standalone
             return current_query

        formatted_history = self._format_chat_history_for_llm(chat_history)

        system_prompt = """You are an expert query assistant for the SatSure LLM. Given a conversation history about SatSure project proposals and related data, and the latest user query, rewrite the latest query to be fully standalone and understandable without the history.
- Incorporate necessary context from the history (e.g., referring to specific projects or topics previously discussed).
- Ensure the refined query is clear and focused for retrieval from proposal documents or data tables.
- If the latest query is already standalone and clear, return it exactly as is.
- Output ONLY the refined, standalone query. No explanations or other text."""

        # Construct prompt for refinement LLM call
        refine_prompt = ChatPromptTemplate.from_messages([
            SystemMessage(content=system_prompt),
            MessagesPlaceholder(variable_name="chat_history"),
            HumanMessage(content=f"Latest user query: \"{current_query}\"\n\nRefined standalone query:")
        ])

        chain = refine_prompt | self._get_llm()
        refined_query = current_query # Default
        try:
            logger.info(f"Refining query '{current_query[:100]}...' using chat history (SatSure context)...")
            response = chain.invoke({"chat_history": formatted_history})
            refined_query_llm = response.content.strip() if hasattr(response, 'content') else str(response).strip()

            # Basic check if refinement happened and is not empty or trivial
            if refined_query_llm and refined_query_llm.strip() and refined_query_llm != current_query and len(refined_query_llm) > 5:
                refined_query = refined_query_llm
                logger.info(f"Query refined to: '{refined_query[:100]}...'")
            else:
                logger.info("Query deemed standalone or refinement failed/empty, using original.")
                refined_query = current_query # Ensure it stays as original if refinement is bad

        except Exception as e:
            logger.error(f"LLM call failed during query refinement: {e}", exc_info=True)
            refined_query = current_query # Fallback to original query on error

        return refined_query


    # ==========================================================================
    # REVISED run_query_with_history FUNCTION (Unchanged Core Logic)
    # Calls the modified helper functions above.
    # ==========================================================================
    def run_query_with_history(self, query, chat_history, context_strategy="top_k"):
        """Complete query pipeline: refine query, classify, (SQL query OR text RAG), synthesize. Incorporates chat history & proposal structure awareness."""
        logger.info(f"--- Starting query execution (SQL+Chat+PDFStruct): '{query[:100]}...' ---")

        # Limit history
        history_limit = self.config.max_chat_history_turns * 2 # x2 for user/ai pair
        limited_history = chat_history[-history_limit:] if len(chat_history) > history_limit else chat_history
        logger.debug(f"Using last {len(limited_history)} messages ({len(limited_history)//2} turns) for context.")

        final_results = {
            "query": query, "refined_query": None,
            "classification": None, "classification_confidence": None,
            "sub_queries": None, "retrieval_results": {},
            "target_table": None, "target_table_metadata": None,
            "sql_query_generated": None, "sql_execution_result": None,
            "aggregated_context_data": {},
            "answer": "", "answer_source": None, # 'SQL Query' or 'Text RAG'
            "status": "Started", "error": None
        }
        is_streamlit = self._is_streamlit
        if is_streamlit: import streamlit as st # Import for spinner

        try:
            # === 0. Refine Query with History ===
            spinner_msg_refine = "Understanding query context..."
            refined_query = query
            if limited_history:
                 if is_streamlit:
                     with st.spinner(spinner_msg_refine): refined_query = self._refine_query_with_history(query, limited_history)
                 else: logger.info(spinner_msg_refine); refined_query = self._refine_query_with_history(query, limited_history)
            final_results["refined_query"] = refined_query


            # === Run Faiss search EARLY (using refined query) ===
            # This helps decide if text RAG is even possible if SQL fails
            spinner_msg_text = "Pre-searching proposal documents (if any)..."
            text_retrieval_results = {}
            if self.faiss_indexes: # Check if any text indexes exist
                 if is_streamlit:
                     with st.spinner(spinner_msg_text): text_retrieval_results = self.query_files(refined_query)
                 else: logger.info(spinner_msg_text); text_retrieval_results = self.query_files(refined_query)
                 final_results["retrieval_results"] = text_retrieval_results # Store pre-search results
            else:
                 logger.info("No text indexes (PDF/PPTX) loaded. Text RAG path will be unavailable.")


            # === 1. Classify Refined Query ===
            spinner_msg_classify = "Analyzing query intent..."
            classification, confidence = "Simple Retrieval", 0.5
            if is_streamlit:
                with st.spinner(spinner_msg_classify): classification, confidence = self._classify_query(refined_query)
            else: logger.info(spinner_msg_classify); classification, confidence = self._classify_query(refined_query)

            final_results["classification"] = classification
            final_results["classification_confidence"] = confidence
            logger.info(f"Refined query classified as: {classification} (Confidence: {confidence:.2f})")

            # === Attempt SQL Query Path ===
            sql_attempted = False; sql_succeeded = False; target_table_name = None
            can_attempt_sql = classification == "Structured Query (SQL)" and confidence >= self.config.dataframe_query_confidence_threshold and self.table_metadata

            if can_attempt_sql:
                sql_attempted = True
                logger.info("Attempting SQL query path based on classification and confidence.")

                # === 2a. Identify Target Table ===
                spinner_msg_identify = "Identifying relevant data table..."
                if is_streamlit:
                    with st.spinner(spinner_msg_identify): target_table_name = self._identify_target_sql_table(refined_query)
                else: logger.info(spinner_msg_identify); target_table_name = self._identify_target_sql_table(refined_query)

                if target_table_name and target_table_name in self.table_metadata:
                    final_results["target_table"] = target_table_name; target_meta = self.table_metadata[target_table_name]; final_results["target_table_metadata"] = target_meta
                    logger.info(f"Proceeding with SQL query on identified table: '{target_table_name}'")

                    # === 3a. Generate SQL Query ===
                    spinner_msg_gen = f"Generating SQL query for '{target_table_name}'..."
                    sql_query = None
                    if is_streamlit:
                        with st.spinner(spinner_msg_gen): sql_query = self._generate_sql_query(refined_query, target_table_name, target_meta)
                    else: logger.info(spinner_msg_gen); sql_query = self._generate_sql_query(refined_query, target_table_name, target_meta)

                    if sql_query:
                        final_results["sql_query_generated"] = sql_query

                        # === 4a. Execute SQL Query ===
                        spinner_msg_exec = f"Executing SQL query on '{target_table_name}'..."
                        sql_result_str, sql_error = None, None
                        if is_streamlit:
                             with st.spinner(spinner_msg_exec): sql_result_str, sql_error = self._execute_sql_query(sql_query)
                        else: logger.info(spinner_msg_exec); sql_result_str, sql_error = self._execute_sql_query(sql_query)

                        if sql_error:
                            logger.error(f"SQL execution failed: {sql_error}"); final_results["status"] = "Failed during SQL execution"; final_results["error"] = f"SQL query execution failed: {sql_error}"
                            # Fall through to text RAG
                        elif sql_result_str is None:
                             logger.error("SQL execution returned None result without error flag. Treating as failure.")
                             final_results["status"] = "Failed during SQL execution"; final_results["error"] = "SQL query execution returned no result."
                             # Fall through
                        else:
                            final_results["sql_execution_result"] = sql_result_str
                            logger.info(f"SQL execution successful. Raw result: {str(sql_result_str)[:200]}...")

                            # === 5a. Synthesize Final Answer from SQL Result (with history) ===
                            spinner_msg_synth = f"Generating answer from table '{target_table_name}' (considering history)..."
                            final_answer = None
                            if is_streamlit:
                                with st.spinner(spinner_msg_synth): final_answer = self._synthesize_answer_from_sql_with_history(refined_query, sql_result_str, target_table_name, target_meta, limited_history)
                            else: logger.info(spinner_msg_synth); final_answer = self._synthesize_answer_from_sql_with_history(refined_query, sql_result_str, target_table_name, target_meta, limited_history)

                            # Check if synthesis was successful (not empty, not generic error, not just the raw data fallback)
                            is_synthesis_successful = bool(final_answer) and final_answer.strip() != "" and not final_answer.lower().startswith("error:")
                            is_fallback_answer = ("had trouble summarizing it naturally" in final_answer or "synthesis failed" in final_answer) if final_answer else True


                            if is_synthesis_successful and not is_fallback_answer:
                                 final_results["answer"] = final_answer; final_results["status"] = "Completed Successfully"; final_results["answer_source"] = "SQL Query"; sql_succeeded = True; final_results["error"] = None
                            else:
                                 logger.error(f"SQL answer synthesis failed or returned fallback. Answer: {final_answer}. Falling back to text.")
                                 final_results["status"] = "Failed SQL answer synthesis";
                                 # Preserve SQL execution result but indicate synthesis failure
                                 final_results["error"] = final_results.get("error", "") + (" Failed to synthesize answer from SQL result." if not is_fallback_answer else f" SQL synthesis returned fallback: {final_answer}")
                                 # Fall through to Text RAG

                    else: # SQL Generation Failed
                        logger.warning(f"Failed to generate SQL query for table '{target_table_name}'. Falling back to text search."); final_results["status"] = "Failed SQL query generation"; final_results["error"] = f"Could not generate SQL query for table '{target_table_name}'."
                        # Fall through
                else: # Target Table Identification Failed
                     if not target_table_name: logger.warning("LLM could not identify a suitable SQL table. Falling back to text search."); final_results["error"] = "Could not identify a relevant SQL table for the query."
                     else: logger.error(f"Identified table '{target_table_name}' not found in metadata. Inconsistency?"); final_results["error"] = f"Identified table '{target_table_name}' not found in metadata."
                     final_results["status"] = "Failed SQL table identification"
                     # Fall through

            # === Fallback to Text RAG Path ===
            # Condition: SQL not attempted OR SQL attempted but failed/yielded no answer
            if not sql_succeeded:
                 if sql_attempted: logger.info("SQL query path failed or was inconclusive. Proceeding with text RAG using proposal documents.")
                 else: logger.info("Query classified for text retrieval/reasoning or no suitable SQL table found. Proceeding with text RAG.")

                 # Check if text indexes are actually available (PDF/PPTX were processed)
                 if not self.faiss_indexes:
                      logger.warning("No text indexes available (PDF/PPTX). Cannot perform text RAG.")
                      final_results["status"] = "Failed: No suitable processing path"
                      # Combine with previous error if SQL was attempted
                      prev_error = final_results.get('error')
                      final_results["error"] = "Query requires text analysis, but no text documents (PDF/PPTX) are indexed."
                      if prev_error: final_results["error"] += f" (Previous SQL attempt error: {prev_error})"
                      # Construct answer indicating failure reason using SatSure format
                      final_results["answer"] = f"I don't really know but my next versions will be able to answer this for sure! (Reason: {final_results['error']})"
                      return final_results # Cannot proceed

                 final_results["answer_source"] = "Text RAG"

                 # === 2b. Decompose if Complex Text Query ===
                 queries_to_retrieve = [refined_query]
                 if classification == "Complex/Reasoning (Text)":
                     spinner_msg_decompose = "Decomposing complex text query..."
                     sub_queries = []
                     if is_streamlit:
                         with st.spinner(spinner_msg_decompose): sub_queries = self._decompose_query(refined_query)
                     else: logger.info(spinner_msg_decompose); sub_queries = self._decompose_query(refined_query)
                     if sub_queries and sub_queries != [refined_query]: queries_to_retrieve = sub_queries; final_results["sub_queries"] = sub_queries; logger.info(f"Using sub-queries for text retrieval: {sub_queries}")

                 # === 3b. Retrieve relevant text chunks ===
                 # Use pre-search results if available and non-empty, otherwise fetch now
                 current_retrieval_results = final_results.get("retrieval_results")
                 needs_fresh_retrieval = not current_retrieval_results or not any(current_retrieval_results.values())

                 if needs_fresh_retrieval:
                     logger.info("Pre-search was empty or skipped, retrieving text now...")
                     spinner_msg_retrieve = f"Searching proposal documents for {len(queries_to_retrieve)} query part(s)..."
                     all_query_results = {}
                     fetch_now_func = lambda q: self.query_files(q) # Use self method
                     if is_streamlit:
                         with st.spinner(spinner_msg_retrieve):
                             for i, q_part in enumerate(queries_to_retrieve): # Renamed q to q_part for clarity
                                 logger.info(f"Retrieving text for query part {i+1}/{len(queries_to_retrieve)}: '{q_part[:100]}...'")
                                 results_for_q = fetch_now_func(q_part);
                                 # Combine results, handling potential overlaps later
                                 for file, res_list in results_for_q.items():
                                     if file not in all_query_results: all_query_results[file] = []
                                     all_query_results[file].extend(res_list)
                     else:
                         logger.info(spinner_msg_retrieve)
                         for i, q_part in enumerate(queries_to_retrieve): # Renamed q to q_part
                             logger.info(f"Retrieving text for query part {i+1}/{len(queries_to_retrieve)}: '{q_part[:100]}...'")
                             results_for_q = fetch_now_func(q_part);
                             for file, res_list in results_for_q.items():
                                    if file not in all_query_results: all_query_results[file] = []
                                    all_query_results[file].extend(res_list)

                     # De-duplicate results based on content and source_info, keep best score
                     final_retrieval_results = {};
                     for file, res_list in all_query_results.items():
                          unique_chunks = {};
                          for res in res_list:
                              # Key based on content hash? or first N chars? Let's use content + page/slide
                              source_key_tuple = tuple(sorted(res.get('source_info',{}).items())) # Make source_info hashable
                              key = (res['content'], source_key_tuple, res.get('potential_section'))
                              if key not in unique_chunks or res['score'] < unique_chunks[key]['score']:
                                  unique_chunks[key] = res
                          # Sort unique results by score
                          sorted_unique_results = sorted(unique_chunks.values(), key=lambda x: x['score'])
                          # Apply K limit *per file* after deduplication
                          final_retrieval_results[file] = sorted_unique_results[:self.config.k_retrieval]

                     final_results["retrieval_results"] = final_retrieval_results;
                     current_retrieval_results = final_retrieval_results # Update for aggregation step
                 else:
                      logger.info("Reusing text retrieval results from pre-search step.")
                      # Ensure results format is correct {file: [results]} if reused
                      # (Pre-search already returns this format)
                      current_retrieval_results = final_results["retrieval_results"]


                 # Check again if retrieval found anything after potential fresh fetch
                 if not current_retrieval_results or not any(current_retrieval_results.values()):
                     logger.warning("No relevant text chunks found for the query/sub-queries.")
                     final_results["status"] = "Completed: No relevant information found."
                     prev_error = final_results.get("error")
                     final_results["error"] = "No relevant information found in the text documents (PDF/PPTX)."
                     if prev_error: final_results["error"] += f" (Previous SQL error: {prev_error})"
                     # Use SatSure "don't know" response
                     final_results["answer"] = f"I don't really know but my next versions will be able to answer this for sure! (Reason: {final_results['error']})"
                     return final_results

                 # === 4b. Aggregate text context ===
                 spinner_msg_agg = "Gathering context from proposals..."
                 aggregated_context_data = {}
                 if is_streamlit:
                     with st.spinner(spinner_msg_agg): aggregated_context_data = self.aggregate_context(current_retrieval_results, strategy=context_strategy)
                 else: logger.info(spinner_msg_agg); aggregated_context_data = self.aggregate_context(current_retrieval_results, strategy=context_strategy)
                 final_results["aggregated_context_data"] = aggregated_context_data

                 if not aggregated_context_data or not aggregated_context_data.get("combined_context"):
                      logger.warning(f"Text context aggregation failed or yielded empty context for query: '{refined_query[:100]}...'")
                      final_results["status"] = "Completed: Text context aggregation failed or empty."
                      prev_error = final_results.get("error")
                      final_results["error"] = "Text context aggregation failed or empty."
                      if prev_error: final_results["error"] += f" (Previous error: {prev_error})"
                      # Use SatSure "don't know" response
                      final_results["answer"] = f"I don't really know but my next versions will be able to answer this for sure! (Reason: Failed to gather relevant text context. {('Previous error: ' + prev_error) if prev_error else ''})"
                      return final_results

                 # === 5b. Query LLM with aggregated text context (using refined query AND history) ===
                 spinner_msg_llm = "Generating answer from proposals (considering history)..."
                 answer = ""
                 if is_streamlit:
                      with st.spinner(spinner_msg_llm): answer = self.query_llm_with_history(refined_query, aggregated_context_data, limited_history)
                 else: logger.info(spinner_msg_llm); answer = self.query_llm_with_history(refined_query, aggregated_context_data, limited_history)

                 # Check if answer is valid (not empty, not generic error)
                 is_valid_answer = bool(answer) and answer.strip() != "" and not answer.lower().startswith("error:")

                 if is_valid_answer:
                      final_results["answer"] = answer # Can be the "don't know" response, which is valid
                      # Only mark as fully successful if there wasn't a prior SQL error
                      if not final_results.get("error") or "Failed SQL" not in final_results.get("status", "") : # Check if SQL path wasn't attempted or didn't fail critically before fallback
                         final_results["status"] = "Completed Successfully"
                      else:
                         # Status reflects the fallback nature
                         final_results["status"] = "Completed via Text Fallback after SQL Issue"
                         logger.info("Text RAG fallback successful, but noting prior SQL path error.")
                 else:
                      # Text RAG LLM call failed
                      logger.error(f"Text RAG LLM call failed or returned empty/error: {answer}")
                      final_results["status"] = "Failed during Text RAG LLM call"
                      prev_error = final_results.get("error")
                      final_results["error"] = f"Failed to generate answer using text context. LLM Response: {answer}"
                      if prev_error: final_results["error"] += f" (Previous SQL error: {prev_error})"
                      # Use SatSure "don't know" response as final fallback
                      final_results["answer"] = f"I don't really know but my next versions will be able to answer this for sure! (Reason: LLM failed to generate answer from text. {('Details: ' + final_results['error']) if final_results['error'] else ''})"


            # Final logging
            logger.info(f"--- Finished query execution (SQL+Chat+PDFStruct) for: '{query[:100]}...'. Answer Source: {final_results.get('answer_source')} --- Status: {final_results.get('status')} ---")

        except Exception as e:
            logger.error(f"Unexpected error during run_query_with_history: {e}", exc_info=True)
            final_results["status"] = "Failed"; final_results["error"] = f"Unexpected runtime error: {str(e)}"
            # Use SatSure "don't know" response for unexpected errors too
            final_results["answer"] = f"I don't really know but my next versions will be able to answer this for sure! (Reason: An unexpected error occurred: {e})"

        return final_results
    # ==========================================================================
    # END OF REVISED run_query_with_history FUNCTION
    # ==========================================================================


# --- START OF STREAMLIT UI CODE ---
if __name__ == "__main__":
    try:
        import streamlit as st
        from packaging.version import parse as parse_version # For version comparison
    except ImportError:
        logger.error("Streamlit or packaging not installed.")
        sys.exit(1)

    # Informational version check for st.dialog context manager,
    # but the code now uses st.sidebar.expander as the primary UI for document management.
    MIN_STREAMLIT_VERSION_FOR_DIALOG_CONTEXT_MANAGER = "1.28.0"
    if parse_version(st.__version__) < parse_version(MIN_STREAMLIT_VERSION_FOR_DIALOG_CONTEXT_MANAGER):
        logger.warning(
            f"Current Streamlit version {st.__version__} is older than {MIN_STREAMLIT_VERSION_FOR_DIALOG_CONTEXT_MANAGER}. "
            f"st.dialog as a context manager might not be available. This app uses st.sidebar.expander as a fallback."
        )

    st.set_page_config(layout="wide", page_title="SatSure Conversational RAG Q&A")

    st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

        :root {
            --primary-color: #4A90E2; /* SatSure Blue */
            --secondary-color: #50E3C2; /* Teal Accent */
            --background-color: #121212; /* Even Darker Background */
            --surface-color: #1E1E1E; /* Dark Surface */
            --text-color: #E0E0E0;
            --text-muted-color: #A0A0A0;
            --border-color: #3A3A3A;
            --success-color: #50E3C2;
            --warning-color: #F5A623;
            --error-color: #D0021B;
            --font-family: 'Inter', sans-serif;
            --primary-color-rgb: 74, 144, 226; 
            --success-color-rgb: 80, 227, 194;
            --warning-color-rgb: 245, 166, 35;
            --error-color-rgb: 208, 2, 27;
            --text-muted-color-rgb: 160, 160, 160;
        }

        body {
            font-family: var(--font-family);
            color: var(--text-color);
            background-color: var(--background-color);
        }
        
        .main .block-container {
            padding-top: 2rem; padding-bottom: 2rem;
            padding-left: 1rem; padding-right: 1rem;
            max-width: 1200px; margin: 0 auto;
        }

        h1, .st-emotion-cache-10trblm { 
            color: var(--primary-color); font-weight: 700; letter-spacing: -1px;
            text-align: center; margin-bottom: 0.3em; font-size: 2.8em;
            text-shadow: 0 2px 4px rgba(0,0,0,0.2);
        }
        .st-caption, .st-emotion-cache-t3byc2 {
            color: var(--text-muted-color); text-align: center;
            font-size: 1em; margin-bottom: 2.5em;
        }
        
        /* Sidebar */
        .st-emotion-cache-16txtl3 { 
            background-color: var(--surface-color);
            border-right: 1px solid var(--border-color);
            box-shadow: 3px 0 15px rgba(0,0,0,0.2);
        }
        .st-emotion-cache-16txtl3 .st-emotion-cache-1cypcdb { /* Sidebar Title / Expander Header */
            color: var(--primary-color); font-weight: 600; padding-top: 1rem;
        }
        .st-emotion-cache-16txtl3 .stButton button { /* General buttons in sidebar */
            background-color: var(--primary-color); color: white;
            border: none; border-radius: 8px; width: 100%;
            transition: background-color 0.3s ease, transform 0.2s ease;
            padding: 10px 0; font-weight: 500;
        }
        .st-emotion-cache-16txtl3 .stButton button:hover {
            background-color: #357ABD; transform: translateY(-2px);
        }
        .st-emotion-cache-16txtl3 .stAlert { /* Alerts in sidebar */
            border-radius: 6px; margin-top: 0.5rem;
        }
        .st-emotion-cache-16txtl3 .stAlert[data-testid="stNotificationContentSuccess"] {
            background-color: rgba(var(--success-color-rgb), 0.1);
            border-left: 4px solid var(--success-color); color: var(--success-color);
        }
        .st-emotion-cache-16txtl3 .stAlert[data-testid="stNotificationContentWarning"] {
            background-color: rgba(var(--warning-color-rgb), 0.1);
            border-left: 4px solid var(--warning-color); color: var(--warning-color);
        }
        .st-emotion-cache-16txtl3 .stAlert[data-testid="stNotificationContentError"] {
            background-color: rgba(var(--error-color-rgb), 0.1);
            border-left: 4px solid var(--error-color); color: var(--error-color);
        }
        .st-emotion-cache-16txtl3 .stInfo { /* Info boxes in sidebar */
            background-color: rgba(var(--primary-color-rgb), 0.08);
            border-left: 4px solid var(--primary-color); color: var(--text-muted-color);
        }
         /* Expander in Sidebar specific styling */
        .st-emotion-cache-16txtl3 .stExpander {
            border: 1px solid var(--border-color);
            border-radius: 8px;
            margin-top: 1rem; /* Space above expander */
            background-color: rgba(var(--surface-color), 0.5); /* Slightly different background for expander */
        }
        .st-emotion-cache-16txtl3 .stExpander summary {
            color: var(--primary-color);
            font-weight: 600; /* Make expander header bold */
            padding: 0.75rem 1rem;
        }
        .st-emotion-cache-16txtl3 .stExpander summary:hover {
            background-color: rgba(var(--primary-color-rgb), 0.1);
        }
        .st-emotion-cache-16txtl3 .stExpander > div { /* Content of expander */
            padding: 0.5rem 1rem 1rem 1rem;
        }


        /* Chat Input */
        .stChatInputContainer {
            background-color: var(--surface-color);
            border-top: 1px solid var(--border-color);
            padding: 1rem 1.5rem; margin: 0 -1rem -2rem -1rem; 
        }
        .stChatInputContainer > div > div > textarea {
            background-color: #282828; color: var(--text-color);
            border: 1px solid #444; border-radius: 10px;
            padding: 12px 15px; box-shadow: inset 0 1px 3px rgba(0,0,0,0.2);
            transition: border-color 0.3s ease, box-shadow 0.3s ease;
        }
        .stChatInputContainer > div > div > textarea:focus {
            border-color: var(--primary-color);
            box-shadow: inset 0 1px 3px rgba(0,0,0,0.2), 0 0 0 3px rgba(var(--primary-color-rgb), 0.3);
        }
        .stChatInputContainer button {
            background-color: var(--primary-color) !important; color: white !important;
            border-radius: 8px !important; transition: background-color 0.3s ease;
            padding: 0.6rem 1rem !important;
        }
        .stChatInputContainer button:hover { background-color: #357ABD !important; }

        /* Chat Messages */
        .stChatMessage {
            background-color: var(--surface-color);
            border-radius: 12px; padding: 0.8em 1.2em; margin-bottom: 1em;
            box-shadow: 0 4px 10px rgba(0,0,0,0.2);
            border: 1px solid var(--border-color);
            opacity: 0; transform: translateY(20px);
            animation: slideInFade 0.6s cubic-bezier(0.25, 0.46, 0.45, 0.94) forwards;
            max-width: 80%; 
        }
        .stChatMessage[data-testid="stChatMessageContent"] { color: var(--text-color); }
        .stChatMessage p, .stChatMessage ul, .stChatMessage ol { line-height: 1.6; }
        .stChatMessage strong { color: var(--secondary-color); }
        .stChatMessage code {
            background-color: rgba(var(--primary-color-rgb), 0.15); color: var(--secondary-color);
            padding: 0.2em 0.4em; border-radius: 4px; font-size: 0.9em;
        }
        .stChatMessage pre {
            background-color: #111 !important; padding: 0.8em; border-radius: 6px;
            border: 1px solid #333;
        }
        .stChatMessage pre code { background-color: transparent !important; color: var(--text-color) !important; }
        
        .stChatMessage[avatar_style="user"] {
            background: linear-gradient(135deg, var(--primary-color) 0%, #357ABD 100%);
            color: white; margin-left: auto;
            animation-name: slideInFadeUser;
            border: none;
        }
        .stChatMessage[avatar_style="user"] p,
        .stChatMessage[avatar_style="user"] ul,
        .stChatMessage[avatar_style="user"] ol,
        .stChatMessage[avatar_style="user"] strong { color: white !important; }
        .stChatMessage[avatar_style="user"] code {
            background-color: rgba(255,255,255, 0.2); color: var(--secondary-color);
        }
        .stChatMessage[avatar_style="user"] pre { background-color: rgba(0,0,0, 0.25) !important; }
        .stChatMessage[avatar_style="user"] pre code { color: #f0f0f0 !important; }
        
        .message-footer {
            font-size: 0.75em;
            color: var(--text-muted-color);
            margin-top: 0.8em;
            padding-top: 0.5em;
            border-top: 1px dashed rgba(var(--text-muted-color-rgb), 0.3);
        }
        .stChatMessage[avatar_style="user"] .message-footer {
            color: rgba(255,255,255,0.7);
            border-top: 1px dashed rgba(255,255,255,0.3);
        }

        @keyframes slideInFade { 0% { opacity: 0; transform: translateY(20px); } 100% { opacity: 1; transform: translateY(0); } }
        @keyframes slideInFadeUser { 0% { opacity: 0; transform: translateY(20px); } 100% { opacity: 1; transform: translateY(0); } }

        /* Expander for Analysis (in chat message) */
        .stChatMessage .stExpander { /* Target expanders specifically within chat messages */
            border: 1px solid var(--border-color); border-radius: 8px;
            margin-top: 0.5rem; background-color: var(--surface-color);
        }
        .stChatMessage .stExpander summary {
            color: var(--primary-color); font-weight: 500; padding: 0.5rem 1rem;
        }
        .stChatMessage .stExpander summary:hover { background-color: #2a2a2a; }
        .stChatMessage .stExpander > div { padding: 0.5rem 1rem 1rem 1rem; } 
        .stChatMessage .stExpander .stCaption { text-align: left; margin-bottom: 0.3em; font-size: 0.85em;}
        .stChatMessage .stExpander .stCodeBlock, .stExpander pre { background-color: #111 !important; }
        
        /* Scrollbars */
        ::-webkit-scrollbar { width: 10px; height: 10px; }
        ::-webkit-scrollbar-track { background: var(--surface-color); border-radius: 10px; }
        ::-webkit-scrollbar-thumb { background: #444; border-radius: 10px; border: 2px solid var(--surface-color); }
        ::-webkit-scrollbar-thumb:hover { background: var(--primary-color); }

        .stSpinner > div { 
            border-top-color: var(--primary-color) !important;
            border-right-color: var(--primary-color) !important;
        }
    </style>
    <script>
        function setRGBColorVar(cssVarName, rgbVarName) {
            const colorVal = getComputedStyle(document.documentElement).getPropertyValue(cssVarName).trim();
            if (colorVal) {
                let r, g, b;
                if (colorVal.startsWith('#')) {
                    const hex = colorVal.substring(1);
                    if (hex.length === 3) { 
                        r = parseInt(hex.substring(0, 1) + hex.substring(0, 1), 16);
                        g = parseInt(hex.substring(1, 2) + hex.substring(1, 2), 16);
                        b = parseInt(hex.substring(2, 3) + hex.substring(2, 3), 16);
                    } else if (hex.length === 6) {
                        r = parseInt(hex.substring(0, 2), 16);
                        g = parseInt(hex.substring(2, 4), 16);
                        b = parseInt(hex.substring(4, 6), 16);
                    }
                } else if (colorVal.startsWith('rgb')) { 
                    const parts = colorVal.match(/\\d+/g);
                    if (parts && parts.length >= 3) { r = parseInt(parts[0]); g = parseInt(parts[1]); b = parseInt(parts[2]); }
                }
                if (r !== undefined && g !== undefined && b !== undefined) { 
                    document.documentElement.style.setProperty(rgbVarName, `${r}, ${g}, ${b}`);
                }
            }
        }
        document.addEventListener('DOMContentLoaded', function() {
            setRGBColorVar('--primary-color', '--primary-color-rgb');
            setRGBColorVar('--success-color', '--success-color-rgb');
            setRGBColorVar('--warning-color', '--warning-color-rgb');
            setRGBColorVar('--error-color', '--error-color-rgb');
            setRGBColorVar('--text-muted-color', '--text-muted-color-rgb');
        });
    </script>
    """, 
    unsafe_allow_html=True)

    st.title("💬 SatSure Proposal Q&A System")
    st.caption("Query Project Proposals (PDFs/PPTX) & Supporting Data (XLSX/CSV) with Chat History (v9-UIFix-Expander)")

    @st.cache_resource
    def load_rag_system(config_path="config.ini"):
        groq_key = os.getenv("GROQ_API_KEY")
        if not groq_key:
            st.error("🚨 GROQ_API_KEY not set in environment variables! The application cannot run effectively."); 
            st.info("Please set your GROQ_API_KEY and restart the application.")
            st.stop()
        
        system = None
        with st.spinner("Initializing RAG System & loading embedding model... This might take a moment."):
            try:
                system = RAGSystem(config_path=config_path)
            except Exception as e:
                logger.error(f"Fatal RAG Initialization Error during RAGSystem instantiation: {e}", exc_info=True)
                raise # Re-raise to be caught by the main app's try-except
        
        if not system or not system.db_conn: 
            st.error("🚨 RAG System or SQLite connection failed critically during initialization! Check logs for details."); 
            st.stop()
        
        if not system.deepseek_api_key: 
            st.sidebar.warning("⚠️ DEEPSEEK_API_KEY missing. Spreadsheet metadata quality may be affected.")
        
        st.sidebar.success("✅ RAG System Initialized.")
        return system
    
    # MODIFIED FUNCTION: Uses st.sidebar.expander instead of st.dialog
    def display_document_management_ui(_rag_system):
        # This UI is now part of the sidebar, within an expander.
        with st.sidebar.expander("📂 Document Management & Processing", expanded=False): # `expanded=True` if you want it open by default
            st.markdown(f"**Data Folder:** `{os.path.abspath(_rag_system.config.data_dir)}`")
            st.markdown(f"**Index Folder:** `{os.path.abspath(_rag_system.config.index_dir)}`")
            st.markdown(f"**DB Location:** `{_rag_system.config.sqlite_db_path}`")
            st.markdown("---")

            status_placeholder_sidebar = st.empty() 

            def sidebar_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
                log_prefix = "✅ " if is_done else "❌ " if is_error else "⚠️ " if is_warning else f"⏳ ({stage}) " if stage else "⏳ "
                progress_val = 1.0 if is_done else 0.0
                if total_steps and current_step is not None:
                    stage_prog = {"Extracting Text": 0.1, "Chunking Text": 0.15, "Verifying Index": 0.2, "Embedding Text": 0.3, "Indexing Text": 0.4,
                                "Loading Sheet": 0.5, "Analyzing Columns": 0.7, "Loading SQL": 0.9}.get(stage, 0.0)
                    progress_val = min(1.0, (current_step + stage_prog) / total_steps) if total_steps > 0 else 0.0
                
                full_message = f"{log_prefix}{message}"
                try:
                    if 0 < progress_val < 1.0 and not is_error and not is_warning:
                        status_placeholder_sidebar.progress(progress_val, text=full_message)
                    else: # Final status or non-progress update
                        if is_error: status_placeholder_sidebar.error(full_message)
                        elif is_warning: status_placeholder_sidebar.warning(full_message)
                        elif is_done: status_placeholder_sidebar.success(full_message)
                        else: status_placeholder_sidebar.info(full_message) 
                except Exception as e: logger.warning(f"Streamlit UI update error in sidebar callback: {e}")

            if st.button("🔄 Process/Re-process Documents", key="process_docs_sidebar_expander", use_container_width=True):
                st.session_state.sidebar_processed_files_list = [] 
                st.session_state.sidebar_processing_status_message = "Starting processing..." 
                st.session_state.is_ready_for_chat = False 
                status_placeholder_sidebar.info(st.session_state.sidebar_processing_status_message)
                
                with st.spinner("Processing documents... This may take a while."):
                    try:
                        processing_successful = _rag_system.process_files(progress_callback=sidebar_progress_callback)
                        st.session_state.sidebar_processed_files_list = sorted(list(_rag_system.processed_files))
                        text_idx_count = len(_rag_system.faiss_indexes); sql_tbl_count = len(_rag_system.table_metadata)
                        processed_count = len(st.session_state.sidebar_processed_files_list)
                        final_msg = f"Processed {processed_count} file(s). Text Indices: {text_idx_count}, SQL Tables: {sql_tbl_count}."
                        
                        if processed_count > 0 and processing_successful: 
                            st.session_state.sidebar_processing_status_message = f"✅ Ready! {final_msg}"
                            st.session_state.is_ready_for_chat = True 
                        elif processing_successful: 
                            st.session_state.sidebar_processing_status_message = f"⚠️ Finished, but no supported documents processed. {final_msg}"
                            st.session_state.is_ready_for_chat = False 
                        else: 
                            st.session_state.sidebar_processing_status_message = f"❌ Processing function reported failure. Check logs."
                            st.session_state.is_ready_for_chat = False
                        
                        # Update the placeholder immediately after processing finishes
                        if st.session_state.is_ready_for_chat and "✅ Ready!" in st.session_state.sidebar_processing_status_message :
                             status_placeholder_sidebar.success(st.session_state.sidebar_processing_status_message)
                        elif "⚠️" in st.session_state.sidebar_processing_status_message:
                             status_placeholder_sidebar.warning(st.session_state.sidebar_processing_status_message)
                        elif "❌" in st.session_state.sidebar_processing_status_message:
                             status_placeholder_sidebar.error(st.session_state.sidebar_processing_status_message)
                        else: # Default to info if no clear status indicator in message (should ideally not happen here)
                            status_placeholder_sidebar.info(st.session_state.sidebar_processing_status_message)


                        st.rerun() # Rerun to update main page based on is_ready_for_chat and other states
                    except Exception as e:
                        logger.error(f"Fatal error during document processing (sidebar): {e}", exc_info=True)
                        st.session_state.sidebar_processing_status_message = f"❌ Fatal error: {e}."
                        status_placeholder_sidebar.error(st.session_state.sidebar_processing_status_message)
                        st.session_state.sidebar_processed_files_list = []
                        st.session_state.is_ready_for_chat = False
                        st.rerun() # Rerun to reflect error state
            else:
                # Display stored status message when not actively processing
                current_msg = st.session_state.get("sidebar_processing_status_message", "Click 'Process/Re-process Documents' to load or update data.")
                if "✅ Ready!" in current_msg: status_placeholder_sidebar.success(current_msg)
                elif "❌" in current_msg: status_placeholder_sidebar.error(current_msg)
                elif "⚠️" in current_msg: status_placeholder_sidebar.warning(current_msg)
                else: status_placeholder_sidebar.info(current_msg)

            if st.session_state.get("sidebar_processed_files_list"): 
                st.markdown("--- \n**Loaded Assets:**")
                assets_df_data = []
                for fname in st.session_state.sidebar_processed_files_list:
                    asset_types = []
                    if fname in _rag_system.faiss_indexes: asset_types.append("Text Index")
                    if fname in _rag_system.file_to_table_map and _rag_system.file_to_table_map.get(fname):
                        tables = _rag_system.file_to_table_map.get(fname, [])
                        if tables: asset_types.append(f"SQL ({len(tables)} table{'s' if len(tables) > 1 else ''})")
                    assets_df_data.append({"File": fname, "Type(s)": ", ".join(asset_types) or "N/A"})
                if assets_df_data: st.dataframe(pd.DataFrame(assets_df_data), use_container_width=True, hide_index=True)
            
        return st.session_state.get("is_ready_for_chat", False)

    try:
        rag_sys = load_rag_system() 

        # Initialize session states if they don't exist
        if "messages" not in st.session_state: 
            st.session_state.messages = []
        if "analysis_details" not in st.session_state: 
            st.session_state.analysis_details = {}
        if "is_ready_for_chat" not in st.session_state:
            st.session_state.is_ready_for_chat = False 
        if "sidebar_processed_files_list" not in st.session_state: # For expander UI
            st.session_state.sidebar_processed_files_list = []
        if "sidebar_processing_status_message" not in st.session_state: # For expander UI
            st.session_state.sidebar_processing_status_message = "Click 'Process/Re-process Documents' within 'Document Management' to load or update data."
        if "query_start_time" not in st.session_state: 
            st.session_state.query_start_time = time.time()

        st.sidebar.title("Controls & Info")
        if st.sidebar.button("✨ New Chat", use_container_width=True, key="new_chat_button"):
            st.session_state.messages = []
            st.session_state.analysis_details = {}
            st.rerun()

        # Call the UI function for document management (now an expander in sidebar)
        is_ready = display_document_management_ui(rag_sys) 

        st.sidebar.markdown("---")
        st.sidebar.subheader("System Configuration")
        config_details = {
            "LLM": rag_sys.config.llm_model,
            "Encoder": rag_sys.config.encoder_model,
            "Metadata LLM": rag_sys.config.deepseek_model,
            "History Turns": rag_sys.config.max_chat_history_turns,
            "Retr. K (Text)": rag_sys.config.k_retrieval,
            "SQL Conf.": f"{rag_sys.config.dataframe_query_confidence_threshold:.2f}",
            "Chunk Size": rag_sys.config.chunk_size,
            "Overlap": rag_sys.config.overlap,
        }
        for key, value in config_details.items():
            st.sidebar.markdown(f"<small><b>{key}:</b> {value}</small>", unsafe_allow_html=True)

        if not is_ready:
             st.info("⬅️ Please process documents using 'Document Management & Processing' in the sidebar to enable chat.")
             st.stop() 
        
        chat_container = st.container() 
        with chat_container:
            for i, message in enumerate(st.session_state.messages):
                avatar_icon = "👤" if message["role"] == "user" else "🛰️"
                with st.chat_message(message["role"], avatar=avatar_icon): 
                    content_part = message["content"]
                    footer_part = ""
                    
                    if message["role"] == "assistant" and "\n\n---\n*" in message["content"]:
                        parts = message["content"].split("\n\n---\n*", 1)
                        content_part = parts[0]
                        if len(parts) > 1:
                            footer_part = f"<div class='message-footer'>*{parts[1]}</div>"
                    
                    st.markdown(content_part, unsafe_allow_html=True)
                    if footer_part:
                        st.markdown(footer_part, unsafe_allow_html=True)

                    if message["role"] == "assistant" and i == len(st.session_state.messages) - 1 and st.session_state.get('analysis_details', {}).get('has_details'):
                        with st.expander("Show Analysis & Context", expanded=False):
                            details = st.session_state.analysis_details
                            if details.get("Original Query"): st.caption(f"**Original Query:** {details['Original Query']}")
                            if details.get("Refined Query") and details["Refined Query"] != details.get("Original Query"): st.caption(f"**Refined Query:** {details['Refined Query']}")
                            if details.get("Classification"): st.caption(f"**Classification:** {details['Classification']}")
                            if details.get("Answer Source"): st.caption(f"**Answer Source:** {details['Answer Source']}")
                            if details.get("Status"): st.caption(f"**Status:** {details['Status']}")
                            if details.get("SQL Table"):
                                st.caption(f"**SQL Table:** `{details['SQL Table']}`")
                                st.caption("**SQL Query Generated:**"); st.code(details.get("SQL Query", "N/A"), language='sql')
                            # Ensure "Text Sources Used" is handled correctly if it's None or empty
                            text_sources = details.get("Text Sources Used")
                            if text_sources: 
                                st.caption(f"**Text Sources Used:** {', '.join(text_sources)}")
                            elif details.get("Answer Source") == "Text RAG": # If source was Text RAG but no files listed (e.g. aggregation failed)
                                st.caption(f"**Text Sources Used:** None or aggregation failed.")

                            if details.get("Error"): st.caption(f"**Note/Error:** {details['Error']}")

        if prompt := st.chat_input("Ask a question about the SatSure proposals..."):
            st.session_state.messages.append({"role": "user", "content": prompt})
            st.session_state.query_start_time = time.time() 
            
            with st.spinner("🛰️ SatSure AI is thinking..."):
                history_for_rag = [msg for msg in st.session_state.messages[:-1] if isinstance(msg, dict) and "role" in msg and "content" in msg]
                results_data = rag_sys.run_query_with_history(prompt, history_for_rag, context_strategy="top_k")
                
                query_end_time = time.time()
                elapsed_time = query_end_time - st.session_state.query_start_time
                elapsed_time_str = f"{elapsed_time:.2f}"

                answer = results_data.get("answer", "Sorry, I encountered an error.")
                answer_source = results_data.get('answer_source', 'N/A')
                query_status = results_data.get("status", "Unknown")
                query_error = results_data.get("error")
                
                full_response = answer
                response_footer_parts = [f"Answer generated via {answer_source} in {elapsed_time_str}s."]
                is_satsure_dont_know = answer.strip().startswith("I don't really know but my next versions")

                if "Completed Successfully" not in query_status or (query_error and not is_satsure_dont_know):
                    if query_status != "Completed Successfully": 
                        response_footer_parts.append(f"Status: {query_status}")
                    if query_error and query_error not in full_response and not is_satsure_dont_know : 
                        response_footer_parts.append(f"Note: {query_error}")
                
                if not is_satsure_dont_know or len(response_footer_parts) > 1 : 
                     full_response += f"\n\n---\n*{' | '.join(response_footer_parts)}"
            
            st.session_state.messages.append({"role": "assistant", "content": full_response})
            
            analysis_info = {
                "has_details": True, "Original Query": results_data.get("query"),
                "Refined Query": results_data.get("refined_query"),
                "Classification": f"{results_data.get('classification')} (Conf: {results_data.get('classification_confidence', 0):.2f})",
                "Answer Source": answer_source, "Status": query_status,
                "SQL Table": results_data.get("target_table"), "SQL Query": results_data.get("sql_query_generated"),
                "Text Sources Used": results_data.get("aggregated_context_data", {}).get("source_files"),
                "Error": query_error
            }
            st.session_state.analysis_details = {k: v for k, v in analysis_info.items() if v is not None and str(v).strip() != ""}
            
            st.rerun()

    except Exception as e:
        st.error(f"Streamlit App Error: {e}")
        logger.exception("Streamlit application error:")
        log_file_path = 'rag_system_sql_chat_pdfstruct.log' 
        try:
            if 'rag_sys' in locals() and rag_sys and hasattr(rag_sys, 'config'): 
                log_file_path = rag_sys.config.log_file
        except Exception: pass 
        st.info(f"Check console or `{os.path.abspath(log_file_path)}` for details.")
        st.stop()

# --- END OF FILE ragtest_c2-7_modified.py ---
# --- END OF FILE ragtest_c2-8.py ---