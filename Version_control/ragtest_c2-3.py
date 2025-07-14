# --- START OF FILE ragtest_c3_query_adapt.py --- # Renamed for clarity

# --- START OF FILE ragtest_c.py --- # <-- Original comment

import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd # Added for CSV/XLSX
import faiss
# import streamlit as st # Import only if needed / within UI code
import time
from sentence_transformers import SentenceTransformer
# from transformers import T5ForConditionalGeneration, T5Tokenizer # Keep if needed elsewhere
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage # Added SystemMessage
from dotenv import load_dotenv
from pptx import Presentation # Added for PPTX
from pptx.enum.shapes import MSO_SHAPE_TYPE # Needed for placeholder check
import sys # Import sys for checking streamlit context

# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted from RAG2.py) ---

# Configure logging (optional for Streamlit, but good for debugging)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system.log"),
        # logging.StreamHandler() # Can be noisy in Streamlit console
    ]
)
logger = logging.getLogger(__name__)


class RAGConfig:
    """Configuration class for RAG system."""
    # (This class remains unchanged from ragtest_c2_merged.py)
    def __init__(self, config_path="config.ini"):
        """Initialize configuration from config file or defaults."""
        self.config = configparser.ConfigParser()
        self.defaults = {
            "PATHS": {"data_dir": "./Data/", "index_dir": "./Faiss_index/", "log_file": "rag_system.log"},
            "MODELS": {"encoder_model": "sentence-transformers/all-MiniLM-L6-v2", "llm_model": "meta-llama/llama-4-scout-17b-16e-instruct", "device": "auto"},
            "PARAMETERS": {"chunk_size": "200", "overlap": "50", "k_retrieval": "5", "temperature": "0.2", "max_context_tokens": "4000", "max_chars_per_element": "1000", "pptx_merge_threshold_words": "50"},
            "SUPPORTED_EXTENSIONS": {"extensions": ".pdf, .xlsx, .csv, .pptx"}
        }
        if os.path.exists(config_path):
            try: self.config.read(config_path); logger.info(f"Loaded configuration from {config_path}"); self._ensure_defaults()
            except Exception as e: logger.error(f"Error reading config file {config_path}: {e}"); self._set_defaults()
        else: logger.warning(f"Config file {config_path} not found, using defaults"); self._set_defaults()
    def _set_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section)
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value)
    def _ensure_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section); logger.info(f"Added missing section [{section}] from defaults.")
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value); logger.info(f"Added missing option '{option} = {value}' to section [{section}] from defaults.")
    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    @property
    def llm_model(self): return self.config.get("MODELS", "llm_model")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto")
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=200)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=50)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=5)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.2)
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=4000)
    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1000)
    @property
    def pptx_merge_threshold_words(self): return self.config.getint("PARAMETERS", "pptx_merge_threshold_words", fallback=50)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf")
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])


class RAGSystem:
    """Retrieval-Augmented Generation system for various document types."""
    # (Init remains unchanged from ragtest_c2_merged.py)
    def __init__(self, config_path="config.ini"):
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system. Device: {self.config.device}, Supported files: {self.config.supported_extensions}")
        self._is_streamlit = "streamlit" in sys.modules
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)
        try:
            if self._is_streamlit:
                import streamlit as st
                with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."): self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            else: logger.info(f"Loading embedding model ({self.config.encoder_model})..."); self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            logger.info("Local embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}"); st.stop()
            else: raise
        self.file_chunks = {}
        self.faiss_indexes = {}
        self.processed_files = set()
        # Initialize Groq LLM client once if needed frequently
        self._llm = None


    def _get_llm(self):
        """Initializes and returns the Groq LLM client."""
        if self._llm is None:
            api_key = os.getenv("GROQ_API_KEY")
            if not api_key:
                logger.error("GROQ_API_KEY not found.")
                if self._is_streamlit:
                    import streamlit as st
                    st.error("Error: GROQ_API_KEY is not set.")
                raise ValueError("GROQ_API_KEY not configured.")
            try:
                self._llm = ChatGroq(
                    temperature=self.config.temperature,
                    model_name=self.config.llm_model,
                    groq_api_key=api_key,
                )
                logger.info(f"Groq LLM client initialized ({self.config.llm_model})")
            except Exception as e:
                 logger.error(f"Failed to initialize Groq client: {e}", exc_info=True)
                 if self._is_streamlit:
                     import streamlit as st
                     st.error(f"Failed to initialize LLM. Error: {e}")
                 raise
        return self._llm


    def clean_text(self, text):
        """Clean and normalize text."""
        # (This function remains unchanged)
        if not isinstance(text, str): text = str(text)
        text = re.sub(r"\(cid:.*?\)", "", text)
        text = re.sub(r"\s+", " ", text).strip()
        text = text.replace('\n', ' ').replace('\r', '')
        return text

    # --- Specific Content Extractors ---
    # (_extract_pdf, _extract_xlsx, _extract_csv, _extract_pptx remain unchanged from ragtest_c2_merged.py)
    def _extract_pdf(self, file_path):
        all_content = []
        base_filename = os.path.basename(file_path)
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Extracting from PDF: {base_filename} ({total_pages} pages)")
                page_bar = None
                if self._is_streamlit and total_pages > 10: import streamlit as st; page_bar = st.progress(0, text=f"Extracting pages from {base_filename}...")
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text(layout="normal") or ""; cleaned_text = self.clean_text(text)
                    if cleaned_text: all_content.append({"type": "text", "content": cleaned_text, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                    try:
                        for table in page.extract_tables():
                             if table:
                                 table_df = pd.DataFrame(table).fillna(''); table_string = table_df.to_string(index=False, header=True); cleaned_table_string = self.clean_text(table_string)
                                 if cleaned_table_string: all_content.append({"type": "table", "content": cleaned_table_string, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                    except Exception as e_table: logger.warning(f"Could not extract table on page {page_num + 1} in {base_filename}: {e_table}")
                    if page_bar: progress_percent = min(1.0, (page_num + 1) / total_pages); page_bar.progress(progress_percent, text=f"Extracting page {page_num + 1}/{total_pages} from {base_filename}...")
                if page_bar: page_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PDF: {base_filename}"); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PDF {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def _extract_xlsx(self, file_path):
        all_content = []
        base_filename = os.path.basename(file_path)
        max_chars = self.config.max_chars_per_element
        try:
            excel_file = pd.ExcelFile(file_path); logger.info(f"Extracting from XLSX: {base_filename} (Sheets: {', '.join(excel_file.sheet_names)})")
            sheet_bar = None; total_sheets = len(excel_file.sheet_names)
            if self._is_streamlit and total_sheets > 1: import streamlit as st; sheet_bar = st.progress(0, text=f"Extracting sheets from {base_filename}...")
            for i, sheet_name in enumerate(excel_file.sheet_names):
                try:
                    df = excel_file.parse(sheet_name); df = df.fillna(''); sheet_string = df.to_string(index=False, header=True, na_rep=''); cleaned_content = self.clean_text(sheet_string)
                    if cleaned_content:
                        if len(cleaned_content) > max_chars * 10: logger.warning(f"Sheet '{sheet_name}' in {base_filename} is very large, truncating."); cleaned_content = cleaned_content[:max_chars*10] + "... (truncated sheet)"
                        all_content.append({"type": "sheet_data", "content": cleaned_content, "source_info": {"sheet": sheet_name, "rows": f"1-{len(df)}"}, "file_type": "xlsx"})
                    logger.debug(f"Extracted content from sheet '{sheet_name}' in {base_filename}")
                except Exception as e_sheet: logger.warning(f"Could not process sheet '{sheet_name}' in {base_filename}: {e_sheet}")
                if sheet_bar: sheet_bar.progress(min(1.0, (i + 1) / total_sheets), text=f"Extracting sheet {i+1}/{total_sheets} ('{sheet_name}') from {base_filename}...")
            if sheet_bar: sheet_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks (sheets) from XLSX: {base_filename}"); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from XLSX {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def _extract_csv(self, file_path):
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element
        try:
            encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']; df = None
            for enc in encodings_to_try:
                 try: df = pd.read_csv(file_path, encoding=enc, low_memory=False); logger.info(f"Read CSV {base_filename} using encoding: {enc}"); break
                 except UnicodeDecodeError: continue
                 except Exception as e_read: logger.warning(f"Error reading CSV {base_filename} with encoding {enc}: {e_read}")
            if df is None: logger.error(f"Could not read CSV {base_filename}"); return []
            df = df.fillna(''); csv_string = df.to_string(index=False, header=True, na_rep=''); cleaned_content = self.clean_text(csv_string)
            if cleaned_content:
                if len(cleaned_content) > max_chars * 20: logger.warning(f"CSV {base_filename} is very large, truncating."); cleaned_content = cleaned_content[:max_chars*20] + "... (truncated CSV)"
                all_content.append({"type": "csv_data", "content": cleaned_content, "source_info": {"rows": f"1-{len(df)}"}, "file_type": "csv"})
            logger.info(f"Extracted content block from CSV: {base_filename} ({len(df)} rows)"); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from CSV {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def _extract_pptx(self, file_path):
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element; merge_threshold_words = self.config.pptx_merge_threshold_words
        try:
            prs = Presentation(file_path); logger.info(f"Extracting from PPTX: {base_filename} ({len(prs.slides)} slides) using merge strategy (threshold: {merge_threshold_words} words).")
            slide_bar = None; total_slides = len(prs.slides)
            if self._is_streamlit and total_slides > 1: import streamlit as st; slide_bar = st.progress(0, text=f"Extracting slides from {base_filename}...")
            pending_title_slide_data = None
            for i, slide in enumerate(prs.slides):
                current_slide_number = i + 1; current_slide_title_text = ""; current_slide_other_texts = []; current_slide_has_title_placeholder = False; title_shape = None
                try:
                    if slide.shapes.title: title_shape = slide.shapes.title; current_slide_has_title_placeholder = True
                except AttributeError: pass
                if title_shape and title_shape.has_text_frame: cleaned_title = self.clean_text(title_shape.text_frame.text); current_slide_title_text = cleaned_title if cleaned_title else ""
                for shape in slide.shapes:
                    if shape == title_shape: continue
                    is_placeholder = shape.is_placeholder; is_body_placeholder = False
                    if is_placeholder:
                        try: ph_type = shape.placeholder_format.type; is_body_placeholder = ph_type in [MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT, MSO_SHAPE_TYPE.SUBTITLE]
                        except AttributeError: pass
                    if shape.has_text_frame:
                        text = shape.text_frame.text; cleaned = self.clean_text(text)
                        if cleaned:
                             if len(cleaned) > max_chars: cleaned = cleaned[:max_chars] + "...(truncated shape)"
                             prefix = "[Body]: " if is_body_placeholder else ""; current_slide_other_texts.append(prefix + cleaned)
                if slide.has_notes_slide:
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text; cleaned_notes = self.clean_text(notes_text)
                        if cleaned_notes:
                            if len(cleaned_notes) > max_chars * 2: cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated notes)"
                            current_slide_other_texts.append(f"[Notes]: {cleaned_notes}")
                    except Exception as e_notes: logger.warning(f"Could not extract notes from slide {current_slide_number}: {e_notes}")
                current_slide_full_content = current_slide_title_text
                if current_slide_other_texts: current_slide_full_content += ("\n" if current_slide_full_content else "") + "\n".join(current_slide_other_texts)
                current_slide_full_content = current_slide_full_content.strip(); other_text_word_count = sum(len(s.split()) for s in current_slide_other_texts)
                merged_content_block = None; should_merge = False
                if pending_title_slide_data:
                    if pending_title_slide_data['has_title'] and pending_title_slide_data['other_words'] <= merge_threshold_words and current_slide_full_content:
                        should_merge = True; logger.info(f"Merging slide {pending_title_slide_data['number']} with slide {current_slide_number}.")
                if should_merge:
                    merged_text = f"[Content from Slide {pending_title_slide_data['number']}]:\n{pending_title_slide_data['content']}\n\n---\n\n[Content from Slide {current_slide_number}]:\n{current_slide_full_content}"
                    merged_content_block = {"type": "slide_text_merged", "content": merged_text.strip(), "source_info": {"slide": current_slide_number, "merged_from": pending_title_slide_data['number']}, "file_type": "pptx"}
                    all_content.append(merged_content_block); pending_title_slide_data = None
                else:
                    if pending_title_slide_data:
                         if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
                         pending_title_slide_data = None
                    if current_slide_has_title_placeholder and other_text_word_count <= merge_threshold_words and current_slide_full_content:
                        logger.debug(f"Slide {current_slide_number} is potential title slide. Holding."); pending_title_slide_data = {"content": current_slide_full_content, "number": current_slide_number, "has_title": current_slide_has_title_placeholder, "other_words": other_text_word_count, "title": current_slide_title_text}
                    else:
                        if current_slide_full_content: all_content.append({"type": "slide_text", "content": current_slide_full_content, "source_info": {"slide": current_slide_number}, "file_type": "pptx"})
                if slide_bar: slide_bar.progress(min(1.0, (i + 1) / total_slides), text=f"Extracting slide {current_slide_number}/{total_slides}...")
            if pending_title_slide_data:
                logger.debug(f"Processing pending title slide {pending_title_slide_data['number']} at end.")
                if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
            if slide_bar: slide_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PPTX {base_filename} (Merge strategy applied)."); return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PPTX {base_filename}: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.warning(f"Could not extract content from {base_filename}. Error: {e}")
            return []

    def extract_content(self, file_path):
        """Extract content based on file extension."""
        # (This function remains unchanged)
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf': return self._extract_pdf(file_path)
        elif extension == '.xlsx': return self._extract_xlsx(file_path)
        elif extension == '.csv': return self._extract_csv(file_path)
        elif extension == '.pptx': return self._extract_pptx(file_path)
        else: logger.warning(f"Unsupported file type skipped: {file_path}"); return []

    def chunk_content(self, all_content):
        """Split extracted content into chunks with overlap, preserving metadata."""
        # (This function remains unchanged)
        chunks = []; total_words_estimate = 0; chunk_bar = None
        if not all_content: return chunks
        try: total_words_estimate = sum(len(str(item.get('content', '')).split()) for item in all_content)
        except Exception: pass
        if self._is_streamlit and total_words_estimate > 5000: import streamlit as st; chunk_bar = st.progress(0, text=f"Chunking content...")
        words_processed = 0
        for item_index, item in enumerate(all_content):
            content = item.get('content', ''); source_info = item.get('source_info', {}); file_type = item.get('file_type', 'unknown'); content_type = item.get('type', 'unknown')
            if not isinstance(content, str): content = str(content)
            words = content.split()
            if not words: logger.debug(f"Skipping empty content block {item_index}"); continue
            item_chunks_created = 0
            for i in range(0, len(words), self.config.chunk_size - self.config.overlap):
                chunk_words = words[i:i + self.config.chunk_size]; chunk_text = " ".join(chunk_words)
                if chunk_text: chunks.append({"content": chunk_text, "source_info": source_info, "file_type": file_type, "type": content_type}); item_chunks_created += 1
            if chunk_bar: words_processed += len(words); progress_percent = min(1.0, words_processed / total_words_estimate) if total_words_estimate > 0 else 0; chunk_bar.progress(progress_percent, text=f"Chunking content... ({len(chunks)} chunks created)")
        if chunk_bar: chunk_bar.empty()
        logger.info(f"Created {len(chunks)} chunks from {len(all_content)} content blocks."); return chunks

    # --- Path Helpers --- (Unchanged)
    def _get_safe_filename(self, file_name): base_name = os.path.splitext(file_name)[0]; return re.sub(r'[^\w\.-]', '_', base_name)
    def get_index_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.index")
    def get_embedding_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.npy")
    def get_chunks_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.json")

    # --- FAISS Index and Chunk Loading/Saving --- (Unchanged)
    def load_faiss_index(self, file_name, embedding_dim):
        index_path = self.get_index_path(file_name)
        if os.path.exists(index_path):
            try: index = faiss.read_index(index_path); logger.info(f"Loaded FAISS index for {file_name} ({index.ntotal} vectors)"); return index
            except Exception as e: logger.error(f"Error reading FAISS index {index_path}: {e}. Creating new index.")
        logger.info(f"FAISS index not found at {index_path}. Will create new."); return faiss.IndexFlatL2(embedding_dim)
    def save_chunks(self, file_name, chunks):
        chunks_path = self.get_chunks_path(file_name)
        try:
            with open(chunks_path, 'w', encoding='utf-8') as f: json.dump(chunks, f, indent=2); logger.info(f"Saved {len(chunks)} chunks to {chunks_path}")
        except Exception as e: logger.error(f"Error saving chunks for {file_name} to {chunks_path}: {e}");
    def load_chunks(self, file_name):
        chunks_path = self.get_chunks_path(file_name)
        if os.path.exists(chunks_path):
            try:
                with open(chunks_path, 'r', encoding='utf-8') as f: chunks = json.load(f); logger.info(f"Loaded {len(chunks)} chunks from {chunks_path}"); return chunks
            except Exception as e: logger.error(f"Error loading/decoding chunks from {chunks_path}: {e}. Will re-process."); return None
        logger.info(f"Chunks file not found at {chunks_path}"); return None

    # --- Document Processing --- (Unchanged)
    def process_files(self, progress_callback=None):
        self.file_chunks = {}; self.faiss_indexes = {}; self.processed_files = set(); data_dir = self.config.data_dir; supported_ext = self.config.supported_extensions
        try: all_files = os.listdir(data_dir); process_list = sorted([f for f in all_files if f.lower().endswith(supported_ext)])
        except Exception as e: logger.error(f"Error listing files in data directory {data_dir}: {e}", exc_info=True); return False
        if not process_list: logger.warning(f"No supported files found in {data_dir}"); return True
        logger.info(f"Processing {len(process_list)} supported files from {data_dir}"); embedding_dim = self.encoder_model.get_sentence_embedding_dimension(); total_files = len(process_list)
        if progress_callback: progress_callback(f"Found {total_files} file(s). Starting processing...", current_step=0, total_steps=total_files)
        for idx, file_name in enumerate(process_list):
            current_file_msg = f"Processing ({idx+1}/{total_files}): {file_name}"; logger.info(f"--- {current_file_msg} ---")
            if progress_callback: progress_callback(current_file_msg, current_step=idx, total_steps=total_files)
            file_path = os.path.join(data_dir, file_name); index_path = self.get_index_path(file_name); emb_path = self.get_embedding_path(file_name); chunks_path = self.get_chunks_path(file_name)
            try:
                chunks = self.load_chunks(file_name)
                if chunks is None:
                    if progress_callback: progress_callback(f"{current_file_msg} - Extracting...", stage="Extracting")
                    all_content = self.extract_content(file_path);
                    if not all_content: logger.warning(f"No content from {file_name}."); continue
                    if progress_callback: progress_callback(f"{current_file_msg} - Chunking...", stage="Chunking")
                    chunks = self.chunk_content(all_content);
                    if not chunks: logger.warning(f"No chunks for {file_name}."); continue
                    self.save_chunks(file_name, chunks)
                self.file_chunks[file_name] = chunks; logger.debug(f"Stored {len(chunks)} chunks for {file_name}")
                faiss_index = None; regenerate_embeddings = False
                if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                    if progress_callback: progress_callback(f"{current_file_msg} - Verifying...", stage="Verifying")
                    try:
                        embeddings = np.load(emb_path)
                        if embeddings.ndim != 2 or embeddings.shape[1] != embedding_dim or embeddings.shape[0] != len(chunks): logger.warning("Mismatch. Regenerating."); regenerate_embeddings = True
                        else: faiss_index = self.load_faiss_index(file_name, embedding_dim); logger.info(f"Verified existing data for {file_name}")
                    except Exception as e: logger.error(f"Error verifying: {e}. Regenerating...", exc_info=True); regenerate_embeddings = True; faiss_index = None
                else: regenerate_embeddings = True
                if regenerate_embeddings or faiss_index is None:
                    logger.info(f"Generating embeddings/index for {file_name}...")
                    if progress_callback: progress_callback(f"{current_file_msg} - Embedding...", stage="Embedding")
                    content_list = [chunk['content'] for chunk in chunks];
                    if not content_list: logger.warning("No content to embed."); continue
                    embeddings = self.encoder_model.encode(content_list, batch_size=64, show_progress_bar=False, convert_to_numpy=True).astype('float32')
                    if embeddings.shape[0] == 0: logger.warning("Embedding yielded no vectors."); continue
                    np.save(emb_path, embeddings); logger.info(f"Saved {embeddings.shape[0]} embeddings")
                    if progress_callback: progress_callback(f"{current_file_msg} - Indexing...", stage="Indexing")
                    faiss_index = faiss.IndexFlatL2(embedding_dim); faiss_index.add(embeddings); faiss.write_index(faiss_index, index_path); logger.info(f"Saved FAISS index ({faiss_index.ntotal} vectors)")
                if faiss_index is not None and faiss_index.ntotal > 0: self.faiss_indexes[file_name] = faiss_index; self.processed_files.add(file_name); logger.debug("Stored FAISS index")
                else: logger.warning(f"No valid index for {file_name}.")
            except Exception as e:
                logger.error(f"Failed to process {file_name}: {e}", exc_info=True)
                if progress_callback: progress_callback(f"❌ Error processing {file_name}: {e}", is_error=True)
                try:
                    for p in [index_path, emb_path, chunks_path]:
                        if os.path.exists(p): os.remove(p); logger.debug(f"Removed {p}")
                except OSError as e_clean: logger.error(f"Error cleaning up: {e_clean}")
                if file_name in self.file_chunks: del self.file_chunks[file_name]
                if file_name in self.faiss_indexes: del self.faiss_indexes[file_name]
        final_indexed_count = len(self.faiss_indexes)
        if final_indexed_count > 0: logger.info(f"--- Processing Complete. Indexed {final_indexed_count}/{total_files} documents. ---"); return True
        else: logger.warning(f"--- Processing Complete. No documents indexed. ---"); return True

    # --- Querying Logic ---

    def query_files(self, query):
        """Query FAISS index for a single query string."""
        # (Modified slightly to handle a single query string robustly)
        if not self.faiss_indexes:
            logger.warning("No indexes available for querying.")
            return {} # Return empty results, let caller handle

        query_results = {}
        try:
            logger.info(f"Encoding query for Faiss: '{query[:100]}...'")
            query_embedding = self.encoder_model.encode(query, convert_to_numpy=True).astype("float32")
            query_embedding = np.array([query_embedding])
            if query_embedding.ndim != 2: raise ValueError("Query embedding shape error")

            # Search across all indexes for this single query
            for file_name, index in self.faiss_indexes.items():
                if index is None or index.ntotal == 0: continue
                try:
                    k_search = min(self.config.k_retrieval, index.ntotal)
                    D, I = index.search(query_embedding, k=k_search)
                    indices, distances = I[0], D[0]

                    current_file_chunks = self.file_chunks.get(file_name)
                    if not current_file_chunks: logger.error(f"Chunks missing for {file_name} during query!"); continue

                    file_results = []
                    processed_indices = set()
                    for i, idx in enumerate(indices):
                        if idx == -1 or idx in processed_indices: continue
                        if not (0 <= idx < len(current_file_chunks)): continue
                        processed_indices.add(idx)
                        chunk = current_file_chunks[idx]
                        file_results.append({
                            "source_info": chunk.get('source_info', {}),
                            "file_type": chunk.get('file_type', 'unknown'),
                            "content": chunk.get('content', ''),
                            "score": round(float(distances[i]), 4),
                            "type": chunk.get('type', 'unknown')
                        })

                    file_results.sort(key=lambda x: x['score']) # Lower L2 is better
                    if file_results:
                        # Store results under the filename
                        if file_name not in query_results: query_results[file_name] = []
                        query_results[file_name].extend(file_results)
                        # Optional: Limit per file after extending if needed, or sort globally later
                        # query_results[file_name] = sorted(query_results[file_name], key=lambda x: x['score'])[:self.config.k_retrieval]


                except Exception as e_search:
                    logger.error(f"Error searching {file_name} for query '{query[:50]}...': {e_search}", exc_info=True)

            logger.debug(f"Faiss search complete for query '{query[:50]}...'. Found results in {len(query_results)} files.")
            return query_results

        except Exception as e_query:
            logger.error(f"Error during Faiss query processing for '{query[:50]}...': {e_query}", exc_info=True)
            return {} # Return empty on error

    # (aggregate_context remains unchanged from ragtest_c2_merged.py)
    def aggregate_context(self, query_results, strategy="top_k"):
        all_context = {}; max_chars = self.config.max_context_tokens * 3
        logger.info(f"Aggregating context using strategy '{strategy}', max_chars ~{max_chars}")
        if not query_results: return all_context
        flat_results = []
        for file_name, results in query_results.items():
            for res in results: flat_results.append({**res, "file_name": file_name})
        flat_results.sort(key=lambda x: x['score'])
        aggregated_context_str = ""; total_aggregated_chars = 0; added_chunks_count = 0; context_sources = set()
        limit = self.config.k_retrieval if strategy == "top_k" else len(flat_results)
        for i, res in enumerate(flat_results):
            if i >= limit: logger.debug(f"Reached top-k limit ({limit})"); break
            source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown'); content_type = res.get('type', 'unknown'); score = res['score']; content_body = res['content']
            source_parts = [f"Source: {file_name}"]
            slide_info = source_info.get('slide', 'N/A'); merged_from = source_info.get('merged_from')
            if file_type == 'pptx': slide_display = f"Slide: {slide_info}" + (f" (merged from {merged_from})" if merged_from else ""); source_parts.append(slide_display)
            elif file_type == 'pdf': source_parts.append(f"Page: {source_info.get('page', 'N/A')}")
            elif file_type == 'xlsx': source_parts.append(f"Sheet: {source_info.get('sheet', 'N/A')}")
            elif file_type == 'csv': source_parts.append(f"CSV (Rows: {source_info.get('rows', 'all')})")
            source_parts.append(f"Type: {content_type}"); source_parts.append(f"Score: {score:.4f}")
            source_str = ", ".join(source_parts); content_header = f"--- Context from {source_str} ---\n"; content_to_add = content_header + content_body + "\n\n"; content_chars = len(content_to_add)
            if total_aggregated_chars + content_chars <= max_chars: aggregated_context_str += content_to_add; total_aggregated_chars += content_chars; added_chunks_count += 1; context_sources.add(file_name)
            else:
                 if added_chunks_count == 0:
                     remaining_chars = max_chars - total_aggregated_chars - len(content_header) - 20
                     if remaining_chars > 50: truncated_body = content_body[:remaining_chars]; aggregated_context_str += content_header + truncated_body + "\n[...TRUNCATED CONTEXT...]\n\n"; total_aggregated_chars += len(aggregated_context_str); added_chunks_count += 1; context_sources.add(file_name); logger.warning("Context truncated (first chunk too large).")
                     else: logger.warning("First chunk too large to fit, skipping.")
                 logger.info(f"Stopping context aggregation at {added_chunks_count} chunks ({total_aggregated_chars}/{max_chars} chars). Limit reached."); break
        final_context = aggregated_context_str.strip()
        if final_context: all_context = {"combined_context": final_context, "source_files": sorted(list(context_sources))}; logger.info(f"Aggregated {total_aggregated_chars} chars from {added_chunks_count} chunks across {len(context_sources)} files.")
        else: logger.warning("No context aggregated within limits.")
        return all_context

    # (query_llm remains unchanged from ragtest_c2_merged.py)
    def query_llm(self, query, context_data, retry_count=1):
        combined_context = context_data.get("combined_context", ""); source_files = context_data.get("source_files", []); source_file_str = ", ".join(source_files) if source_files else "the provided documents"
        if not combined_context: logger.warning("No context for LLM query."); return f"Could not generate answer: No relevant context found in {source_file_str}."
        try:
            llm = self._get_llm() # Use helper to get/init LLM
            system_prompt = f"""You are an AI assistant answering questions based ONLY on the provided context from document(s): '{source_file_str}'. Use ONLY information in the 'START CONTEXT'/'END CONTEXT' section. If the answer isn't there, state: "Based on the provided context from {source_file_str}, I cannot answer this question." Do NOT use external knowledge."""
            human_prompt = f"""Context from document(s) '{source_file_str}':\n--- START CONTEXT ---\n{combined_context}\n--- END CONTEXT ---\n\nUser Question: {query}\n\nAnswer based ONLY on the provided context:"""
            full_prompt_messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
            logger.info(f"Querying Groq model {self.config.llm_model} using context from {source_file_str}...")
            answer = f"Error: LLM query failed after retries."
            for attempt in range(retry_count + 1):
                try:
                    response = llm.invoke(full_prompt_messages)
                    answer = response.content.strip() if hasattr(response, 'content') else str(response).strip()
                    logger.info(f"Groq response received (attempt {attempt+1}). Length: {len(answer)}")
                    if not answer and attempt < retry_count: logger.warning("Empty response, retrying..."); time.sleep(1); continue
                    return answer or f"Received empty response from LLM for context from {source_file_str}."
                except Exception as e_api:
                    logger.warning(f"Groq API attempt {attempt+1} failed: {e_api}")
                    if attempt < retry_count: time.sleep(1.5 ** attempt); logger.info("Retrying Groq query...")
                    else: answer = f"Error: Failed to get answer from LLM after {retry_count+1} attempts. (API Error: {e_api})"
            return answer
        except Exception as e_setup:
            logger.error(f"Error setting up/calling Groq API: {e_setup}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"LLM query failed: {e_setup}")
            return f"Error: Could not query LLM due to setup error: {e_setup}"


    # ==========================================================================
    # START OF NEW HELPER METHODS FOR QUERY ANALYSIS
    # ==========================================================================

    def _call_llm_for_analysis(self, prompt_messages, task_description):
        """Helper to call LLM for internal analysis tasks like classification/decomposition."""
        try:
            llm = self._get_llm() # Get initialized LLM
            logger.info(f"Calling LLM for {task_description}...")
            response = llm.invoke(prompt_messages)
            content = response.content.strip() if hasattr(response, 'content') else str(response).strip()
            logger.info(f"LLM response for {task_description}: '{content[:100]}...'")
            return content
        except Exception as e:
            logger.error(f"LLM call failed during {task_description}: {e}", exc_info=True)
            return None # Indicate failure

    def _classify_query(self, query):
        """Classify the query type using the configured LLM."""
        system_prompt = """You are an expert query analyzer. Your task is to classify a user query based on whether it requires complex reasoning/aggregation/calculation typically needed for structured data (like Excel tables) or if it's a simple lookup.
Consider queries asking for averages, sums, comparisons across rows/columns, or multi-step lookups as 'Complex/Reasoning'.
Consider queries asking for specific facts, definitions, or summaries of text as 'Simple Retrieval'.
Respond ONLY with 'Simple Retrieval' or 'Complex/Reasoning'."""
        human_prompt = f"Classify the following user query: \"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        classification = self._call_llm_for_analysis(messages, "query classification")

        if classification in ["Simple Retrieval", "Complex/Reasoning"]:
            return classification
        else:
            logger.warning(f"Unexpected classification result: '{classification}'. Defaulting to 'Simple Retrieval'.")
            return "Simple Retrieval" # Default fallback

    def _decompose_query(self, query):
        """Decompose a complex query into simpler sub-queries using the configured LLM."""
        system_prompt = """You are an expert query decomposer. Your task is to break down a complex user query, potentially requiring aggregation or multi-step lookups in data (like Excel), into a series of simpler, factual sub-queries. Each sub-query should aim to retrieve a specific piece of information needed to answer the original query.
Focus on extracting the necessary raw data points. Do not perform calculations.
Format the output as a numbered list, with each sub-query on a new line. If the query cannot be decomposed or is already simple, return the original query prefixed with '1. '."""
        human_prompt = f"Decompose the following complex query into simple, factual sub-queries:\n\"{query}\""
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]

        decomposition = self._call_llm_for_analysis(messages, "query decomposition")

        if decomposition:
            # Basic parsing: split by newline, remove numbering like "1. ", strip whitespace
            sub_queries = [re.sub(r"^\d+\.\s*", "", line).strip() for line in decomposition.split('\n') if line.strip()]
            if sub_queries:
                logger.info(f"Decomposed query into: {sub_queries}")
                return sub_queries
            else:
                 logger.warning(f"Decomposition resulted in empty list. Using original query.")
                 return [query] # Fallback to original query
        else:
            logger.warning(f"Decomposition failed. Using original query.")
            return [query] # Fallback to original query

    # ==========================================================================
    # END OF NEW HELPER METHODS
    # ==========================================================================


    # ==========================================================================
    # START OF MODIFIED run_query FUNCTION
    # ==========================================================================
    def run_query(self, query, context_strategy="top_k"):
        """Complete query pipeline: classify, (decompose), retrieve, aggregate, query LLM."""
        logger.info(f"--- Starting query execution (v3): '{query[:100]}...' ---")
        final_results = {
            "query": query,
            "classification": None,
            "sub_queries": None,
            "retrieval_results": {},
            "aggregated_context_data": {},
            "answer": "",
            "status": "Started",
            "error": None
        }
        is_streamlit = self._is_streamlit
        if is_streamlit: import streamlit as st

        try:
            # === 1. Classify Query ===
            spinner_msg = "Analyzing query..."
            if is_streamlit:
                with st.spinner(spinner_msg): classification = self._classify_query(query)
            else: logger.info(spinner_msg); classification = self._classify_query(query)
            final_results["classification"] = classification
            logger.info(f"Query classified as: {classification}")

            # === 2. Decompose if Complex ===
            queries_to_retrieve = [query] # Default: use original query
            if classification == "Complex/Reasoning":
                spinner_msg = "Decomposing complex query..."
                if is_streamlit:
                     with st.spinner(spinner_msg): sub_queries = self._decompose_query(query)
                else: logger.info(spinner_msg); sub_queries = self._decompose_query(query)

                if sub_queries and sub_queries != [query]: # Check if decomposition actually happened
                     queries_to_retrieve = sub_queries
                     final_results["sub_queries"] = sub_queries
                     logger.info(f"Using sub-queries for retrieval: {sub_queries}")
                else:
                     logger.info("Decomposition didn't yield distinct sub-queries or failed, using original query for retrieval.")
                     # Keep queries_to_retrieve = [query]

            # === 3. Retrieve relevant chunks for determined queries ===
            spinner_msg = f"Searching documents for {len(queries_to_retrieve)} query part(s)..."
            all_query_results = {}
            if is_streamlit:
                with st.spinner(spinner_msg):
                    for i, q in enumerate(queries_to_retrieve):
                        logger.info(f"Retrieving for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                        results_for_q = self.query_files(q) # Calls Faiss search
                        # Merge results, simple aggregation for now
                        for file, res_list in results_for_q.items():
                            if file not in all_query_results: all_query_results[file] = []
                            all_query_results[file].extend(res_list)
            else:
                logger.info(spinner_msg)
                for i, q in enumerate(queries_to_retrieve):
                     logger.info(f"Retrieving for query part {i+1}/{len(queries_to_retrieve)}: '{q[:100]}...'")
                     results_for_q = self.query_files(q)
                     for file, res_list in results_for_q.items():
                            if file not in all_query_results: all_query_results[file] = []
                            all_query_results[file].extend(res_list)

            # De-duplicate and sort results per file *after* merging from sub-queries
            final_retrieval_results = {}
            for file, res_list in all_query_results.items():
                 unique_content = {}
                 for res in res_list:
                     # Use content and maybe slide/page as key for uniqueness
                     key = (res['content'], tuple(sorted(res.get('source_info',{}).items())))
                     if key not in unique_content or res['score'] < unique_content[key]['score']:
                         unique_content[key] = res
                 sorted_unique_results = sorted(unique_content.values(), key=lambda x: x['score'])
                 final_retrieval_results[file] = sorted_unique_results

            final_results["retrieval_results"] = final_retrieval_results # Store final, de-duped results

            if not final_retrieval_results:
                logger.warning("No relevant chunks found across any documents for the query/sub-queries.")
                final_results["status"] = "Completed: No relevant information found."
                if is_streamlit: st.info("Could not find relevant information.")
                return final_results

            # === 4. Aggregate context (uses combined results) ===
            spinner_msg = "Gathering context..."
            if is_streamlit:
                with st.spinner(spinner_msg): aggregated_context_data = self.aggregate_context(final_retrieval_results, strategy=context_strategy)
            else: logger.info(spinner_msg); aggregated_context_data = self.aggregate_context(final_retrieval_results, strategy=context_strategy)
            final_results["aggregated_context_data"] = aggregated_context_data

            if not aggregated_context_data or not aggregated_context_data.get("combined_context"):
                 logger.warning(f"Context aggregation failed or yielded empty context for query: '{query[:100]}...'")
                 final_results["status"] = "Completed: Context aggregation failed or empty."
                 final_results["error"] = "Context aggregation failed or empty."
                 if is_streamlit: st.warning("Couldn't prepare context for answering.")
                 return final_results

            # === 5. Query LLM with aggregated context BUT use the ORIGINAL query ===
            spinner_msg = "Generating final answer..."
            if is_streamlit:
                 with st.spinner(spinner_msg): answer = self.query_llm(query, aggregated_context_data) # Pass ORIGINAL query here
            else: logger.info(spinner_msg); answer = self.query_llm(query, aggregated_context_data) # Pass ORIGINAL query here

            final_results["answer"] = answer
            final_results["status"] = "Completed Successfully."
            logger.info(f"--- Finished query execution (v3) for: '{query[:100]}...' ---")

        except Exception as e:
            logger.error(f"Unexpected error during run_query (v3): {e}", exc_info=True)
            final_results["status"] = "Failed"; final_results["error"] = str(e)
            if is_streamlit: st.error(f"An unexpected error occurred: {e}")

        return final_results
    # ==========================================================================
    # END OF MODIFIED run_query FUNCTION
    # ==========================================================================


# --- START OF STREAMLIT UI CODE ---
# (Streamlit UI code remains identical to the previous version - ragtest_c2_merged.py)
if __name__ == "__main__":
    try: import streamlit as st
    except ImportError: logger.error("Streamlit not installed."); sys.exit(1)

    st.set_page_config(layout="wide", page_title="Multi-Document Q&A with RAG v3")
    st.title("📄 Multi-Document Question Answering System (v3)")
    st.caption("Query PDFs, XLSX, CSV, PPTX using Embeddings, FAISS, Groq LLM (with Query Analysis)")

    @st.cache_resource
    def load_rag_system(config_path="config.ini"):
        if not os.getenv("GROQ_API_KEY"): st.error("🚨 GROQ_API_KEY not set!"); st.stop()
        try: return RAGSystem(config_path=config_path)
        except Exception as e: st.error(f"Fatal RAG Init Error: {e}"); logger.error("RAG Init Error", exc_info=True); st.stop()

    @st.cache_data
    def process_documents(_rag_system):
        # (This cached function remains unchanged - relies on RAGSystem.process_files)
        status_messages = []; status_placeholder = st.empty()
        def streamlit_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
            log_prefix = "✅ " if is_done else "❌ " if is_error else "⚠️ " if is_warning else f"⏳ ({stage}) " if stage else "⏳ "; progress_val = 1.0 if is_done else 0.0
            if total_steps and current_step is not None: stage_prog = {"Extracting": 0.1, "Chunking": 0.3, "Verifying": 0.5, "Embedding": 0.6, "Indexing": 0.9}.get(stage, 0.0); progress_val = min(1.0, (current_step + stage_prog) / total_steps) if total_steps > 0 else 0.0
            full_message = f"{log_prefix}{message}"; status_messages.append(full_message)
            try:
                if progress_val > 0 or is_done: status_placeholder.progress(progress_val)
                status_placeholder.caption(full_message)
            except Exception as e: logger.warning(f"Streamlit UI update error: {e}")
        st.header("📚 Document Processing")
        st.caption(f"Data: `{os.path.abspath(_rag_system.config.data_dir)}` | Index: `{os.path.abspath(_rag_system.config.index_dir)}`")
        processing_successful = False
        with st.spinner("Initializing and processing documents..."):
            try:
                processing_successful = _rag_system.process_files(progress_callback=streamlit_progress_callback)
                indexed_count = len(_rag_system.faiss_indexes)
                if indexed_count > 0: status_placeholder.success(f"✅ Ready! Indexed {indexed_count} document(s).")
                elif processing_successful: status_placeholder.warning(f"⚠️ Processing finished, but no documents indexed.")
                elif not status_messages or not any("❌" in msg for msg in status_messages): status_placeholder.error(f"❌ Processing failed. Check logs.")
            except Exception as e: logger.error(f"Fatal error during processing call: {e}", exc_info=True); status_placeholder.error(f"❌ Fatal error: {e}. Check logs.")
        return len(_rag_system.faiss_indexes) > 0, sorted(list(_rag_system.faiss_indexes.keys()))

    try:
        rag_sys = load_rag_system()
        if st.button("🔄 Re-process Documents"): st.cache_data.clear(); st.cache_resource.clear(); st.rerun()
        is_ready, indexed_files = process_documents(rag_sys)

        if is_ready and indexed_files:
            st.sidebar.success(f"Indexed Documents ({len(indexed_files)}):")
            with st.sidebar.expander("Show Indexed Files"):
                 for fname in indexed_files: st.caption(f"- {fname}")
            st.sidebar.info(f"LLM: `{rag_sys.config.llm_model}`")
            st.sidebar.info(f"Retrieval K: `{rag_sys.config.k_retrieval}`")
            st.sidebar.info(f"Max Context Tokens: `{rag_sys.config.max_context_tokens}`")

            st.header("💬 Ask a Question")
            user_query = st.text_input("Enter your query:", key="query_input")
            if user_query and st.button("Get Answer", key="submit_query"):
                query_start_time = time.time()
                results_data = rag_sys.run_query(user_query, context_strategy="top_k") # Calls the modified run_query
                query_end_time = time.time()

                st.subheader("💡 Answer")
                answer = results_data.get("answer")
                if answer: st.markdown(answer)
                else: st.warning("Could not generate an answer."); st.error(f"Error: {results_data['error']}") if results_data.get("error") else None

                # Display Query Analysis Info if available
                st.subheader("📊 Query Analysis")
                classification = results_data.get("classification")
                sub_queries = results_data.get("sub_queries")
                st.caption(f"Query Classification: `{classification or 'N/A'}`")
                if sub_queries:
                    with st.expander("Sub-queries used for retrieval"):
                        for i, sq in enumerate(sub_queries):
                            st.code(f"{i+1}. {sq}")

                # Display Supporting Evidence (Context + Retrieval)
                st.subheader("🔍 Supporting Evidence")
                agg_context_data = results_data.get("aggregated_context_data", {})
                combined_context = agg_context_data.get("combined_context", "")
                context_source_files = agg_context_data.get("source_files", [])
                if combined_context:
                     with st.expander(f"Context Used from: {', '.join(context_source_files)}"):
                         st.text_area("Combined Context Sent to LLM", combined_context, height=250, key="context_display")
                else: st.info("No context was aggregated for the LLM.")

                retrieval_results = results_data.get("retrieval_results", {})
                if retrieval_results:
                     # Displaying retrieved chunks (unchanged)
                     with st.expander("Retrieved Chunks (De-duplicated, Top K overall)"):
                         # Flatten results for display, respecting original file source
                         display_chunks = []
                         for file_name, res_list in retrieval_results.items():
                              for res in res_list:
                                   display_chunks.append({**res, 'file_name': file_name})
                         # Sort by score globally for display order
                         display_chunks.sort(key=lambda x: x['score'])

                         for i, res in enumerate(display_chunks[:rag_sys.config.k_retrieval * 2]): # Show a bit more than K for inspection
                             source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown'); content_type = res.get('type', 'unknown'); score = res['score']
                             source_display_parts = [f"File: {file_name}"]
                             slide_info = source_info.get('slide', 'N/A'); merged_from = source_info.get('merged_from')
                             if file_type == 'pptx': slide_display = f"Slide: {slide_info}" + (f" (merged from {merged_from})" if merged_from else ""); source_display_parts.append(slide_display)
                             elif file_type == 'pdf': source_display_parts.append(f"Page: {source_info.get('page', 'N/A')}")
                             elif file_type == 'xlsx': source_display_parts.append(f"Sheet: {source_info.get('sheet', 'N/A')}")
                             elif file_type == 'csv': source_display_parts.append(f"CSV (Rows: {source_info.get('rows', 'all')})")
                             source_display_parts.append(f"Type: {content_type}")
                             st.markdown(f"**Chunk {i+1} (Score: {score:.4f})**")
                             st.caption(" | ".join(source_display_parts))
                             st.text(f"{res['content'][:350]}...")
                             if i < len(display_chunks) -1 : st.divider() # Don't add divider after last chunk
                else: st.info("No specific document chunks were retrieved.")
                st.caption(f"Query processed in {query_end_time - query_start_time:.2f} seconds.")
        elif not indexed_files:
            st.warning(f"No documents indexed.")
            st.info("Check 'Document Processing' status or add files and click 'Re-process Documents'.")
        else: st.error("System not ready. Check 'Document Processing' status and logs.")
    except Exception as e:
        st.error(f"Streamlit App Error: {e}")
        logger.exception("Streamlit application error:")
        st.info("Check console or `rag_system.log`.")

# --- END OF STREAMLIT UI CODE ---

# --- END OF FILE ragtest_c3_query_adapt.py ---