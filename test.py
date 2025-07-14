# --- START OF FILE ragtest_c2-3.py ---

# --- START OF FILE ragtest_c3_query_adapt_excel_count.py --- # Renamed

# --- START OF FILE ragtest_c.py --- # <-- Original comment

import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd # Added for CSV/XLSX
import faiss
# import streamlit as st # Import only if needed / within UI code
import time
from sentence_transformers import SentenceTransformer
# from transformers import T5ForConditionalGeneration, T5Tokenizer # Keep if needed elsewhere
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage # Added SystemMessage
from dotenv import load_dotenv
from pptx import Presentation # Added for PPTX
from pptx.enum.shapes import MSO_SHAPE_TYPE # Needed for placeholder check
import sys # Import sys for checking streamlit context
import string # Import string for keyword cleaning

# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted from RAG2.py) ---

# Configure logging (optional for Streamlit, but good for debugging)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system.log"),
        # logging.StreamHandler() # Can be noisy in Streamlit console
    ]
)
logger = logging.getLogger(__name__)


class RAGConfig:
    """Configuration class for RAG system."""
    # (Unchanged from previous version)
    def __init__(self, config_path="config.ini"):
        self.config = configparser.ConfigParser()
        self.defaults = {
            "PATHS": {"data_dir": "./Data/", "index_dir": "./Faiss_index/", "log_file": "rag_system.log"},
            "MODELS": {"encoder_model": "sentence-transformers/all-MiniLM-L6-v2", "llm_model": "mixtral-8x7b-32768", "device": "auto"},
            "PARAMETERS": {
                "chunk_size": "200", "overlap": "50", "k_retrieval": "5", "k_retrieval_large": "100",
                "temperature": "0.1", "max_context_tokens": "4000", "max_chars_per_element": "1000",
                "pptx_merge_threshold_words": "50"
            },
            "SUPPORTED_EXTENSIONS": {"extensions": ".pdf, .xlsx, .csv, .pptx"}
        }
        self.config_path = config_path
        if os.path.exists(config_path):
            try:
                self.config.read(config_path); logger.info(f"Loaded configuration from {config_path}"); self._ensure_defaults()
            except Exception as e: logger.error(f"Error reading config file {config_path}: {e}"); self._set_defaults()
        else: logger.warning(f"Config file {config_path} not found, using defaults"); self._set_defaults()
    def _set_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section)
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value)
    def _ensure_defaults(self):
        changes_made = False
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section); logger.info(f"Added missing section [{section}]"); changes_made = True
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value); logger.info(f"Added missing option '{option}' to [{section}]"); changes_made = True
    # Properties (Unchanged from previous version)
    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    @property
    def llm_model(self): return self.config.get("MODELS", "llm_model")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto")
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=200)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=50)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=5)
    @property
    def k_retrieval_large(self): return self.config.getint("PARAMETERS", "k_retrieval_large", fallback=100)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.1)
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=4000)
    @property
    def max_context_chars(self): return self.max_context_tokens * 3
    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1000)
    @property
    def pptx_merge_threshold_words(self): return self.config.getint("PARAMETERS", "pptx_merge_threshold_words", fallback=50)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf")
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])


class RAGSystem:
    """Retrieval-Augmented Generation system for various document types."""
    # (Init remains unchanged)
    def __init__(self, config_path="config.ini"):
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system. Device: {self.config.device}, Supported files: {self.config.supported_extensions}")
        logger.info(f"Config: k_retrieval={self.config.k_retrieval}, k_retrieval_large={self.config.k_retrieval_large}, max_context_tokens={self.config.max_context_tokens}")
        self._is_streamlit = "streamlit" in sys.modules
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)
        try:
            # Model loading (unchanged)
            if self._is_streamlit:
                import streamlit as st
                with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."): self.encoder_model = SentenceTransformer(self.config.encoder_model, device=self.config.device)
            else: logger.info(f"Loading embedding model ({self.config.encoder_model})..."); self.encoder_model = SentenceTransformer(self.config.encoder_model, device=self.config.device)
            logger.info("Local embedding model loaded successfully.")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}"); st.stop()
            else: raise
        self.file_chunks = {}
        self.faiss_indexes = {}
        self.processed_files = set()
        self._llm = None

    # (_get_llm, clean_text, extractors, chunk_content, path helpers, load/save, process_files, query_files, aggregate_context unchanged)
    # --- OMITTED FOR BREVITY ---
    # (Keep all these methods exactly as they were in the previous 'filtered_agg' version)
    def _get_llm(self):
        # (Unchanged)
        if self._llm is None:
            api_key = os.getenv("GROQ_API_KEY")
            if not api_key: logger.error("GROQ_API_KEY not found."); raise ValueError("GROQ_API_KEY not configured.")
            try:
                self._llm = ChatGroq(temperature=self.config.temperature, model_name=self.config.llm_model, groq_api_key=api_key)
                logger.info(f"Groq LLM client initialized ({self.config.llm_model})")
            except Exception as e: logger.error(f"Failed to initialize Groq client: {e}", exc_info=True); raise
        return self._llm

    def clean_text(self, text): # (Unchanged)
        if not isinstance(text, str): text = str(text)
        text = re.sub(r"\(cid:.*?\)", "", text); text = re.sub(r"\s+", " ", text).strip(); return text

    def _extract_pdf(self, file_path): # (Unchanged)
        all_content = []; base_filename = os.path.basename(file_path)
        try:
            with pdfplumber.open(file_path) as pdf:
                # ... (rest of extraction logic) ...
                for page_num, page in enumerate(pdf.pages):
                     text = page.extract_text(layout="normal") or ""; cleaned_text = self.clean_text(text)
                     if cleaned_text: all_content.append({"type": "text", "content": cleaned_text, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                     try:
                         for table in page.extract_tables():
                              if table: table_df=pd.DataFrame(table).fillna(''); table_string=table_df.to_string(index=False,header=True,na_rep=''); cleaned_table_string=self.clean_text(table_string);
                              if cleaned_table_string: all_content.append({"type": "table", "content": cleaned_table_string, "source_info": {"page": page_num + 1}, "file_type": "pdf"})
                     except Exception as e_table: logger.warning(f"Table extraction error PDF page {page_num+1}: {e_table}")
            logger.info(f"Extracted {len(all_content)} blocks from PDF: {base_filename}"); return all_content
        except Exception as e: logger.error(f"Error extracting PDF {base_filename}: {e}", exc_info=True); return []

    def _extract_xlsx(self, file_path): # (Unchanged)
        all_content = []; base_filename = os.path.basename(file_path)
        try:
            excel_file = pd.ExcelFile(file_path); logger.info(f"Extracting XLSX: {base_filename}")
            for i, sheet_name in enumerate(excel_file.sheet_names):
                try:
                    df = excel_file.parse(sheet_name); df = df.fillna(''); sheet_string = df.to_string(index=False, header=True, na_rep=''); cleaned_content = self.clean_text(sheet_string)
                    if cleaned_content: all_content.append({"type": "excel_sheet", "content": cleaned_content, "source_info": {"sheet": sheet_name, "rows": f"1-{len(df)}"}, "file_type": "xlsx"})
                except Exception as e_sheet: logger.warning(f"Sheet processing error XLSX sheet '{sheet_name}': {e_sheet}")
            logger.info(f"Extracted {len(all_content)} sheets from XLSX: {base_filename}"); return all_content
        except Exception as e: logger.error(f"Error extracting XLSX {base_filename}: {e}", exc_info=True); return []

    def _extract_csv(self, file_path): # (Unchanged)
        all_content = []; base_filename = os.path.basename(file_path)
        try:
            encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']; df = None
            for enc in encodings_to_try:
                 try: df = pd.read_csv(file_path, encoding=enc, low_memory=False); logger.info(f"Read CSV {base_filename} with {enc}"); break
                 except UnicodeDecodeError: continue
                 except Exception as e_read: logger.warning(f"CSV read error {base_filename} with {enc}: {e_read}")
            if df is None: logger.error(f"Could not read CSV {base_filename}"); return []
            df = df.fillna(''); csv_string = df.to_string(index=False, header=True, na_rep=''); cleaned_content = self.clean_text(csv_string)
            if cleaned_content: all_content.append({"type": "csv_data", "content": cleaned_content, "source_info": {"rows": f"1-{len(df)}"}, "file_type": "csv"})
            logger.info(f"Extracted content from CSV: {base_filename}"); return all_content
        except Exception as e: logger.error(f"Error extracting CSV {base_filename}: {e}", exc_info=True); return []

    def _extract_pptx(self, file_path): # (Unchanged)
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element; merge_threshold = self.config.pptx_merge_threshold_words
        try:
            prs = Presentation(file_path); logger.info(f"Extracting PPTX: {base_filename}")
            pending_title_slide_data = None
            for i, slide in enumerate(prs.slides):
                 current_slide_number = i + 1; current_slide_title_text = ""; current_slide_other_texts = []; current_slide_has_title = False; title_shape = None
                 try:
                     if slide.shapes.title: title_shape = slide.shapes.title; current_slide_has_title = True
                 except AttributeError: pass
                 if title_shape and title_shape.has_text_frame: current_slide_title_text = self.clean_text(title_shape.text_frame.text)
                 for shape in slide.shapes:
                     if shape == title_shape: continue
                     is_body = False
                     if shape.is_placeholder:
                          try: ph_type = shape.placeholder_format.type; is_body = ph_type in [MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT, MSO_SHAPE_TYPE.SUBTITLE, MSO_SHAPE_TYPE.CONTENT]
                          except AttributeError: pass
                     if shape.has_text_frame:
                          text = shape.text_frame.text; cleaned = self.clean_text(text)
                          if cleaned: cleaned = cleaned[:max_chars] + "...(truncated)" if len(cleaned)>max_chars else cleaned; prefix = "[Body]: " if is_body else ""; current_slide_other_texts.append(prefix + cleaned)
                 if slide.has_notes_slide:
                     try: 
                        notes_text = slide.notes_slide.notes_text_frame.text; cleaned_notes = self.clean_text(notes_text);
                        if cleaned_notes: cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated)" if len(cleaned_notes)>max_chars*2 else cleaned_notes; current_slide_other_texts.append(f"[Notes]: {cleaned_notes}")
                     except Exception as e_notes: logger.warning(f"PPTX notes extraction error slide {current_slide_number}: {e_notes}")
                 current_full_content = current_slide_title_text + ("\n" if current_slide_title_text and current_slide_other_texts else "") + "\n".join(current_slide_other_texts)
                 current_full_content = current_full_content.strip(); other_words = sum(len(s.split()) for s in current_slide_other_texts)
                 should_merge = False
                 if pending_title_slide_data and pending_title_slide_data['has_title'] and pending_title_slide_data['other_words'] <= merge_threshold and current_full_content: should_merge = True
                 if should_merge:
                     merged_text = f"[Slide {pending_title_slide_data['number']}]:\n{pending_title_slide_data['content']}\n---\n[Slide {current_slide_number}]:\n{current_full_content}";
                     all_content.append({"type": "slide_text_merged", "content": merged_text.strip(), "source_info": {"slide": current_slide_number, "merged_from": pending_title_slide_data['number']}, "file_type": "pptx"})
                     pending_title_slide_data = None
                 else:
                     if pending_title_slide_data and pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"}); pending_title_slide_data = None
                     if current_slide_has_title and other_words <= merge_threshold and current_full_content: pending_title_slide_data = {"content": current_full_content, "number": current_slide_number, "has_title": current_slide_has_title, "other_words": other_words}
                     elif current_full_content: all_content.append({"type": "slide_text", "content": current_full_content, "source_info": {"slide": current_slide_number}, "file_type": "pptx"})
            if pending_title_slide_data and pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
            logger.info(f"Extracted {len(all_content)} blocks from PPTX: {base_filename}"); return all_content
        except Exception as e: logger.error(f"Error extracting PPTX {base_filename}: {e}", exc_info=True); return []

    def extract_content(self, file_path): # (Unchanged)
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf': return self._extract_pdf(file_path)
        elif ext == '.xlsx': return self._extract_xlsx(file_path)
        elif ext == '.csv': return self._extract_csv(file_path)
        elif ext == '.pptx': return self._extract_pptx(file_path)
        else: logger.warning(f"Unsupported file type skipped: {file_path}"); return []

    def chunk_content(self, all_content): # (Unchanged)
        chunks = []; stride = self.config.chunk_size - self.config.overlap; stride = max(1, stride)
        for item in all_content:
            content = item.get('content', ''); words = content.split()
            if not words: continue
            for i in range(0, len(words), stride):
                chunk_words = words[i : i + self.config.chunk_size]; chunk_text = " ".join(chunk_words)
                if chunk_text: chunks.append({"content": chunk_text, "source_info": item.get('source_info', {}).copy(), "file_type": item.get('file_type', 'unknown'), "type": item.get('type', 'unknown')})
        logger.info(f"Created {len(chunks)} chunks from {len(all_content)} content blocks."); return chunks

    def _get_safe_filename(self, file_name): # (Unchanged)
        base = os.path.splitext(file_name)[0]; safe = re.sub(r'[^\w\.\-]+', '_', base); safe = safe.strip('._'); return safe or "invalid_filename"

    def get_index_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.index")
    def get_embedding_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.npy")
    def get_chunks_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.json")

    def load_faiss_index(self, file_name, embedding_dim): # (Unchanged)
        path = self.get_index_path(file_name)
        if os.path.exists(path):
            try: 
                index = faiss.read_index(path)
                if index.d != embedding_dim: logger.warning(f"Index dim mismatch {file_name}"); return faiss.IndexFlatL2(embedding_dim)
                logger.info(f"Loaded FAISS index {file_name} ({index.ntotal} vectors)"); return index
            except Exception as e: logger.error(f"Error reading index {path}: {e}")
        return faiss.IndexFlatL2(embedding_dim)

    def save_chunks(self, file_name, chunks): # (Unchanged)
        path = self.get_chunks_path(file_name)
        try:
            with open(path, 'w', encoding='utf-8') as f: json.dump(chunks, f, indent=2, ensure_ascii=False)
            logger.info(f"Saved {len(chunks)} chunks to {path}")
        except Exception as e: logger.error(f"Error saving chunks {path}: {e}")

    def load_chunks(self, file_name): # (Unchanged)
        path = self.get_chunks_path(file_name)
        if os.path.exists(path):
            try:
                with open(path, 'r', encoding='utf-8') as f: chunks = json.load(f)
                logger.info(f"Loaded {len(chunks)} chunks from {path}"); return chunks
            except Exception as e: logger.error(f"Error loading/decoding chunks {path}: {e}"); return None
        return None

    def process_files(self, progress_callback=None): # (Unchanged)
        # ... (identical to previous version) ...
        self.file_chunks = {}; self.faiss_indexes = {}; self.processed_files = set(); data_dir = self.config.data_dir; supported_ext = self.config.supported_extensions
        try: all_files = os.listdir(data_dir); process_list = sorted([f for f in all_files if f.lower().endswith(supported_ext)])
        except Exception as e: logger.error(f"Error listing files in {data_dir}: {e}"); return False
        if not process_list: logger.warning(f"No supported files in {data_dir}"); return True
        logger.info(f"Processing {len(process_list)} files from {data_dir}"); embedding_dim = self.encoder_model.get_sentence_embedding_dimension(); total_files = len(process_list)
        files_processed_successfully = 0
        for idx, file_name in enumerate(process_list):
            # ... (inner loop identical to previous version - extract, chunk, embed, index, verify) ...
             try: # Wrap file processing
                # ... (load chunks or extract/chunk/save) ...
                chunks = self.load_chunks(file_name)
                if chunks is None:
                    # ... (extract/chunk/save logic) ...
                    all_content = self.extract_content(os.path.join(data_dir, file_name))
                    if not all_content: logger.warning(f"No content from {file_name}"); continue
                    chunks = self.chunk_content(all_content)
                    if not chunks: logger.warning(f"No chunks for {file_name}"); continue
                    self.save_chunks(file_name, chunks)
                self.file_chunks[file_name] = chunks

                # ... (verification and regeneration logic for embeddings/index) ...
                faiss_index = None; regenerate = False; embeddings = None
                index_path=self.get_index_path(file_name); emb_path=self.get_embedding_path(file_name); chunks_path=self.get_chunks_path(file_name)
                if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                    try:
                        embeddings = np.load(emb_path)
                        if embeddings.ndim!=2 or embeddings.shape[1]!=embedding_dim or embeddings.shape[0]!=len(chunks): regenerate=True; logger.warning(f"Verification mismatch {file_name}")
                        else: faiss_index = self.load_faiss_index(file_name, embedding_dim);
                        if faiss_index.ntotal!=embeddings.shape[0]: regenerate=True; faiss_index=None; logger.warning(f"Index/Embedding count mismatch {file_name}")
                    except Exception as e_v: regenerate=True; faiss_index=None; embeddings=None; logger.error(f"Verification error {file_name}: {e_v}")
                else: regenerate=True

                if regenerate or faiss_index is None:
                    if not chunks: continue
                    content_list = [c['content'] for c in chunks]
                    if not content_list: continue
                    should_gen_new_emb = regenerate and (embeddings is None or embeddings.shape[0]!=len(chunks) or embeddings.shape[1]!=embedding_dim)
                    if should_gen_new_emb:
                         embeddings = self.encoder_model.encode(content_list, batch_size=64, show_progress_bar=False, convert_to_numpy=True).astype('float32')
                         if embeddings.shape[0] == 0: continue
                         np.save(emb_path, embeddings); logger.info(f"Saved {embeddings.shape[0]} new embeddings {file_name}")
                    elif embeddings is None: continue # Error state
                    # Generate index
                    faiss_index = faiss.IndexFlatL2(embedding_dim); faiss_index.add(embeddings); faiss.write_index(faiss_index, index_path); logger.info(f"Saved new index {file_name}")

                # ... (store index) ...
                if faiss_index is not None and faiss_index.ntotal > 0: self.faiss_indexes[file_name] = faiss_index; self.processed_files.add(file_name); files_processed_successfully += 1
                else: logger.warning(f"No valid index for {file_name}")
             except Exception as e_file: logger.error(f"Error processing {file_name}: {e_file}", exc_info=True); # ... (cleanup) ...

        logger.info(f"--- Processing Complete. Indexed {files_processed_successfully}/{total_files} documents. ---")
        return True

    def query_files(self, query, k_override=None): # (Unchanged)
        # ... (identical to previous version) ...
        if not self.faiss_indexes: logger.warning("No indexes available"); return {}
        k = k_override if k_override is not None else self.config.k_retrieval
        logger.info(f"Retrieving top {k} candidates per file for query: '{query[:100]}...'")
        results = {}
        try:
            q_emb = self.encoder_model.encode([query], convert_to_numpy=True).astype("float32")
            if q_emb.ndim != 2 or q_emb.shape[0] != 1: raise ValueError("Bad query embedding shape")
            for fname, index in self.faiss_indexes.items():
                 if index is None or index.ntotal == 0: continue
                 chunks = self.file_chunks.get(fname)
                 if not chunks: logger.error(f"Chunks missing for {fname}!"); continue
                 k_search = min(k, index.ntotal);
                 if k_search <= 0: continue
                 try:
                     D, I = index.search(q_emb, k=k_search); indices, dists = I[0], D[0]
                     file_res = []
                     processed = set()
                     for i, idx in enumerate(indices):
                          if idx == -1 or idx in processed or not (0 <= idx < len(chunks)): continue
                          processed.add(idx); chunk = chunks[idx]
                          file_res.append({"source_info": chunk.get('source_info',{}), "file_type": chunk.get('file_type','?'), "content": chunk.get('content',''), "score": round(float(dists[i]),4), "type": chunk.get('type','?')})
                     file_res.sort(key=lambda x: x['score'])
                     if file_res:
                          if fname not in results: results[fname] = []
                          results[fname].extend(file_res)
                 except Exception as e_s: logger.error(f"Error searching {fname}: {e_s}")
            return results
        except Exception as e_q: logger.error(f"Error in query_files setup: {e_q}"); return {}

    def aggregate_context(self, query_results, strategy="top_k"): # (Unchanged)
        # ... (identical to previous version - Standard Top K by score) ...
        all_context = {}; max_chars = self.config.max_context_chars
        logger.info(f"Standard Aggregation: strategy '{strategy}', max_chars ~{max_chars}")
        if not query_results: return {"combined_context": "", "source_files": [], "aggregated_chunks_count": 0}
        flat_results = []
        for fname, res_list in query_results.items():
             for res in res_list:
                  if 'content' in res and 'score' in res: flat_results.append({**res, "file_name": fname})
        if not flat_results: return {"combined_context": "", "source_files": [], "aggregated_chunks_count": 0}
        flat_results.sort(key=lambda x: x['score'])
        context_str = ""; chars = 0; count = 0; sources = set(); limit = self.config.k_retrieval if strategy=="top_k" else len(flat_results)
        processed_keys = set()
        for i, res in enumerate(flat_results):
            if count >= limit: break
            key = (res['file_name'], tuple(sorted(res.get('source_info',{}).items())), res['content'])
            if key in processed_keys: continue
            # ... (build source string) ...
            source_parts = [f"Source: {res['file_name']}"]; file_type = res.get('file_type','?'); source_info = res.get('source_info',{})
            if file_type=='pptx': source_parts.append(f"Slide: {source_info.get('slide','N/A')}" + (f" (merged from {source_info.get('merged_from')})" if source_info.get('merged_from') else ""))
            elif file_type=='pdf': source_parts.append(f"Page: {source_info.get('page','N/A')}")
            elif file_type=='xlsx': source_parts.append(f"Sheet: {source_info.get('sheet','N/A')}")
            elif file_type=='csv': source_parts.append(f"CSV Rows: ~{source_info.get('rows','all')}")
            source_parts.append(f"Score: {res['score']:.4f}")
            source_str = ", ".join(source_parts)
            header = f"--- Context from {source_str} ---\n"; body = res['content']; to_add = header + body + "\n\n"; n_chars = len(to_add)
            if chars + n_chars <= max_chars: context_str += to_add; chars += n_chars; count += 1; sources.add(res['file_name']); processed_keys.add(key)
            else:
                 if count==0: # First chunk too large
                      remain = max_chars - len(header) - 20
                      if remain > 50: context_str += header + body[:remain] + "\n[...TRUNCATED...]\n\n"; chars=len(context_str); count+=1; sources.add(res['file_name']); processed_keys.add(key); logger.warning("Standard Agg: Truncated first chunk")
                 break # Stop aggregation
        final_context = context_str.strip()
        result = {"combined_context": final_context, "source_files": sorted(list(sources)), "aggregated_chunks_count": count}
        if count > 0: logger.info(f"Standard Aggregation: Aggregated {chars} chars from {count} chunks.")
        else: logger.warning("Standard Aggregation: No context aggregated.")
        return result

    def _call_llm_for_analysis(self, prompt_messages, task_description): # (Unchanged)
        try:
            llm = self._get_llm(); logger.info(f"Calling LLM for {task_description}...")
            response = llm.invoke(prompt_messages)
            content = response.content.strip() if hasattr(response, 'content') else str(response).strip()
            if not content or content.lower().startswith("i cannot"): logger.warning(f"LLM analysis {task_description} gave empty/refusal"); return None
            logger.info(f"LLM analysis response: '{content[:100]}...'"); return content
        except Exception as e: logger.error(f"LLM call failed {task_description}: {e}"); return None

    def _classify_query(self, query): # (Unchanged)
        # ... (uses _call_llm_for_analysis) ...
        system_prompt = """You are an expert query analyzer... Classify as 'Simple Retrieval' or 'Complex/Reasoning'. Respond ONLY with the classification."""
        human_prompt = f"Classify:\n---\n{query}\n---\nClassification:"
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        try:
            classification = self._call_llm_for_analysis(messages, "query classification")
            if classification in ["Simple Retrieval", "Complex/Reasoning"]: return classification
            cleaned = str(classification).strip().replace("'", "").replace('"', '')
            if cleaned in ["Simple Retrieval", "Complex/Reasoning"]: logger.warning(f"Cleaned class: '{classification}'->'{cleaned}'"); return cleaned
            logger.warning(f"Unexpected class: '{classification}'. Defaulting Simple."); return "Simple Retrieval"
        except Exception as e: logger.error(f"Classify error: {e}"); return "Simple Retrieval"

    def _decompose_query(self, query): # (Unchanged)
        # ... (uses _call_llm_for_analysis) ...
        system_prompt = """You are an expert query decomposer... Format ONLY as a numbered list... If simple, return '1. <original query>'."""
        human_prompt = f"Decompose:\n---\n{query}\n---\nDecomposition:"
        messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
        decomposition = self._call_llm_for_analysis(messages, "query decomposition")
        if decomposition:
            sub_queries = [re.sub(r"^\s*\d+\.\s*", "", line).strip() for line in decomposition.split('\n') if line.strip()]
            sub_queries = [q for q in sub_queries if q]
            if sub_queries:
                if len(sub_queries) == 1 and sub_queries[0].lower() == query.lower(): logger.info("Query simple, no decomp."); return [query]
                logger.info(f"Decomposed '{query[:50]}...' into: {sub_queries}"); return sub_queries
            else: logger.warning(f"Decomp empty list for '{query[:50]}...'"); return [query]
        else: logger.warning(f"Decomp failed for '{query[:50]}...'"); return [query]

    def _extract_keywords(self, text): # (Unchanged)
        # ... (simple keyword extraction logic) ...
        if not text: return []
        words = text.split(); translator = str.maketrans('', '', string.punctuation); keywords = set()
        for word in words:
             cleaned = word.translate(translator)
             if len(cleaned) > 1:
                 # Keep title case, all caps (potential acronyms), or digits
                 if cleaned.istitle() or cleaned.isupper() or cleaned.isdigit(): keywords.add(cleaned.lower())
        logger.info(f"Keywords: {list(keywords)} from '{text[:50]}...'"); return list(keywords) if keywords else []


    # ==========================================================================
    # START OF MODIFIED query_llm FUNCTION (Specific Excel/CSV Counting Prompt)
    # ==========================================================================
    def query_llm(self, query, context_data, retry_count=1, is_structured_query_hint=False, is_excel_csv_counting_hint=False):
        """Generates an answer using the LLM based on the query and context.
           Uses specialized prompts based on hints.
        """
        combined_context = context_data.get("combined_context", ""); source_files = context_data.get("source_files", []); source_file_str = ", ".join(source_files) if source_files else "the provided documents"

        if not combined_context:
            logger.warning("No context provided for LLM query.")
            no_context_msg = f"Could not generate answer: No relevant context was found or aggregated from {source_file_str}."
            # Check if retrieval/filtering happened but aggregation failed
            if context_data.get("retrieval_results_after_filter_count", 0) > 0 and not combined_context:
                 no_context_msg = f"Could not generate answer: Relevant chunks might have been found in {source_file_str}, but could not be aggregated into the context (likely due to size limits)."
            return no_context_msg

        try:
            llm = self._get_llm()
            system_prompt = ""

            # === CHOOSE PROMPT BASED ON HINTS (MOST SPECIFIC FIRST) ===
            if is_excel_csv_counting_hint:
                # NEW: Highly specific prompt for counting in Excel/CSV context
                system_prompt = f"""You are an AI assistant specialized in analyzing data from Excel or CSV files, presented as text in the context below (from file(s): '{source_file_str}').
Your task is to answer the user's question which specifically asks for a COUNT or QUANTITY based on the provided context.
1. Carefully analyze the text context, paying attention to rows, columns, and specific values mentioned.
2. Identify the exact criteria from the User Question (e.g., the region 'South USA', specific values, etc.).
3. Count the number of rows/entries/instances in the context that match ALL the criteria from the User Question.
4. Respond ONLY with the final numerical count. If the count is 0, respond with '0'.
5. If the context does not contain enough information to perform the count accurately (e.g., the relevant column or data is missing), respond ONLY with: "Based on the provided context from {source_file_str}, I cannot determine the count."
Do NOT provide explanations, summaries, or list examples. Just the number or the inability statement. Do not use external knowledge."""
                logger.info("Using SPECIFIC system prompt for Excel/CSV counting query.")

            elif is_structured_query_hint:
                # Existing prompt for general structured data reasoning
                system_prompt = f"""You are an AI assistant specialized in analyzing data presented in text format, often extracted from tables in documents like Excel or CSV files found in: '{source_file_str}'.
Your task is to answer the user's question based ONLY on the provided context below.
The user's question likely requires calculation, aggregation, comparison, filtering, or specific data lookup within the text representation of the tables provided in the context.
Analyze the structure and data within the 'START CONTEXT'/'END CONTEXT' section carefully. Pay attention to headers, columns, and rows implied by the text formatting.
Perform the necessary reasoning or calculation based *strictly* on the information presented in the context. Do not invent data or perform operations not supported by the context. Present results clearly. If calculations are involved, briefly mention the result.
If the answer cannot be determined or calculated from the provided context, state clearly: "Based on the provided context from {source_file_str}, I cannot answer this question." And briefly explain why (e.g., "the relevant data for [specific region/column] is missing").
Do NOT use any external knowledge or perform web searches."""
                logger.info("Using GENERALIZED system prompt for structured data query.")

            else:
                # Standard RAG prompt for simple retrieval / other file types
                system_prompt = f"""You are an AI assistant answering questions based ONLY on the provided context from document(s): '{source_file_str}'.
Use ONLY information contained within the 'START CONTEXT'/'END CONTEXT' section below.
Do not make assumptions or use external knowledge. Base your entire answer on the text provided.
If the answer isn't explicitly stated or cannot be inferred solely from the context, state: "Based on the provided context from {source_file_str}, I cannot answer this question." """
                logger.info("Using standard RAG system prompt.")
            # === END OF PROMPT CHOICE ===

            human_prompt = f"""Context from document(s) '{source_file_str}':
--- START CONTEXT ---
{combined_context}
--- END CONTEXT ---

User Question: {query}

Answer based ONLY on the provided context:"""

            full_prompt_messages = [SystemMessage(content=system_prompt), HumanMessage(content=human_prompt)]
            logger.info(f"Querying Groq model {self.config.llm_model} using context from {source_file_str}...")

            answer = f"Error: LLM query failed after retries."
            for attempt in range(retry_count + 1):
                try:
                    response = llm.invoke(full_prompt_messages)
                    answer = response.content.strip() if hasattr(response, 'content') else str(response).strip()
                    logger.info(f"Groq response received (attempt {attempt+1}). Length: {len(answer)}")
                    if not answer and attempt < retry_count: logger.warning("Empty response received, retrying..."); time.sleep(1); continue
                    # For counting prompt, we expect a number or the specific refusal.
                    # We could add validation here if needed, but let's trust the LLM first.
                    return answer or f"Received empty response from LLM for context from {source_file_str} after {attempt+1} attempts."
                except Exception as e_api:
                    logger.warning(f"Groq API call attempt {attempt+1} failed: {e_api}")
                    if attempt < retry_count: sleep_time = 1.5 ** attempt; logger.info(f"Retrying Groq query in {sleep_time:.2f}s..."); time.sleep(sleep_time)
                    else: logger.error(f"LLM query failed after {retry_count+1} attempts. Error: {e_api}", exc_info=True); answer = f"Error: Failed to get answer from LLM. Please check logs. (API Error: {e_api})"
            return answer
        except Exception as e_setup:
            logger.error(f"Error setting up or calling Groq API: {e_setup}", exc_info=True)
            if self._is_streamlit: import streamlit as st; st.error(f"LLM query failed: {e_setup}")
            return f"Error: Could not query LLM due to setup error: {e_setup}"
    # ==========================================================================
    # END OF MODIFIED query_llm FUNCTION
    # ==========================================================================


    # ==========================================================================
    # START OF MODIFIED run_query FUNCTION (Adds Excel/CSV Counting Hint)
    # ==========================================================================
    def run_query(self, query, context_strategy="top_k"):
        """Complete query pipeline: classify, (decompose), retrieve MORE, filter, aggregate, determine hints, query LLM."""
        logger.info(f"--- Starting query execution (v3+excel_count): '{query[:100]}...' ---")
        final_results = {
            "query": query, "classification": None, "sub_queries": None, "keywords_used_for_filter": None,
            "retrieval_results_before_filter_count": 0, "retrieval_results_after_filter_count": 0,
            "chunks_considered_for_context": [], "chunks_included_in_context": [],
            "aggregated_context_data": {}, "answer": "", "status": "Started", "error": None,
            "query_type_hint": "standard", # standard, structured, excel_csv_counting
            "aggregation_method": "standard" # standard_top_k, filtered_all
        }
        is_streamlit = self._is_streamlit
        if is_streamlit: import streamlit as st

        # Define keywords indicating a counting/aggregation query
        COUNTING_KEYWORDS = ['how many', 'count', 'total', 'number of', 'sum', 'average', 'list all', 'show all']

        try:
            # === 1. Classify Query === (Unchanged)
            spinner_msg = "Analyzing query..."
            if is_streamlit: 
                with st.spinner(spinner_msg): classification = self._classify_query(query)
            else: logger.info(spinner_msg); classification = self._classify_query(query)
            final_results["classification"] = classification; logger.info(f"Query classified as: {classification}")

            # === 2. Decompose if Complex === (Unchanged)
            queries_to_retrieve = [query]
            if classification == "Complex/Reasoning":
                # ... (decomposition logic unchanged) ...
                spinner_msg = "Decomposing complex query..."
                if is_streamlit: 
                    with st.spinner(spinner_msg): sub_queries = self._decompose_query(query)
                else: logger.info(spinner_msg); sub_queries = self._decompose_query(query)
                if sub_queries and (len(sub_queries) > 1 or sub_queries[0] != query):
                     queries_to_retrieve = sub_queries; final_results["sub_queries"] = sub_queries; logger.info(f"Using sub-queries: {sub_queries}")
                else: logger.info("No distinct sub-queries generated.")

            # === 3. Determine Retrieval K & Extract Keywords === (Unchanged)
            k_to_use = self.config.k_retrieval; keywords_for_filter = []; needs_filtering = False
            keywords_for_filter = self._extract_keywords(query)
            if classification == "Complex/Reasoning" or keywords_for_filter:
                 k_to_use = self.config.k_retrieval_large; needs_filtering = True
                 logger.info(f"Using K={k_to_use} and keyword filtering.")
                 if keywords_for_filter: final_results["keywords_used_for_filter"] = keywords_for_filter
                 else: # Complex but no keywords in main query, try sub-queries
                     if final_results["sub_queries"]:
                          all_sub_q_keywords=set().union(*(self._extract_keywords(sq) for sq in final_results["sub_queries"]))
                          if all_sub_q_keywords: keywords_for_filter=list(all_sub_q_keywords); final_results["keywords_used_for_filter"]=keywords_for_filter; logger.info(f"Using keywords from sub-queries: {keywords_for_filter}")
                          else: needs_filtering = False # Revert if still no keywords
            if not needs_filtering: logger.info(f"Using standard K={k_to_use}, no keyword filtering.")


            # === 4. Retrieve relevant chunks === (Unchanged)
            spinner_msg = f"Searching documents (k={k_to_use})..."; all_query_results_raw = {}; retrieval_spinner = None
            if is_streamlit: retrieval_spinner = st.spinner(spinner_msg)
            # ... (retrieval loop unchanged) ...
            try:
                if retrieval_spinner: retrieval_spinner.__enter__()
                for i, q in enumerate(queries_to_retrieve):
                    logger.info(f"Retrieving part {i+1}/{len(queries_to_retrieve)} (k={k_to_use}): '{q[:50]}...'")
                    results_for_q = self.query_files(q, k_override=k_to_use)
                    for file, res_list in results_for_q.items():
                        if file not in all_query_results_raw: all_query_results_raw[file] = []
                        all_query_results_raw[file].extend(res_list)
            finally:
                if retrieval_spinner: retrieval_spinner.__exit__(None, None, None)
            raw_count = sum(len(v) for v in all_query_results_raw.values()); final_results["retrieval_results_before_filter_count"] = raw_count; logger.info(f"Retrieved {raw_count} raw candidates.")


            # === 5. Filter Retrieval Results === (Unchanged)
            filtered_query_results_dict = {}; filtered_count = 0
            if needs_filtering and keywords_for_filter and raw_count > 0:
                 # ... (filtering logic unchanged) ...
                 logger.info(f"Applying keyword filter: {keywords_for_filter}")
                 kept_results_by_file = {}
                 flat_raw_results = [ {**res, "file_name": file_name} for file_name, res_list in all_query_results_raw.items() for res in res_list ]
                 for res in flat_raw_results:
                     content_lower = res.get('content', '').lower()
                     if any(keyword.lower() in content_lower for keyword in keywords_for_filter):
                         file_name = res['file_name']
                         if file_name not in kept_results_by_file: kept_results_by_file[file_name] = []
                         kept_results_by_file[file_name].append(res); filtered_count += 1
                 if filtered_count == 0: logger.warning("Keyword filtering removed all chunks!")
                 filtered_query_results_dict = kept_results_by_file; logger.info(f"Filtered down to {filtered_count} chunks.")
            elif needs_filtering: logger.warning("Filtering needed but no keywords; using raw results."); filtered_query_results_dict = all_query_results_raw; needs_filtering = False
            else: logger.info("No filtering applied."); filtered_query_results_dict = all_query_results_raw


            # === 6. De-duplicate and Sort === (Unchanged)
            unique_results_list = []
            # ... (de-duplication logic identical, operates on filtered_query_results_dict) ...
            for file, res_list in filtered_query_results_dict.items():
                 unique_content_per_file = {}
                 for res in res_list:
                     key = (res['content'], tuple(sorted(res.get('source_info',{}).items())))
                     if key not in unique_content_per_file or res['score'] < unique_content_per_file[key]['score']: unique_content_per_file[key] = {**res, "file_name": file}
                 unique_results_list.extend(unique_content_per_file.values())
            unique_results_list.sort(key=lambda x: x['score']) # Sort unique results by score
            final_results["retrieval_results_after_filter_count"] = len(unique_results_list); logger.info(f"Prepared {len(unique_results_list)} unique chunks.")
            final_results["chunks_considered_for_context"] = unique_results_list # Store candidates
            if not unique_results_list: # Handle no results after filtering/dedup
                logger.warning("No relevant chunks found after filtering/de-duplication."); final_results["status"] = "Completed: No relevant information found."; final_results["answer"] = f"Could not find relevant information for: '{query}'"; return final_results


            # === 7. Aggregate Context (Conditional Logic) === (Unchanged)
            spinner_msg = "Gathering context..."; aggregated_context_data = {}; max_chars = self.config.max_context_chars
            if needs_filtering and final_results["retrieval_results_after_filter_count"] > 0:
                # --- AGGREGATE ALL FILTERED RESULTS --- (Unchanged)
                final_results["aggregation_method"] = "filtered_all"
                logger.info(f"Using Filtered Aggregation: Including all {len(unique_results_list)} chunks (limit {max_chars} chars).")
                # ... (build context string from unique_results_list, respecting max_chars) ...
                agg_context_str = ""; agg_chars = 0; agg_count = 0; agg_sources = set(); included_chunks = []
                for i, res in enumerate(unique_results_list):
                    # ... (build source header, content_to_add) ...
                    source_parts = [f"Source: {res['file_name']}"]; file_type=res.get('file_type','?'); info=res.get('source_info',{}); score=res.get('score',0.0) # simplified
                    if file_type=='xlsx': source_parts.append(f"Sheet: {info.get('sheet','?')}")
                    elif file_type=='csv': source_parts.append(f"CSV: ~{info.get('rows','?')}")
                    # ... (other types) ...
                    source_parts.append(f"Score: {score:.4f}")
                    source_str = ", ".join(source_parts); header = f"--- Context from {source_str} ---\n"; body = res['content']; to_add = header + body + "\n\n"; n_chars = len(to_add)
                    if agg_chars + n_chars <= max_chars: agg_context_str += to_add; agg_chars += n_chars; agg_count += 1; agg_sources.add(res['file_name']); included_chunks.append(res)
                    else: logger.info(f"Filtered Agg: Max chars reached at chunk {agg_count+1}."); break
                final_context = agg_context_str.strip()
                aggregated_context_data = {"combined_context": final_context, "source_files": sorted(list(agg_sources)), "aggregated_chunks_count": agg_count}
                final_results["chunks_included_in_context"] = included_chunks
                logger.info(f"Filtered Aggregation: Aggregated {agg_chars} chars from {agg_count} chunks.")

            else:
                # --- USE STANDARD AGGREGATION --- (Unchanged)
                final_results["aggregation_method"] = "standard_top_k"
                logger.info(f"Using Standard Aggregation: Top {self.config.k_retrieval} by score.")
                # ... (call self.aggregate_context) ...
                results_for_standard_agg = {}
                for res in unique_results_list: fname = res['file_name']; results_for_standard_agg.setdefault(fname, []).append(res)
                if is_streamlit: 
                    with st.spinner(spinner_msg): aggregated_context_data = self.aggregate_context(results_for_standard_agg, strategy=context_strategy)
                else: logger.info(spinner_msg); aggregated_context_data = self.aggregate_context(results_for_standard_agg, strategy=context_strategy)
                # Approx included chunks for standard method
                final_results["chunks_included_in_context"] = unique_results_list[:aggregated_context_data.get("aggregated_chunks_count", 0)]

            final_results["aggregated_context_data"] = aggregated_context_data


            # === 8. Determine LLM Prompt Hints (ADDED Excel/CSV Counting Check) ===
            is_structured_query_hint = False
            is_excel_csv_counting_hint = False # New hint

            if aggregated_context_data and aggregated_context_data.get("source_files"):
                source_files = aggregated_context_data["source_files"]
                # Check if primary sources are Excel or CSV
                is_excel_csv_source = any(f.lower().endswith(('.xlsx', '.csv')) for f in source_files)

                if classification == "Complex/Reasoning" and is_excel_csv_source:
                    # It's complex and involves Excel/CSV. Is it specifically a counting query?
                    query_lower = query.lower()
                    is_counting_query = any(keyword in query_lower for keyword in COUNTING_KEYWORDS)

                    if is_counting_query:
                        is_excel_csv_counting_hint = True
                        final_results["query_type_hint"] = "excel_csv_counting"
                        logger.info("Query is Complex, involves Excel/CSV, and contains counting keywords. Using EXCEL/CSV COUNTING prompt.")
                    else:
                        # Complex, involves Excel/CSV, but not explicitly counting
                        is_structured_query_hint = True
                        final_results["query_type_hint"] = "structured"
                        logger.info("Query is Complex, involves Excel/CSV, but no counting keywords detected. Using GENERAL STRUCTURED prompt.")
                elif classification == "Complex/Reasoning": # Complex but not Excel/CSV source
                     is_structured_query_hint = True # Use general structured prompt? Or standard? Let's try structured.
                     final_results["query_type_hint"] = "structured"
                     logger.info("Query is Complex, but primary sources are not Excel/CSV. Using GENERAL STRUCTURED prompt.")
                else: # Simple Retrieval
                    final_results["query_type_hint"] = "standard"
                    logger.info("Query is Simple Retrieval. Using STANDARD RAG prompt.")
            else: # No context / sources
                 final_results["query_type_hint"] = "standard"
                 logger.info("No context/sources for LLM. Using STANDARD RAG prompt by default.")


            # === 9. Query LLM with Aggregated Context and Hints ===
            if not aggregated_context_data or not aggregated_context_data.get("combined_context"):
                 logger.warning(f"Context is empty for LLM query.")
                 final_results["status"] = "Completed: Context empty."; final_results["answer"] = self.query_llm(query, {}, False, False); return final_results # Pass empty context

            spinner_msg = "Generating final answer..."
            llm_call_params = {
                "query": query,
                "context_data": aggregated_context_data,
                "is_structured_query_hint": is_structured_query_hint,
                "is_excel_csv_counting_hint": is_excel_csv_counting_hint # Pass the new hint
            }
            if is_streamlit:
                 with st.spinner(spinner_msg): answer = self.query_llm(**llm_call_params)
            else: logger.info(spinner_msg); answer = self.query_llm(**llm_call_params)

            final_results["answer"] = answer; final_results["status"] = "Completed Successfully."
            logger.info(f"--- Finished query execution (v3+excel_count) for: '{query[:100]}...' ---")

        except Exception as e:
            logger.error(f"Unexpected error during run_query (v3+excel_count): {e}", exc_info=True)
            final_results["status"] = "Failed"; final_results["error"] = str(e); final_results["answer"] = f"Error processing query: {e}"
            if is_streamlit: st.error(f"Unexpected error: {e}")

        return final_results
    # ==========================================================================
    # END OF MODIFIED run_query FUNCTION
    # ==========================================================================


# --- START OF STREAMLIT UI CODE ---
# (Adjust display for new query type hint)
if __name__ == "__main__":
    try: import streamlit as st
    except ImportError: logger.error("Streamlit not installed."); sys.exit(1)

    st.set_page_config(layout="wide", page_title="Multi-Document Q&A with RAG v3.4")
    st.title("📄 Multi-Document Question Answering System (v3.4 - Excel/CSV Counting)")
    st.caption("Query PDFs, XLSX, CSV, PPTX (with enhanced handling for Excel/CSV counting)")

    # --- (load_rag_system, process_documents_once functions remain unchanged) ---
    @st.cache_resource
    def load_rag_system(config_path="config.ini"):
        # (Unchanged)
        if not os.getenv("GROQ_API_KEY"): st.error("🚨 GROQ_API_KEY not set!"); st.stop()
        try: return RAGSystem(config_path=config_path)
        except Exception as e: st.error(f"Fatal RAG Init Error: {e}"); logger.error("RAG Init Error", exc_info=True); st.stop()

    if 'docs_processed' not in st.session_state: st.session_state.docs_processed = False
    if 'indexed_files_list' not in st.session_state: st.session_state.indexed_files_list = []

    def process_documents_once(_rag_system):
        # (Unchanged)
        if st.session_state.docs_processed: st.success(f"✅ Docs processed. Indexed {len(st.session_state.indexed_files_list)} file(s)."); return True, st.session_state.indexed_files_list
        st.header("📚 Document Processing"); st.caption(f"Data: `{os.path.abspath(_rag_system.config.data_dir)}` | Index: `{os.path.abspath(_rag_system.config.index_dir)}`")
        status_placeholder = st.empty(); progress_bar = st.progress(0.0); status_messages = []
        def streamlit_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
             log_prefix = "✅ " if is_done else "❌ " if is_error else "⚠️ " if is_warning else f"⏳ ({stage}) " if stage else "⏳ "; progress_val = 0.0
             if total_steps and current_step is not None and total_steps > 0: stage_prog = {"Starting":0.0,"Extracting":0.1,"Chunking":0.3,"Verifying":0.5,"Embedding":0.6,"Indexing":0.9}.get(stage,0.0); progress_val=min(1.0,(current_step+min(stage_prog,0.95))/total_steps)
             elif is_done: progress_val = 1.0
             full_message=f"{log_prefix}{message}";
             if is_error or is_warning: status_messages.append(full_message)
             try: progress_bar.progress(progress_val); status_placeholder.caption(full_message)
             except Exception as e: logger.warning(f"UI update error: {e}")
        processing_successful = False; indexed_files = []
        with st.spinner("Processing documents..."):
            try:
                processing_successful = _rag_system.process_files(progress_callback=streamlit_progress_callback); indexed_files = sorted(list(_rag_system.faiss_indexes.keys())); indexed_count = len(indexed_files)
                if indexed_count > 0: final_msg = f"✅ Ready! Indexed {indexed_count} document(s)."; streamlit_progress_callback(final_msg, is_done=True); st.session_state.docs_processed = True; st.session_state.indexed_files_list = indexed_files; status_placeholder.success(final_msg)
                elif processing_successful: final_msg = f"⚠️ Processing finished, but no documents indexed."; streamlit_progress_callback(final_msg, is_warning=True, is_done=True); status_placeholder.warning(final_msg)
                else: final_msg = f"❌ Processing failed."; error_details = ". ".join([msg for msg in status_messages if "❌" in msg]); final_msg += f" Details: {error_details}" if error_details else " Check logs."; streamlit_progress_callback(final_msg, is_error=True, is_done=True); status_placeholder.error(final_msg)
            except Exception as e: logger.error(f"Fatal processing error: {e}", exc_info=True); final_msg = f"❌ Fatal error: {e}. Check logs."; status_placeholder.error(final_msg); 
            try: progress_bar.progress(1.0)
            except: pass
        return st.session_state.docs_processed, st.session_state.indexed_files_list

    try:
        rag_sys = load_rag_system()

        if st.button("🔄 Re-process Documents"):
            st.session_state.docs_processed = False; st.session_state.indexed_files_list = []; st.cache_data.clear(); st.cache_resource.clear(); st.rerun()

        is_ready, indexed_files = process_documents_once(rag_sys)

        if is_ready and indexed_files:
            # --- Sidebar (Unchanged) ---
            st.sidebar.success(f"Indexed Documents ({len(indexed_files)}):")
            with st.sidebar.expander("Show Indexed Files"): st.caption("\n".join(f"- {fname}" for fname in indexed_files))
            st.sidebar.info(f"LLM: `{rag_sys.config.llm_model}`"); st.sidebar.info(f"Embedding: `{os.path.basename(rag_sys.config.encoder_model)}`")
            st.sidebar.info(f"K (Final): `{rag_sys.config.k_retrieval}`"); st.sidebar.info(f"K (Initial): `{rag_sys.config.k_retrieval_large}`"); st.sidebar.info(f"Max Tokens: `{rag_sys.config.max_context_tokens}`")

            st.header("💬 Ask a Question")
            user_query = st.text_input("Enter your query:", key="query_input", placeholder="e.g., How many rows have 'South USA' in the region column? or Summarize project alpha.")

            if user_query and st.button("Get Answer", key="submit_query"):
                query_start_time = time.time()
                results_data = rag_sys.run_query(user_query) # Calls the latest run_query
                query_end_time = time.time()

                st.subheader("💡 Answer")
                answer = results_data.get("answer")
                if answer: st.markdown(answer)
                else: st.warning("Could not generate an answer."); st.error(f"Error: {results_data['error']}") if results_data.get("error") else None

                # === Display Query Analysis & Filtering Info ===
                st.subheader("📊 Query Analysis & Filtering")
                col1, col2, col3 = st.columns(3)
                with col1:
                    classification = results_data.get("classification", 'N/A')
                    st.metric("Query Classification", classification)
                    sub_queries = results_data.get("sub_queries")
                    if sub_queries and sub_queries != [user_query]: st.caption(f"{len(sub_queries)} sub-queries used")
                with col2:
                    keywords = results_data.get("keywords_used_for_filter")
                    filter_status = "Active" if keywords else "Inactive"
                    st.metric("Keyword Filter", filter_status)
                    if keywords: st.caption(f"Keywords: {', '.join(keywords)}")
                with col3:
                    # Display more specific prompt type
                    query_type_hint = results_data.get("query_type_hint", "standard")
                    prompt_map = { "standard": "Standard RAG", "structured": "General Structured", "excel_csv_counting": "Excel/CSV Counting" }
                    prompt_display = prompt_map.get(query_type_hint, "Unknown")
                    st.metric("LLM Prompt Used", prompt_display)
                    agg_method = results_data.get("aggregation_method", "N/A").replace("_", " ").title()
                    st.caption(f"Aggregation: {agg_method}")


                # --- Display Sub-queries if they exist ---
                if sub_queries and sub_queries != [user_query]:
                    with st.expander("Sub-queries used for retrieval"): st.code("\n".join(f"{i+1}. {sq}" for i, sq in enumerate(sub_queries)), language=None)

                # === Display Supporting Evidence ===
                st.subheader("🔍 Supporting Evidence")
                agg_context_data = results_data.get("aggregated_context_data", {})
                combined_context = agg_context_data.get("combined_context", "")
                context_source_files = agg_context_data.get("source_files", [])
                chunks_in_context = results_data.get("chunks_included_in_context", [])

                # --- Context Used Expander ---
                if combined_context:
                     context_title = f"Context Used ({len(chunks_in_context)} chunks from {len(context_source_files)} file(s))"
                     with st.expander(context_title):
                         st.text_area("Combined Context Sent to LLM", combined_context, height=300, key="context_display", help="Final text provided to LLM.")
                else: st.info("No context was aggregated or provided to the LLM.")

                # --- Included Chunks Expander ---
                if chunks_in_context:
                     raw_count = results_data.get("retrieval_results_before_filter_count", "?")
                     filtered_count = results_data.get("retrieval_results_after_filter_count", len(chunks_in_context))
                     filter_info = f" (Initial: {raw_count} -> Filtered: {filtered_count})" if results_data.get("keywords_used_for_filter") else ""
                     chunks_title = f"Chunks Included in Final Context ({len(chunks_in_context)} chunks{filter_info})"
                     with st.expander(chunks_title):
                         st.caption(f"These {len(chunks_in_context)} chunks (sorted by relevance score) formed the context. Score may be less relevant if 'Filtered All' aggregation was used.")
                         for i, res in enumerate(chunks_in_context):
                             # (Display logic unchanged)
                             info = res.get('source_info', {}); fname = res.get('file_name', 'N/A'); ftype = res.get('file_type', '?'); ctype = res.get('type', '?'); score = res.get('score', 0.0)
                             parts = [f"**File:** {fname}"]
                             if ftype=='pptx': parts.append(f"Slide: {info.get('slide','?')}" + (f" (merged from {info.get('merged_from')})" if info.get('merged_from') else ""))
                             elif ftype=='pdf': parts.append(f"Page: {info.get('page','?')}")
                             elif ftype=='xlsx': parts.append(f"Sheet: {info.get('sheet','?')}")
                             elif ftype=='csv': parts.append(f"CSV Rows: ~{info.get('rows','?')}")
                             parts.append(f"Type: {ctype}"); parts.append(f"Score: {score:.4f}")
                             st.markdown(f"**Chunk {i+1}**"); st.caption(" | ".join(parts))
                             st.markdown(f"> {res.get('content', '')[:400].replace('\n', ' ')}...")
                             if i < len(chunks_in_context) - 1: st.divider()
                else: st.info("No specific document chunks were included in the context.")

                st.caption(f"Query processed in {query_end_time - query_start_time:.2f} seconds.")

        elif not is_ready and not st.session_state.docs_processed: st.warning("Document processing needed.")
        elif not indexed_files: st.warning(f"No documents indexed in '{rag_sys.config.data_dir}'.")

    except Exception as e: st.error(f"Streamlit App Error: {e}"); logger.exception("Streamlit App Error:")

# --- END OF STREAMLIT UI CODE ---

# --- END OF FILE ragtest_c3_query_adapt_excel_count.py ---