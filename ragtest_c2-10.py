
import os
import re
import logging
import json
import configparser
import pdfplumber
import torch
import numpy as np
import pandas as pd
import faiss
import time
from sentence_transformers import SentenceTransformer
# Langchain Core imports are still useful for structuring prompts even if not directly passed to LLM
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder # Still useful for templating
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import io
from contextlib import redirect_stdout
import sqlite3
import requests # No longer needed for DeepSeek, but kept if other HTTP requests are made.
import openai # For SambaNova

# Load environment variables from .env file
load_dotenv()

# --- START OF RAG SYSTEM CODE (Adapted for SambaNova) ---

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("rag_system_sambanova.log"),
        # logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


class RAGConfig:
    """Configuration class for RAG system with SambaNova."""
    def __init__(self, config_path="config1.ini"):
        self.config = configparser.ConfigParser()
        self.defaults = {
            "PATHS": {
                "data_dir": "./Data/",
                "index_dir": "./Faiss_index/",
                "log_file": "rag_system_sambanova.log",
                "sqlite_db_path": ":memory:"
            },
            "MODELS": {
                "encoder_model": "sentence-transformers/all-MiniLM-L6-v2",
                "sambanova_model": "Meta-Llama-3.3-70B-Instruct", # SambaNova model
                "sambanova_base_url": "https://api.sambanova.ai/v1",
                "device": "auto",
            },
            "PARAMETERS": {
                "chunk_size": "250",
                "overlap": "60",
                "k_retrieval": "6",
                "temperature": "0.1",
                "top_p": "0.1", # Added for SambaNova
                "max_context_tokens": "20000", # This is more for context window, less directly for SambaNova's max_tokens param
                "max_output_tokens_sambanova": "4096", # Max tokens for SambaNova response
                "max_chars_per_element": "1200",
                "pptx_merge_threshold_words": "50",
                "dataframe_query_confidence_threshold": "0.70",
                "max_chat_history_turns": "5"
            },
            "SUPPORTED_EXTENSIONS": {
                "extensions": ".pdf, .xlsx, .csv, .pptx"
            },
            "API_KEYS": {
                 "sambanova_api_key_config": "" # Can be set here or via .env
            }
        }
        if os.path.exists(config_path):
            try: self.config.read(config_path); logger.info(f"Loaded configuration from {config_path}"); self._ensure_defaults()
            except Exception as e: logger.error(f"Error reading config file {config_path}: {e}"); self._set_defaults()
        else: logger.warning(f"Config file {config_path} not found, using defaults"); self._set_defaults()

    def _set_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section)
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value)

    def _ensure_defaults(self):
        for section, options in self.defaults.items():
            if not self.config.has_section(section): self.config.add_section(section); logger.info(f"Added missing section [{section}] from defaults.")
            for option, value in options.items():
                if not self.config.has_option(section, option): self.config.set(section, option, value); logger.info(f"Added missing option '{option} = {value}' to section [{section}] from defaults.")

    @property
    def data_dir(self): return self.config.get("PATHS", "data_dir")
    @property
    def index_dir(self): return self.config.get("PATHS", "index_dir")
    @property
    def log_file(self): return self.config.get("PATHS", "log_file")
    @property
    def sqlite_db_path(self): return self.config.get("PATHS", "sqlite_db_path", fallback=":memory:")
    @property
    def encoder_model(self): return self.config.get("MODELS", "encoder_model")
    @property
    def sambanova_model(self): return self.config.get("MODELS", "sambanova_model")
    @property
    def sambanova_base_url(self): return self.config.get("MODELS", "sambanova_base_url")
    @property
    def device(self):
        device_setting = self.config.get("MODELS", "device", fallback="auto")
        if device_setting == "auto": return "cuda" if torch.cuda.is_available() else "cpu"
        return device_setting
    @property
    def chunk_size(self): return self.config.getint("PARAMETERS", "chunk_size", fallback=250)
    @property
    def overlap(self): return self.config.getint("PARAMETERS", "overlap", fallback=60)
    @property
    def k_retrieval(self): return self.config.getint("PARAMETERS", "k_retrieval", fallback=6)
    @property
    def temperature(self): return self.config.getfloat("PARAMETERS", "temperature", fallback=0.1)
    @property
    def top_p(self): return self.config.getfloat("PARAMETERS", "top_p", fallback=0.1)
    @property
    def max_context_tokens(self): return self.config.getint("PARAMETERS", "max_context_tokens", fallback=6000) # Context window guidance
    @property
    def max_output_tokens_sambanova(self): return self.config.getint("PARAMETERS", "max_output_tokens_sambanova", fallback=1024)

    @property
    def max_chars_per_element(self): return self.config.getint("PARAMETERS", "max_chars_per_element", fallback=1200)
    @property
    def pptx_merge_threshold_words(self): return self.config.getint("PARAMETERS", "pptx_merge_threshold_words", fallback=50)
    @property
    def supported_extensions(self):
        ext_str = self.config.get("SUPPORTED_EXTENSIONS", "extensions", fallback=".pdf, .xlsx, .csv, .pptx")
        return tuple([e.strip() for e in ext_str.lower().split(',') if e.strip()])
    @property
    def dataframe_query_confidence_threshold(self):
        return self.config.getfloat("PARAMETERS", "dataframe_query_confidence_threshold", fallback=0.70)
    @property
    def max_chat_history_turns(self):
        return self.config.getint("PARAMETERS", "max_chat_history_turns", fallback=5)
    @property
    def sambanova_api_key(self):
        key = os.getenv("SAMBANOVA_API_KEY")
        if not key:
            key = self.config.get("API_KEYS", "sambanova_api_key_config", fallback=None)
            if key: logger.info("Using SambaNova API key from config file.")
        return key


class RAGSystem:
    """Retrieval-Augmented Generation system adapted for SambaNova LLM."""
    def __init__(self, config_path="config1.ini"):
        self.config = RAGConfig(config_path)
        logger.info(f"Initializing RAG system (SambaNova Version). Device: {self.config.device}, Supported files: {self.config.supported_extensions}")
        self._is_streamlit = "streamlit" in sys.modules
        os.makedirs(self.config.data_dir, exist_ok=True)
        os.makedirs(self.config.index_dir, exist_ok=True)

        try:
            if self._is_streamlit:
                import streamlit as st
                with st.spinner(f"Loading embedding model ({self.config.encoder_model})..."):
                    self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            else:
                logger.info(f"Loading embedding model ({self.config.encoder_model})...")
                self.encoder_model = SentenceTransformer(self.config.encoder_model).to(self.config.device)
            logger.info("Local embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading embedding model: {e}", exc_info=True)
            if self._is_streamlit:
                import streamlit as st
                st.error(f"Fatal Error: Failed to load embedding model. Check logs. Error: {e}")
                st.stop()
            else: raise

        self.file_chunks = {}
        self.faiss_indexes = {}
        self.processed_files = set()
        self.db_conn = None
        self.db_cursor = None
        self.table_metadata = {}
        self.file_to_table_map = {}
        self._llm_client = None # SambaNova client cache

        self.sambanova_api_key_val = self.config.sambanova_api_key
        if not self.sambanova_api_key_val:
            logger.warning("SAMBANOVA_API_KEY not found in environment variables or config file. LLM calls will fail.")
        
        self._connect_db()

    def _connect_db(self):
        try:
            db_path = self.config.sqlite_db_path
            self.db_conn = sqlite3.connect(db_path, check_same_thread=False)
            self.db_cursor = self.db_conn.cursor()
            logger.info(f"Connected to SQLite database at '{db_path}'")
        except sqlite3.Error as e:
            logger.error(f"Error connecting to SQLite database at '{self.config.sqlite_db_path}': {e}", exc_info=True)
            self.db_conn = None; self.db_cursor = None
            if self._is_streamlit: import streamlit as st; st.error(f"Fatal: Failed to connect to SQLite DB: {e}"); st.stop()
            else: raise ConnectionError(f"Failed to connect to SQLite DB: {e}") from e

    def _close_db(self):
        if self.db_conn:
            try: self.db_conn.commit(); self.db_conn.close(); logger.info("Closed SQLite database connection.")
            except sqlite3.Error as e: logger.error(f"Error closing SQLite connection: {e}", exc_info=True)
            finally: self.db_conn = None; self.db_cursor = None
    
    def __del__(self): self._close_db()

    def _get_llm_client(self):
        """Initializes and returns the SambaNova OpenAI client."""
        if self._llm_client is None:
            if not self.sambanova_api_key_val:
                logger.error("SAMBANOVA_API_KEY not found.")
                if self._is_streamlit: import streamlit as st; st.error("Error: SAMBANOVA_API_KEY is not set.")
                raise ValueError("SAMBANOVA_API_KEY not configured.")
            try:
                self._llm_client = openai.OpenAI(
                    api_key=self.sambanova_api_key_val,
                    base_url=self.config.sambanova_base_url,
                )
                logger.info(f"SambaNova OpenAI client initialized (Model: {self.config.sambanova_model}, Base URL: {self.config.sambanova_base_url})")
            except Exception as e:
                 logger.error(f"Failed to initialize SambaNova client: {e}", exc_info=True)
                 if self._is_streamlit: import streamlit as st; st.error(f"Failed to initialize LLM (SambaNova). Error: {e}")
                 raise
        return self._llm_client

    def _convert_langchain_messages_to_samba_format(self, langchain_messages):
        """Converts Langchain message objects to SambaNova's expected list of dicts."""
        samba_messages = []
        for msg in langchain_messages:
            role = "user" # Default
            if isinstance(msg, SystemMessage):
                role = "system"
            elif isinstance(msg, HumanMessage):
                role = "user"
            elif isinstance(msg, AIMessage):
                role = "assistant"
            
            if hasattr(msg, 'content'):
                samba_messages.append({"role": role, "content": msg.content})
            else:
                logger.warning(f"Message object of type {type(msg)} has no 'content' attribute, skipping.")
        return samba_messages

    # --- Text Cleaning, File/Table Naming, Index Paths (Largely Unchanged) ---
    def clean_text(self, text):
        if not isinstance(text, str): text = str(text)
        text = re.sub(r"\(cid:.*?\)", "", text) 
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text) 
        text = re.sub(r"\s+", " ", text).strip() 
        return text

    def _get_safe_filename(self, file_name):
        base_name = os.path.splitext(file_name)[0]
        return re.sub(r'[^\w\.-]', '_', base_name)

    def _get_safe_table_name(self, file_name, sheet_name=None):
        base = self._get_safe_filename(file_name)
        if sheet_name:
             safe_sheet = re.sub(r'[^\w]', '_', sheet_name)
             name = f"{base}_{safe_sheet}"
        else: name = base
        if not re.match(r"^[a-zA-Z_]", name): name = "_" + name
        return name[:60].lower()

    def _table_exists(self, table_name):
        if not self.db_cursor: return False
        try: self.db_cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name=?;", (table_name,)); return self.db_cursor.fetchone() is not None
        except sqlite3.Error as e: logger.error(f"Error checking if table '{table_name}' exists: {e}"); return False

    def get_index_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.index")
    def get_embedding_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.npy")
    def get_chunks_path(self, file_name): return os.path.join(self.config.index_dir, f"{self._get_safe_filename(file_name)}.json")


    # --- Content Extractors (PDF, PPTX - Unchanged) ---
    # _extract_pdf, _extract_pptx methods remain the same as in ragtest_c2-9.py

    # MODIFIED _extract_pdf for Structure Awareness
    def _extract_pdf(self, file_path):
        all_content = []
        base_filename = os.path.basename(file_path)
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Extracting text/tables from PDF: {base_filename} ({total_pages} pages) - Applying structural heuristics.")
                page_bar = None
                if self._is_streamlit and total_pages > 5: 
                    import streamlit as st
                    page_bar = st.progress(0, text=f"Extracting pages from {base_filename}...")

                toc_pattern = re.compile(r"(\.|\s){4,}\s*\d+\s*$") 
                intro_keywords = re.compile(r'\b(introduction|about sat ?sure|company profile|executive summary|overview|problem statement|client need|challenge)\b', re.IGNORECASE)
                closing_keywords = re.compile(r'\b(appendix|annexure|pricing|cost|timeline|schedule|next steps|conclusion)\b', re.IGNORECASE)

                for page_num, page in enumerate(pdf.pages):
                    current_page = page_num + 1
                    text = page.extract_text(layout="normal") or "" 
                    cleaned_text = self.clean_text(text) 
                    page_word_count = len(cleaned_text.split())

                    potential_section = "main_content" 
                    is_first_page = (page_num == 0)
                    is_likely_toc_page = (page_num == 1 or page_num == 2) 
                    is_near_end = (page_num >= total_pages - 3) 

                    if is_first_page:
                        potential_section = "title_page"
                        potential_title = ""
                        lines = [line.strip() for line in text.split('\n') if line.strip()]
                        if lines:
                            first_few_lines = lines[:7] 
                            if first_few_lines:
                                potential_title = max(first_few_lines, key=len)
                                if len(potential_title) > 80 or potential_title.islower():
                                    potential_title = first_few_lines[0] 
                        logger.debug(f"PDF Pg {current_page}: Marked as 'title_page'. Potential Title: '{potential_title[:50]}...'")

                    elif is_likely_toc_page:
                        lines = text.split('\n')
                        toc_lines_count = sum(1 for line in lines if toc_pattern.search(line.strip()))
                        if toc_lines_count >= 3 and page_word_count < 300 : 
                            potential_section = "table_of_contents"
                            logger.debug(f"PDF Pg {current_page}: Marked as 'table_of_contents' (Found {toc_lines_count} matching lines).")

                    if potential_section == "main_content" and page_num < 5:
                        if intro_keywords.search(cleaned_text):
                            potential_section = "introduction_overview"
                            logger.debug(f"PDF Pg {current_page}: Marked as 'introduction_overview' based on keywords.")
                        elif page_num < 3: 
                             potential_section = "front_matter"

                    if potential_section == "main_content" and is_near_end:
                         if closing_keywords.search(cleaned_text):
                             potential_section = "closing_appendix"
                             logger.debug(f"PDF Pg {current_page}: Marked as 'closing_appendix' based on keywords.")
                         elif page_num == total_pages -1: 
                             potential_section = "final_page"

                    source_info = {"page": current_page}
                    if is_first_page and potential_title:
                         source_info["potential_doc_title"] = potential_title

                    if cleaned_text:
                        all_content.append({
                            "type": "text", "content": cleaned_text, 
                            "source_info": source_info, "file_type": "pdf",
                            "potential_section": potential_section 
                        })

                    try:
                        for table_idx, table_content in enumerate(page.extract_tables()): 
                             if table_content:
                                 table_df = pd.DataFrame(table_content)
                                 if not table_df.empty and not pd.api.types.is_numeric_dtype(table_df.iloc[0].dropna()):
                                     try:
                                         table_df.columns = table_df.iloc[0].fillna('').astype(str)
                                         table_df = table_df[1:]
                                     except Exception: pass 
                                 table_df = table_df.fillna('')
                                 table_string = table_df.to_string(index=False, header=True)
                                 cleaned_table_string = self.clean_text(table_string)
                                 if cleaned_table_string:
                                     all_content.append({
                                        "type": "table", "content": cleaned_table_string,
                                        "source_info": {**source_info, "table_index_on_page": table_idx + 1},
                                        "file_type": "pdf", "potential_section": potential_section
                                     })
                    except Exception as e_table:
                        logger.warning(f"Could not extract/process table on page {current_page} in {base_filename}: {e_table}")

                    if page_bar:
                        progress_percent = min(1.0, (page_num + 1) / total_pages)
                        page_bar.progress(progress_percent, text=f"Extracting page {current_page}/{total_pages} from {base_filename}...")
                if page_bar: page_bar.empty()
            logger.info(f"Extracted {len(all_content)} text/table content blocks from PDF: {base_filename} with structural tagging.")
            return all_content
        except Exception as e:
            logger.error(f"Error extracting content from PDF {base_filename}: {e}", exc_info=True)
            return []

    def _extract_pptx(self, file_path):
        all_content = []; base_filename = os.path.basename(file_path); max_chars = self.config.max_chars_per_element; merge_threshold_words = self.config.pptx_merge_threshold_words
        try:
            prs = Presentation(file_path); logger.info(f"Extracting text from PPTX: {base_filename} ({len(prs.slides)} slides)")
            slide_bar = None; total_slides = len(prs.slides)
            if self._is_streamlit and total_slides > 1: import streamlit as st; slide_bar = st.progress(0, text=f"Extracting slides from {base_filename}...")
            pending_title_slide_data = None
            for i, slide in enumerate(prs.slides):
                 current_slide_number = i + 1; current_slide_title_text = ""; current_slide_other_texts = []; current_slide_has_title_placeholder = False; title_shape = None
                 try:
                    if slide.shapes.title: title_shape = slide.shapes.title; current_slide_has_title_placeholder = True
                 except AttributeError: pass 
                 if title_shape and title_shape.has_text_frame: cleaned_title = self.clean_text(title_shape.text_frame.text); current_slide_title_text = cleaned_title if cleaned_title else ""
                 for shape in slide.shapes:
                     if shape == title_shape: continue 
                     is_placeholder = shape.is_placeholder; is_body_placeholder = False
                     if is_placeholder:
                         try: ph_type = shape.placeholder_format.type; is_body_placeholder = ph_type in [MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT, MSO_SHAPE_TYPE.SUBTITLE, MSO_SHAPE_TYPE.CONTENT, MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.PICTURE] 
                         except AttributeError: pass
                     if shape.has_text_frame:
                         text = shape.text_frame.text; cleaned = self.clean_text(text)
                         if cleaned:
                              if len(cleaned) > max_chars: cleaned = cleaned[:max_chars] + "...(truncated shape)"
                              prefix = "[Body]: " if is_body_placeholder and not current_slide_title_text and not current_slide_other_texts else "" 
                              current_slide_other_texts.append(prefix + cleaned)
                 if slide.has_notes_slide:
                    try:
                         notes_text = slide.notes_slide.notes_text_frame.text; cleaned_notes = self.clean_text(notes_text)
                         if cleaned_notes:
                             if len(cleaned_notes) > max_chars * 2: cleaned_notes = cleaned_notes[:max_chars*2] + "...(truncated notes)"
                             current_slide_other_texts.append(f"[Notes]: {cleaned_notes}")
                    except Exception as e_notes: logger.warning(f"Could not extract notes from slide {current_slide_number}: {e_notes}")

                 current_slide_full_content = current_slide_title_text
                 if current_slide_other_texts: current_slide_full_content += ("\n" if current_slide_full_content else "") + "\n".join(current_slide_other_texts)
                 current_slide_full_content = current_slide_full_content.strip(); other_text_word_count = sum(len(s.split()) for s in current_slide_other_texts) 

                 merged_content_block = None; should_merge = False
                 if pending_title_slide_data:
                     if pending_title_slide_data['has_title'] and pending_title_slide_data['other_words'] <= merge_threshold_words and current_slide_full_content:
                         should_merge = True; logger.info(f"Merging slide {pending_title_slide_data['number']} (title-like) with content from slide {current_slide_number}.")

                 if should_merge:
                    merged_text = f"[Title from Slide {pending_title_slide_data['number']}]: {pending_title_slide_data['title']}\n\n[Content from Slide {pending_title_slide_data['number']} (if any)]:\n{pending_title_slide_data['content_without_title']}\n\n---\n\n[Content from Slide {current_slide_number}]:\n{current_slide_full_content}"
                    merged_content_block = {"type": "slide_text_merged", "content": merged_text.strip(), "source_info": {"slide_title": pending_title_slide_data['number'], "slide_content": current_slide_number}, "file_type": "pptx"}
                    all_content.append(merged_content_block); pending_title_slide_data = None 
                 else:
                    if pending_title_slide_data:
                         if pending_title_slide_data['content']: 
                             all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})
                         pending_title_slide_data = None 

                    if current_slide_has_title_placeholder and current_slide_title_text and other_text_word_count <= merge_threshold_words and current_slide_full_content:
                        logger.debug(f"Slide {current_slide_number} ('{current_slide_title_text[:30]}...') is potential title slide for next content. Holding.");
                        content_without_title_parts = [txt for txt in current_slide_other_texts if not txt.startswith(current_slide_title_text)] 
                        content_w_o_title_str = "\n".join(content_without_title_parts).strip()
                        pending_title_slide_data = {
                            "content": current_slide_full_content, "content_without_title": content_w_o_title_str,
                            "number": current_slide_number, "has_title": current_slide_has_title_placeholder,
                            "other_words": other_text_word_count, "title": current_slide_title_text
                        }
                    else: 
                        if current_slide_full_content:
                            all_content.append({"type": "slide_text", "content": current_slide_full_content, "source_info": {"slide": current_slide_number}, "file_type": "pptx"})

                 if slide_bar: slide_bar.progress(min(1.0, (i + 1) / total_slides), text=f"Extracting slide {current_slide_number}/{total_slides}...")
            if pending_title_slide_data:
                 logger.debug(f"Processing final pending title slide {pending_title_slide_data['number']} at end.")
                 if pending_title_slide_data['content']: all_content.append({"type": "slide_text", "content": pending_title_slide_data['content'], "source_info": {"slide": pending_title_slide_data['number']}, "file_type": "pptx"})

            if slide_bar: slide_bar.empty()
            logger.info(f"Extracted {len(all_content)} content blocks from PPTX {base_filename} (Merge strategy applied)."); return all_content
        except Exception as e: logger.error(f"Error extracting content from PPTX {base_filename}: {e}", exc_info=True); return []


    # --- XLSX/CSV Extraction and Loading (Metadata generation now uses SambaNova) ---
    def _extract_and_load_xlsx(self, file_path, file_name, progress_callback=None):
        base_filename = os.path.basename(file_path)
        sheet_load_success_count = 0; sheet_metadata_success_count = 0; sheet_sql_load_success_count = 0; processed_sheet_tables = []
        try:
            excel_file = pd.ExcelFile(file_path); sheet_names = excel_file.sheet_names
            logger.info(f"Processing XLSX: {base_filename} (Sheets: {len(sheet_names)})")
            if not sheet_names: logger.warning(f"No sheets found in {base_filename}"); return False, []
            sheet_bar = None; total_sheets = len(sheet_names)
            if self._is_streamlit and total_sheets > 1: import streamlit as st; sheet_bar = st.progress(0, text=f"Processing sheets in {base_filename}...")
            for i, sheet_name in enumerate(sheet_names):
                stage_msg = f"Processing sheet '{sheet_name}' ({i+1}/{total_sheets}) in {base_filename}"
                logger.info(stage_msg)
                if progress_callback: progress_callback(f"{stage_msg} - Loading...", stage="Loading Sheet")
                if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.1) / total_sheets), text=f"Loading '{sheet_name}'...")
                try:
                    df = excel_file.parse(sheet_name); df = df.fillna('')
                    df.columns = [re.sub(r'[^a-zA-Z0-9_]', '_', str(col)).strip('_') for col in df.columns]
                    df.columns = [f"col_{j}" if not name else name for j, name in enumerate(df.columns)]
                    cols = pd.Series(df.columns)
                    for dup_idx, dup_val in cols[cols.duplicated()].items():
                         indices = cols[cols == dup_val].index.tolist()
                         for k_idx, col_idx in enumerate(indices):
                             if k_idx > 0: cols[col_idx] = f"{dup_val}_{k_idx}"
                    df.columns = cols

                    if df.empty: logger.warning(f"Sheet '{sheet_name}' in {base_filename} is empty. Skipping."); continue
                    sheet_load_success_count += 1; logger.info(f"Loaded DataFrame from sheet '{sheet_name}' (Rows: {len(df)}, Cols: {len(df.columns)})")
                    table_name = self._get_safe_table_name(file_name, sheet_name)
                    if not table_name: logger.error(f"Could not generate table name for sheet '{sheet_name}'. Skipping SQL load."); continue
                    
                    if progress_callback: progress_callback(f"{stage_msg} - Analyzing Columns (SambaNova)...", stage="Analyzing Columns")
                    if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.4) / total_sheets), text=f"Analyzing '{sheet_name}' (SambaNova)...")
                    column_metadata = self._get_column_metadata_sambanova(df, file_name, sheet_name) # MODIFIED CALL
                    
                    if not column_metadata: 
                        logger.warning(f"Failed to get column metadata for sheet '{sheet_name}' via SambaNova. Proceeding without descriptions.")
                        column_metadata = [{"name": col, "type": self._infer_sql_type(df[col]), "description": "N/A - Metadata generation failed."} for col in df.columns]
                    else: 
                        sheet_metadata_success_count +=1
                    
                    if progress_callback: progress_callback(f"{stage_msg} - Loading to SQL Database...", stage="Loading SQL")
                    if sheet_bar: sheet_bar.progress(min(1.0, (i + 0.8) / total_sheets), text=f"Storing '{sheet_name}'...")
                    sql_loaded = self._load_df_to_sql(df, table_name, column_metadata, file_name, sheet_name)
                    if sql_loaded: sheet_sql_load_success_count += 1; processed_sheet_tables.append(table_name); logger.info(f"Successfully loaded sheet '{sheet_name}' into SQL table '{table_name}'")
                    else: logger.error(f"Failed to load sheet '{sheet_name}' into SQL.")
                except Exception as e_sheet:
                    logger.error(f"Error processing sheet '{sheet_name}' in {base_filename}: {e_sheet}", exc_info=True);
                    if progress_callback: progress_callback(f"⚠️ Error processing sheet '{sheet_name}': {e_sheet}", is_warning=True)
                if sheet_bar: sheet_bar.progress(min(1.0, (i + 1) / total_sheets), text=f"Finished sheet {i+1}/{total_sheets}...")
            if sheet_bar: sheet_bar.empty()
            logger.info(f"Finished processing XLSX {base_filename}. Sheets Loaded: {sheet_load_success_count}/{total_sheets}, Metadata OK: {sheet_metadata_success_count}, SQL OK: {sheet_sql_load_success_count}")
            return sheet_sql_load_success_count > 0, processed_sheet_tables
        except Exception as e:
            logger.error(f"Error opening or processing XLSX file {base_filename}: {e}", exc_info=True);
            if progress_callback: progress_callback(f"❌ Error processing XLSX {base_filename}: {e}", is_error=True)
            return False, []

    def _extract_and_load_csv(self, file_path, file_name, progress_callback=None):
        base_filename = os.path.basename(file_path); df = None; load_success = False; metadata_success = False; sql_load_success = False; table_name = None
        try:
            stage_msg = f"Processing CSV: {base_filename}"; logger.info(stage_msg)
            if progress_callback: progress_callback(f"{stage_msg} - Loading...", stage="Loading CSV")
            encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
            for enc in encodings_to_try:
                try:
                    df = pd.read_csv(file_path, encoding=enc, low_memory=False, on_bad_lines='warn'); df = df.fillna('')
                    df.columns = [re.sub(r'[^a-zA-Z0-9_]', '_', str(col)).strip('_') for col in df.columns]
                    df.columns = [f"col_{j}" if not name else name for j, name in enumerate(df.columns)]
                    cols = pd.Series(df.columns)
                    for dup_idx, dup_val in cols[cols.duplicated()].items():
                         indices = cols[cols == dup_val].index.tolist()
                         for k_idx, col_idx in enumerate(indices):
                             if k_idx > 0: cols[col_idx] = f"{dup_val}_{k_idx}"
                    df.columns = cols
                    logger.info(f"Read CSV {base_filename} using encoding: {enc}. (Rows: {len(df)}, Cols: {len(df.columns)})"); load_success = True; break
                except UnicodeDecodeError: continue
                except Exception as e_read: logger.warning(f"Pandas read_csv error for {base_filename} with encoding {enc}: {e_read}")
            
            if df is None: logger.error(f"Could not read CSV {base_filename} after trying encodings."); return False, []
            if df.empty: logger.warning(f"CSV {base_filename} is empty."); return False, []
            if not load_success: return False, []
            
            table_name = self._get_safe_table_name(file_name)
            if not table_name: logger.error("Could not generate table name for CSV. Skipping SQL load."); return False, []
            
            if progress_callback: progress_callback(f"{stage_msg} - Analyzing Columns (SambaNova)...", stage="Analyzing Columns")
            column_metadata = self._get_column_metadata_sambanova(df, file_name) # MODIFIED CALL
            
            if not column_metadata: 
                logger.warning("Failed to get column metadata for CSV via SambaNova. Proceeding without descriptions.")
                column_metadata = [{"name": col, "type": self._infer_sql_type(df[col]), "description": "N/A - Metadata generation failed."} for col in df.columns]
            else: 
                metadata_success = True
            
            if progress_callback: progress_callback(f"{stage_msg} - Loading to SQL Database...", stage="Loading SQL")
            sql_loaded = self._load_df_to_sql(df, table_name, column_metadata, file_name)
            if sql_loaded: sql_load_success = True; logger.info(f"Successfully loaded CSV {base_filename} into SQL table '{table_name}'")
            else: logger.error(f"Failed to load CSV {base_filename} into SQL.")
            return sql_load_success, ([table_name] if sql_load_success else [])
        except Exception as e:
            logger.error(f"Error processing CSV file {base_filename}: {e}", exc_info=True);
            if progress_callback: progress_callback(f"❌ Error processing CSV {base_filename}: {e}", is_error=True)
            return False, []

    # --- SQL Helper Methods (_infer_sql_type unchanged, _load_df_to_sql unchanged) ---
    def _infer_sql_type(self, series):
        num_numeric = pd.to_numeric(series, errors='coerce').notna().sum()
        num_total = len(series)
        if num_total > 0 and num_numeric / num_total > 0.9:
            numeric_series = pd.to_numeric(series, errors='coerce').dropna()
            if not numeric_series.empty and (numeric_series == numeric_series.astype(np.int64)).all():
                 return "INTEGER"
            return "REAL"
        if pd.api.types.is_integer_dtype(series.dtype): return "INTEGER"
        if pd.api.types.is_float_dtype(series.dtype): return "REAL"
        if pd.api.types.is_bool_dtype(series.dtype): return "INTEGER"
        if pd.api.types.is_datetime64_any_dtype(series.dtype): return "TEXT"
        return "TEXT"

    def _get_column_metadata_sambanova(self, df, file_name, sheet_name=None):
        """Generates column metadata (descriptions and SQL types) using SambaNova."""
        if not self.sambanova_api_key_val:
            logger.warning("SambaNova API key not available. Cannot generate column metadata.")
            return None

        client = self._get_llm_client()
        sample_data_str = df.head(3).to_string(index=False) # Smaller sample for token limits
        column_names = list(df.columns)

        prompt_content = f"""Analyze the following table columns from a spreadsheet ('{file_name}'{f", sheet '{sheet_name}'" if sheet_name else ""}).
The column names are: {column_names}
Here are the first few rows of data:
{sample_data_str}

For each column name, provide:
1. A concise one-sentence description of what the data in the column represents.
2. The most appropriate SQLite-compatible SQL data type (choose from: TEXT, INTEGER, REAL - use TEXT for dates/datetimes or mixed types).

Respond ONLY with a valid JSON object where keys are the exact column names and values are objects containing 'description' and 'sql_type'. Ensure the JSON is well-formed and the response contains ONLY the JSON object itself. Example:
{{
  "ColumnA": {{ "description": "Unique identifier for each record.", "sql_type": "INTEGER" }},
  "ColumnB": {{ "description": "Name of the customer.", "sql_type": "TEXT" }}
}}
"""
        messages = [
            {"role": "system", "content": "You are an expert data analyst specializing in understanding spreadsheet columns and assigning SQLite-compatible SQL types. You respond only with valid JSON."},
            {"role": "user", "content": prompt_content}
        ]

        try:
            logger.info(f"Requesting column metadata from SambaNova for {file_name}{f'/{sheet_name}' if sheet_name else ''}...")
            response = client.chat.completions.create(
                model=self.config.sambanova_model,
                messages=messages,
                temperature=0.0, # Low temp for factual task
                top_p=self.config.top_p,
                max_tokens=self.config.max_output_tokens_sambanova # Adjust as needed
            )
            raw_content = response.choices[0].message.content
            logger.debug(f"Raw SambaNova response for metadata: {raw_content}")

            json_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw_content, re.DOTALL | re.IGNORECASE)
            if not json_match: json_match = re.search(r"(\{.*?\})", raw_content, re.DOTALL)

            if json_match:
                json_str = json_match.group(1)
                metadata_dict = json.loads(json_str)
                
                formatted_metadata = []
                processed_cols = set()
                for expected_col_name in column_names:
                    found_meta = None
                    # Case-insensitive matching for keys from LLM response
                    for resp_col_name_key in list(metadata_dict.keys()):
                        if resp_col_name_key.lower() == expected_col_name.lower():
                            found_meta = metadata_dict.pop(resp_col_name_key) # Remove to avoid re-processing
                            break
                    
                    if found_meta and isinstance(found_meta, dict) and 'description' in found_meta and 'sql_type' in found_meta:
                        sql_type = str(found_meta['sql_type']).upper()
                        if sql_type not in ["TEXT", "INTEGER", "REAL"]:
                            logger.warning(f"SambaNova proposed invalid SQL type '{sql_type}' for column '{expected_col_name}'. Defaulting to TEXT.")
                            sql_type = "TEXT"
                        formatted_metadata.append({"name": expected_col_name, "description": str(found_meta['description']).strip(), "type": sql_type})
                        processed_cols.add(expected_col_name)
                    else:
                        logger.warning(f"Metadata missing or invalid format for column '{expected_col_name}' in SambaNova response. Adding basic entry.")
                
                # Add fallbacks for any columns missed by the LLM
                missing_cols = set(column_names) - processed_cols
                if missing_cols:
                     logger.warning(f"SambaNova metadata generation missed columns: {missing_cols}. Adding basic entries.")
                     final_meta_ordered = []
                     current_meta_map = {m['name']: m for m in formatted_metadata}
                     for col_name_in_order in column_names:
                         if col_name_in_order in current_meta_map:
                             final_meta_ordered.append(current_meta_map[col_name_in_order])
                         else:
                             final_meta_ordered.append({"name": col_name_in_order, "description": "N/A - Metadata generation by LLM failed or incomplete.", "type": self._infer_sql_type(df[col_name_in_order])})
                     formatted_metadata = final_meta_ordered

                logger.info(f"Successfully generated/handled metadata for {len(formatted_metadata)} columns using SambaNova for '{file_name}'{f'/{sheet_name}' if sheet_name else ''}.")
                return formatted_metadata
            else:
                logger.error(f"Could not find JSON object in SambaNova response for metadata: {raw_content}")
                return None
        except openai.APIError as e_api:
            logger.error(f"SambaNova API error during metadata generation: {e_api}", exc_info=True)
            return None
        except json.JSONDecodeError as e_json:
            logger.error(f"Failed to parse JSON metadata from SambaNova: {e_json}. Raw response: {raw_content}")
            return None
        except Exception as e:
            logger.error(f"Error getting column metadata from SambaNova: {e}", exc_info=True)
            return None

    def _load_df_to_sql(self, df, table_name, column_metadata, file_name, sheet_name=None):
        if not self.db_conn or not self.db_cursor: logger.error("Database connection not available. Cannot load table."); return False
        if not table_name: logger.error("Invalid table name provided. Cannot load table."); return False
        if self._table_exists(table_name):
             logger.warning(f"Table '{table_name}' already exists. Replacing.")
             try: self.db_cursor.execute(f'DROP TABLE "{table_name}"')
             except sqlite3.Error as e_drop: logger.error(f"Error dropping existing table {table_name}: {e_drop}"); return False
        try:
            logger.info(f"Loading DataFrame into SQL table '{table_name}'...")
            df.columns = [meta['name'] for meta in column_metadata] # Assume metadata names are final
            temp_df = df.copy()
            for meta in column_metadata:
                col, sql_type = meta['name'], meta['type']
                if sql_type in ["INTEGER", "REAL"]:
                    temp_df[col] = pd.to_numeric(temp_df[col], errors='coerce')
                    if sql_type == "INTEGER":
                         try: temp_df[col] = temp_df[col].astype('Int64')
                         except Exception: pass 
                elif sql_type == "TEXT":
                     temp_df[col] = temp_df[col].astype(str)
            
            temp_df.to_sql(table_name, self.db_conn, if_exists='replace', index=False); self.db_conn.commit()
            logger.info(f"DataFrame loaded into table '{table_name}' ({len(df)} rows).")
            self.table_metadata[table_name] = {"columns": column_metadata, "source_file": file_name, "source_sheet": sheet_name, "row_count": len(df)}
            if file_name not in self.file_to_table_map: self.file_to_table_map[file_name] = []
            if table_name not in self.file_to_table_map[file_name]: self.file_to_table_map[file_name].append(table_name)
            return True
        except Exception as e:
            logger.error(f"Error loading DataFrame to table '{table_name}': {e}", exc_info=True);
            try: self.db_conn.rollback()
            except: pass
            return False

    # --- Text Chunking, FAISS Indexing, Document Processing (Largely Unchanged) ---
    # extract_content, chunk_content, load_faiss_index, save_chunks, load_chunks, process_files
    # These methods are mostly independent of the LLM choice, so they remain similar to ragtest_c2-9.py
    # Small logging changes in process_files might be needed if progress_callback uses LLM-specific terms.
    # For now, assume process_files structure is fine.
    def extract_content(self, file_path):
        extension = os.path.splitext(file_path)[1].lower()
        if extension == '.pdf': return self._extract_pdf(file_path) 
        elif extension == '.pptx': return self._extract_pptx(file_path)
        elif extension in ['.csv', '.xlsx']: logger.debug(f"Skipping text extraction for {os.path.basename(file_path)}, handled for SQL."); return []
        else: logger.warning(f"Unsupported file type for text extraction: {file_path}"); return []

    def chunk_content(self, all_content):
        chunks = []; 
        if not all_content: return chunks
        # ... (rest of chunk_content is identical to ragtest_c2-9.py) ...
        for item_index, item in enumerate(all_content):
            content = item.get('content', '');
            source_info = item.get('source_info', {}); 
            file_type = item.get('file_type', 'unknown');
            content_type = item.get('type', 'unknown')
            potential_section = item.get('potential_section', 'unknown') 
            if not isinstance(content, str): content = str(content)
            words = content.split();
            if not words: continue
            start_index = 0
            while start_index < len(words):
                end_index = start_index + self.config.chunk_size
                chunk_text = " ".join(words[start_index:end_index])
                if chunk_text:
                     chunks.append({
                         "content": chunk_text, "source_info": source_info, 
                         "file_type": file_type, "type": content_type,
                         "potential_section": potential_section
                     })
                start_index += (self.config.chunk_size - self.config.overlap)
                if start_index >= len(words): break
        logger.info(f"Created {len(chunks)} text chunks from {len(all_content)} content blocks.");
        return chunks

    def load_faiss_index(self, file_name, embedding_dim):
        index_path = self.get_index_path(file_name)
        if os.path.exists(index_path):
            try: index = faiss.read_index(index_path); logger.info(f"Loaded FAISS index for {file_name} ({index.ntotal} vectors)"); return index
            except Exception as e: logger.error(f"Error reading FAISS index {index_path}: {e}. Creating new index.")
        return faiss.IndexFlatL2(embedding_dim)

    def save_chunks(self, file_name, chunks):
        chunks_path = self.get_chunks_path(file_name)
        try:
            with open(chunks_path, 'w', encoding='utf-8') as f: json.dump(chunks, f, indent=2); logger.info(f"Saved {len(chunks)} text chunks to {chunks_path}")
        except Exception as e: logger.error(f"Error saving text chunks for {file_name} to {chunks_path}: {e}");

    def load_chunks(self, file_name):
        chunks_path = self.get_chunks_path(file_name)
        if os.path.exists(chunks_path):
            try:
                with open(chunks_path, 'r', encoding='utf-8') as f: chunks = json.load(f); logger.info(f"Loaded {len(chunks)} text chunks from {chunks_path}"); return chunks
            except Exception as e: logger.error(f"Error loading/decoding text chunks from {chunks_path}: {e}. Will re-process."); return None
        return None

    def process_files(self, progress_callback=None):
        # ... (process_files logic is largely the same, but calls to metadata generation are updated)
        # This method orchestrates the calls to _extract_and_load_xlsx/_csv which now use _get_column_metadata_sambanova
        # and text extraction/indexing parts remain similar.
        self.file_chunks = {}; self.faiss_indexes = {}; self.processed_files = set()
        self.table_metadata = {}; self.file_to_table_map = {}
        if self.config.sqlite_db_path == ":memory:" and self.db_cursor:
            logger.info("Clearing existing tables from in-memory database...")
            try:
                self.db_cursor.execute("SELECT name FROM sqlite_master WHERE type='table';"); tables = self.db_cursor.fetchall()
                for table_tuple in tables: self.db_cursor.execute(f'DROP TABLE IF EXISTS "{table_tuple[0]}"')
                self.db_conn.commit(); logger.info(f"Dropped {len(tables)} tables.")
            except sqlite3.Error as e_drop: logger.error(f"Error dropping tables: {e_drop}")

        data_dir = self.config.data_dir; supported_ext = self.config.supported_extensions
        try: all_files = os.listdir(data_dir); process_list = sorted([f for f in all_files if f.lower().endswith(supported_ext)])
        except Exception as e: logger.error(f"Error listing files in data directory {data_dir}: {e}", exc_info=True); return False
        if not process_list: logger.warning(f"No supported files found in {data_dir}"); return True

        logger.info(f"Processing {len(process_list)} supported files from {data_dir}");
        embedding_dim = self.encoder_model.get_sentence_embedding_dimension();
        total_files = len(process_list)
        if progress_callback: progress_callback(f"Found {total_files} file(s). Starting processing...", current_step=0, total_steps=total_files)

        any_success = False
        for idx, file_name in enumerate(process_list):
            current_file_msg = f"Processing ({idx+1}/{total_files}): {file_name}"; logger.info(f"--- {current_file_msg} ---")
            if progress_callback: progress_callback(current_file_msg, current_step=idx, total_steps=total_files)
            file_path = os.path.join(data_dir, file_name); index_path = self.get_index_path(file_name); emb_path = self.get_embedding_path(file_name); chunks_path = self.get_chunks_path(file_name)
            file_extension = os.path.splitext(file_name)[1].lower()

            if file_extension in ['.csv', '.xlsx']:
                sql_success, _ = self._extract_and_load_xlsx(file_path, file_name, progress_callback) if file_extension == '.xlsx' else self._extract_and_load_csv(file_path, file_name, progress_callback)
                if sql_success: any_success = True; self.processed_files.add(file_name)
                else: logger.error(f"Failed SQL processing for {file_name}")
            
            elif file_extension in ['.pdf', '.pptx']:
                try:
                    chunks = self.load_chunks(file_name)
                    if chunks is None:
                        all_content = self.extract_content(file_path);
                        if not all_content: logger.warning(f"No text content from {file_name}."); continue
                        chunks = self.chunk_content(all_content);
                        if not chunks: logger.warning(f"No chunks for {file_name}."); continue
                        self.save_chunks(file_name, chunks)
                    
                    if chunks:
                        self.file_chunks[file_name] = chunks
                        faiss_index = None; regenerate_embeddings = True # Simplified logic to always regen if index not perfect
                        if os.path.exists(index_path) and os.path.exists(emb_path) and os.path.exists(chunks_path):
                            try:
                                embeddings_on_disk = np.load(emb_path)
                                if embeddings_on_disk.ndim == 2 and embeddings_on_disk.shape[1] == embedding_dim and embeddings_on_disk.shape[0] == len(chunks):
                                    loaded_index = self.load_faiss_index(file_name, embedding_dim)
                                    if loaded_index.ntotal == len(chunks):
                                        faiss_index = loaded_index; regenerate_embeddings = False
                                        logger.info(f"Verified existing text index for {file_name}")
                                    else: logger.warning(f"Index size mismatch for {file_name}. Regenerating.")
                                else: logger.warning(f"Embedding/chunk mismatch for {file_name}. Regenerating.")
                            except Exception as e_verify: logger.error(f"Error verifying index for {file_name}: {e_verify}. Regenerating.")
                        
                        if regenerate_embeddings:
                            content_list = [chunk['content'] for chunk in chunks];
                            if content_list:
                                embeddings = self.encoder_model.encode(content_list, batch_size=32, show_progress_bar=False, convert_to_numpy=True).astype('float32')
                                if embeddings.shape[0] > 0:
                                    np.save(emb_path, embeddings)
                                    faiss_index = faiss.IndexFlatL2(embedding_dim); faiss_index.add(embeddings); faiss.write_index(faiss_index, index_path)
                                    logger.info(f"Generated and saved text index for {file_name} ({faiss_index.ntotal} vectors)")
                                else: logger.warning(f"Embedding yielded no vectors for {file_name}.")
                            else: logger.warning(f"No content to embed for {file_name}.")

                        if faiss_index and faiss_index.ntotal > 0:
                             self.faiss_indexes[file_name] = faiss_index; self.processed_files.add(file_name); any_success = True
                        else: logger.warning(f"No valid text index for {file_name}.")
                except Exception as e_text:
                    logger.error(f"Failed text processing for {file_name}: {e_text}", exc_info=True)
        
        logger.info(f"--- Processing Complete. Processed {len(self.processed_files)} files. ---")
        return True


    # --- Querying Logic (Faiss search, Aggregate context - Unchanged) ---
    # query_files, aggregate_context methods remain the same as in ragtest_c2-9.py
    def query_files(self, query):
        if not self.faiss_indexes: logger.warning("No text indexes available for Faiss querying."); return {}
        query_results = {};
        try:
            query_embedding = self.encoder_model.encode(query, convert_to_numpy=True).astype("float32"); query_embedding = np.array([query_embedding])
            if query_embedding.ndim != 2: raise ValueError("Query embedding shape error")
            for file_name, index in self.faiss_indexes.items():
                if index is None or index.ntotal == 0: continue
                try:
                    k_search = min(self.config.k_retrieval, index.ntotal);
                    D, I = index.search(query_embedding, k=k_search); 
                    indices, distances = I[0], D[0]
                    current_file_chunks = self.file_chunks.get(file_name)
                    if not current_file_chunks: logger.error(f"Text chunks missing for {file_name}!"); continue
                    file_results = []
                    for i, idx in enumerate(indices):
                        if idx != -1 and 0 <= idx < len(current_file_chunks):
                            chunk = current_file_chunks[idx]
                            file_results.append({
                                "source_info": chunk.get('source_info', {}), "file_type": chunk.get('file_type', 'unknown'),
                                "content": chunk.get('content', ''), "score": round(float(distances[i]), 4),
                                "type": chunk.get('type', 'unknown'), "potential_section": chunk.get('potential_section', 'unknown')
                            })
                    if file_results:
                        query_results.setdefault(file_name, []).extend(sorted(file_results, key=lambda x: x['score']))
                except Exception as e_search: logger.error(f"Error searching text index {file_name}: {e_search}", exc_info=True)
            return query_results
        except Exception as e_query: logger.error(f"Error during Faiss text query: {e_query}", exc_info=True); return {}

    def aggregate_context(self, query_results, strategy="top_k"):
        # ... (aggregate_context is identical to ragtest_c2-9.py) ...
        all_context = {}; max_chars = self.config.max_context_tokens * 3 
        if not query_results: return all_context
        flat_results = []
        for file_name, results in query_results.items():
            for res in results: flat_results.append({**res, "file_name": file_name})
        flat_results.sort(key=lambda x: x['score']) 
        aggregated_context_str = ""; total_aggregated_chars = 0; added_chunks_count = 0; context_sources = set()
        limit_k = self.config.k_retrieval 
        for res in flat_results:
            if added_chunks_count >= limit_k: break
            source_info = res.get('source_info', {}); file_name = res['file_name']; file_type = res.get('file_type', 'unknown');
            content_body = res['content']; potential_section = res.get('potential_section', 'N/A') 
            source_parts = [f"Source: {file_name}"]
            if file_type == 'pptx':
                 slide_info_parts = []
                 if 'slide' in source_info: slide_info_parts.append(f"Slide: {source_info['slide']}")
                 if 'slide_title' in source_info : slide_info_parts.append(f"Title Slide: {source_info['slide_title']}")
                 if 'slide_content' in source_info : slide_info_parts.append(f"Content Slide: {source_info['slide_content']}")
                 slide_combined_info = ", ".join(filter(None, slide_info_parts))
                 if slide_combined_info: source_parts.append(slide_combined_info)
            elif file_type == 'pdf':
                 page_info = f"Page: {source_info.get('page', 'N/A')}"
                 if 'table_index_on_page' in source_info: page_info += f", Table {source_info['table_index_on_page']}"
                 source_parts.append(page_info)
            source_parts.append(f"Section: {potential_section}")
            source_str = ", ".join(filter(None, source_parts));
            content_header = f"--- Context from {source_str} ---\n";
            content_to_add = content_header + content_body + "\n\n";
            content_chars = len(content_to_add)
            if total_aggregated_chars + content_chars <= max_chars:
                aggregated_context_str += content_to_add; total_aggregated_chars += content_chars;
                added_chunks_count += 1; context_sources.add(file_name)
            elif added_chunks_count == 0: # Truncate if first chunk is too large
                 remaining_chars = max_chars - total_aggregated_chars - len(content_header) - 20 
                 if remaining_chars > 50:
                     aggregated_context_str += content_header + content_body[:remaining_chars] + "\n[...TRUNCATED CONTEXT...]\n\n";
                     added_chunks_count += 1; context_sources.add(file_name); break
                 else: break 
            else: break
        if aggregated_context_str.strip():
            all_context = {"combined_context": aggregated_context_str.strip(), "source_files": sorted(list(context_sources))};
        return all_context

    # --- LLM Call Helpers (Adapted for SambaNova) ---
    def _call_llm_for_analysis(self, samba_messages, task_description):
        """Calls the SambaNova LLM for internal analysis tasks."""
        client = self._get_llm_client()
        if not client: return None
        try:
            logger.info(f"Calling SambaNova LLM ({self.config.sambanova_model}) for task: {task_description}...")
            response = client.chat.completions.create(
                model=self.config.sambanova_model,
                messages=samba_messages,
                temperature=0.0, # Low temp for analysis
                top_p=self.config.top_p,
                max_tokens=self.config.max_output_tokens_sambanova 
            )
            content = response.choices[0].message.content.strip()
            logger.info(f"SambaNova response for {task_description}: '{content[:150]}...'")
            return content
        except openai.APIError as e:
            logger.error(f"SambaNova API error during {task_description}: {e}", exc_info=True)
        except Exception as e:
            logger.error(f"LLM call failed during {task_description}: {e}", exc_info=True)
        return None

    def query_llm_with_history(self, query, context_data, chat_history, retry_count=0): # Simplified retry for now
        combined_context = context_data.get("combined_context", "")
        source_files = context_data.get("source_files", [])
        source_file_str = ", ".join(source_files) if source_files else "the provided text documents"

        if not combined_context:
             return "I don't really know but my next versions will be able to answer this for sure! (Reason: No relevant text context found)"

        client = self._get_llm_client()
        if not client: return "Error: LLM client not available."

        # Prepare messages for SambaNova
        samba_history = []
        for msg in chat_history: # chat_history is already in {"role": ..., "content": ...}
            if msg["role"] == "assistant":
                clean_content = re.sub(r"\n\n---\n\*Answer generated.*", "", msg["content"], flags=re.DOTALL).strip()
                samba_history.append({"role": "assistant", "content": clean_content})
            else:
                samba_history.append(msg)
        
        system_prompt_content = f"""You are the SatSure LLM, a helpful AI assistant answering questions based ONLY on the provided text context from SatSure project proposal documents and the ongoing conversation history. The proposals outline solutions for clients.
The context below comes from document(s): '{source_file_str}'. Context sections are marked with 'Source:', 'Page:', 'Table Index', 'Slide:', and 'Section:' (e.g., 'Section: title_page', 'Section: main_content').
Use ONLY information presented between '--- START CONTEXT ---' and '--- END CONTEXT ---' or previous messages.
This is sensitive data, there must be no hallucination. Do not use any prior knowledge.
If the answer cannot be found in the text or history, state EXACTLY: "I don't really know but my next versions will be able to answer this for sure!"
Respond conversationally, considering the flow of the chat. Do NOT perform calculations unless explicitly shown in the text."""

        user_prompt_content = f"""Retrieved Context from SatSure proposal(s):\n--- START CONTEXT ---\n{combined_context}\n--- END CONTEXT ---\n\nLatest User Question: {query}\n\nAnswer based ONLY on the provided text context and conversation history:"""

        messages_for_samba = [{"role": "system", "content": system_prompt_content}]
        messages_for_samba.extend(samba_history)
        messages_for_samba.append({"role": "user", "content": user_prompt_content})

        try:
            logger.info(f"Querying SambaNova model {self.config.sambanova_model} with text context and history...")
            response = client.chat.completions.create(
                model=self.config.sambanova_model,
                messages=messages_for_samba,
                temperature=self.config.temperature,
                top_p=self.config.top_p,
                max_tokens=self.config.max_output_tokens_sambanova
            )
            answer = response.choices[0].message.content.strip()
            logger.info(f"SambaNova response received. Length: {len(answer)}")
            return answer
        except openai.APIError as e_api:
            logger.error(f"SambaNova API error during text RAG query: {e_api}", exc_info=True)
            return f"Error: Failed to get answer from LLM (SambaNova API Error: {e_api})"
        except Exception as e_gen:
            logger.error(f"Error calling SambaNova API for text RAG: {e_gen}", exc_info=True)
            return f"Error: Could not query LLM (SambaNova) due to: {e_gen}"

    # --- Query Analysis Helpers (Adapted for SambaNova) ---
    def _classify_query(self, query): # Takes refined query
        system_prompt = """You are the SatSure LLM, an expert query analyzer. Your task is to classify the user query related to SatSure project proposal documents (PDFs) and potentially related data (XLSX/CSV). This is sensitive data so there must be no hallucination during retrieval. If you don't know an answer, simple answer with "I don't really know but my next versions will be able to answer this for sure!"

Classify the user query into ONE of the following categories:
1.  'Simple Retrieval': Asking for specific facts, definitions, or simple summaries directly extractable from text within the proposals.
2.  'Complex/Reasoning (Text)': Requires combining information from multiple text passages, summarization, comparison, or reasoning based *only* on text.
3.  'Structured Query (SQL)': Requires operations on structured data (tables), such as calculations, filtering, grouping, or sorting.

Respond ONLY in JSON format with "classification" (category string) and "confidence" (float 0.0-1.0).
Examples:
Query: "What is Project SkyWatch constellation?" -> {"classification": "Simple Retrieval", "confidence": 0.99}
Query: "Compare AquaMonitor and TerraScan methodologies." -> {"classification": "Complex/Reasoning (Text)", "confidence": 0.9}
Query: "Total cost for Phase 1 deliverables in spreadsheet?" -> {"classification": "Structured Query (SQL)", "confidence": 0.98}"""
        
        samba_messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Classify the following user query about SatSure proposals/data:\n\"{query}\""}
        ]
        raw_response = self._call_llm_for_analysis(samba_messages, "query classification")
        if not raw_response: return "Simple Retrieval", 0.5
        try:
            json_match = re.search(r"\{.*\}", raw_response, re.DOTALL)
            if json_match:
                result = json.loads(json_match.group(0))
                classification = result.get("classification"); confidence = float(result.get("confidence", 0.0))
                if classification in ["Simple Retrieval", "Complex/Reasoning (Text)", "Structured Query (SQL)"]:
                    logger.info(f"Query classified as '{classification}' with confidence {confidence:.2f}"); return classification, confidence
            logger.warning(f"Invalid classification JSON from SambaNova: {raw_response}")
        except Exception as e: logger.warning(f"Failed to parse classification from SambaNova: '{raw_response}'. Error: {e}.")
        return "Simple Retrieval", 0.5 # Default

    def _decompose_query(self, query):
        system_prompt = """You are an expert query decomposer. Break down a complex user query needing information potentially spanning multiple SatSure project proposal documents into simpler, factual sub-queries. Each sub-query should aim to retrieve a specific piece of information from the text content of the proposals. Do NOT decompose queries asking for calculations or filtering on tables. Format as a numbered list. If the query is simple or cannot be decomposed, return the original query prefixed with '1. '."""
        samba_messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Decompose the following complex query about SatSure proposals:\n\"{query}\""}
        ]
        decomposition = self._call_llm_for_analysis(samba_messages, "text query decomposition")
        if decomposition:
            sub_queries = [re.sub(r"^\s*[\*\-\d]+\.?\s*", "", line).strip() for line in decomposition.split('\n') if line.strip()]
            if sub_queries and not (len(sub_queries) == 1 and sub_queries[0].strip() == query.strip()):
                logger.info(f"Decomposed text query into: {sub_queries}"); return sub_queries
        logger.info("Decomposition didn't yield distinct sub-queries for text. Using original."); return [query]

    def _identify_target_sql_table(self, query):
        if not self.table_metadata: logger.warning("No table metadata to identify target SQL table."); return None
        schema_overview = "Available SQL Tables:\n"; table_options = []
        for table_name, meta in self.table_metadata.items():
            table_options.append(table_name); col_descs = []
            for col in meta.get('columns', []):
                desc_short = (col.get('description', 'N/A')[:70] + '...') if len(col.get('description', 'N/A')) > 73 else col.get('description', 'N/A')
                col_descs.append(f"- `{col.get('name', 'UNKNOWN_COL')}` (Type: {col.get('type', 'UNK')}, Desc: {desc_short})")
            sheet_info = f" (Sheet: '{meta.get('source_sheet')}')" if meta.get('source_sheet') else ""
            schema_overview += f"\nTable: `{table_name}` (From: '{meta.get('source_file', 'N/A')}'{sheet_info})\n" + "\n".join(col_descs) + "\n"
        if not table_options: return None
        
        system_prompt = f"""You are an expert data analyst. Choose the single most relevant SQL table to answer the user's query, based on schemas.
Available Table Schemas:\n{schema_overview}
Respond ONLY with the exact name of the single most appropriate table from: {table_options}. If none, respond "NONE"."""
        samba_messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"User Query: \"{query}\"\n\nMost relevant table name:"}
        ]
        target_table_name = self._call_llm_for_analysis(samba_messages, "SQL target table identification")
        if not target_table_name: return None
        target_table_name = target_table_name.strip().strip('`"')
        if target_table_name == "NONE": return None
        if target_table_name in self.table_metadata: return target_table_name
        logger.warning(f"SambaNova returned invalid table name '{target_table_name}'. Options: {table_options}"); return None

    def _generate_sql_query(self, query, table_name, table_meta):
        if not table_meta or 'columns' not in table_meta: return None
        schema_description = f"Table: `{table_name}`\nColumns:\n"
        for col in table_meta['columns']:
             schema_description += f"- \"{col['name']}\" (Type: {col.get('type', 'TEXT')}, Desc: {col.get('description', 'N/A')})\n"
        system_prompt = f"""You are an expert SQLite query writer. Given a user query and table schema, write a valid SQLite query.
Schema:\n{schema_description}
Instructions:
- Write ONLY a single, valid SQLite SELECT query.
- Quote column names with double quotes.
- Use CAST if needed for comparisons on TEXT columns.
- If impossible, output ONLY: QUERY_CANNOT_BE_ANSWERED"""
        samba_messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"User Query: \"{query}\"\n\nSQLite Query:"}
        ]
        sql_query = self._call_llm_for_analysis(samba_messages, f"SQL query generation for {table_name}")
        if not sql_query or sql_query == "QUERY_CANNOT_BE_ANSWERED": return None
        sql_query = re.sub(r"^```(?:sql)?\s*|\s*```$", "", sql_query.strip(), flags=re.DOTALL)
        if not sql_query.upper().startswith("SELECT"): return None
        if sql_query.endswith(';'): sql_query = sql_query[:-1]
        logger.info(f"Generated SQL for '{table_name}': {sql_query}"); return sql_query

    def _execute_sql_query(self, sql_query):
        # ... (This method remains the same as it only interacts with SQLite)
        if not self.db_cursor: return None, "Database connection error."
        if not sql_query: return None, "No SQL query generated."
        try:
            logger.info(f"Executing SQL: {sql_query}"); 
            df_results = pd.read_sql_query(sql_query, self.db_conn)
            if df_results.empty: return "No results found.", None
            output_str = df_results.to_string(index=False, max_rows=25, na_rep='NULL')
            if len(df_results) > 25: output_str += f"\n... (truncated, {len(df_results)} total rows)"
            return f"Query Result ({len(df_results)} row(s)):\n" + output_str, None
        except Exception as e_sql:
            err_msg = f"SQLite execution error: {e_sql}"
            logger.error(err_msg, exc_info=False)
            return None, f"Database error: {e_sql}"


    def _synthesize_answer_from_sql_with_history(self, user_query, sql_result_str, table_name, table_meta, chat_history):
        source_desc = f"table '{table_name}'"
        if table_meta: source_desc += f" from file '{table_meta.get('source_file', 'N/A')}'"
        
        samba_history = []
        for msg in chat_history:
            if msg["role"] == "assistant":
                clean_content = re.sub(r"\n\n---\n\*Answer generated.*", "", msg["content"], flags=re.DOTALL).strip()
                samba_history.append({"role": "assistant", "content": clean_content})
            else:
                samba_history.append(msg)

        system_content_template = """You are the SatSure LLM. Data was retrieved from {source_desc}.
Formulate a natural language answer to the LATEST user question based ONLY on the provided SQL query results and conversation history.
Do NOT repeat the raw data. Summarize or present key information. If data is incomplete, mention it.
If data doesn't answer, say so (e.g., "I don't really know..."). Start answer directly. No hallucination."""
        
        if sql_result_str == "No results found.":
            human_content = f"Original User Question: \"{user_query}\"\nData from {source_desc}: No results found.\nFormulate a 'not found' answer, considering history. Use \"I looked in the data but couldn't find...\" or \"I don't really know...\" if it fits."
        else:
            human_content = f"Original User Question: \"{user_query}\"\nData from {source_desc}:\n--- START SQL RESULT ---\n{sql_result_str}\n--- END SQL RESULT ---\nFormulate final answer based ONLY on this data and history:"

        messages_for_samba = [{"role": "system", "content": system_content_template.format(source_desc=source_desc)}]
        messages_for_samba.extend(samba_history)
        messages_for_samba.append({"role": "user", "content": human_content})
        
        final_answer = self._call_llm_for_analysis(messages_for_samba, "SQL result synthesis")

        if not final_answer:
            return "I had trouble summarizing the findings. " + (f"SQL query found no results from {source_desc}." if sql_result_str == "No results found." else f"Here's the raw data from {source_desc}:\n```\n{sql_result_str}\n```")
        return final_answer

    def _refine_query_with_history(self, current_query, chat_history):
        if not chat_history: return current_query
        samba_history = [] # Already in correct format from st.session_state.messages
        for msg in chat_history:
            if msg["role"] == "assistant":
                clean_content = re.sub(r"\n\n---\n\*Answer generated.*", "", msg["content"], flags=re.DOTALL).strip()
                samba_history.append({"role": "assistant", "content": clean_content})
            else:
                samba_history.append(msg)

        system_prompt = """You are an expert query assistant for SatSure LLM. Rewrite the latest user query to be standalone, incorporating context from history. If already clear, return it as is. Output ONLY the refined query."""
        
        messages_for_samba = [{"role": "system", "content": system_prompt}]
        messages_for_samba.extend(samba_history)
        messages_for_samba.append({"role": "user", "content": f"Latest user query: \"{current_query}\"\n\nRefined standalone query:"})

        refined_query_llm = self._call_llm_for_analysis(messages_for_samba, "query refinement")
        
        if refined_query_llm and refined_query_llm.strip() and refined_query_llm != current_query and len(refined_query_llm) > 5:
            logger.info(f"Query refined by SambaNova to: '{refined_query_llm[:100]}...'")
            return refined_query_llm
        logger.info("Query deemed standalone or SambaNova refinement failed/empty, using original.")
        return current_query

    # ==========================================================================
    # run_query_with_history FUNCTION (Calls adapted helpers)
    # ==========================================================================
    def run_query_with_history(self, query, chat_history, context_strategy="top_k"):
        # ... (This method's high-level logic remains similar, but it calls the SambaNova-adapted sub-methods) ...
        # The core flow: refine -> classify -> (SQL or Text RAG) -> synthesize will use the new SambaNova calls.
        logger.info(f"--- Starting SambaNova query: '{query[:100]}...' ---")
        history_limit = self.config.max_chat_history_turns * 2 
        limited_history = chat_history[-history_limit:]

        final_results = {
            "query": query, "refined_query": None, "classification": None, "classification_confidence": None,
            "answer": "", "answer_source": None, "status": "Started", "error": None
        }
        is_streamlit = self._is_streamlit
        if is_streamlit: import streamlit as st

        try:
            refined_query = query
            if limited_history:
                spinner_msg = "Understanding query context (SambaNova)..."
                if is_streamlit: 
                    with st.spinner(spinner_msg): refined_query = self._refine_query_with_history(query, limited_history)
                else: refined_query = self._refine_query_with_history(query, limited_history)
            final_results["refined_query"] = refined_query

            text_retrieval_results = {}
            if self.faiss_indexes:
                 if is_streamlit: 
                     with st.spinner("Pre-searching documents..."): text_retrieval_results = self.query_files(refined_query)
                 else: text_retrieval_results = self.query_files(refined_query)
                 final_results["retrieval_results"] = text_retrieval_results
            
            classification, confidence = "Simple Retrieval", 0.5
            spinner_msg = "Analyzing query (SambaNova)..."
            if is_streamlit: 
                with st.spinner(spinner_msg): classification, confidence = self._classify_query(refined_query)
            else: classification, confidence = self._classify_query(refined_query)
            final_results["classification"], final_results["classification_confidence"] = classification, confidence

            sql_succeeded = False
            if classification == "Structured Query (SQL)" and confidence >= self.config.dataframe_query_confidence_threshold and self.table_metadata:
                logger.info("Attempting SQL path with SambaNova assistance.")
                target_table_name = self._identify_target_sql_table(refined_query)
                if target_table_name and target_table_name in self.table_metadata:
                    target_meta = self.table_metadata[target_table_name]
                    sql_query = self._generate_sql_query(refined_query, target_table_name, target_meta)
                    if sql_query:
                        sql_result_str, sql_error = self._execute_sql_query(sql_query)
                        if not sql_error and sql_result_str:
                            final_answer = self._synthesize_answer_from_sql_with_history(refined_query, sql_result_str, target_table_name, target_meta, limited_history)
                            if final_answer and not final_answer.lower().startswith("error:"):
                                final_results.update({"answer": final_answer, "status": "Completed Successfully", "answer_source": "SQL Query (SambaNova assisted)"})
                                sql_succeeded = True
                            else: final_results["error"] = "Failed to synthesize SQL answer with SambaNova."
                        else: final_results["error"] = sql_error or "SQL execution returned no result."
                    else: final_results["error"] = "Failed to generate SQL query with SambaNova."
                else: final_results["error"] = "Could not identify SQL table with SambaNova."
            
            if not sql_succeeded:
                logger.info("Proceeding with Text RAG using SambaNova.")
                if not self.faiss_indexes:
                    final_results.update({"answer": "I don't know (no text documents indexed).", "status": "Failed", "error": "No text documents indexed."})
                    return final_results
                
                current_retrieval_results = final_results.get("retrieval_results", {})
                if not current_retrieval_results or not any(current_retrieval_results.values()): # If pre-search was empty
                    current_retrieval_results = self.query_files(refined_query)
                    final_results["retrieval_results"] = current_retrieval_results


                if not current_retrieval_results or not any(current_retrieval_results.values()):
                    final_results.update({"answer": "I don't know (no relevant text found).", "status": "Completed: No relevant info."})
                    return final_results

                aggregated_context_data = self.aggregate_context(current_retrieval_results, strategy=context_strategy)
                final_results["aggregated_context_data"] = aggregated_context_data

                if not aggregated_context_data or not aggregated_context_data.get("combined_context"):
                    final_results.update({"answer": "I don't know (failed to aggregate text).", "status": "Failed: Context aggregation."})
                    return final_results

                spinner_msg = "Generating answer from proposals (SambaNova)..."
                answer = ""
                if is_streamlit: 
                    with st.spinner(spinner_msg): answer = self.query_llm_with_history(refined_query, aggregated_context_data, limited_history)
                else: answer = self.query_llm_with_history(refined_query, aggregated_context_data, limited_history)

                if answer and not answer.lower().startswith("error:"):
                    final_results.update({"answer": answer, "status": "Completed Successfully", "answer_source": "Text RAG (SambaNova)"})
                else:
                    final_results.update({"answer": answer or "Failed to get response from SambaNova for text RAG.", "status": "Failed Text RAG LLM Call"})
            
            logger.info(f"--- SambaNova query finished. Source: {final_results.get('answer_source')} --- Status: {final_results.get('status')} ---")

        except Exception as e:
            logger.error(f"Error in run_query_with_history (SambaNova): {e}", exc_info=True)
            final_results.update({"status": "Failed", "error": str(e), "answer": f"An unexpected error occurred: {e}"})
        return final_results

# --- Streamlit UI Code (largely unchanged, but config display needs update) ---
# The Streamlit UI part from ragtest_c2-9.py can be used, but you'll need to modify
# the `config_details` in the sidebar to show `sambanova_model` instead of Groq/DeepSeek.
# And ensure SAMBANOVA_API_KEY is checked in `load_rag_system`.

if __name__ == "__main__":
    try:
        import streamlit as st
        from packaging.version import parse as parse_version 
    except ImportError:
        logger.error("Streamlit or packaging not installed.")
        sys.exit(1)
    
    st.set_page_config(layout="wide", page_title="SatSure Conversational RAG Q&A (SambaNova)")

    st.markdown("""
    <style>
    /* ... (Your existing CSS from previous response) ... */
    </style>
    <script>
    /* ... (Your existing JS from previous response) ... */
    </script>
    """, unsafe_allow_html=True)

    st.title("SatSureLLM")
    st.caption("Query Project Proposals & Supporting Data")

    @st.cache_resource
    def load_rag_system(config_path="config1.ini"):
        # Check for SAMBANOVA_API_KEY early
        sambanova_key = os.getenv("SAMBANOVA_API_KEY")
        # Also check from config as a fallback, though .env is preferred
        if not sambanova_key:
            temp_config_parser = configparser.ConfigParser()
            if os.path.exists(config_path):
                temp_config_parser.read(config_path)
                sambanova_key = temp_config_parser.get("API_KEYS", "sambanova_api_key_config", fallback=None)
        
        if not sambanova_key:
            st.error("🚨 SAMBANOVA_API_KEY not set! The application cannot run effectively."); 
            st.info("Please set your SAMBANOVA_API_KEY in .env or config1.ini and restart.")
            st.stop()
        
        system = None
        with st.spinner("Initializing RAG System & loading models..."):
            try:
                system = RAGSystem(config_path=config_path) # RAGSystem now handles its own key check message
            except Exception as e:
                logger.error(f"Fatal RAG Initialization Error: {e}", exc_info=True)
                st.error(f"Fatal RAG Init Error: {e}")
                st.stop()
        
        if not system or not system.db_conn: 
            st.error("🚨 RAG System or SQLite connection failed critically during initialization! Check logs."); 
            st.stop()
        
        st.sidebar.success("✅ RAG System Initialized.")
        return system
    
    def display_document_management_ui(_rag_system):
        with st.sidebar.expander("📂 Document Management & Processing", expanded=False):
            st.markdown(f"**Data Folder:** `{os.path.abspath(_rag_system.config.data_dir)}`")
            st.markdown(f"**Index Folder:** `{os.path.abspath(_rag_system.config.index_dir)}`")
            st.markdown(f"**DB Location:** `{_rag_system.config.sqlite_db_path}`")
            st.markdown("---")
            status_placeholder_sidebar = st.empty() 

            def sidebar_progress_callback(message, is_error=False, is_warning=False, is_done=False, current_step=None, total_steps=None, stage=None):
                log_prefix = "✅ " if is_done else "❌ " if is_error else "⚠️ " if is_warning else f"⏳ ({stage}) " if stage else "⏳ "
                progress_val = 1.0 if is_done else 0.0
                if total_steps and current_step is not None:
                    stage_prog = {"Extracting Text": 0.1, "Chunking Text": 0.15, "Verifying Index": 0.2, "Embedding Text": 0.3, "Indexing Text": 0.4,
                                "Loading Sheet": 0.5, "Analyzing Columns": 0.7, "Loading SQL": 0.9}.get(stage, 0.0)
                    progress_val = min(1.0, (current_step + stage_prog) / total_steps) if total_steps > 0 else 0.0
                full_message = f"{log_prefix}{message}"
                try:
                    if 0 < progress_val < 1.0 and not is_error and not is_warning: status_placeholder_sidebar.progress(progress_val, text=full_message)
                    else:
                        if is_error: status_placeholder_sidebar.error(full_message)
                        elif is_warning: status_placeholder_sidebar.warning(full_message)
                        elif is_done: status_placeholder_sidebar.success(full_message)
                        else: status_placeholder_sidebar.info(full_message) 
                except Exception as e: logger.warning(f"Streamlit UI update error in sidebar callback: {e}")

            if st.button("🔄 Process/Re-process Documents", key="process_docs_sidebar_expander", use_container_width=True):
                st.session_state.sidebar_processed_files_list = [] 
                st.session_state.sidebar_processing_status_message = "Starting processing..." 
                st.session_state.is_ready_for_chat = False 
                status_placeholder_sidebar.info(st.session_state.sidebar_processing_status_message)
                with st.spinner("Processing documents..."):
                    try:
                        processing_successful = _rag_system.process_files(progress_callback=sidebar_progress_callback)
                        st.session_state.sidebar_processed_files_list = sorted(list(_rag_system.processed_files))
                        text_idx_count = len(_rag_system.faiss_indexes); sql_tbl_count = len(_rag_system.table_metadata)
                        processed_count = len(st.session_state.sidebar_processed_files_list)
                        final_msg = f"Processed {processed_count} file(s). Text Indices: {text_idx_count}, SQL Tables: {sql_tbl_count}."
                        if processed_count > 0 and processing_successful: 
                            st.session_state.sidebar_processing_status_message = f"✅ Ready! {final_msg}"; st.session_state.is_ready_for_chat = True 
                        elif processing_successful: st.session_state.sidebar_processing_status_message = f"⚠️ Finished, but no supported documents processed. {final_msg}"; st.session_state.is_ready_for_chat = False 
                        else: st.session_state.sidebar_processing_status_message = f"❌ Processing function reported failure."; st.session_state.is_ready_for_chat = False
                        
                        if "✅" in st.session_state.sidebar_processing_status_message: status_placeholder_sidebar.success(st.session_state.sidebar_processing_status_message)
                        elif "⚠️" in st.session_state.sidebar_processing_status_message: status_placeholder_sidebar.warning(st.session_state.sidebar_processing_status_message)
                        else: status_placeholder_sidebar.error(st.session_state.sidebar_processing_status_message)
                        st.rerun() 
                    except Exception as e:
                        st.session_state.sidebar_processing_status_message = f"❌ Fatal error during processing: {e}."
                        status_placeholder_sidebar.error(st.session_state.sidebar_processing_status_message)
                        st.session_state.is_ready_for_chat = False; st.rerun()
            else:
                current_msg = st.session_state.get("sidebar_processing_status_message", "Click 'Process/Re-process' to load data.")
                if "✅" in current_msg: status_placeholder_sidebar.success(current_msg)
                elif "❌" in current_msg: status_placeholder_sidebar.error(current_msg)
                elif "⚠️" in current_msg: status_placeholder_sidebar.warning(current_msg)
                else: status_placeholder_sidebar.info(current_msg)
            if st.session_state.get("sidebar_processed_files_list"): 
                st.markdown("--- \n**Loaded Assets:**")
                assets_df_data = [{"File": fn, "Type(s)": ", ".join(filter(None, [("Text Index" if fn in _rag_system.faiss_indexes else None), (f"SQL ({len(_rag_system.file_to_table_map.get(fn,[]))} tables)" if fn in _rag_system.file_to_table_map and _rag_system.file_to_table_map.get(fn) else None)])) or "N/A"} for fn in st.session_state.sidebar_processed_files_list]
                if assets_df_data: st.dataframe(pd.DataFrame(assets_df_data), use_container_width=True, hide_index=True)
        return st.session_state.get("is_ready_for_chat", False)

    try:
        rag_sys = load_rag_system() 

        if "messages" not in st.session_state: st.session_state.messages = []
        if "analysis_details" not in st.session_state: st.session_state.analysis_details = {}
        if "is_ready_for_chat" not in st.session_state: st.session_state.is_ready_for_chat = False 
        if "sidebar_processed_files_list" not in st.session_state: st.session_state.sidebar_processed_files_list = []
        if "sidebar_processing_status_message" not in st.session_state: st.session_state.sidebar_processing_status_message = "Click 'Process/Re-process Documents' to load data."
        if "query_start_time" not in st.session_state: st.session_state.query_start_time = time.time()

        st.sidebar.title("Controls & Info")
        if st.sidebar.button("New Chat", use_container_width=True, key="new_chat_button"):
            st.session_state.messages = []; st.session_state.analysis_details = {}; st.rerun()

        is_ready = display_document_management_ui(rag_sys) 

        st.sidebar.markdown("---")
        st.sidebar.subheader("System Configuration")
        config_details = {
            "LLM (LLama 3.3 Instruct 70B)": rag_sys.config.sambanova_model, # UPDATED
            "Encoder": rag_sys.config.encoder_model,
            "History Turns": rag_sys.config.max_chat_history_turns,
            "Retr. K (Text)": rag_sys.config.k_retrieval,
            "SQL Conf.": f"{rag_sys.config.dataframe_query_confidence_threshold:.2f}",
        }
        for key, value in config_details.items():
            st.sidebar.markdown(f"<small><b>{key}:</b> {value}</small>", unsafe_allow_html=True)

        if not is_ready:
             st.info("⬅️ Please process documents using 'Document Management & Processing' in the sidebar to enable chat.")
             st.stop() 
        
        chat_container = st.container() 
        with chat_container:
            for i, message in enumerate(st.session_state.messages):
                avatar_icon = "👤" if message["role"] == "user" else "🛰️" # SambaNova icon maybe?
                with st.chat_message(message["role"], avatar=avatar_icon): 
                    content_part = message["content"]
                    footer_part = ""
                    if message["role"] == "assistant" and "\n\n---\n*" in message["content"]:
                        parts = message["content"].split("\n\n---\n*", 1)
                        content_part = parts[0]
                        if len(parts) > 1: footer_part = f"<div class='message-footer'>*{parts[1]}</div>"
                    st.markdown(content_part, unsafe_allow_html=True)
                    if footer_part: st.markdown(footer_part, unsafe_allow_html=True)

                    if message["role"] == "assistant" and i == len(st.session_state.messages) - 1 and st.session_state.get('analysis_details', {}).get('has_details'):
                        with st.expander("Show Analysis & Context", expanded=False):
                            # ... (analysis details display, mostly unchanged logic but source would be SambaNova)
                            details = st.session_state.analysis_details
                            if details.get("Original Query"): st.caption(f"**Original Query:** {details['Original Query']}")
                            if details.get("Refined Query") and details["Refined Query"] != details.get("Original Query"): st.caption(f"**Refined Query:** {details['Refined Query']}")
                            if details.get("Classification"): st.caption(f"**Classification:** {details['Classification']}")
                            if details.get("Answer Source"): st.caption(f"**Answer Source:** {details['Answer Source']}") # Will now show SambaNova
                            if details.get("Status"): st.caption(f"**Status:** {details['Status']}")
                            if details.get("SQL Table"):
                                st.caption(f"**SQL Table:** `{details['SQL Table']}`")
                                st.caption("**SQL Query Generated:**"); st.code(details.get("SQL Query", "N/A"), language='sql')
                            text_sources = details.get("Text Sources Used")
                            if text_sources: st.caption(f"**Text Sources Used:** {', '.join(text_sources)}")
                            if details.get("Error"): st.caption(f"**Note/Error:** {details['Error']}")


        if prompt := st.chat_input("Ask a question about SatSure"):
            st.session_state.messages.append({"role": "user", "content": prompt})
            st.session_state.query_start_time = time.time() 
            
            with st.spinner("🛰️ SatSureLLM is thinking..."):
                history_for_rag = st.session_state.messages[:-1] # Already in {"role": ..., "content": ...}
                results_data = rag_sys.run_query_with_history(prompt, history_for_rag, context_strategy="top_k")
                
                query_end_time = time.time()
                elapsed_time = query_end_time - st.session_state.query_start_time
                elapsed_time_str = f"{elapsed_time:.2f}"

                answer = results_data.get("answer", "Sorry, I encountered an error with SambaNova.")
                answer_source = results_data.get('answer_source', 'N/A (SambaNova)')
                query_status = results_data.get("status", "Unknown")
                query_error = results_data.get("error")
                
                full_response = answer
                response_footer_parts = [f"Answer via {answer_source} in {elapsed_time_str}s."]
                is_satsure_dont_know = answer.strip().startswith("I don't really know but my next versions")

                if "Completed Successfully" not in query_status or (query_error and not is_satsure_dont_know):
                    if query_status != "Completed Successfully": response_footer_parts.append(f"Status: {query_status}")
                    if query_error and query_error not in full_response and not is_satsure_dont_know : response_footer_parts.append(f"Note: {query_error}")
                
                if not is_satsure_dont_know or len(response_footer_parts) > 1 : 
                     full_response += f"\n\n---\n*{' | '.join(response_footer_parts)}"
            
            st.session_state.messages.append({"role": "assistant", "content": full_response})
            analysis_info = { "has_details": True, **results_data } # Simpler way to pass all results
            st.session_state.analysis_details = {k: v for k, v in analysis_info.items() if v is not None and str(v).strip() != ""}
            st.rerun()

    except Exception as e:
        st.error(f"Streamlit App Error: {e}")
        logger.exception("Streamlit application error:")
        log_file_path = 'rag_system_sambanova.log' 
        try:
            if 'rag_sys' in locals() and rag_sys and hasattr(rag_sys, 'config'): 
                log_file_path = rag_sys.config.log_file
        except Exception: pass 
        st.info(f"Check console or `{os.path.abspath(log_file_path)}` for details.")
        st.stop()
